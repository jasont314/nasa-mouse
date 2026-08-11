"""Build the approved-template ASGSR expiMap scientific poster."""

from __future__ import annotations

import shutil
import subprocess
import tempfile
import textwrap
from io import BytesIO
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.lines import Line2D
import numpy as np
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from .integrate_reassessed_tissues_paper import (
    MAIN_TISSUES,
    RETAINED_TERMS,
    load_retained_evidence,
    project_effects,
)


ROOT = Path(__file__).resolve().parents[2]
PAPER_DIR = ROOT / "paper/asgsr_expimap_hvg"
FIGURE_DIR = PAPER_DIR / "figures"
POSTER_DIR = PAPER_DIR / "poster"
ASSET_DIR = POSTER_DIR / "assets"
TEMPLATE_PATH = (
    ROOT
    / "assets/poster_template/00 Poster Session Student Template_Approved by Legal.pptx"
)

TITLE = (
    "Cross-Mission expiMap Analysis Recovers Established Tissue Responses "
    "and Identifies Complementary Pathway Shifts in Mouse Spaceflight "
    "Transcriptomes"
)
TITLE_LINE_1 = (
    "Cross-Mission expiMap Analysis Recovers Established Tissue Responses"
)
TITLE_LINE_2 = (
    "and Identifies Complementary Pathway Shifts in Mouse Spaceflight Transcriptomes"
)

SLIDE_W = 48.0
SLIDE_H = 27.0
FONT = "Arial"

WHITE = "FFFFFF"
INK = "071C3F"
BODY = "33445D"
MUTED = "697889"
HEADER = "DDECF7"
PANEL = "DCEAF4"
PANEL_LIGHT = "EEF4F8"
RULE = "AEBCC5"
NAVY = "082552"
BLUE = "236DA2"
GOLD = "9A6D1D"
GREEN = "178C6A"
ORANGE = "D97800"
THYMUS = "6A4C93"
SKIN = "C65D37"
LIVER = "13877C"
SPLEEN = "A64D67"

POSTER_TERMS = {
    "thymus": RETAINED_TERMS["thymus"],
    "skin": RETAINED_TERMS["skin"],
    "liver": RETAINED_TERMS["liver"],
    "spleen": RETAINED_TERMS["spleen"],
}


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_fill(shape, color: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)


def set_line(shape, color: str, width: float = 1.0) -> None:
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)


def keep_words_intact(paragraph) -> None:
    """Require PowerPoint to wrap Latin text only at word boundaries."""
    properties = paragraph._p.get_or_add_pPr()
    properties.set("latinLnBrk", "0")
    properties.set("eaLnBrk", "0")
    properties.set("hangingPunct", "0")


def add_text(
    slide,
    value: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float,
    color: str = BODY,
    bold: bool = False,
    italic: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.0,
    word_wrap: bool = True,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = word_wrap
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0
    keep_words_intact(paragraph)
    run = paragraph.add_run()
    run.text = value
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return shape


def add_text_lines(
    slide,
    values: tuple[str, ...],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float,
    color: str = BODY,
    bold: bool = False,
    italic: bool = False,
    align=PP_ALIGN.LEFT,
) -> list:
    """Author one complete, nonwrapping line per PowerPoint text box."""
    if not values or any("\n" in value for value in values):
        raise ValueError("Fixed text lines must be nonempty and contain no line breaks")
    line_height = h / len(values)
    if line_height < size / 72:
        raise ValueError("Fixed text lines do not have enough vertical space")
    shapes = []
    for index, value in enumerate(values):
        shape = add_text(
            slide,
            value,
            x,
            y + index * line_height,
            w,
            line_height,
            size=size,
            color=color,
            bold=bold,
            italic=italic,
            align=align,
            valign=MSO_ANCHOR.MIDDLE,
            word_wrap=False,
        )
        shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        shapes.append(shape)
    return shapes


def add_paragraphs(
    slide,
    values: list[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float,
    color: str = BODY,
    gap: float = 4.0,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    for index, value in enumerate(values):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(gap)
        paragraph.line_spacing = 1.0
        keep_words_intact(paragraph)
        run = paragraph.add_run()
        run.text = value
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = rgb(color)
    return shape


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str,
    *,
    line: str | None = None,
    line_width: float = 1.0,
    rounded: bool = False,
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        kind, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_fill(shape, fill)
    if line:
        set_line(shape, line, line_width)
    else:
        shape.line.fill.background()
    return shape


def add_panel(slide, x: float, y: float, w: float, h: float) -> None:
    add_rect(slide, x, y, w, h, PANEL)


def add_section(
    slide, label: str, x: float, y: float, w: float, *, size: float = 26
) -> None:
    add_rect(slide, x, y, w, 0.72, WHITE, line=RULE, line_width=1.1)
    add_text(
        slide,
        label,
        x + 0.15,
        y + 0.02,
        w - 0.30,
        0.66,
        size=size,
        color="111111",
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        word_wrap=False,
    )


def add_box(
    slide,
    title: str,
    subtitle: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str,
    line: str,
    title_size: float = 18,
    subtitle_size: float = 14.5,
) -> None:
    shape = add_rect(
        slide, x, y, w, h, fill, line=line, line_width=1.4, rounded=True
    )
    shape.adjustments[0] = 0.06
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.margin_left = Inches(0.10)
    frame.margin_right = Inches(0.10)
    frame.margin_top = Inches(0.05)
    frame.margin_bottom = Inches(0.04)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    first = frame.paragraphs[0]
    first.alignment = PP_ALIGN.CENTER
    first.space_after = Pt(2)
    keep_words_intact(first)
    run = first.add_run()
    run.text = title
    run.font.name = FONT
    run.font.size = Pt(title_size)
    run.font.bold = True
    run.font.color.rgb = rgb(INK)
    second = frame.add_paragraph()
    second.alignment = PP_ALIGN.CENTER
    second.space_before = Pt(0)
    second.space_after = Pt(0)
    keep_words_intact(second)
    run = second.add_run()
    run.text = subtitle
    run.font.name = FONT
    run.font.size = Pt(subtitle_size)
    run.font.color.rgb = rgb(MUTED)


def add_down_arrow(slide, center_x: float, y: float, color: str = BLUE) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        Inches(center_x - 0.18),
        Inches(y),
        Inches(0.36),
        Inches(0.34),
    )
    set_fill(shape, color)
    shape.line.fill.background()


def add_right_arrow(
    slide, x: float, y: float, w: float, h: float, color: str = BLUE
) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_fill(shape, color)
    shape.line.fill.background()


def add_connector(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: str,
    width: float = 1.1,
    dashed: bool = False,
):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    set_line(connector, color, width)
    if dashed:
        connector.line.dash_style = 2
    return connector


def set_table_cell(
    cell,
    value: str,
    *,
    size: float,
    color: str,
    bold: bool,
    center: bool,
) -> None:
    cell.text = ""
    cell.margin_left = Inches(0.20 if not center else 0.07)
    cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.03)
    cell.margin_bottom = Inches(0.03)
    cell.text_frame.word_wrap = False
    cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
    keep_words_intact(paragraph)
    run = paragraph.add_run()
    run.text = value
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_scope_table(slide, x: float, y: float, w: float, h: float) -> None:
    rows = [
        ("Thymus", "1,362", "117", "5", "387", THYMUS),
        ("Skin", "2,593", "151", "4", "319", SKIN),
        ("Liver", "5,000", "197", "9", "364", LIVER),
        ("Spleen", "6,289", "100", "5", "360", SPLEEN),
    ]
    table_rows = len(rows) + 1
    shape = slide.shapes.add_table(
        table_rows, 5, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    table = shape.table
    widths = [3.10, 2.75, 2.35, 2.25, 2.65]
    for index, value in enumerate(widths):
        table.columns[index].width = Inches(value)
    for column, value in enumerate(("Tissue", "ARCHS4", "OSDR", "Projects", "Programs")):
        cell = table.cell(0, column)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(NAVY)
        set_table_cell(
            cell,
            value,
            size=14,
            color=WHITE,
            bold=True,
            center=column > 0,
        )
        if column == 0:
            cell.margin_left = Inches(0.65)
    for row_index, row in enumerate(rows, 1):
        for column, value in enumerate(row[:5]):
            cell = table.cell(row_index, column)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(WHITE if row_index % 2 else PANEL_LIGHT)
            set_table_cell(
                cell,
                value,
                size=14.5,
                color=row[5] if column == 0 else BODY,
                bold=column == 0,
                center=column > 0,
            )
            if column == 0:
                cell.margin_left = Inches(0.65)
    for row in table.rows:
        row.height = Inches(h / table_rows)


def add_nasa_logo(slide, template: Presentation) -> None:
    template_slide = template.slides[0]
    candidates = [
        shape
        for shape in template_slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        and shape.left / template.slide_width > 0.85
        and shape.top / template.slide_height < 0.25
    ]
    if not candidates:
        raise RuntimeError("NASA logo was not found in the approved template")
    logo = max(candidates, key=lambda shape: shape.width * shape.height)
    slide.shapes.add_picture(BytesIO(logo.image.blob), Inches(44.20), Inches(0.78), width=Inches(3.05))


def add_challenge_equation(slide, x: float, y: float, w: float) -> None:
    add_text(
        slide,
        "THE CENTRAL CHALLENGE",
        x,
        y,
        w,
        0.40,
        size=19,
        color=GOLD,
        bold=True,
    )
    add_text(
        slide,
        "Observed expression",
        x,
        y + 0.48,
        2.70,
        0.60,
        size=17,
        color=INK,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        word_wrap=False,
    )
    parts = [
        ("Spaceflight biology", ORANGE, 2.75),
        ("Tissue biology", NAVY, 2.20),
        ("Animal variation", MUTED, 2.35),
        ("Study / protocol", GOLD, 2.55),
    ]
    cursor = x + 2.85
    for index, (label, color, width) in enumerate(parts):
        if index:
            add_text(
                slide,
                "+",
                cursor,
                y + 0.48,
                0.38,
                0.60,
                size=22,
                color=MUTED,
                bold=True,
                align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE,
                word_wrap=False,
            )
            cursor += 0.43
        add_rect(slide, cursor, y + 0.48, width, 0.60, color, rounded=True)
        add_text(
            slide,
            label,
            cursor + 0.07,
            y + 0.50,
            width - 0.14,
            0.54,
            size=15.5,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
            word_wrap=False,
        )
        cursor += width


def add_architecture(slide, x: float, y: float, w: float, h: float) -> None:
    add_text(
        slide,
        "Pathway structure is wired into the decoder",
        x,
        y,
        w,
        0.50,
        size=22,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
        word_wrap=False,
    )
    add_text(
        slide,
        "Train a tissue reference, then map flight and ground samples.",
        x + 0.35,
        y + 0.52,
        w - 0.70,
        0.46,
        size=16,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        word_wrap=False,
    )

    add_text(
        slide,
        "1  REFERENCE TRAINING",
        x + 0.15,
        y + 1.03,
        4.00,
        0.35,
        size=17,
        color=BLUE,
        bold=True,
    )
    add_box(
        slide,
        "ARCHS4 reference",
        "non-spaceflight tissue counts",
        x + 0.15,
        y + 1.45,
        3.55,
        1.10,
        fill="E9F3F8",
        line=BLUE,
        title_size=18,
        subtitle_size=14,
    )
    add_right_arrow(slide, x + 3.88, y + 1.79, 0.55, 0.36, BLUE)

    model = add_rect(
        slide,
        x + 4.55,
        y + 1.28,
        10.20,
        2.65,
        WHITE,
        line=BLUE,
        line_width=1.2,
        rounded=True,
    )
    model.adjustments[0] = 0.04
    add_text(
        slide,
        "expiMap tissue model",
        x + 4.82,
        y + 1.42,
        4.40,
        0.36,
        size=17.5,
        color=INK,
        bold=True,
        word_wrap=False,
    )

    input_x = x + 4.95
    encoder_x = x + 5.48
    encoder_w = 1.48
    latent_x = x + 7.42
    decoder_x = x + 10.58
    decoder_w = 1.48
    output_x = x + 13.55
    gene_ys = [y + 2.05 + 0.42 * index for index in range(4)]
    for gene_y in gene_ys:
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(input_x),
            Inches(gene_y),
            Inches(0.24),
            Inches(0.24),
        )
        set_fill(node, "A9BBD0")
        node.line.fill.background()

    encoder = slide.shapes.add_shape(
        MSO_SHAPE.TRAPEZOID,
        Inches(encoder_x),
        Inches(y + 1.92),
        Inches(encoder_w),
        Inches(1.42),
    )
    encoder.rotation = 90
    set_fill(encoder, "DCE6F2")
    set_line(encoder, BLUE, 1.4)
    add_text(
        slide,
        "Encoder",
        encoder_x - 0.05,
        y + 2.42,
        encoder_w + 0.10,
        0.38,
        size=14,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        word_wrap=False,
    )
    for gene_y in gene_ys:
        add_connector(
            slide,
            input_x + 0.24,
            gene_y + 0.12,
            encoder_x + 0.04,
            y + 2.63,
            color="A9BBD0",
            width=0.8,
        )

    latent_rows = [
        ("DNA repair", THYMUS, y + 1.95),
        ("T-cell signaling", SPLEEN, y + 2.43),
        ("Cell junctions", SKIN, y + 2.91),
    ]
    for label, color, node_y in latent_rows:
        add_connector(
            slide,
            encoder_x + encoder_w - 0.04,
            y + 2.63,
            latent_x,
            node_y + 0.16,
            color="A9BBD0",
            width=0.9,
        )
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(latent_x),
            Inches(node_y),
            Inches(0.32),
            Inches(0.32),
        )
        set_fill(node, color)
        node.line.fill.background()
        label_box = add_rect(
            slide,
            latent_x + 0.40,
            node_y - 0.02,
            2.22,
            0.36,
            "F8FAFC",
            line=color,
            line_width=0.9,
            rounded=True,
        )
        label_box.adjustments[0] = 0.08
        add_text(
            slide,
            label,
            latent_x + 0.50,
            node_y - 0.02,
            2.02,
            0.36,
            size=15.5,
            color=color,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
            word_wrap=False,
        )

    output_ys = [y + 2.03 + 0.38 * index for index in range(4)]
    connection_map = [(0, 0), (0, 2), (1, 1), (1, 3), (2, 2), (2, 3)]
    for latent_index, output_index in connection_map:
        _, color, node_y = latent_rows[latent_index]
        add_connector(
            slide,
            latent_x + 2.62,
            node_y + 0.16,
            output_x,
            output_ys[output_index] + 0.12,
            color=color,
            width=1.2,
        )

    decoder = slide.shapes.add_shape(
        MSO_SHAPE.TRAPEZOID,
        Inches(decoder_x),
        Inches(y + 1.92),
        Inches(decoder_w),
        Inches(1.42),
    )
    decoder.rotation = 270
    decoder.fill.background()
    set_line(decoder, GREEN, 1.4)
    add_text(
        slide,
        "Decoder",
        decoder_x + 0.14,
        y + 1.48,
        1.20,
        0.34,
        size=13.5,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        word_wrap=False,
    )

    for out_y in output_ys:
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(output_x),
            Inches(out_y),
            Inches(0.24),
            Inches(0.24),
        )
        set_fill(node, "A9BBD0")
        node.line.fill.background()

    add_text(
        slide,
        "genes",
        input_x - 0.18,
        y + 3.48,
        0.66,
        0.35,
        size=14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        word_wrap=False,
    )
    add_text(
        slide,
        "program scores",
        latent_x - 0.05,
        y + 3.48,
        2.65,
        0.35,
        size=14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        word_wrap=False,
    )
    add_text(
        slide,
        "masked decoder",
        decoder_x - 0.18,
        y + 3.48,
        decoder_w + 0.36,
        0.35,
        size=14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        word_wrap=False,
    )
    add_text(
        slide,
        "output genes",
        x + 12.85,
        y + 3.48,
        1.65,
        0.35,
        size=14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        word_wrap=False,
    )

    add_text(
        slide,
        "2  QUERY MAPPING AND CROSS-MISSION COMPARISON",
        x + 0.15,
        y + 4.08,
        8.10,
        0.36,
        size=17,
        color=ORANGE,
        bold=True,
    )
    add_box(
        slide,
        "NASA OSDR query",
        "FLT, GC, and accession",
        x + 0.15,
        y + 4.48,
        2.55,
        1.08,
        fill="FAEFE7",
        line=ORANGE,
        title_size=16.5,
        subtitle_size=13.5,
    )
    add_right_arrow(slide, x + 2.82, y + 4.82, 0.38, 0.36, ORANGE)
    add_box(
        slide,
        "Reference-query map",
        "scArches adaptation",
        x + 3.32,
        y + 4.48,
        2.62,
        1.08,
        fill=WHITE,
        line=ORANGE,
        title_size=15.5,
        subtitle_size=13.5,
    )
    add_right_arrow(slide, x + 6.06, y + 4.82, 0.38, 0.36, BLUE)
    add_box(
        slide,
        "Program scores",
        "posterior mean per sample",
        x + 6.56,
        y + 4.48,
        2.32,
        1.08,
        fill=WHITE,
        line=NAVY,
        title_size=15.5,
        subtitle_size=13,
    )
    add_right_arrow(slide, x + 9.00, y + 4.82, 0.36, 0.36, BLUE)
    add_box(
        slide,
        "Mean FLT - mean GC",
        "within each project",
        x + 9.48,
        y + 4.48,
        2.62,
        1.08,
        fill=WHITE,
        line=BLUE,
        title_size=15.5,
        subtitle_size=13.5,
    )
    add_right_arrow(slide, x + 12.20, y + 4.82, 0.34, 0.36, BLUE)

    heatmap_x = x + 12.64
    add_rect(
        slide,
        heatmap_x,
        y + 4.48,
        2.11,
        1.08,
        WHITE,
        line=BLUE,
        line_width=1.1,
        rounded=True,
    )
    add_text(
        slide,
        "Shift heatmap",
        heatmap_x + 0.10,
        y + 4.53,
        1.91,
        0.26,
        size=13.5,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
        word_wrap=False,
    )
    heatmap_colors = (
        ("3F7EAE", "C4DCE9", "F2C1A7", "CE5B36"),
        ("78A7C7", "E7EFF3", "DD815A", "F4D6C6"),
        ("2F6F9F", "A9C9DB", "F7E8E1", "D66A43"),
        ("D8E7EF", "4F89B5", "E9A17F", "C64E2A"),
    )
    cell_w, cell_h = 0.31, 0.13
    cell_gap_x, cell_gap_y = 0.06, 0.04
    grid_x, grid_y = heatmap_x + 0.31, y + 4.84
    for row_index, row in enumerate(heatmap_colors):
        for column_index, color in enumerate(row):
            add_rect(
                slide,
                grid_x + column_index * (cell_w + cell_gap_x),
                grid_y + row_index * (cell_h + cell_gap_y),
                cell_w,
                cell_h,
                color,
            )

    add_down_arrow(slide, x + 13.69, y + 5.60, GREEN)
    add_box(
        slide,
        "LLM-assisted literature annotation",
        "literature roles, source verified",
        x + 9.10,
        y + 5.94,
        5.65,
        0.62,
        fill="EBF5EF",
        line=GREEN,
        title_size=14.5,
        subtitle_size=13,
    )

    add_text_lines(
        slide,
        (
            "Validation: held-out projects and three complete trainings; conventional enrichment,",
            "member-gene expression, and composition sensitivity.",
        ),
        x + 0.20,
        y + 5.96,
        8.55,
        0.60,
        size=13,
        color=MUTED,
        italic=True,
        align=PP_ALIGN.LEFT,
    )


def add_reference_query_mapping(slide, x: float, y: float, w: float) -> None:
    """Draw compact reference training and query mapping beneath the scope table."""
    add_text(
        slide,
        "1  REFERENCE TRAINING",
        x,
        y,
        w,
        0.30,
        size=15,
        color=BLUE,
        bold=True,
        word_wrap=False,
    )
    add_box(
        slide,
        "ARCHS4 reference",
        "non-spaceflight tissue counts",
        x,
        y + 0.36,
        3.05,
        1.00,
        fill="E9F3F8",
        line=BLUE,
        title_size=15.5,
        subtitle_size=12.5,
    )
    add_right_arrow(slide, x + 3.18, y + 0.68, 0.38, 0.34, BLUE)

    model_x, model_y = x + 3.70, y + 0.28
    model_w, model_h = w - 3.70, 1.18
    add_rect(
        slide,
        model_x,
        model_y,
        model_w,
        model_h,
        WHITE,
        line=BLUE,
        line_width=1.1,
        rounded=True,
    )
    add_text(
        slide,
        "expiMap reference with Reactome-masked decoder",
        model_x + 0.15,
        model_y + 0.07,
        model_w - 0.30,
        0.26,
        size=13.5,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
        word_wrap=False,
    )
    input_x = model_x + 0.45
    latent_x = model_x + 3.45
    decoder_x = model_x + 5.70
    output_x = model_x + 8.55
    node_ys = [model_y + 0.48 + 0.20 * index for index in range(3)]
    output_ys = [model_y + 0.45 + 0.18 * index for index in range(4)]
    colors = (THYMUS, SPLEEN, SKIN)
    for node_y in node_ys:
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(input_x),
            Inches(node_y),
            Inches(0.16),
            Inches(0.16),
        )
        set_fill(node, "A9BBD0")
        node.line.fill.background()
    for index, node_y in enumerate(node_ys):
        latent = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(latent_x),
            Inches(node_y),
            Inches(0.19),
            Inches(0.19),
        )
        set_fill(latent, colors[index])
        latent.line.fill.background()
        for input_y in node_ys:
            add_connector(
                slide,
                input_x + 0.16,
                input_y + 0.08,
                latent_x,
                node_y + 0.09,
                color="BCC9D3",
                width=0.55,
            )
    connection_map = ((0, 0), (0, 2), (1, 1), (1, 3), (2, 2), (2, 3))
    for latent_index, output_index in connection_map:
        add_connector(
            slide,
            latent_x + 0.19,
            node_ys[latent_index] + 0.09,
            output_x,
            output_ys[output_index] + 0.08,
            color=colors[latent_index],
            width=0.85,
        )
    decoder = slide.shapes.add_shape(
        MSO_SHAPE.TRAPEZOID,
        Inches(decoder_x),
        Inches(model_y + 0.40),
        Inches(0.95),
        Inches(0.67),
    )
    decoder.rotation = 270
    decoder.fill.background()
    set_line(decoder, GREEN, 1.1)
    for output_y in output_ys:
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(output_x),
            Inches(output_y),
            Inches(0.16),
            Inches(0.16),
        )
        set_fill(node, "A9BBD0")
        node.line.fill.background()
    add_text(
        slide,
        "2  QUERY MAPPING",
        x,
        y + 1.68,
        w,
        0.30,
        size=15,
        color=ORANGE,
        bold=True,
        word_wrap=False,
    )
    add_box(
        slide,
        "NASA OSDR query",
        "FLT, GC, and accession",
        x,
        y + 2.02,
        3.25,
        0.92,
        fill="FAEFE7",
        line=ORANGE,
        title_size=15,
        subtitle_size=12.5,
    )
    add_right_arrow(slide, x + 3.40, y + 2.31, 0.38, 0.32, ORANGE)
    add_box(
        slide,
        "scArches adaptation",
        "map into the trained reference",
        x + 3.92,
        y + 2.02,
        4.10,
        0.92,
        fill=WHITE,
        line=ORANGE,
        title_size=14.5,
        subtitle_size=12.5,
    )
    add_right_arrow(slide, x + 8.18, y + 2.31, 0.38, 0.32, BLUE)
    add_box(
        slide,
        "Mapped OSDR samples",
        "posterior pathway scores",
        x + 8.70,
        y + 2.02,
        w - 8.70,
        0.92,
        fill=WHITE,
        line=NAVY,
        title_size=14.5,
        subtitle_size=12.5,
    )
    add_text(
        slide,
        "~2,000 tissue HVGs | 319-387 Reactome programs | accession retained as a query covariate",
        x,
        y + 3.08,
        w,
        0.28,
        size=11.5,
        color=MUTED,
        italic=True,
        align=PP_ALIGN.CENTER,
        word_wrap=False,
    )


def add_program_annotation_workflow(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    heatmap_asset: Path,
) -> float:
    """Show real pathway-shift quantification and the reviewed LLM workflow."""
    split_x = x + 7.35
    add_connector(
        slide,
        split_x,
        y + 0.05,
        split_x,
        y + h - 0.10,
        color=RULE,
        width=1.0,
    )

    add_text(
        slide,
        "PROGRAM-SCORE SUMMARY",
        x,
        y,
        7.05,
        0.34,
        size=17,
        color=BLUE,
        bold=True,
        align=PP_ALIGN.CENTER,
        word_wrap=False,
    )
    add_box(
        slide,
        "Posterior program score",
        "one value per sample and pathway",
        x + 0.15,
        y + 0.42,
        3.05,
        0.82,
        fill=WHITE,
        line=NAVY,
        title_size=13.5,
        subtitle_size=11.5,
    )
    add_right_arrow(slide, x + 3.35, y + 0.67, 0.42, 0.32, BLUE)
    add_box(
        slide,
        "Within-project shift",
        "mean FLT - mean GC",
        x + 3.92,
        y + 0.42,
        2.98,
        0.82,
        fill=WHITE,
        line=BLUE,
        title_size=13.5,
        subtitle_size=12,
    )
    heatmap_ppi = add_picture_contain(
        slide,
        heatmap_asset,
        x + 0.08,
        y + 1.34,
        6.95,
        3.55,
        name="Actual spleen project-shift heatmap",
        description=(
            "Actual primary-seed mean flight-minus-ground latent-score shifts "
            "for three retained spleen pathways across five OSDR mission projects."
        ),
    )
    add_text_lines(
        slide,
        (
            "Heatmap cells are actual seed-2020 project effects; OSD-288 was excluded before spleen means.",
            "Tissue summaries weight projects equally. Retained programs were checked with held-out projects,",
            "three full trainings, enrichment, member genes, and composition sensitivity.",
        ),
        x + 0.10,
        y + 4.98,
        6.85,
        0.82,
        size=11.2,
        color=MUTED,
        italic=True,
        align=PP_ALIGN.LEFT,
    )

    right_x = split_x + 0.30
    right_w = w - 7.65
    add_text(
        slide,
        "LLM-ASSISTED LITERATURE ANNOTATION",
        right_x,
        y,
        right_w,
        0.34,
        size=17,
        color=GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
        word_wrap=False,
    )
    input_y = y + 0.42
    card_w = (right_w - 0.30) / 3
    input_cards = (
        ("Program + effect", "TCR: -0.216"),
        ("Mission context", "5 of 5 lower"),
        ("Primary sources", "references 35-37"),
    )
    for index, (title, subtitle) in enumerate(input_cards):
        add_box(
            slide,
            title,
            subtitle,
            right_x + index * (card_w + 0.15),
            input_y,
            card_w,
            0.78,
            fill=WHITE,
            line=GREEN if index == 2 else BLUE,
            title_size=12,
            subtitle_size=10.5,
        )
    add_down_arrow(slide, right_x + right_w / 2, y + 1.27, BLUE)

    model_x = right_x + 1.15
    model_y = y + 1.67
    model_w = right_w - 2.30
    for offset, fill in ((0.16, "DCE6F2"), (0.08, "E5EDF5"), (0.0, "F7FAFC")):
        add_rect(
            slide,
            model_x + offset,
            model_y - offset,
            model_w,
            1.28,
            fill,
            line=BLUE,
            line_width=1.0,
            rounded=True,
        )
    add_text(
        slide,
        "OpenAI Codex (GPT-5)",
        model_x + 0.18,
        model_y + 0.10,
        model_w - 0.36,
        0.30,
        size=15,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
        word_wrap=False,
    )
    token_colors = (THYMUS, BLUE, ORANGE, GREEN)
    token_x = model_x + 0.45
    token_y = model_y + 0.57
    for index, color in enumerate(token_colors):
        add_rect(
            slide,
            token_x + index * 0.37,
            token_y,
            0.25,
            0.25,
            color,
            rounded=True,
        )
    network_columns = (
        (model_x + 2.25, (model_y + 0.55, model_y + 0.85)),
        (model_x + 3.08, (model_y + 0.48, model_y + 0.70, model_y + 0.92)),
        (model_x + 3.92, (model_y + 0.58, model_y + 0.84)),
    )
    for column_index in range(len(network_columns) - 1):
        source_x, source_ys = network_columns[column_index]
        target_x, target_ys = network_columns[column_index + 1]
        for source_y in source_ys:
            for target_y in target_ys:
                add_connector(
                    slide,
                    source_x + 0.13,
                    source_y + 0.13,
                    target_x,
                    target_y + 0.13,
                    color="9FB3C3",
                    width=0.45,
                )
    for column_x, column_ys in network_columns:
        for node_y in column_ys:
            node = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(column_x),
                Inches(node_y),
                Inches(0.26),
                Inches(0.26),
            )
            set_fill(node, BLUE)
            node.line.fill.background()
    add_text(
        slide,
        "LLM-assisted synthesis",
        model_x + 0.20,
        model_y + 0.92,
        1.95,
        0.22,
        size=10.5,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        word_wrap=False,
    )

    add_down_arrow(slide, right_x + right_w / 2, y + 3.04, GREEN)
    add_box(
        slide,
        "Source verification",
        "direction, citations, and tissue context checked",
        right_x + 1.05,
        y + 3.40,
        right_w - 2.10,
        0.72,
        fill="EBF5EF",
        line=GREEN,
        title_size=13.5,
        subtitle_size=10.5,
    )
    add_down_arrow(slide, right_x + right_w / 2, y + 4.18, GREEN)

    role_y = y + 4.58
    role_gap = 0.15
    role_w = (right_w - 2 * role_gap) / 3
    roles = (
        ("Literature aligned", "T-cell receptor", GREEN),
        ("Complementary", "Neutrophil degranulation", BLUE),
        ("Context sensitive", "Skin DNA repair", ORANGE),
    )
    for index, (role, example, color) in enumerate(roles):
        role_x = right_x + index * (role_w + role_gap)
        add_rect(
            slide,
            role_x,
            role_y,
            role_w,
            0.96,
            WHITE,
            line=color,
            line_width=1.2,
            rounded=True,
        )
        add_text(
            slide,
            role,
            role_x + 0.08,
            role_y + 0.08,
            role_w - 0.16,
            0.28,
            size=11.5,
            color=color,
            bold=True,
            align=PP_ALIGN.CENTER,
            word_wrap=False,
        )
        add_text(
            slide,
            example,
            role_x + 0.08,
            role_y + 0.47,
            role_w - 0.16,
            0.30,
            size=10.5,
            color=BODY,
            align=PP_ALIGN.CENTER,
            word_wrap=False,
        )
    add_text(
        slide,
        "The LLM organized evidence; retained labels were not accepted without source-level review.",
        right_x,
        y + 5.72,
        right_w,
        0.34,
        size=11.5,
        color=MUTED,
        italic=True,
        align=PP_ALIGN.CENTER,
        word_wrap=False,
    )
    return heatmap_ppi


def load_poster_evidence() -> pd.DataFrame:
    """Load the manuscript-retained primary-tissue programs."""
    retained = load_retained_evidence().copy()
    return retained.loc[retained["tissue"].isin(MAIN_TISSUES)].copy()


def render_spleen_shift_heatmap(output: Path) -> None:
    """Render actual spleen project-level latent shifts for the workflow panel."""
    terms = POSTER_TERMS["spleen"]
    projects = ("MHU-3", "RR-1", "RR-23", "RR-6", "RR-9")
    effects = project_effects("spleen", terms)
    matrix = (
        effects.pivot(
            index="term",
            columns="heldout_project",
            values="project_effect",
        )
        .reindex(index=terms, columns=projects)
        .astype(float)
    )
    if matrix.isna().any().any():
        raise ValueError("Spleen project-shift heatmap contains missing values")

    labels = (
        "T-cell receptor",
        "Neutrophil degranulation",
        "C-type lectin receptor",
    )
    style = {
        "font.family": "DejaVu Sans",
        "font.size": 11,
        "axes.titlesize": 15,
        "axes.labelsize": 11,
        "xtick.labelsize": 10.5,
        "ytick.labelsize": 11.5,
        "text.color": "#202629",
        "axes.labelcolor": "#202629",
        "xtick.color": "#202629",
        "ytick.color": "#202629",
    }
    with plt.rc_context(style):
        fig, ax = plt.subplots(figsize=(7.4, 3.45), layout="constrained")
        image = ax.imshow(
            matrix.to_numpy(),
            cmap="RdBu_r",
            vmin=-0.75,
            vmax=0.75,
            aspect="auto",
        )
        ax.set_xticks(np.arange(len(projects)), projects)
        ax.set_yticks(np.arange(len(labels)), labels)
        ax.set_xlabel("OSDR mission project")
        ax.set_title(
            "Actual spleen project shifts\ncell = mean FLT - mean GC latent score",
            fontweight="bold",
            pad=8,
        )
        ax.tick_params(length=0)
        for row_index in range(matrix.shape[0]):
            for column_index in range(matrix.shape[1]):
                value = float(matrix.iloc[row_index, column_index])
                ax.text(
                    column_index,
                    row_index,
                    f"{value:+.2f}",
                    ha="center",
                    va="center",
                    fontsize=10.5,
                    fontweight="bold",
                    color="white" if abs(value) >= 0.38 else "#263238",
                )
        colorbar = fig.colorbar(
            image,
            ax=ax,
            orientation="horizontal",
            pad=0.18,
            fraction=0.10,
            aspect=32,
        )
        colorbar.set_label("Flight - ground latent score")
        colorbar.outline.set_linewidth(0.6)
        output.parent.mkdir(parents=True, exist_ok=True)
        fig.savefig(output, dpi=400, facecolor="white")
        plt.close(fig)


def render_poster_pathway_asset(output: Path) -> None:
    """Render the manuscript-retained primary-tissue pathways."""
    role_colors = {
        "aligned": "#009E73",
        "complementary": "#0072B2",
        "context_sensitive": "#D55E00",
    }
    role_markers = {
        "aligned": "o",
        "complementary": "s",
        "context_sensitive": "^",
    }
    evidence = load_poster_evidence()
    style = {
        "font.family": "DejaVu Sans",
        "font.size": 14,
        "axes.labelsize": 15,
        "xtick.labelsize": 13,
        "ytick.labelsize": 15,
        "text.color": "#202629",
        "axes.labelcolor": "#202629",
        "xtick.color": "#202629",
        "ytick.color": "#202629",
    }
    with plt.rc_context(style):
        fig, axes = plt.subplots(
            2,
            2,
            figsize=(14.4, 8.7),
            layout="constrained",
        )
        for panel, ax, tissue in zip("abcd", axes.flat, MAIN_TISSUES):
            terms = POSTER_TERMS[tissue]
            subset = (
                evidence.loc[evidence["tissue"].eq(tissue)]
                .set_index("term")
                .loc[list(terms)]
                .reset_index()
            )
            points = project_effects(tissue, terms)
            positions = np.arange(len(subset))[::-1]
            for position, row in zip(positions, subset.itertuples(index=False)):
                local = points.loc[
                    points["term"].eq(row.term), "project_effect"
                ].to_numpy(dtype=float)
                ax.scatter(
                    local,
                    np.full(len(local), position),
                    s=64,
                    facecolor="white",
                    edgecolor="#707A7F",
                    linewidth=1.2,
                    zorder=2,
                )
                role = str(row.evidence_role)
                ax.hlines(
                    position,
                    float(row.seed_effect_minimum),
                    float(row.seed_effect_maximum),
                    color=role_colors[role],
                    linewidth=2.4,
                    zorder=3,
                )
                ax.scatter(
                    float(row.seed_effect_median),
                    position,
                    marker=role_markers[role],
                    s=132,
                    facecolor=role_colors[role],
                    edgecolor="white",
                    linewidth=1.1,
                    zorder=4,
                )

            ax.axvline(0, color="#3F494E", linewidth=1.1)
            ax.set_yticks(positions)
            ax.set_yticklabels(
                [
                    textwrap.fill(
                        str(label),
                        width=29,
                        break_long_words=False,
                        break_on_hyphens=False,
                    )
                    for label in subset["display_label"]
                ]
            )
            for tick, row in zip(ax.get_yticklabels(), subset.itertuples(index=False)):
                tick.set_color(role_colors[str(row.evidence_role)])
            project_count = int(subset["expimap_n_projects"].max())
            ax.set_title(
                f"{panel}  {tissue.title()} ({project_count} projects)",
                loc="left",
                fontweight="bold",
                fontsize=18,
                pad=9,
            )
            if panel in "cd":
                ax.set_xlabel("Flight - ground pathway shift")
            ax.grid(axis="x", color="#DCE1E3", linewidth=1.0)
            ax.set_axisbelow(True)
            ax.tick_params(axis="y", length=0, pad=8)
            ax.tick_params(axis="x", width=1.0, length=4)
            for side in ("top", "right", "left"):
                ax.spines[side].set_visible(False)
            ax.spines["bottom"].set_color("#697277")
            ax.spines["bottom"].set_linewidth(1.0)

        handles = [
            Line2D(
                [0],
                [0],
                marker="o",
                linestyle="none",
                markerfacecolor="white",
                markeredgecolor="#707A7F",
                markeredgewidth=1.2,
                markersize=8,
                label="OSDR project",
            ),
            *[
                Line2D(
                    [0],
                    [0],
                    marker=role_markers[role],
                    color=role_colors[role],
                    linewidth=2.4,
                    markersize=8,
                    label=label,
                )
                for role, label in (
                    ("aligned", "Literature aligned"),
                    ("complementary", "Complementary"),
                    ("context_sensitive", "Context sensitive"),
                )
            ],
        ]
        fig.legend(
            handles=handles,
            loc="outside lower center",
            ncol=4,
            frameon=False,
            fontsize=12.5,
            handlelength=1.8,
            columnspacing=1.25,
        )
        fig.get_layout_engine().set(
            w_pad=0.12,
            h_pad=0.12,
            wspace=0.12,
            hspace=0.12,
        )
        output.parent.mkdir(parents=True, exist_ok=True)
        fig.savefig(
            output,
            dpi=400,
            facecolor="white",
            metadata={
                "Title": "Retained expiMap pathway shifts"
            },
        )
        plt.close(fig)


def render_pdf_asset(
    source: Path,
    output: Path,
    dpi: int = 700,
    *,
    crop_bottom_fraction: float = 0.0,
) -> None:
    output.parent.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        [
            "pdftocairo",
            "-singlefile",
            "-png",
            "-r",
            str(dpi),
            str(source),
            str(output.with_suffix("")),
        ],
        check=True,
    )
    if crop_bottom_fraction:
        with Image.open(output) as rendered:
            width, height = rendered.size
            cropped = rendered.crop(
                (0, 0, width, round(height * (1 - crop_bottom_fraction)))
            )
        cropped.save(output)


def add_picture_contain(
    slide,
    path: Path,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    name: str,
    description: str,
) -> float:
    with Image.open(path) as image:
        px_w, px_h = image.size
    image_ratio = px_w / px_h
    region_ratio = w / h
    if image_ratio >= region_ratio:
        placed_w = w
        placed_h = w / image_ratio
        placed_x = x
        placed_y = y + (h - placed_h) / 2
    else:
        placed_h = h
        placed_w = h * image_ratio
        placed_x = x + (w - placed_w) / 2
        placed_y = y
    picture = slide.shapes.add_picture(
        str(path),
        Inches(placed_x),
        Inches(placed_y),
        Inches(placed_w),
        Inches(placed_h),
    )
    properties = picture._element.xpath(".//p:cNvPr")[0]
    properties.set("name", name)
    properties.set("descr", description)
    ppi = min(px_w / placed_w, px_h / placed_h)
    if ppi < 300:
        raise ValueError(f"{path.name} has only {ppi:.0f} effective ppi")
    return ppi


def validate(prs: Presentation) -> None:
    if len(prs.slides) != 1:
        raise ValueError("Poster must contain exactly one slide")
    for shape in prs.slides[0].shapes:
        if shape.left < 0 or shape.top < 0:
            raise ValueError(f"{shape.name!r} starts outside the slide")
        if shape.left + shape.width > prs.slide_width + 1:
            raise ValueError(f"{shape.name!r} exceeds the slide width")
        if shape.top + shape.height > prs.slide_height + 1:
            raise ValueError(f"{shape.name!r} exceeds the slide height")
        if shape.has_text_frame:
            normalized = " ".join(shape.text.split())
            if len(normalized) >= 60 and shape.text_frame.word_wrap:
                raise ValueError(
                    f"{shape.name!r} contains long reflowable text; use fixed lines"
                )
            for paragraph in shape.text_frame.paragraphs:
                if not paragraph.text:
                    continue
                properties = paragraph._p.get_or_add_pPr()
                if properties.get("latinLnBrk") != "0":
                    raise ValueError(
                        f"{shape.name!r} allows a Latin word to break between letters"
                    )
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        if not paragraph.text:
                            continue
                        properties = paragraph._p.get_or_add_pPr()
                        if properties.get("latinLnBrk") != "0":
                            raise ValueError(
                                "A table cell allows a Latin word to break between letters"
                            )


def render_poster(pptx_path: Path) -> tuple[Path | None, Path | None]:
    libreoffice = shutil.which("libreoffice")
    pdftocairo = shutil.which("pdftocairo")
    if not libreoffice:
        return None, None
    pdf_path = pptx_path.with_suffix(".pdf")
    pdf_path.unlink(missing_ok=True)
    with tempfile.TemporaryDirectory(prefix="asgsr_poster_lo_") as profile:
        subprocess.run(
            [
                libreoffice,
                "--headless",
                f"-env:UserInstallation={Path(profile).resolve().as_uri()}",
                "--convert-to",
                "pdf",
                "--outdir",
                str(pptx_path.parent),
                str(pptx_path),
            ],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            text=True,
        )
    if not pdf_path.exists():
        raise FileNotFoundError(f"LibreOffice did not create {pdf_path}")
    if not pdftocairo:
        return pdf_path, None
    preview = pptx_path.with_name(pptx_path.stem + "_preview.png")
    subprocess.run(
        [
            pdftocairo,
            "-singlefile",
            "-png",
            "-r",
            "100",
            str(pdf_path),
            str(preview.with_suffix("")),
        ],
        check=True,
    )
    return pdf_path, preview


def render_architecture_crop(pdf_path: Path) -> Path:
    output = ASSET_DIR / "expimap_architecture_visualization_300dpi.png"
    dpi = 300
    x, y, w, h = 0.70, 19.85, 14.45, 4.20
    subprocess.run(
        [
            "pdftocairo",
            "-singlefile",
            "-png",
            "-r",
            str(dpi),
            "-x",
            str(round(x * dpi)),
            "-y",
            str(round(y * dpi)),
            "-W",
            str(round(w * dpi)),
            "-H",
            str(round(h * dpi)),
            str(pdf_path),
            str(output.with_suffix("")),
        ],
        check=True,
    )
    return output


def build() -> tuple[Path, Path | None, Path | None, Path | None, list[float]]:
    POSTER_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    result_asset = ASSET_DIR / "figure_3_tissue_pathway_shifts_poster_400dpi.png"
    hypothesis_asset = ASSET_DIR / "figure_6_tissue_state_hypotheses_700dpi.png"
    shift_heatmap_asset = ASSET_DIR / "spleen_project_latent_shift_heatmap_400dpi.png"
    render_poster_pathway_asset(result_asset)
    render_spleen_shift_heatmap(shift_heatmap_asset)
    render_pdf_asset(
        FIGURE_DIR / "figure_6_tissue_state_hypotheses.pdf",
        hypothesis_asset,
        crop_bottom_fraction=0.07,
    )

    template = Presentation(TEMPLATE_PATH)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.title = TITLE
    prs.core_properties.author = "Jason Trinh"
    prs.core_properties.subject = (
        "Template-based cross-mission expiMap poster using NASA OSDR mouse RNA-seq"
    )
    prs.core_properties.keywords = (
        "spaceflight, expiMap, OSDR, ARCHS4, Reactome, transcriptomics, NASA"
    )
    prs.core_properties.comments = (
        "Layout and NASA branding reference 00 Poster Session Student "
        "Template_Approved by Legal.pptx."
    )
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(WHITE)

    # Approved-template header
    add_rect(slide, 0, 0, SLIDE_W, 5.05, HEADER)
    add_text(
        slide,
        "National Aeronautics and Space Administration",
        0.80,
        0.28,
        18.0,
        0.48,
        size=16,
        color=MUTED,
    )
    for line, line_y in ((TITLE_LINE_1, 1.00), (TITLE_LINE_2, 1.92)):
        add_text(
            slide,
            line,
            0.60,
            line_y,
            42.80,
            0.82,
            size=50,
            color="111111",
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
            word_wrap=False,
        )
    add_text(
        slide,
        (
            "Jason Trinh / University of California, Berkeley | "
            "NASA Space Life Sciences Training Program, NASA Ames Research Center"
        ),
        1.50,
        3.70,
        41.0,
        0.72,
        size=26,
        color=MUTED,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        word_wrap=False,
    )
    add_nasa_logo(slide, template)

    # Template-style panels
    left_x, left_w = 0.35, 15.15
    center_x, center_w = 15.85, 15.75
    right_x, right_w = 31.95, 15.70

    add_panel(slide, left_x, 5.35, left_w, 3.78)
    add_section(slide, "ABSTRACT", left_x + 0.28, 5.55, left_w - 0.56)
    abstract_lines = (
        "Spaceflight affects multiple organs, but mission differences can obscure recurring biological responses. We mapped",
        "NASA OSDR mouse bulk RNA-sequencing samples into tissue-matched expiMap models trained on non-spaceflight",
        "ARCHS4 data, using about 2,000 highly variable genes linked to mouse Reactome pathways. Project-balanced shifts",
        "were evaluated with enrichment, held-out projects, three complete trainings, composition proxies, and member genes.",
        "Four tissues produced reproducible patterns: lower thymic repair and cytoskeletal programs; lower skin regulatory,",
        "repair, barrier-lipid, and cell-junction programs; lower hepatic antigen-presentation and T-cell signaling; and lower",
        "splenic T-cell, neutrophil-degranulation, and C-type lectin-receptor programs; all three splenic programs had GSEA",
        "FDR < 0.05. These results recover established responses and identify complementary tissue hypotheses to test.",
    )
    add_text_lines(
        slide,
        abstract_lines,
        left_x + 0.50,
        6.48,
        left_w - 1.00,
        2.08,
        size=17.5,
        color=BODY,
    )

    add_panel(slide, left_x, 9.43, left_w, 5.15)
    add_section(slide, "INTRODUCTION", left_x + 0.28, 9.63, left_w - 0.56)
    add_text(
        slide,
        "PROJECT OBJECTIVE",
        left_x + 0.55,
        10.52,
        left_w - 1.10,
        0.42,
        size=19,
        color=GOLD,
        bold=True,
    )
    add_text(
        slide,
        "Identify reproducible tissue-specific gene programs across spaceflight missions.",
        left_x + 0.55,
        10.95,
        left_w - 1.10,
        0.72,
        size=23,
        color=INK,
        bold=True,
        word_wrap=False,
    )
    add_challenge_equation(slide, left_x + 0.55, 11.85, left_w - 1.10)
    add_text_lines(
        slide,
        (
            "Study identity and protocol can dominate an unconstrained expression representation. Tissue-matched reference",
            "mapping, accession conditioning, and equal project weighting reduce this bias without removing all confounding.",
        ),
        left_x + 0.55,
        13.15,
        left_w - 1.10,
        0.76,
        size=16.5,
        color=BODY,
    )

    add_panel(slide, left_x, 14.88, left_w, 9.80)
    add_section(
        slide,
        "DATA AND REFERENCE MAPPING",
        left_x + 0.28,
        15.08,
        left_w - 0.56,
        size=23,
    )
    add_text(
        slide,
        "ANALYSIS SCOPE",
        left_x + 0.55,
        16.03,
        left_w - 1.10,
        0.40,
        size=17,
        color=GOLD,
        bold=True,
    )
    add_scope_table(
        slide, left_x + 0.55, 16.43, left_w - 1.10, 2.75
    )
    add_text_lines(
        slide,
        ("Spleen excludes nine condition-strain-confounded OSD-288 samples from the primary analysis.",),
        left_x + 0.55,
        19.32,
        left_w - 1.10,
        0.34,
        size=12.5,
        color=MUTED,
        italic=True,
    )
    add_text(
        slide,
        "METHODS: REFERENCE TRAINING AND QUERY MAPPING",
        left_x + 0.55,
        19.88,
        left_w - 1.10,
        0.34,
        size=17,
        color=GOLD,
        bold=True,
        word_wrap=False,
    )
    add_reference_query_mapping(
        slide,
        left_x + 0.55,
        20.32,
        left_w - 1.10,
    )

    add_panel(slide, center_x, 5.35, center_w, 8.05)
    add_section(
        slide,
        "PROGRAM SCORES AND LITERATURE ANNOTATION",
        center_x + 0.28,
        5.55,
        center_w - 0.56,
        size=22,
    )
    heatmap_ppi = add_program_annotation_workflow(
        slide,
        center_x + 0.40,
        6.42,
        center_w - 0.80,
        6.65,
        shift_heatmap_asset,
    )

    add_panel(slide, center_x, 13.70, center_w, 10.98)
    add_section(
        slide,
        "RESULTS: RETAINED PROGRAM SHIFTS",
        center_x + 0.28,
        13.90,
        center_w - 0.56,
        size=24,
    )
    result_ppi = add_picture_contain(
        slide,
        result_asset,
        center_x + 0.35,
        14.78,
        center_w - 0.70,
        9.52,
        name="Retained pathway shifts",
        description=(
            "Manuscript-retained expiMap pathway shifts for thymus, skin, "
            "liver, and spleen, with project effects and three-training ranges."
        ),
    )
    add_panel(slide, right_x, 5.35, right_w, 13.28)
    add_section(
        slide,
        "RESULTS: BIOLOGICAL INTERPRETATION",
        right_x + 0.28,
        5.55,
        right_w - 0.56,
        size=24,
    )
    add_text(
        slide,
        "Prior evidence, observed shifts, and hypotheses to test",
        right_x + 0.50,
        6.42,
        right_w - 1.00,
        0.42,
        size=19,
        color=BLUE,
        bold=True,
        align=PP_ALIGN.CENTER,
        word_wrap=False,
    )
    hypothesis_ppi = add_picture_contain(
        slide,
        hypothesis_asset,
        right_x + 0.35,
        6.87,
        right_w - 0.70,
        11.48,
        name="Tissue-state hypotheses",
        description=(
            "Prior-literature phenotypes, observed lower expiMap program scores, "
            "and complementary hypotheses for thymus, skin, liver, and spleen."
        ),
    )

    add_panel(slide, right_x, 18.93, right_w, 3.55)
    add_section(slide, "CONCLUSIONS", right_x + 0.28, 19.13, right_w - 0.56)
    add_text_lines(
        slide,
        (
            "\u2022 Spleen showed the strongest evidence: three immune programs were lower across five projects and",
            "  three trainings, with GSEA FDR < 0.05 for each.",
            "\u2022 Thymus retained lower repair, cytoskeletal, and niche-interaction programs; skin retained lower",
            "  regulatory, repair, regenerative, barrier-lipid, and cell-junction programs.",
            "\u2022 Liver retained lower antigen-presentation and T-cell signaling; kidney showed a reproducible",
            "  structural and growth-factor pattern but remains exploratory.",
            "\u2022 These bulk latent shifts generate tissue-level hypotheses; they do not prove pathway inhibition",
            "  or distinguish altered cell abundance from altered cell state.",
        ),
        right_x + 0.55,
        20.00,
        right_w - 1.10,
        2.12,
        size=16.5,
        color=BODY,
    )

    add_panel(slide, right_x, 22.75, right_w, 1.93)
    add_section(
        slide,
        "REFERENCES",
        right_x + 0.28,
        22.95,
        right_w - 0.56,
        size=22,
    )
    add_text_lines(
        slide,
        (
            "Selected sources; numbering follows the manuscript: [1] Gebre et al., NAR 2025; [2] Lotfollahi et al., Nat Cell Biol 2023.",
            "[4] Lachmann et al., Nat Commun 2018; [5] Milacic et al., NAR 2024; [6] Horie et al., Sci Rep 2019.",
            "[12] Cope et al., Commun Med 2024; [13] Park et al., Nat Commun 2024; [17] Beheshti et al., Sci Rep 2019.",
            "[35] Gridley et al., J Appl Physiol 2009. Complete citations and source tables accompany the manuscript.",
        ),
        right_x + 0.52,
        23.78,
        right_w - 1.04,
        0.78,
        size=12.5,
        color=BODY,
    )

    # Approved-template footer and acknowledgements
    add_rect(slide, center_x, 24.92, center_w + right_w + 0.05, 2.08, HEADER)
    add_text(
        slide,
        "Acknowledgements",
        center_x + 0.35,
        25.02,
        5.0,
        0.45,
        size=22,
        color=MUTED,
        bold=True,
    )
    add_text_lines(
        slide,
        (
            "Mentor: James Casaletto | NASA Space Life Sciences Training Program",
            "NASA OSDR, ARCHS4, and Reactome investigators and curation teams",
        ),
        center_x + 0.35,
        25.52,
        20.60,
        0.66,
        size=16,
        color=BODY,
    )
    add_text(
        slide,
        "jasontrinh@berkeley.edu\ngithub.com/jasont314/nasa-mouse",
        39.25,
        25.32,
        7.92,
        0.90,
        size=16,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.RIGHT,
        word_wrap=False,
    )
    add_text(
        slide,
        "www.nasa.gov",
        0.70,
        26.10,
        3.60,
        0.40,
        size=15,
        color=MUTED,
    )

    validate(prs)
    output = POSTER_DIR / "asgsr_expimap_poster.pptx"
    prs.save(output)
    pdf_path, preview_path = render_poster(output)
    architecture_path = render_architecture_crop(pdf_path) if pdf_path else None
    return (
        output,
        pdf_path,
        preview_path,
        architecture_path,
        [result_ppi, heatmap_ppi, hypothesis_ppi],
    )


def run() -> None:
    output, pdf_path, preview_path, architecture_path, ppi = build()
    print(output)
    if pdf_path:
        print(pdf_path)
    if preview_path:
        print(preview_path)
    if architecture_path:
        print(architecture_path)
    print(
        "Embedded data-figure effective resolution: "
        + ", ".join(f"{value:.0f} ppi" for value in ppi)
    )


if __name__ == "__main__":
    run()
