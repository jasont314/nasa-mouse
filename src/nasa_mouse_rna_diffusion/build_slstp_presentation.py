"""Build the 2026 SLSTP presentation for the generative transcriptomics study."""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
TEMPLATE = ROOT / "presentation/SLSTP_template_2026.pptx"
OUTPUT = ROOT / "presentation/SLSTP_2026_Generative_Transcriptomics.pptx"
PACKAGE_DIR = ROOT / "presentation/generative_slstp_2026"
ASSET_DIR = PACKAGE_DIR / "assets"
NOTES_PATH = PACKAGE_DIR / "speaker_notes.md"
PAPER_DIR = ROOT / "paper/synthetic_guided_spaceflight"
EXPIMAP_PAPER_DIR = ROOT / "paper/asgsr_expimap_hvg"

SLIDE_W = 13.333333
SLIDE_H = 7.5
FONT = "Arial"

NAVY = "082B55"
BLUE = "2D6496"
TEAL = "178681"
GREEN = "478A64"
ORANGE = "D37B00"
GOLD = "9A6B20"
CORAL = "D96552"
PURPLE = "8063A6"
DARK = "263746"
GRAY = "59697A"
MID_GRAY = "8A969E"
LIGHT = "F2F5F7"
PALE_BLUE = "E9F1F7"
PALE_TEAL = "E8F3F1"
PALE_GOLD = "F8F1E2"
PALE_CORAL = "F8ECE9"
WHITE = "FFFFFF"


@dataclass(frozen=True)
class SlideNote:
    number: int
    title: str
    time: str
    text: str


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def _diverging_color(value: float, limit: float) -> str:
    position = np.clip((float(value) / float(limit) + 1.0) / 2.0, 0.0, 1.0)
    colormap = matplotlib.colors.LinearSegmentedColormap.from_list(
        "program_score", [f"#{BLUE}", f"#{WHITE}", f"#{ORANGE}"]
    )
    return matplotlib.colors.to_hex(colormap(position)).lstrip("#").upper()


def _set_fill(shape, color: str, transparency: int | None = None) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    if transparency is not None:
        color_nodes = shape._element.spPr.xpath("./a:solidFill/a:srgbClr")
        if color_nodes:
            alpha = OxmlElement("a:alpha")
            alpha.set("val", str(int((100 - transparency) * 1000)))
            color_nodes[0].append(alpha)


def _set_line(shape, color: str, width: float = 1.0) -> None:
    shape.line.color.rgb = _rgb(color)
    shape.line.width = Pt(width)


def _add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: str = DARK,
    bold: bool = False,
    italic: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.02,
    name: str | None = None,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color)
    return shape


def _add_rich_text(
    slide,
    segments: Iterable[tuple[str, dict]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: str = DARK,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.02,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.space_after = Pt(0)
    for text, style in segments:
        run = paragraph.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(style.get("size", size))
        run.font.bold = style.get("bold", False)
        run.font.italic = style.get("italic", False)
        run.font.color.rgb = _rgb(style.get("color", color))
    return shape


def _add_panel(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = LIGHT,
    line: str = "DDE4E8",
    radius: bool = True,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    _set_fill(shape, fill)
    _set_line(shape, line, 0.8)
    return shape


def _add_rule(slide, x: float, y: float, w: float, color: str, height: float = 0.04):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(height)
    )
    _set_fill(shape, color)
    shape.line.fill.background()
    return shape


def _add_circle(slide, x: float, y: float, diameter: float, color: str):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x),
        Inches(y),
        Inches(diameter),
        Inches(diameter),
    )
    _set_fill(shape, color)
    shape.line.fill.background()
    return shape


def _add_data_badge(slide, label: str, x: float, y: float, color: str, diameter: float = 0.24):
    _add_circle(slide, x, y, diameter, color)
    _add_text(
        slide,
        label,
        x,
        y + 0.005,
        diameter,
        diameter - 0.01,
        size=8.2,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )


def _add_arrow(slide, x: float, y: float, w: float, h: float, color: str = MID_GRAY):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    _set_fill(shape, color)
    shape.line.fill.background()
    return shape


def _add_bullet_rows(
    slide,
    rows: list[str],
    x: float,
    y: float,
    w: float,
    *,
    size: float = 16,
    color: str = DARK,
    bullet_color: str = TEAL,
    row_h: float = 0.48,
):
    for index, row in enumerate(rows):
        top = y + index * row_h
        _add_circle(slide, x, top + 0.10, 0.10, bullet_color)
        _add_text(
            slide,
            row,
            x + 0.20,
            top,
            w - 0.20,
            row_h,
            size=size,
            color=color,
            valign=MSO_ANCHOR.TOP,
            margin=0,
        )


def _add_picture_contain(
    slide,
    path: Path,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    alt: str,
):
    with Image.open(path) as image:
        aspect = image.width / image.height
    box_aspect = w / h
    if aspect >= box_aspect:
        width = w
        height = w / aspect
        left = x
        top = y + (h - height) / 2
    else:
        height = h
        width = h * aspect
        left = x + (w - width) / 2
        top = y
    shape = slide.shapes.add_picture(
        str(path), Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.name = alt
    shape._element.nvPicPr.cNvPr.set("descr", alt)
    return shape


def _add_slide_title(slide, eyebrow: str, title: str, subtitle: str | None = None):
    _add_text(
        slide,
        eyebrow.upper(),
        0.32,
        0.58,
        5.5,
        0.25,
        size=10.5,
        color=GOLD,
        bold=True,
    )
    _add_text(
        slide,
        title,
        0.32,
        0.86,
        12.5,
        0.64,
        size=27,
        color=NAVY,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    if subtitle:
        _add_text(
            slide,
            subtitle,
            0.34,
            1.48,
            12.3,
            0.34,
            size=14.5,
            color=GRAY,
            valign=MSO_ANCHOR.MIDDLE,
        )


def _prepare_content_slide(slide, number: int):
    for placeholder in slide.placeholders:
        kind = placeholder.placeholder_format.type
        if kind in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.OBJECT):
            placeholder.text = ""
        elif kind == PP_PLACEHOLDER.FOOTER:
            placeholder.text = ""
        elif kind == PP_PLACEHOLDER.SLIDE_NUMBER:
            placeholder.text = str(number)
    return slide


def _add_source(slide, text: str):
    _add_text(
        slide,
        text,
        0.36,
        7.12,
        12.0,
        0.19,
        size=7.2,
        color=MID_GRAY,
        valign=MSO_ANCHOR.MIDDLE,
    )


def _build_tissue_utility_chart() -> Path:
    data = pd.read_csv(
        PAPER_DIR / "source_data/table_s18_matched_all_gene_utility.tsv",
        sep="\t",
    )
    order = [
        "eye",
        "retina",
        "skin",
        "thymus",
        "spleen",
        "liver",
        "skeletal_muscle",
        "soleus",
        "gastrocnemius",
        "hippocampus",
    ]
    data = data.loc[data["arm"].eq("real_plus_generated")].copy()
    data = data.set_index("tissue").loc[order].reset_index()
    display = {
        "eye": "Eye",
        "retina": "Retina",
        "skin": "Skin",
        "thymus": "Thymus",
        "skeletal_muscle": "Skeletal muscle",
        "soleus": "Soleus",
        "gastrocnemius": "Gastrocnemius",
        "hippocampus": "Hippocampus",
        "spleen": "Spleen",
        "liver": "Liver",
    }
    real = data["real_mean_balanced_accuracy"].to_numpy(float)
    candidate = data["arm_mean_balanced_accuracy"].to_numpy(float)
    delta = data["mean_delta_balanced_accuracy"].to_numpy(float)
    passed = data["joint_mean_all_metrics_nonworse"].astype(bool).to_numpy()
    y = np.arange(len(data))
    fig, ax = plt.subplots(figsize=(8.1, 4.8))
    fig.patch.set_alpha(0)
    ax.set_facecolor("none")
    for yi, start, end in zip(y, real, candidate):
        ax.plot([start, end], [yi, yi], color="#AEB9BF", lw=2.0, zorder=1)
    candidate_colors = np.where(passed, "#178681", "#D96552")
    ax.scatter(real, y, s=58, color="#7E8A92", label="Real only", zorder=3)
    ax.scatter(candidate, y, s=70, color=candidate_colors, zorder=4)
    for yi, end, change, color in zip(y, candidate, delta, candidate_colors):
        label_x = end + 0.009 if change >= 0 else end - 0.012
        ax.text(
            label_x,
            yi,
            f"{change:+.3f}",
            ha="left" if change >= 0 else "right",
            va="center",
            fontsize=8.7,
            color=color,
            weight="bold",
        )
    ax.set_yticks(y, [display[value] for value in data["tissue"]])
    ax.invert_yaxis()
    ax.set_xlim(0.45, 1.04)
    ax.set_xlabel("Balanced accuracy on held-out real profiles", fontsize=10)
    ax.grid(axis="x", color="#DDE4E8", lw=0.8)
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.tick_params(axis="y", length=0, labelsize=10)
    ax.tick_params(axis="x", labelsize=9)
    ax.scatter([], [], s=70, color="#178681", label="Passed all six metrics")
    ax.scatter([], [], s=70, color="#D96552", label="Did not pass")
    ax.legend(
        frameon=False,
        ncol=3,
        fontsize=8.3,
        loc="lower center",
        bbox_to_anchor=(0.5, 1.01),
    )
    fig.tight_layout()
    path = ASSET_DIR / "tissue_balanced_accuracy_dumbbell.png"
    fig.savefig(path, dpi=220, transparent=True, bbox_inches="tight")
    plt.close(fig)
    return path


def _build_trajectory_crop() -> Path:
    source = PAPER_DIR / "figures/figure_2a_archs4_denoising_trajectory.png"
    with Image.open(source) as image:
        top = int(image.height * 0.07)
        bottom = int(image.height * 0.84)
        cropped = image.crop((0, top, image.width, bottom))
        path = ASSET_DIR / "ddim_trajectory_panels.png"
        cropped.save(path)
    return path


def _build_pca_comparison_chart(primary: str) -> Path:
    if primary not in {"tissue", "condition"}:
        raise ValueError(f"Unsupported PCA comparison: {primary!r}")
    source = pd.read_csv(
        PACKAGE_DIR / "source_data/locked_pca_by_accession.tsv",
        sep="\t",
    )
    required = {"source", primary, "accession", "pc1", "pc2"}
    missing = required.difference(source.columns)
    if missing:
        raise ValueError(
            "PCA source data are missing columns: "
            + ", ".join(sorted(missing))
        )
    accessions = sorted(
        source["accession"].astype(str).unique(),
        key=lambda value: int(value.rsplit("-", 1)[-1]),
    )
    hues = (np.arange(len(accessions)) * 0.61803398875) % 1.0
    palette = {
        accession: matplotlib.colors.hsv_to_rgb((hue, 0.66, 0.82))
        for accession, hue in zip(accessions, hues)
    }
    if primary == "condition":
        primary_levels = ["flight", "ground_control"]
        primary_palette = {
            "flight": "#D96552",
            "ground_control": "#2D6496",
        }
        primary_labels = {
            "flight": "Flight",
            "ground_control": "Ground control",
        }
    else:
        primary_levels = sorted(source[primary].astype(str).unique())
        tissue_cmap = plt.get_cmap("tab20", max(len(primary_levels), 1))
        primary_palette = {
            level: tissue_cmap(index)
            for index, level in enumerate(primary_levels)
        }
        primary_labels = {
            level: level.replace("_", " ").title()
            for level in primary_levels
        }
    x_pad = 0.04 * (source["pc1"].max() - source["pc1"].min())
    y_pad = 0.06 * (source["pc2"].max() - source["pc2"].min())
    x_limits = (source["pc1"].min() - x_pad, source["pc1"].max() + x_pad)
    y_limits = (source["pc2"].min() - y_pad, source["pc2"].max() + y_pad)

    figure, axes = plt.subplots(1, 2, figsize=(12.4, 4.45), sharex=True, sharey=True)
    figure.patch.set_alpha(0)

    def add_grouped_points(
        axis,
        *,
        column: str,
        levels: list[str],
        colors: dict[str, object],
        labels: dict[str, str],
        legend: bool,
    ) -> None:
        for level in levels:
            group = source.loc[source[column].astype(str).eq(level)]
            real = group.loc[group["source"].eq("real")]
            synthetic = group.loc[group["source"].eq("synthetic")]
            axis.scatter(
                real["pc1"],
                real["pc2"],
                s=24,
                marker="o",
                color=colors[level],
                alpha=0.42,
                linewidths=0,
                edgecolors="none",
            )
            axis.scatter(
                synthetic["pc1"],
                synthetic["pc2"],
                s=28,
                marker="x",
                color=colors[level],
                alpha=0.82,
                linewidths=1.0,
                label=labels[level],
            )
        if legend:
            axis.legend(
                frameon=False,
                fontsize=6.0 if column == "tissue" else 8.5,
                ncol=2 if column == "tissue" else 1,
                loc="upper right",
                borderaxespad=0.25,
                handletextpad=0.3,
                columnspacing=0.7,
            )

    add_grouped_points(
        axes[0],
        column=primary,
        levels=primary_levels,
        colors=primary_palette,
        labels=primary_labels,
        legend=True,
    )
    add_grouped_points(
        axes[1],
        column="accession",
        levels=accessions,
        colors=palette,
        labels={accession: accession for accession in accessions},
        legend=False,
    )
    titles = {
        "tissue": "Colored by tissue",
        "condition": "Colored by condition",
    }
    for axis, title in zip(axes, [titles[primary], "Colored by OSDR accession"]):
        axis.set_facecolor("none")
        axis.set_title(title, fontsize=14, color="#082B55", weight="bold", pad=9)
        axis.set_xlabel("PC1", fontsize=11, color="#263746")
        axis.set_xlim(x_limits)
        axis.set_ylim(y_limits)
        axis.grid(color="#DDE4E8", linewidth=0.75, alpha=0.78)
        axis.tick_params(labelsize=8.5, colors="#59697A")
        axis.spines[["top", "right"]].set_visible(False)
        axis.spines[["left", "bottom"]].set_color("#AEB9BF")
    axes[0].set_ylabel("PC2", fontsize=11, color="#263746")
    figure.subplots_adjust(left=0.06, right=0.99, top=0.90, bottom=0.13, wspace=0.16)
    path = ASSET_DIR / f"locked_pca_{primary}_vs_accession.png"
    figure.savefig(path, dpi=240, transparent=True)
    plt.close(figure)
    return path


def _set_title_slide(slide):
    main = next(
        placeholder
        for placeholder in slide.placeholders
        if placeholder.placeholder_format.type == PP_PLACEHOLDER.BODY
    )
    main.text_frame.clear()
    main.text_frame.word_wrap = True
    main.text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
    paragraph = main.text_frame.paragraphs[0]
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0
    run = paragraph.add_run()
    run.text = "Interpretable and generative models\nfor mouse spaceflight"
    run.font.name = FONT
    run.font.size = Pt(29)
    run.font.bold = True
    run.font.color.rgb = _rgb(WHITE)

    subtitle = next(
        placeholder
        for placeholder in slide.placeholders
        if placeholder.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE
    )
    subtitle.text_frame.clear()
    subtitle.text_frame.word_wrap = True
    entries = [
        ("Jason Trinh", 21, True, WHITE),
        ("EECS | UC Berkeley", 14, False, WHITE),
        ("", 6, False, WHITE),
        ("Mentor: James Casaletto | August 2026", 13, False, WHITE),
    ]
    for index, (text, size, bold, color) in enumerate(entries):
        paragraph = subtitle.text_frame.paragraphs[0] if index == 0 else subtitle.text_frame.add_paragraph()
        paragraph.space_after = Pt(2)
        run = paragraph.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color)


def _slide_project_scope(slide):
    _add_slide_title(
        slide,
        "Project",
        "One dataset, two machine-learning questions",
        "We used mouse bulk RNA-seq to study how spaceflight changes genes, pathways, and tissue state.",
    )

    _add_panel(slide, 0.55, 2.05, 12.24, 0.76, fill=NAVY, line=NAVY, radius=False)
    _add_text(slide, "NASA OSDR API", 0.86, 2.20, 1.78, 0.31, size=16, color=WHITE, bold=True, margin=0)
    _add_text(
        slide,
        "Mus musculus bulk RNA-seq | spaceflight (FLT) and ground control (GC) | multiple missions and tissues",
        2.72,
        2.18,
        9.58,
        0.34,
        size=13.2,
        color="DCE7F2",
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )

    panels = [
        (
            0.55,
            "01",
            "Which biological programs change?",
            "expiMap",
            "Train a pathway-constrained reference on ARCHS4, then map OSDR samples into mouse Reactome programs.",
            "Output: tissue-specific pathway shifts",
            BLUE,
            PALE_BLUE,
        ),
        (
            6.83,
            "02",
            "Can synthetic profiles sharpen the analysis?",
            "Conditional generation",
            "Generate matched expression profiles by tissue, study context, and FLT or GC condition.",
            "Output: validated profiles and gene ranking",
            TEAL,
            PALE_TEAL,
        ),
    ]
    for x, number, question, method, body, output, color, fill in panels:
        _add_panel(slide, x, 3.22, 5.96, 3.20, fill=fill, line=fill, radius=False)
        _add_text(slide, number, x + 0.24, 3.48, 0.45, 0.25, size=11, color=color, bold=True, margin=0)
        _add_text(slide, question, x + 0.82, 3.40, 4.76, 0.50, size=17.2, color=NAVY, bold=True, margin=0)
        _add_rule(slide, x + 0.25, 4.12, 5.43, color, 0.025)
        _add_text(slide, method, x + 0.25, 4.37, 5.30, 0.34, size=15.4, color=color, bold=True, margin=0)
        _add_text(slide, body, x + 0.25, 4.82, 5.30, 0.66, size=13.0, color=DARK, margin=0)
        _add_text(slide, output, x + 0.25, 5.78, 5.30, 0.31, size=12.3, color=NAVY, bold=True, margin=0)
    _add_source(slide, "Data source: NASA Open Science Data Repository Biological Data API.")


def _slide_autoencoder_foundation(slide):
    _add_slide_title(
        slide,
        "Neural networks",
        "Autoencoders compress expression into a latent space",
        "The encoder summarizes a gene profile; the decoder learns enough structure to reconstruct it.",
    )

    _add_text(slide, "AUTOENCODER", 0.58, 2.05, 1.44, 0.24, size=10.2, color=BLUE, bold=True, margin=0)
    _add_panel(slide, 0.50, 2.36, 7.38, 3.88, fill="F7F9FA", line="D9E1E5", radius=False)
    encoder_background = slide.shapes.add_shape(
        MSO_SHAPE.TRAPEZOID,
        Inches(3.18),
        Inches(2.91),
        Inches(1.72),
        Inches(2.46),
    )
    encoder_background.rotation = 270
    _set_fill(encoder_background, PALE_BLUE, transparency=18)
    _set_line(encoder_background, "D5E2EC", 0.8)
    decoder_background = slide.shapes.add_shape(
        MSO_SHAPE.TRAPEZOID,
        Inches(5.10),
        Inches(2.91),
        Inches(1.72),
        Inches(2.46),
    )
    decoder_background.rotation = 90
    _set_fill(decoder_background, PALE_TEAL, transparency=18)
    _set_line(decoder_background, "CEE2DE", 0.8)
    _add_text(slide, "gene\nprofile", 0.75, 3.82, 0.72, 0.52, size=11.2, color=GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    gene_values = [0.58, 1.07, 0.73, 1.32, 0.90, 0.46]
    for index, height in enumerate(gene_values):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.64 + index * 0.19),
            Inches(4.72 - height),
            Inches(0.11),
            Inches(height),
        )
        _set_fill(bar, BLUE if index % 2 == 0 else "8FB2CC")
        bar.line.fill.background()
    _add_arrow(slide, 2.86, 3.86, 0.34, 0.25, MID_GRAY)

    encoder_positions = [
        [(3.43, 3.08 + row * (2.02 / 4)) for row in range(5)],
        [(3.93, 3.08 + row * (2.02 / 2)) for row in range(3)],
    ]
    latent_positions = [(4.98, y + 0.10) for y in [3.44, 3.88, 4.32]]
    decoder_positions = [
        [(5.96, 3.08 + row * (2.02 / 2)) for row in range(3)],
        [(6.46, 3.08 + row * (2.02 / 4)) for row in range(5)],
    ]

    def connect_layers(source, target):
        for x1, y1 in source:
            for x2, y2 in target:
                line = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    Inches(x1),
                    Inches(y1),
                    Inches(x2),
                    Inches(y2),
                )
                _set_line(line, "C5D1D8", 0.65)

    connect_layers(encoder_positions[0], encoder_positions[1])
    connect_layers(encoder_positions[1], latent_positions)
    connect_layers(latent_positions, decoder_positions[0])
    connect_layers(decoder_positions[0], decoder_positions[1])

    for col_index, positions in enumerate(encoder_positions):
        for center_x, center_y in positions:
            _add_circle(slide, center_x - 0.08, center_y - 0.08, 0.16, "6D97B5" if col_index == 0 else TEAL)
    _add_text(slide, "encoder", 3.18, 5.35, 1.05, 0.24, size=10.2, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)

    for index, (center_x, center_y) in enumerate(latent_positions):
        _add_circle(slide, center_x - 0.10, center_y - 0.10, 0.20, ORANGE if index == 1 else GOLD)
    _add_text(slide, "latent\nvariables", 4.56, 4.88, 0.84, 0.48, size=10.2, color=ORANGE, bold=True, align=PP_ALIGN.CENTER, margin=0)

    for col_index, positions in enumerate(decoder_positions):
        for center_x, center_y in positions:
            _add_circle(slide, center_x - 0.08, center_y - 0.08, 0.16, TEAL if col_index == 0 else "6D97B5")
    _add_text(slide, "decoder", 5.70, 5.35, 1.05, 0.24, size=10.2, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_arrow(slide, 6.82, 3.86, 0.30, 0.25, MID_GRAY)
    _add_text(slide, "output\nprofile", 7.30, 3.72, 0.48, 0.56, size=10.0, color=GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)

    _add_text(slide, "LATENT SPACE", 8.30, 2.05, 1.34, 0.24, size=10.2, color=TEAL, bold=True, margin=0)
    _add_panel(slide, 8.22, 2.36, 4.58, 3.88, fill=PALE_TEAL, line="C9DFDB", radius=False)
    _add_rule(slide, 8.72, 5.55, 3.43, MID_GRAY, 0.018)
    _add_rule(slide, 8.72, 2.86, 0.018, MID_GRAY, 2.71)
    clusters = [
        (9.33, 4.63, BLUE),
        (10.47, 3.66, ORANGE),
        (11.48, 4.78, TEAL),
    ]
    offsets = [(-0.19, -0.07), (0.08, -0.16), (0.20, 0.05), (-0.06, 0.18), (0.02, 0.02)]
    for cx, cy, color in clusters:
        for dx, dy in offsets:
            _add_circle(slide, cx + dx, cy + dy, 0.13, color)
    _add_text(slide, "profiles with similar expression\noccupy nearby regions", 8.93, 5.72, 3.10, 0.42, size=10.5, color=GRAY, align=PP_ALIGN.CENTER, margin=0)

    _add_panel(slide, 0.50, 6.40, 12.30, 0.52, fill=NAVY, line=NAVY, radius=False)
    _add_text(slide, "expiMap:", 0.82, 6.54, 0.92, 0.23, size=12.2, color=WHITE, bold=True, margin=0)
    _add_text(slide, "each constrained latent variable corresponds to a Reactome gene program.", 1.81, 6.54, 10.40, 0.23, size=12.2, color="DCE7F2", margin=0)
    _add_source(slide, "Concept adapted from the midpoint presentation. expiMap: Lotfollahi et al., Nature Cell Biology (2023).")


def _slide_expimap_program_scores(slide):
    _add_slide_title(
        slide,
        "expiMap scores",
        "Program scores turn expression into pathway-level changes",
        "Each sample receives a Reactome program score; FLT minus GC gives the direction of change.",
    )
    scores = pd.read_csv(
        EXPIMAP_PAPER_DIR
        / "source_data/table_s33_representative_program_sample_scores.tsv.gz",
        sep="\t",
    )
    project_means = (
        scores.groupby(
            ["tissue", "display_label", "project", "condition"], observed=True
        )["project_centered_pathway_score"]
        .mean()
        .reset_index()
    )
    condition_means = (
        project_means.groupby(
            ["tissue", "display_label", "condition"], observed=True
        )["project_centered_pathway_score"]
        .mean()
        .unstack("condition")
    )
    condition_means["change"] = (
        condition_means["flight"] - condition_means["ground_control"]
    )
    order = ["thymus", "skin", "liver", "spleen"]
    condition_means = condition_means.reindex(order, level="tissue")

    _add_panel(slide, 0.47, 2.04, 8.82, 4.55, fill="F7F9FA", line="DDE4E8", radius=False)
    _add_text(slide, "REPRESENTATIVE OBSERVED OSDR PROGRAMS", 0.77, 2.25, 3.60, 0.23, size=9.3, color=BLUE, bold=True, margin=0)
    headers = [
        (4.52, "FLT score", ORANGE),
        (6.18, "GC score", BLUE),
        (7.87, "FLT - GC", TEAL),
    ]
    for x, label, color in headers:
        _add_text(slide, label, x, 2.47, 1.08, 0.26, size=10.5, color=color, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "-", 5.80, 2.47, 0.25, 0.26, size=14, color=MID_GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "=", 7.46, 2.47, 0.25, 0.26, size=13, color=MID_GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)

    tissue_colors = {
        "thymus": PURPLE,
        "skin": ORANGE,
        "liver": TEAL,
        "spleen": CORAL,
    }
    for index, ((tissue, program), row) in enumerate(condition_means.iterrows()):
        y = 2.92 + index * 0.76
        if index:
            _add_rule(slide, 0.77, y - 0.12, 8.17, "E1E6E9", 0.012)
        _add_rule(slide, 0.77, y + 0.05, 0.06, tissue_colors[tissue], 0.48)
        _add_text(slide, tissue.replace("_", " ").title(), 0.98, y - 0.01, 0.94, 0.24, size=11.0, color=tissue_colors[tissue], bold=True, margin=0)
        _add_text(slide, program, 1.91, y - 0.01, 2.35, 0.48, size=10.8, color=DARK, bold=True, valign=MSO_ANCHOR.MIDDLE, margin=0)
        values = [float(row["flight"]), float(row["ground_control"]), float(row["change"])]
        for x, value in zip((4.52, 6.18, 7.87), values):
            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x),
                Inches(y - 0.03),
                Inches(1.08),
                Inches(0.48),
            )
            _set_fill(cell, _diverging_color(value, 1.05))
            _set_line(cell, "D5DDE2", 0.6)
            _add_text(
                slide,
                f"{value:+.2f}",
                x,
                y + 0.06,
                1.08,
                0.24,
                size=11.0,
                color=WHITE if abs(value) > 0.58 else DARK,
                bold=True,
                align=PP_ALIGN.CENTER,
                margin=0,
            )
        _add_text(slide, "-", 5.80, y + 0.04, 0.25, 0.26, size=14, color=MID_GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)
        _add_text(slide, "=", 7.46, y + 0.04, 0.25, 0.26, size=13, color=MID_GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)

    legend_values = [(-1.0, "lower"), (0.0, "no shift"), (1.0, "higher")]
    for index, (value, label) in enumerate(legend_values):
        x = 3.30 + index * 1.20
        square = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(6.18), Inches(0.25), Inches(0.18)
        )
        _set_fill(square, _diverging_color(value, 1.0))
        _set_line(square, "D5DDE2", 0.6)
        _add_text(slide, label, x + 0.31, 6.15, 0.82, 0.23, size=8.8, color=GRAY, margin=0)

    _add_rule(slide, 9.62, 2.12, 0.015, "D5DDE2", 4.37)
    _add_text(slide, "HOW TO READ IT", 9.92, 2.24, 2.08, 0.24, size=10.2, color=BLUE, bold=True, margin=0)
    steps = [
        ("1", "Score each sample", "The constrained latent value represents one Reactome program."),
        ("2", "Compare within project", "Flight samples are compared with ground controls from the same study."),
        ("3", "Combine project shifts", "Project-level changes are integrated across the available missions."),
    ]
    for index, (number, heading, body) in enumerate(steps):
        y = 2.75 + index * 1.03
        _add_circle(slide, 9.93, y + 0.02, 0.30, TEAL)
        _add_text(slide, number, 9.93, y + 0.06, 0.30, 0.18, size=9.0, color=WHITE, bold=True, align=PP_ALIGN.CENTER, margin=0)
        _add_text(slide, heading, 10.39, y, 2.18, 0.27, size=12.0, color=NAVY, bold=True, margin=0)
        _add_text(slide, body, 10.39, y + 0.34, 2.20, 0.54, size=10.2, color=GRAY, margin=0)
    _add_panel(slide, 9.91, 5.92, 2.66, 0.53, fill=PALE_BLUE, line="CBDDE9", radius=False)
    _add_text(slide, "Negative change = lower program score in flight", 10.08, 6.03, 2.32, 0.28, size=10.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_source(slide, "Observed example scores: expiMap manuscript Table S33. Formal effects are estimated within project; lower scores do not prove biochemical inhibition.")


def _slide_expimap_tissue_results(slide):
    _add_slide_title(
        slide,
        "expiMap results",
        "Four tissues produced the clearest pathway patterns",
        "Effects were estimated within each project before the project-level shifts were combined.",
    )
    rows = [
        (
            "Thymus",
            "117 samples | 5 projects",
            "DNA repair; RHOA cytoskeletal cycle; lymphoid-stromal interactions",
            "Known involution may also involve lower repair, motility, and niche coordination.",
            PURPLE,
            "F1EDF7",
        ),
        (
            "Skin",
            "151 samples | 4 projects",
            "Chromatin regulation; DNA repair; Hedgehog; sphingolipids; cell junctions",
            "Barrier injury may include a broader loss of tissue maintenance and coordination.",
            ORANGE,
            PALE_GOLD,
        ),
        (
            "Liver",
            "197 samples | 9 projects",
            "MHC class II antigen presentation; T-cell receptor signaling",
            "Metabolic heterogeneity coexisted with lower adaptive immune communication.",
            TEAL,
            PALE_TEAL,
        ),
        (
            "Spleen",
            "100 samples | 5 projects",
            "T-cell receptor signaling; neutrophil degranulation; C-type lectin signaling",
            "The most consistent multi-pathway result joined adaptive and innate immune changes.",
            CORAL,
            PALE_CORAL,
        ),
    ]
    for index, (tissue, scope, pathways, interpretation, color, fill) in enumerate(rows):
        y = 2.02 + index * 1.17
        _add_panel(slide, 0.48, y, 12.38, 0.98, fill=fill, line=fill, radius=False)
        _add_rule(slide, 0.48, y, 0.075, color, 0.98)
        _add_text(slide, tissue, 0.76, y + 0.18, 1.42, 0.32, size=17.0, color=color, bold=True, margin=0)
        _add_text(slide, scope, 0.76, y + 0.56, 1.58, 0.22, size=9.5, color=GRAY, margin=0)
        _add_text(slide, "LOWER IN FLIGHT", 2.54, y + 0.16, 1.34, 0.22, size=9.0, color=BLUE, bold=True, margin=0)
        _add_text(slide, pathways, 2.54, y + 0.43, 4.31, 0.40, size=11.7, color=DARK, bold=True, valign=MSO_ANCHOR.MIDDLE, margin=0)
        _add_rule(slide, 7.02, y + 0.16, 0.012, "D3DBDF", 0.66)
        _add_text(slide, interpretation, 7.30, y + 0.19, 5.18, 0.58, size=12.2, color=NAVY, valign=MSO_ANCHOR.MIDDLE, margin=0)
    _add_source(slide, "Source: expiMap manuscript Figures 3 and 6. Lower latent scores do not by themselves prove biochemical inhibition.")


def _slide_expimap_evidence(slide):
    _add_slide_title(
        slide,
        "expiMap evidence",
        "Cross-checks narrowed the biological interpretation",
        "Blue cells mark directional support. Conventional GSEA FDR is reported separately.",
    )
    retained = pd.read_csv(
        EXPIMAP_PAPER_DIR / "source_data/table_2_retained_pathway_evidence.tsv",
        sep="\t",
    )
    evidence = pd.read_csv(
        EXPIMAP_PAPER_DIR / "source_data/table_s24_pathway_robustness_evidence.tsv",
        sep="\t",
    )
    kidney_spleen_evidence = pd.read_csv(
        EXPIMAP_PAPER_DIR / "source_data/table_s27_kidney_spleen_pathway_evidence.tsv",
        sep="\t",
    )
    retained = retained.loc[retained["analysis_role"].eq("main")]
    columns = [
        "tissue",
        "term",
        "ssgsea_direction_support",
        "preranked_gsea_direction_support",
        "heldout_direction_support",
        "seed_direction_support",
        "composition_proxy_support",
    ]
    evidence = pd.concat(
        [evidence[columns], kidney_spleen_evidence[columns]],
        ignore_index=True,
    ).drop_duplicates(["tissue", "term"], keep="last")
    data = retained.merge(evidence, on=["tissue", "term"], how="left", validate="one_to_one")
    order = {"thymus": 0, "skin": 1, "liver": 2, "spleen": 3}
    data = data.assign(_order=data["tissue"].map(order)).sort_values(["_order", "display_label"])
    if len(data) != 13 or data[columns[2:]].isna().any().any():
        raise ValueError("Unexpected retained expiMap evidence table")

    label_x = 0.48
    label_w = 3.10
    check_x = 3.78
    cell_w = 0.60
    fdr_x = 6.93
    headers = ["ssGSEA", "GSEA", "Held-out", "3 seeds", "Composition"]
    _add_text(slide, "Tissue and pathway", label_x, 2.00, label_w, 0.24, size=9.5, color=GRAY, bold=True, margin=0)
    for index, header in enumerate(headers):
        _add_text(slide, header, check_x + index * cell_w - 0.05, 1.98, cell_w + 0.10, 0.32, size=8.2, color=GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "GSEA FDR", fdr_x, 2.00, 0.80, 0.24, size=8.8, color=GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_rule(slide, label_x, 2.34, 7.28, NAVY, 0.022)

    tissue_colors = {"thymus": PURPLE, "skin": ORANGE, "liver": TEAL, "spleen": CORAL}
    row_h = 0.31
    start_y = 2.46
    last_tissue = None
    for index, row in enumerate(data.itertuples(index=False)):
        y = start_y + index * row_h
        if index % 2 == 0:
            shade = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(label_x), Inches(y - 0.015), Inches(7.28), Inches(row_h - 0.01))
            _set_fill(shade, "F6F8F9")
            shade.line.fill.background()
        if last_tissue is not None and row.tissue != last_tissue:
            _add_rule(slide, label_x, y - 0.045, 7.28, "BBC7CE", 0.018)
        color = tissue_colors[row.tissue]
        label = f"{row.tissue.title()}: {row.display_label}"
        _add_rule(slide, label_x + 0.03, y + 0.105, 0.12, color, 0.045)
        _add_text(slide, label, label_x + 0.23, y + 0.035, label_w - 0.20, 0.23, size=8.8, color=DARK, margin=0)
        values = [
            row.ssgsea_direction_support,
            row.preranked_gsea_direction_support,
            row.heldout_direction_support,
            row.seed_direction_support,
            row.composition_proxy_support,
        ]
        for check_index, supported in enumerate(values):
            x = check_x + check_index * cell_w + 0.13
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + 0.025), Inches(0.34), Inches(0.24))
            _set_fill(box, BLUE if bool(supported) else "E7ECEF")
            box.line.fill.background()
            _add_text(slide, "+" if bool(supported) else "-", x, y + 0.025, 0.34, 0.24, size=9.0, color=WHITE if bool(supported) else MID_GRAY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
        fdr = float(row.gsea_fdr)
        fdr_label = "<0.001" if fdr < 0.001 else f"{fdr:.3f}"
        _add_text(slide, fdr_label, fdr_x, y + 0.035, 0.80, 0.22, size=8.8, color=NAVY if fdr < 0.05 else GRAY, bold=fdr < 0.05, align=PP_ALIGN.CENTER, margin=0)
        last_tissue = row.tissue

    _add_rule(slide, 8.02, 2.08, 0.015, "D5DDE2", 4.38)
    _add_text(slide, "WHAT REMAINED", 8.38, 2.10, 1.62, 0.24, size=10.2, color=TEAL, bold=True, margin=0)
    summaries = [
        ("Spleen", "All three programs passed GSEA FDR < 0.05 and all five directional checks.", CORAL),
        ("Skin", "Chromatin, DNA repair, and cell-junction programs passed GSEA FDR < 0.05.", ORANGE),
        ("Thymus", "DNA repair passed GSEA FDR < 0.001; the cytoskeletal and niche programs had directional support.", PURPLE),
        ("Liver", "MHC II and T-cell receptor scores were lower, but conventional FDR was 0.121 and 0.051.", TEAL),
    ]
    for index, (tissue, text, color) in enumerate(summaries):
        y = 2.59 + index * 0.90
        if index:
            _add_rule(slide, 8.38, y - 0.16, 4.02, "E0E5E8", 0.012)
        _add_text(slide, tissue, 8.38, y, 1.12, 0.26, size=13.5, color=color, bold=True, margin=0)
        _add_text(slide, text, 9.56, y - 0.02, 2.84, 0.63, size=10.8, color=DARK, margin=0)
    _add_panel(slide, 8.37, 6.19, 4.05, 0.46, fill=PALE_BLUE, line="C9DCE9", radius=False)
    _add_text(slide, "Spleen was the strongest multi-pathway result.", 8.58, 6.29, 3.62, 0.24, size=11.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_source(slide, "Source: expiMap manuscript Table 2 and Table S24. Held-out checks are internal project-wise validation.")


def _slide_why_synthetic(slide):
    _add_slide_title(
        slide,
        "Background",
        "What is synthetic transcriptomics?",
        "A generator learns patterns in measured RNA-seq data, then samples new expression vectors under chosen conditions.",
    )

    columns = [
        (
            0.50,
            "01",
            "Observed profiles",
            "Measured gene expression from real mouse tissue",
            BLUE,
            PALE_BLUE,
        ),
        (
            4.62,
            "02",
            "Learn the distribution",
            "Capture gene relationships and variation by tissue and FLT/GC condition",
            ORANGE,
            PALE_GOLD,
        ),
        (
            8.74,
            "03",
            "Synthetic profiles",
            "Sample new numeric expression vectors for a selected context",
            TEAL,
            PALE_TEAL,
        ),
    ]
    for index, (x, number, heading, body, color, fill) in enumerate(columns):
        if index:
            _add_arrow(slide, x - 0.47, 4.00, 0.34, 0.28, MID_GRAY)
        _add_text(slide, number, x, 2.10, 0.42, 0.25, size=11, color=color, bold=True, margin=0)
        _add_text(slide, heading, x + 0.50, 2.04, 3.08, 0.38, size=17.2, color=NAVY, bold=True, margin=0)
        _add_panel(slide, x, 2.58, 3.58, 3.72, fill=fill, line=fill, radius=False)

        if index in (0, 2):
            cell_colors = [color, "AFC3D4", color, "D8E2E8", color]
            for row in range(5):
                for col in range(6):
                    cell_color = cell_colors[(row * 2 + col + index) % len(cell_colors)]
                    square = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(x + 0.36 + col * 0.40),
                        Inches(3.03 + row * 0.31),
                        Inches(0.29),
                        Inches(0.20),
                    )
                    _set_fill(square, cell_color)
                    square.line.fill.background()
            _add_text(slide, "genes", x + 0.36, 4.72, 2.29, 0.20, size=9.5, color=GRAY, align=PP_ALIGN.CENTER, margin=0)
            _add_text(slide, "profiles", x + 2.74, 3.62, 0.54, 0.22, size=9.5, color=GRAY, align=PP_ALIGN.CENTER, margin=0)
        else:
            node_positions = [
                (x + 0.36, 3.72),
                (x + 1.03, 3.20),
                (x + 1.03, 4.28),
                (x + 1.76, 3.74),
                (x + 2.48, 3.20),
                (x + 2.48, 4.28),
                (x + 3.08, 3.74),
            ]
            connections = [(0, 1), (0, 2), (1, 3), (2, 3), (3, 4), (3, 5), (4, 6), (5, 6)]
            for start, end in connections:
                x1, y1 = node_positions[start]
                x2, y2 = node_positions[end]
                line = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    Inches(x1 + 0.08),
                    Inches(y1 + 0.08),
                    Inches(x2 + 0.08),
                    Inches(y2 + 0.08),
                )
                _set_line(line, "C6B07A", 1.5)
            for node_index, (node_x, node_y) in enumerate(node_positions):
                _add_circle(slide, node_x, node_y, 0.20, ORANGE if node_index in (0, 3, 6) else GOLD)

        _add_rule(slide, x + 0.28, 5.08, 3.02, "D3DCE1", 0.018)
        _add_text(slide, body, x + 0.30, 5.30, 2.98, 0.72, size=14.2, color=DARK, valign=MSO_ANCHOR.MIDDLE, margin=0)


def _slide_scientific_objective(slide):
    _add_slide_title(
        slide,
        "Study goals",
        "Match tissue distributions, then test FLT versus GC biology",
        "The first goal validates the generator. The second asks whether synthetic data improves a real-data analysis.",
    )

    _add_text(slide, "01", 0.58, 2.08, 0.42, 0.26, size=11, color=TEAL, bold=True, margin=0)
    _add_text(slide, "Reproduce tissue structure", 1.08, 2.02, 4.50, 0.38, size=18, color=NAVY, bold=True, margin=0)
    _add_text(
        slide,
        "Synthetic bulk RNA-seq should occupy the same tissue-defined expression space as real profiles.",
        0.58,
        2.49,
        5.58,
        0.56,
        size=13.4,
        color=DARK,
        margin=0,
    )

    _add_rule(slide, 0.83, 4.58, 4.92, MID_GRAY, 0.018)
    _add_rule(slide, 0.83, 3.14, 0.018, MID_GRAY, 1.46)

    def add_cross(x, y, color):
        first = _add_rule(slide, x, y + 0.06, 0.18, color, 0.025)
        first.rotation = 45
        second = _add_rule(slide, x, y + 0.06, 0.18, color, 0.025)
        second.rotation = -45

    clusters = [
        (1.52, 3.98, TEAL, "Liver"),
        (3.12, 3.43, ORANGE, "Muscle"),
        (4.77, 4.03, PURPLE, "Thymus"),
    ]
    offsets = [(-0.23, -0.10), (0.05, -0.19), (0.22, 0.02), (-0.08, 0.18)]
    for center_x, center_y, color, label in clusters:
        for point_index, (dx, dy) in enumerate(offsets):
            _add_circle(slide, center_x + dx, center_y + dy, 0.13, color)
            add_cross(center_x + dx + 0.10, center_y + dy + (0.08 if point_index % 2 else -0.03), color)
        _add_text(slide, label, center_x - 0.37, center_y + 0.40, 0.98, 0.22, size=9.2, color=GRAY, align=PP_ALIGN.CENTER, margin=0)
    _add_circle(slide, 0.93, 3.05, 0.12, NAVY)
    _add_text(slide, "Real", 1.10, 3.02, 0.55, 0.20, size=9.2, color=GRAY, margin=0)
    add_cross(1.70, 3.03, NAVY)
    _add_text(slide, "Generated", 1.95, 3.02, 0.90, 0.20, size=9.2, color=GRAY, margin=0)

    _add_rule(slide, 6.50, 2.05, 0.015, "D6DEE3", 2.72)

    _add_text(slide, "02", 6.86, 2.08, 0.42, 0.26, size=11, color=CORAL, bold=True, margin=0)
    _add_text(slide, "Preserve the FLT/GC difference", 7.36, 2.02, 4.96, 0.38, size=18, color=NAVY, bold=True, margin=0)
    _add_text(
        slide,
        "Within each tissue, generated profiles should retain the smaller condition signal.",
        6.86,
        2.49,
        5.60,
        0.56,
        size=13.4,
        color=DARK,
        margin=0,
    )
    _add_text(slide, "FLT", 7.18, 3.24, 0.72, 0.28, size=14, color=CORAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "GC", 9.08, 3.24, 0.72, 0.28, size=14, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    flt_points = [(7.18, 3.73), (7.52, 3.58), (7.79, 3.81), (7.37, 4.04)]
    gc_points = [(9.09, 3.73), (9.43, 3.58), (9.70, 3.81), (9.28, 4.04)]
    for x, y in flt_points:
        _add_circle(slide, x, y, 0.17, CORAL)
    for x, y in gc_points:
        _add_circle(slide, x, y, 0.17, BLUE)
    _add_text(slide, "vs", 8.38, 3.75, 0.42, 0.24, size=11.5, color=MID_GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_arrow(slide, 9.98, 3.67, 0.31, 0.28, MID_GRAY)
    _add_panel(slide, 10.49, 3.30, 2.00, 1.14, fill=PALE_BLUE, line="D4E2EC", radius=False)
    _add_text(slide, "OUTPUT", 10.68, 3.47, 1.63, 0.18, size=9.2, color=BLUE, bold=True, margin=0)
    _add_text(
        slide,
        "Tissue-specific\ngenes and pathways",
        10.68,
        3.75,
        1.63,
        0.54,
        size=12.2,
        color=NAVY,
        bold=True,
        margin=0,
    )

    _add_rule(slide, 0.60, 5.12, 12.08, TEAL, 0.035)
    _add_text(slide, "Model choice", 0.62, 5.43, 1.55, 0.30, size=16.5, color=TEAL, bold=True, margin=0)
    _add_text(
        slide,
        "Use the pipeline that matches tissue structure and performs best on held-out real FLT/GC samples.",
        2.18,
        5.34,
        7.40,
        0.58,
        size=15,
        color=NAVY,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )
    _add_text(slide, "Gene effects and BH FDR use observed OSDR data.", 9.88, 5.42, 2.57, 0.46, size=11.5, color=GRAY, margin=0)


def _slide_2(slide):
    _add_slide_title(
        slide,
        "Question",
        "Small studies and study effects complicate tissue comparisons",
        "Can a generator help rank spaceflight signal without pretending it creates new animals?",
    )
    _add_text(slide, "NASA OSDR", 0.62, 2.04, 2.20, 0.30, size=17, color=BLUE, bold=True)
    _add_text(slide, "Observed spaceflight cohort", 0.62, 2.37, 2.80, 0.24, size=11.5, color=GRAY)
    _add_text(slide, "1,610", 0.58, 2.64, 2.25, 0.69, size=42, color=NAVY, bold=True)
    _add_text(slide, "profiles", 0.64, 3.29, 1.15, 0.25, size=12.5, color=GRAY)
    _add_text(slide, "75", 3.01, 2.72, 1.20, 0.54, size=31, color=NAVY, bold=True)
    _add_text(slide, "accessions", 3.04, 3.26, 1.40, 0.25, size=12.5, color=GRAY)

    bar_x = 0.64
    bar_y = 4.03
    bar_w = 5.24
    flt_w = bar_w * 835 / 1610
    _add_text(slide, "835 FLT", bar_x, 3.65, 1.10, 0.26, size=13.5, color=CORAL, bold=True)
    _add_text(slide, "775 GC", bar_x + bar_w - 1.10, 3.65, 1.10, 0.26, size=13.5, color=BLUE, bold=True, align=PP_ALIGN.RIGHT)
    _add_rule(slide, bar_x, bar_y, flt_w, CORAL, 0.12)
    _add_rule(slide, bar_x + flt_w, bar_y, bar_w - flt_w, BLUE, 0.12)

    _add_rule(slide, 6.49, 2.02, 0.018, "CBD4DA", 2.21)
    _add_text(slide, "ARCHS4 mouse", 6.88, 2.04, 2.40, 0.30, size=17, color=TEAL, bold=True)
    _add_text(slide, "Reference pretraining cohort", 6.88, 2.37, 2.80, 0.24, size=11.5, color=GRAY)
    _add_text(slide, "997,515", 6.84, 2.64, 2.55, 0.69, size=39, color=NAVY, bold=True)
    _add_text(slide, "profiles screened", 6.90, 3.29, 1.75, 0.25, size=12.5, color=GRAY)
    _add_arrow(slide, 9.26, 2.88, 0.42, 0.30, MID_GRAY)
    _add_text(slide, "17,244", 9.86, 2.66, 1.70, 0.54, size=29, color=NAVY, bold=True)
    _add_text(slide, "selected", 9.89, 3.24, 1.10, 0.25, size=12.5, color=GRAY)
    _add_text(slide, "20 tissues", 11.42, 2.73, 1.25, 0.36, size=17, color=TEAL, bold=True)

    _add_rule(slide, 0.62, 4.57, 12.08, "D6DEE3", 0.018)
    _add_text(slide, "Study effects can look like spaceflight biology.", 0.62, 4.82, 11.80, 0.42, size=22, color=NAVY, bold=True)
    challenge_rows = [
        ("01", "Preserve tissue and FLT/GC structure."),
        ("02", "Avoid memorizing individual profiles."),
        ("03", "Test biological effects in observed OSDR samples."),
    ]
    for index, (number, text) in enumerate(challenge_rows):
        y = 5.42 + index * 0.42
        _add_text(slide, number, 0.64, y, 0.42, 0.26, size=10.5, color=ORANGE, bold=True, margin=0)
        _add_text(slide, text, 1.17, y - 0.03, 10.95, 0.31, size=14.5, color=DARK, margin=0)
    _add_source(slide, "Sources: NASA OSDR Biological Data API; ARCHS4 mouse v2.5.")


def _slide_wgan_explainer(slide, wgan_gif: Path):
    _add_slide_title(
        slide,
        "WGAN-GP",
        "A generator learns by competing with a critic",
        "The generator makes conditioned expression profiles; the critic compares them with measured RNA-seq.",
    )
    _add_panel(slide, 0.48, 1.94, 5.40, 4.88, fill="EEF2F6", line="D9E1E5", radius=False)
    _add_picture_contain(
        slide,
        wgan_gif,
        0.73,
        2.03,
        4.90,
        4.67,
        alt="Animated teaching diagram of a WGAN generator competing with a critic",
    )

    _add_text(slide, "TRAINING LOOP", 6.30, 2.03, 1.55, 0.24, size=10.2, color=BLUE, bold=True, margin=0)
    stages = [
        (2.50, "01", "Generate", "Noise and a requested condition become one synthetic expression profile.", ORANGE),
        (3.78, "02", "Score", "The critic scores measured and generated profiles without seeing their labels.", BLUE),
        (5.06, "03", "Update", "The critic sharpens the comparison; the generator learns to close the gap.", TEAL),
    ]
    for y, number, heading, body, color in stages:
        _add_text(slide, number, 6.30, y + 0.02, 0.40, 0.25, size=10.5, color=color, bold=True, margin=0)
        _add_rule(slide, 6.82, y + 0.07, 0.055, color, 0.72)
        _add_text(slide, heading, 7.10, y, 1.55, 0.30, size=16.0, color=NAVY, bold=True, margin=0)
        _add_text(slide, body, 8.66, y - 0.02, 3.78, 0.68, size=11.4, color=DARK, valign=MSO_ANCHOR.MIDDLE, margin=0)
        if y < 5.0:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Inches(7.54),
                Inches(y + 0.83),
                Inches(0.22),
                Inches(0.28),
            )
            _set_fill(arrow, MID_GRAY)
            arrow.line.fill.background()

    _add_rule(slide, 6.28, 6.18, 6.22, "D5DDE2", 0.015)
    _add_text(slide, "CONDITION", 6.31, 6.39, 0.92, 0.18, size=8.5, color=TEAL, bold=True, margin=0)
    condition_tiles = [
        (7.34, 1.03, "Tissue"),
        (8.48, 1.12, "FLT / GC"),
        (9.71, 1.15, "Accession"),
        (10.97, 1.42, "Material type"),
    ]
    for x, width, label in condition_tiles:
        _add_panel(slide, x, 6.29, width, 0.38, fill=PALE_TEAL, line="C9DFDB", radius=False)
        _add_text(slide, label, x, 6.39, width, 0.16, size=8.5, color=TEAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_source(slide, "Animation reused from the project midpoint presentation. WGAN-GP follows the adversarial RNA-seq framework of Vinas et al. (2022).")


def _slide_3(slide):
    _add_slide_title(
        slide,
        "Method",
        "We built a configurable bulk RNA-seq generation pipeline",
        "The downstream analysis used the outlined options; the remaining choices stay configurable.",
    )

    def option_tile(label, x, y, w, color, selected=False, *, size=9.1):
        if selected:
            tile = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x),
                Inches(y),
                Inches(w),
                Inches(0.36),
            )
            _set_fill(tile, WHITE)
            _set_line(tile, color, 1.8)
            _add_rule(slide, x, y, 0.035, color, 0.36)
        _add_text(
            slide,
            label,
            x + 0.05,
            y + 0.055,
            w - 0.10,
            0.23,
            size=size,
            color=NAVY if selected else MID_GRAY,
            bold=selected,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
            margin=0,
        )

    def segmented_row(x, y, label, options, color):
        _add_text(slide, label, x + 0.16, y, 2.02, 0.18, size=8.4, color=GRAY, bold=True, margin=0)
        gap = 0.05
        available = 2.03
        tile_width = (available - gap * (len(options) - 1)) / len(options)
        for index, (option, selected) in enumerate(options):
            option_tile(
                option,
                x + 0.16 + index * (tile_width + gap),
                y + 0.22,
                tile_width,
                color,
                selected,
                size=8.2 if len(options) == 3 else 8.8,
            )

    def stage_panel(x, number, title, color, fill):
        _add_text(slide, number, x + 0.14, 2.15, 0.34, 0.24, size=11.5, color=MID_GRAY, bold=True, margin=0)
        _add_text(slide, title, x + 0.56, 2.11, 1.64, 0.32, size=14.2, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)
        _add_rule(slide, x + 0.14, 2.55, 2.05, color, 0.022)

    xs = [0.35, 2.90, 5.45, 8.00, 10.55]
    for separator_x in [2.79, 5.34, 7.89, 10.44]:
        _add_rule(slide, separator_x, 2.03, 0.012, "D9E0E4", 3.24)

    stage_panel(xs[0], "01", "Data scope", BLUE, WHITE)
    segmented_row(xs[0], 2.72, "Sources", [("OSDR API", True), ("ARCHS4", True)], BLUE)
    segmented_row(xs[0], 3.48, "Studies", [("Single", False), ("Multiple", True)], BLUE)
    segmented_row(xs[0], 4.24, "Tissues", [("Per tissue", False), ("All tissues", True)], BLUE)

    stage_panel(xs[1], "02", "Transform", TEAL, WHITE)
    segmented_row(xs[1], 2.72, "Expression", [("Raw", False), ("CPM", False), ("TPM", True)], TEAL)
    segmented_row(xs[1], 3.48, "Scaling", [("None", False), ("Z-score", False), ("MaxAbs", True)], TEAL)
    segmented_row(xs[1], 4.24, "Features", [("All", False), ("HVG", False), ("L1000", True)], TEAL)

    stage_panel(xs[2], "03", "Harmonization", TEAL, WHITE)
    _add_text(slide, "Global or study correction", xs[2] + 0.16, 2.72, 2.02, 0.18, size=8.4, color=GRAY, bold=True, margin=0)
    for index, (label, selected) in enumerate([
        ("None", True),
        ("Within-study z-score", False),
        ("ComBat / MBatch", False),
        ("MOBER", False),
    ]):
        option_tile(label, xs[2] + 0.16, 2.94 + index * 0.50, 2.03, TEAL, selected, size=9.2)

    stage_panel(xs[3], "04", "Model training", ORANGE, WHITE)
    segmented_row(xs[3], 2.72, "Generator", [("WGAN-GP", False), ("DDIM", True)], ORANGE)
    _add_text(slide, "Training source", xs[3] + 0.16, 3.48, 2.02, 0.18, size=8.4, color=GRAY, bold=True, margin=0)
    for index, (label, selected) in enumerate([
        ("OSDR only", False),
        ("ARCHS4 only", False),
        ("Pretrain + adapt", True),
    ]):
        option_tile(label, xs[3] + 0.16, 3.70 + index * 0.50, 2.03, ORANGE, selected, size=9.2)

    stage_panel(xs[4], "05", "Conditioning", TEAL, WHITE)
    _add_text(slide, "Model inputs", xs[4] + 0.16, 2.72, 2.02, 0.18, size=8.4, color=GRAY, bold=True, margin=0)
    for index, (label, selected) in enumerate([
        ("Tissue", True),
        ("FLT / GC", True),
        ("Accession", True),
        ("Material type", True),
        ("Sex / age", False),
    ]):
        option_tile(label, xs[4] + 0.16, 2.92 + index * 0.42, 2.03, TEAL, selected, size=8.9)

    _add_panel(slide, 0.35, 5.57, 12.55, 1.23, fill="F7F9FA", line="F7F9FA", radius=False)
    _add_rule(slide, 0.35, 5.57, 12.55, "C9D3D9", 0.018)
    _add_text(slide, "Downstream", 0.64, 5.80, 1.42, 0.25, size=10.8, color=NAVY, bold=True)
    _add_text(slide, "configuration", 0.64, 6.17, 1.35, 0.24, size=9.6, color=GRAY)

    def selected_step(x, width, heading, value, color):
        _add_panel(slide, x, 5.72, width, 0.49, fill=WHITE, line=color, radius=False)
        _add_text(slide, heading, x + 0.10, 5.78, width - 0.20, 0.12, size=7.1, color=color, bold=True, margin=0)
        _add_text(slide, value, x + 0.10, 5.94, width - 0.20, 0.19, size=9.1, color=DARK, bold=True, align=PP_ALIGN.CENTER, margin=0)

    selected_path = [
        (2.10, 2.15, "PREPROCESS", "TPM / MaxAbs / 974 landmarks", TEAL),
        (4.65, 1.85, "PRETRAIN", "ARCHS4", TEAL),
        (6.90, 1.85, "ADAPT", "OSDR", BLUE),
        (9.15, 2.45, "GENERATE", "Conditional DDIM", ORANGE),
    ]
    for index, (x, width, heading, value, color) in enumerate(selected_path):
        selected_step(x, width, heading, value, color)
        if index < len(selected_path) - 1:
            _add_arrow(slide, x + width + 0.07, 5.86, 0.23, 0.20, MID_GRAY)

    _add_text(slide, "CONDITION", 2.10, 6.37, 0.72, 0.16, size=7.3, color=TEAL, bold=True, margin=0)
    condition_tiles = [
        (2.88, 0.90, "Tissue"),
        (3.85, 1.00, "FLT / GC"),
        (4.92, 1.14, "Accession"),
        (6.13, 1.38, "Material type"),
    ]
    for x, width, label in condition_tiles:
        option_tile(label, x, 6.27, width, TEAL, True, size=7.8)

    _add_text(slide, "HARMONIZE", 7.76, 6.37, 0.82, 0.16, size=7.3, color=TEAL, bold=True, margin=0)
    option_tile("None", 8.65, 6.27, 0.78, TEAL, True, size=7.8)
    _add_text(slide, "SCOPE", 9.70, 6.37, 0.46, 0.16, size=7.3, color=BLUE, bold=True, margin=0)
    option_tile("All tissues", 10.24, 6.27, 1.20, BLUE, True, size=7.8)
    _add_source(slide, "Screen: 463 planned configurations and nine matched liver harmonization arms. Models: Vinas et al. (2022); Lacan et al. (2026).")


def _slide_5(slide, trajectory: Path):
    _add_slide_title(
        slide,
        "Diffusion",
        "Diffusion learns tissue structure from noise",
        "The same PCA axes follow 1,024 generated profiles through reverse diffusion.",
    )
    _add_picture_contain(
        slide,
        trajectory,
        0.34,
        1.78,
        12.64,
        4.72,
        alt="DDIM reverse trajectory at timesteps 1000, 200 and 0",
    )
    _add_text(slide, "t = 1000: noise", 0.78, 6.50, 2.5, 0.28, size=13, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "t = 200: partial structure", 5.23, 6.50, 2.9, 0.28, size=13, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "t = 0: tissue-conditioned profiles", 9.15, 6.50, 3.4, 0.28, size=13, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    _add_source(slide, "Gray points are real ARCHS4 profiles; colors identify generated tissue conditions. PC1 and PC2 limits are shared across panels.")


def _slide_diffusion_explainer(slide, diffusion_gif: Path):
    _add_slide_title(
        slide,
        "Diffusion",
        "Diffusion turns noise into a conditioned sample",
        "The model repeatedly removes predicted noise while following the requested biological context.",
    )
    _add_panel(slide, 0.48, 1.98, 5.30, 4.76, fill="EEF2F6", line="D9E1E5", radius=False)
    _add_picture_contain(
        slide,
        diffusion_gif,
        0.76,
        2.07,
        4.74,
        4.55,
        alt="Animated two-dimensional teaching example of FLT- and GC-conditioned reverse diffusion",
    )
    _add_text(slide, "FOR RNA-SEQ", 6.18, 2.06, 1.52, 0.24, size=10.2, color=BLUE, bold=True, margin=0)
    stages = [
        (2.54, "t = 1000", "Random noise", "974 unconstrained gene values", MID_GRAY, "noise"),
        (3.78, "reverse steps", "Conditional denoising", "Tissue + FLT/GC + accession + material", TEAL, "denoise"),
        (5.15, "t = 0", "Synthetic expression", "One generated 974-gene profile", ORANGE, "profile"),
    ]
    for y, tag, heading, body, color, kind in stages:
        _add_rule(slide, 6.29, y + 0.07, 0.07, color, 0.68)
        _add_text(slide, tag.upper(), 6.58, y, 1.34, 0.20, size=8.5, color=color, bold=True, margin=0)
        _add_text(slide, heading, 6.58, y + 0.26, 2.32, 0.30, size=15.0, color=NAVY, bold=True, margin=0)
        if body and kind != "denoise":
            _add_text(slide, body, 9.05, y + 0.21, 2.10, 0.43, size=11.2, color=DARK, valign=MSO_ANCHOR.MIDDLE, margin=0)
        if kind == "noise":
            rng = np.random.default_rng(2026)
            for dx, dy in rng.uniform((0.0, 0.0), (1.18, 0.45), size=(13, 2)):
                _add_circle(slide, 11.07 + dx, y + 0.16 + dy, 0.06, MID_GRAY)
        elif kind == "denoise":
            for index, label in enumerate(("Tissue", "FLT / GC", "Study", "Material")):
                x = 9.23 + (index % 2) * 1.52
                yy = y + 0.08 + (index // 2) * 0.34
                _add_panel(slide, x, yy, 1.35, 0.26, fill=PALE_TEAL, line="C9DFDB", radius=False)
                _add_text(slide, label, x, yy + 0.05, 1.35, 0.15, size=8.0, color=TEAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
        else:
            profile_values = [0.28, 0.51, 0.36, 0.62, 0.44, 0.22, 0.55, 0.31]
            for index, height in enumerate(profile_values):
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(11.45 + index * 0.16),
                    Inches(y + 0.64 - height),
                    Inches(0.10),
                    Inches(height),
                )
                _set_fill(bar, ORANGE if index % 2 == 0 else "E0A24D")
                bar.line.fill.background()
    first_arrow = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, Inches(6.72), Inches(3.37), Inches(0.22), Inches(0.31)
    )
    _set_fill(first_arrow, MID_GRAY)
    first_arrow.line.fill.background()
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, Inches(6.72), Inches(4.73), Inches(0.22), Inches(0.31)
    )
    _set_fill(arrow, MID_GRAY)
    arrow.line.fill.background()
    _add_source(slide, "Animation reused from the project midpoint presentation. The 2-D spiral is explanatory; the trained model generates 974-gene vectors.")


def _slide_6(slide, tissue_accession_pca: Path):
    _add_slide_title(
        slide,
        "Diffusion output",
        "Tissue and study structure dominate the PCA space",
        "Left: tissue. Right: OSDR accession. Circles are observed profiles; crosses are matched DDIM profiles.",
    )
    _add_picture_contain(
        slide,
        tissue_accession_pca,
        0.38,
        1.79,
        12.52,
        4.66,
        alt="PCA of observed and matched DDIM profiles colored by tissue and OSDR accession",
    )
    _add_panel(slide, 0.48, 6.47, 12.37, 0.47, fill=PALE_BLUE, line="C9DCE9", radius=False)
    _add_text(
        slide,
        "Tissue and accession clusters are visible across both observed and generated profiles.",
        0.76,
        6.56,
        11.80,
        0.26,
        size=13.0,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_source(slide, "Same 293 observed and 293 seed-5020 generated profiles; PCA was fitted once to the combined locked expression space.")


def _slide_condition_accession_pca(slide, condition_accession_pca: Path):
    _add_slide_title(
        slide,
        "Diffusion output",
        "Flight condition is subtler than study structure",
        "Left: FLT or GC. Right: OSDR accession. Circles are observed profiles; crosses are matched DDIM profiles.",
    )
    _add_picture_contain(
        slide,
        condition_accession_pca,
        0.38,
        1.79,
        12.52,
        4.66,
        alt="PCA of observed and matched DDIM profiles colored by condition and OSDR accession",
    )
    _add_panel(slide, 0.48, 6.47, 12.37, 0.47, fill=PALE_BLUE, line="C9DCE9", radius=False)
    _add_text(
        slide,
        "Flight and ground-control profiles overlap within the larger accession-defined structure.",
        0.76,
        6.56,
        11.80,
        0.26,
        size=13.0,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_source(slide, "The accession colors are shared across the two PCA slides; FLT-GC effects are evaluated within study.")


def _slide_4(slide, architecture_figure: Path):
    _add_slide_title(
        slide,
        "Validation",
        "DDIM had lower separability and distributional distance",
        "Both generators matched expression. DDIM was harder to distinguish from real profiles and had higher F1.",
    )
    _add_text(slide, "Architecture from Lacan et al.", 0.52, 2.08, 7.12, 0.32, size=16.5, color=NAVY, bold=True)
    _add_picture_contain(
        slide,
        architecture_figure,
        0.20,
        2.39,
        7.48,
        2.73,
        alt="Lacan et al. Figure 1C residual diffusion generator architecture",
    )
    _add_text(slide, "Residual dense denoiser conditioned on timestep and tissue.", 0.56, 5.17, 6.90, 0.28, size=11.3, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "OSDR adaptation adds FLT/GC, accession and material context through LoRA.", 0.56, 5.51, 6.90, 0.34, size=10.6, color=GRAY, align=PP_ALIGN.CENTER)

    _add_rule(slide, 7.83, 2.05, 0.015, "D5DDE2", 4.32)
    _add_text(slide, "Metric", 8.04, 2.17, 1.42, 0.28, size=10.0, color=GRAY, bold=True)
    _add_text(slide, "WGAN-GP", 9.57, 2.17, 0.85, 0.28, size=10.0, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "DDIM", 10.52, 2.17, 0.76, 0.28, size=10.0, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "Read", 11.43, 2.17, 1.10, 0.28, size=10.0, color=GRAY, bold=True)
    _add_rule(slide, 8.04, 2.56, 4.58, NAVY, 0.022)
    metrics = [
        ("Correlation", "0.976", "0.974", "similar"),
        ("F1", "0.985", "0.997", "higher"),
        ("Adversarial acc.", "0.636", "0.475", "near 0.5"),
        ("FD / real P95", "0.144", "0.074", "lower"),
    ]
    for index, (label, wgan, ddim, reading) in enumerate(metrics):
        y = 2.80 + index * 0.59
        if index % 2 == 0:
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.02), Inches(y - 0.07), Inches(4.61), Inches(0.50))
            _set_fill(shape, "F5F7F8")
            shape.line.fill.background()
        _add_text(slide, label, 8.10, y, 1.39, 0.32, size=10.3, color=DARK, bold=index >= 2)
        _add_text(slide, wgan, 9.57, y, 0.85, 0.32, size=12.6, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, ddim, 10.52, y, 0.76, 0.32, size=12.6, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, reading, 11.43, y, 1.10, 0.32, size=10.2, color=GRAY)

    _add_rule(slide, 0.55, 5.91, 7.03, "D5DDE2", 0.015)
    _add_text(slide, "ARCHS4 tissue probe", 0.55, 6.08, 1.75, 0.24, size=10.4, color=BLUE, bold=True, margin=0)
    _add_text(slide, "Real 0.781  |  Synthetic 0.781", 4.30, 6.06, 3.12, 0.27, size=10.4, color=DARK, align=PP_ALIGN.RIGHT, margin=0)
    _add_text(slide, "Balanced accuracy is preserved in generated tissue labels.", 0.55, 6.39, 6.87, 0.25, size=9.7, color=GRAY, align=PP_ALIGN.CENTER, margin=0)
    _add_rule(slide, 8.04, 5.36, 0.055, TEAL, 0.78)
    _add_text(slide, "Use DDIM", 8.29, 5.40, 1.32, 0.31, size=16, color=NAVY, bold=True, margin=0)
    _add_text(slide, "Higher F1 with lower separability and FD.", 9.73, 5.38, 2.78, 0.47, size=11.3, color=DARK, valign=MSO_ANCHOR.MIDDLE, margin=0)
    _add_text(slide, "AA near 0.5 means real and generated profiles are difficult to separate.", 8.29, 5.97, 4.22, 0.40, size=9.5, color=GRAY, margin=0)
    _add_source(slide, "Architecture excerpt: Lacan et al. (2026), Fig. 1C, doi:10.1186/s12859-026-06470-8. Metrics use each model's stated split and are not paired.")


def _slide_7(slide):
    _add_slide_title(
        slide,
        "Analysis",
        "Five arms separate gene ranking from classifier fitting",
        "Guided arms use synthetic profiles to help choose genes; they do not treat generated profiles as new animals.",
    )
    _add_text(slide, "INPUTS", 0.47, 2.00, 0.66, 0.22, size=8.8, color=GRAY, bold=True, margin=0)
    _add_data_badge(slide, "R", 1.20, 1.98, BLUE)
    _add_text(slide, "real OSDR", 1.52, 1.98, 1.18, 0.24, size=10.8, color=BLUE, bold=True, margin=0)
    _add_data_badge(slide, "S", 2.91, 1.98, CORAL)
    _add_text(slide, "matched DDIM", 3.23, 1.98, 1.42, 0.24, size=10.8, color=CORAL, bold=True, margin=0)
    _add_text(slide, "GUIDED", 6.10, 1.99, 0.70, 0.22, size=8.8, color=TEAL, bold=True, margin=0)
    _add_text(slide, "R + S help rank genes", 6.85, 1.97, 2.05, 0.25, size=10.5, color=DARK, margin=0)
    _add_text(slide, "5%", 9.41, 1.99, 0.38, 0.22, size=8.8, color=ORANGE, bold=True, margin=0)
    _add_text(slide, "S carries 5% of total fit weight", 9.84, 1.97, 2.78, 0.25, size=10.5, color=DARK, margin=0)

    headers = [
        ("ANALYSIS ARM", 0.66, 2.18),
        ("GENE RANKING", 3.07, 2.90),
        ("CLASSIFIER FIT", 6.43, 2.95),
        ("WHAT IT TESTS", 9.82, 2.72),
    ]
    for heading, x, width in headers:
        _add_text(slide, heading, x, 2.38, width, 0.22, size=9.0, color=GRAY, bold=True, margin=0)
    _add_rule(slide, 0.47, 2.64, 12.18, NAVY, 0.022)

    arms = [
        ("Real only", BLUE, ("R",), "Rank from real", ("R",), "Fit real", "Baseline"),
        ("Generated only", CORAL, ("S",), "Rank from generated", ("S",), "Fit generated", "Synthetic-only stress test"),
        ("Real + generated", TEAL, ("R", "S"), "Consensus rank", ("R", "S"), "Equal total weight", "Direct augmentation"),
        ("Guided: real fit", TEAL, ("R", "S"), "Consensus rank", ("R",), "Fit real only", "Feature guidance only"),
        ("Guided: 5% synthetic", ORANGE, ("R", "S"), "Consensus rank", ("R", "S"), "S = 5% weight", "Guidance + light regularization"),
    ]

    def add_sources(tokens: tuple[str, ...], x: float, y: float) -> float:
        if len(tokens) == 1:
            color = BLUE if tokens[0] == "R" else CORAL
            _add_data_badge(slide, tokens[0], x, y, color)
            return x + 0.34
        _add_data_badge(slide, "R", x, y, BLUE)
        _add_text(slide, "+", x + 0.27, y + 0.01, 0.18, 0.20, size=9.5, color=GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)
        _add_data_badge(slide, "S", x + 0.47, y, CORAL)
        return x + 0.82

    for index, (arm, color, rank_tokens, rank_label, fit_tokens, fit_label, purpose) in enumerate(arms):
        top = 2.73 + index * 0.57
        if index % 2 == 0:
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.47), Inches(top), Inches(12.18), Inches(0.50))
            _set_fill(shape, "F4F6F7")
            shape.line.fill.background()
        _add_rule(slide, 0.47, top, 0.055, color, 0.50)
        _add_text(slide, arm, 0.67, top + 0.08, 2.17, 0.33, size=11.2, color=color, bold=True, valign=MSO_ANCHOR.MIDDLE, margin=0)

        rank_end = add_sources(rank_tokens, 3.08, top + 0.13)
        _add_arrow(slide, rank_end, top + 0.17, 0.26, 0.14, MID_GRAY)
        _add_text(slide, rank_label, rank_end + 0.37, top + 0.08, 1.78, 0.33, size=10.4, color=DARK, bold=True, valign=MSO_ANCHOR.MIDDLE, margin=0)

        fit_end = add_sources(fit_tokens, 6.45, top + 0.13)
        _add_arrow(slide, fit_end, top + 0.17, 0.26, 0.14, MID_GRAY)
        _add_text(slide, fit_label, fit_end + 0.37, top + 0.08, 1.71, 0.33, size=10.4, color=DARK, bold=True, valign=MSO_ANCHOR.MIDDLE, margin=0)
        _add_text(slide, purpose, 9.83, top + 0.08, 2.70, 0.33, size=10.5, color=DARK, valign=MSO_ANCHOR.MIDDLE, margin=0)

    _add_rule(slide, 0.47, 5.64, 12.18, "CDD6DC", 0.018)
    _add_text(slide, "All five arms", 0.62, 5.83, 1.36, 0.27, size=12.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_arrow(slide, 2.13, 5.87, 0.34, 0.19, MID_GRAY)
    _add_text(slide, "Held-out real profiles", 2.66, 5.82, 2.00, 0.30, size=12.5, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_arrow(slide, 4.83, 5.87, 0.34, 0.19, MID_GRAY)
    _add_text(slide, "BA  |  AUROC  |  AP", 5.37, 5.83, 1.88, 0.27, size=12.2, color=DARK, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_arrow(slide, 7.41, 5.87, 0.34, 0.19, MID_GRAY)
    _add_text(slide, "Choose an eligible arm within each tissue", 7.92, 5.80, 4.40, 0.34, size=12.8, color=TEAL, bold=True, align=PP_ALIGN.CENTER, margin=0)

    _add_panel(slide, 0.46, 6.25, 12.22, 0.58, fill=NAVY, line=NAVY, radius=False)
    _add_text(slide, "Association test", 0.74, 6.40, 1.56, 0.25, size=13.5, color=WHITE, bold=True, margin=0)
    _add_text(slide, "FLT vs GC effects and BH FDR use real OSDR profiles only.", 2.47, 6.39, 6.40, 0.27, size=12.9, color="DCE7F2", margin=0)
    _add_text(slide, "Animal n stays unchanged.", 9.19, 6.39, 3.13, 0.27, size=11.5, color="FFD69A", bold=True, align=PP_ALIGN.RIGHT, margin=0)


def _slide_matched_classifier_design(slide):
    _add_slide_title(
        slide,
        "Primary analysis",
        "Change the training source; hold the classifier fixed",
        "All arms use the same 974 genes, data splits, preprocessing, and ridge regularization.",
    )

    _add_panel(slide, 0.47, 1.98, 12.18, 0.72, fill="F2F5F7", line="D3DDE2", radius=False)
    shared = [
        (0.73, "974 genes"),
        (3.27, "Same outer splits"),
        (6.14, "Real-fitted scaler"),
        (9.18, "Same ridge penalty"),
    ]
    for index, (x, label) in enumerate(shared):
        if index:
            _add_rule(slide, x - 0.28, 2.13, 0.012, "C9D3D9", 0.42)
        _add_text(
            slide,
            label,
            x,
            2.17,
            2.48,
            0.28,
            size=13.2,
            color=NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
            margin=0,
        )

    arms = [
        (0.48, "REAL ONLY", BLUE, ("R",), "Reference classifier"),
        (4.50, "SYNTHETIC ONLY", CORAL, ("S",), "Generator transfer test"),
        (8.52, "REAL + SYNTHETIC", TEAL, ("R", "S"), "Synthetic augmentation"),
    ]
    for x, heading, color, tokens, purpose in arms:
        _add_panel(slide, x, 2.93, 3.64, 1.38, fill=WHITE, line=color, radius=False)
        _add_rule(slide, x, 2.93, 3.64, color, 0.055)
        _add_text(
            slide,
            heading,
            x + 0.19,
            3.17,
            3.26,
            0.26,
            size=13.0,
            color=color,
            bold=True,
            align=PP_ALIGN.CENTER,
            margin=0,
        )
        if len(tokens) == 1:
            _add_data_badge(slide, tokens[0], x + 1.55, 3.57, color)
        else:
            _add_data_badge(slide, "R", x + 1.25, 3.57, BLUE)
            _add_text(slide, "+", x + 1.57, 3.58, 0.23, 0.22, size=11, color=GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)
            _add_data_badge(slide, "S", x + 1.84, 3.57, CORAL)
        _add_text(
            slide,
            purpose,
            x + 0.20,
            3.91,
            3.24,
            0.24,
            size=10.5,
            color=DARK,
            align=PP_ALIGN.CENTER,
            margin=0,
        )
        _add_arrow(slide, x + 1.56, 4.38, 0.42, 0.24, MID_GRAY)

    _add_panel(slide, 0.47, 4.76, 12.18, 0.86, fill=PALE_BLUE, line="C9DCE9", radius=False)
    _add_text(slide, "Same held-out real profiles", 0.80, 4.97, 3.30, 0.30, size=16.0, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_rule(slide, 4.27, 4.93, 0.012, "AFC0CB", 0.50)
    _add_text(slide, "Pooled metrics", 4.61, 4.92, 2.25, 0.25, size=12.6, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "BA  |  AUROC  |  AP", 4.61, 5.20, 2.25, 0.22, size=10.4, color=GRAY, align=PP_ALIGN.CENTER, margin=0)
    _add_rule(slide, 7.07, 4.93, 0.012, "AFC0CB", 0.50)
    _add_text(slide, "Accession-macro metrics", 7.39, 4.92, 2.65, 0.25, size=12.6, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "BA  |  AUROC  |  AP", 7.39, 5.20, 2.65, 0.22, size=10.4, color=GRAY, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "Pass only if all six means are nonworse", 10.13, 4.92, 2.22, 0.48, size=11.2, color=TEAL, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)

    _add_panel(slide, 0.47, 5.91, 12.18, 0.80, fill=NAVY, line=NAVY, radius=False)
    _add_text(slide, "Gene-level follow-up", 0.75, 6.12, 2.15, 0.28, size=14.2, color=WHITE, bold=True, margin=0)
    _add_text(slide, "Blocked permutation + linear SHAP on held-out real profiles", 3.03, 6.10, 5.42, 0.30, size=12.4, color="DCE7F2", margin=0)
    _add_text(slide, "BH FDR still uses real OSDR only", 8.71, 6.10, 3.58, 0.30, size=12.4, color="FFD69A", bold=True, align=PP_ALIGN.RIGHT, margin=0)
    _add_source(slide, "Primary matched analysis: 27 tissue units, eight outer splits, 648 fitted classifiers; generated profiles do not increase animal n.")


def _slide_guidance_mechanism(slide):
    _add_slide_title(
        slide,
        "Analysis",
        "Consensus ranking chooses the classifier input genes",
        "Real and generated rankings combine, then only the top k genes are used for FLT-GC classification.",
    )

    def add_rank_track(x: float, y: float, label: str, position: float, color: str):
        track_x = x + 0.90
        track_w = 1.92
        _add_text(slide, label, x, y - 0.055, 0.78, 0.23, size=10.4, color=DARK, bold=True, margin=0)
        _add_rule(slide, track_x, y + 0.035, track_w, "C9D2D8", 0.025)
        _add_circle(slide, track_x + track_w * position - 0.075, y - 0.025, 0.15, color)

    panel_y = 2.05
    panel_h = 2.60
    _add_panel(slide, 0.48, panel_y, 3.61, panel_h, fill=PALE_BLUE, line="C9DCE9", radius=False)
    _add_text(slide, "1  REAL RANKING", 0.72, 2.27, 1.86, 0.25, size=14.0, color=BLUE, bold=True, margin=0)
    _add_text(slide, "Rank all 974 genes using real OSDR", 0.72, 2.61, 3.02, 0.23, size=10.4, color=DARK, margin=0)
    _add_text(slide, "lower-ranked", 1.59, 2.96, 0.82, 0.18, size=8.5, color=GRAY, margin=0)
    _add_text(slide, "top-ranked", 3.02, 2.96, 0.72, 0.18, size=8.5, color=GRAY, align=PP_ALIGN.RIGHT, margin=0)
    example_genes = ["Gene 1", "Gene 2", "Gene 3", "Gene 4", "Gene 5"]
    real_positions = [0.90, 0.82, 0.72, 0.56, 0.35]
    for index, (gene, position) in enumerate(zip(example_genes, real_positions)):
        add_rank_track(0.72, 3.18 + index * 0.25, gene, position, BLUE)

    _add_arrow(slide, 4.16, 2.93, 0.27, 0.20, MID_GRAY)
    _add_panel(slide, 4.50, panel_y, 3.61, panel_h, fill=PALE_CORAL, line="E6CEC8", radius=False)
    _add_text(slide, "2  GENERATED RANKING", 4.74, 2.27, 2.48, 0.25, size=14.0, color=CORAL, bold=True, margin=0)
    _add_text(slide, "Rank the same 974 genes using generated profiles", 4.74, 2.61, 3.07, 0.23, size=10.4, color=DARK, margin=0)
    _add_text(slide, "lower-ranked", 5.61, 2.96, 0.82, 0.18, size=8.5, color=GRAY, margin=0)
    _add_text(slide, "top-ranked", 7.04, 2.96, 0.72, 0.18, size=8.5, color=GRAY, align=PP_ALIGN.RIGHT, margin=0)
    generated_positions = [0.88, 0.45, 0.86, 0.68, 0.32]
    for index, (gene, position) in enumerate(zip(example_genes, generated_positions)):
        add_rank_track(4.74, 3.18 + index * 0.25, gene, position, CORAL)
    _add_arrow(slide, 8.18, 2.93, 0.27, 0.20, MID_GRAY)
    _add_panel(slide, 8.52, panel_y, 4.13, panel_h, fill=PALE_TEAL, line="C9DFDB", radius=False)
    _add_text(slide, "3  CONSENSUS RANKING", 8.76, 2.27, 2.55, 0.25, size=14.0, color=TEAL, bold=True, margin=0)
    _add_text(slide, "Genes supported by both move upward", 8.76, 2.61, 2.82, 0.23, size=10.4, color=DARK, margin=0)
    _add_text(slide, "lower-ranked", 9.66, 2.96, 0.82, 0.18, size=8.5, color=GRAY, margin=0)
    _add_text(slide, "top-ranked", 11.61, 2.96, 0.72, 0.18, size=8.5, color=GRAY, align=PP_ALIGN.RIGHT, margin=0)
    shared_rows = [
        ("Gene 1", 0.89, TEAL, "selected"),
        ("Gene 2", 0.61, MID_GRAY, "not selected"),
        ("Gene 3", 0.79, TEAL, "selected"),
        ("Gene 4", 0.62, MID_GRAY, "not selected"),
        ("Gene 5", 0.33, MID_GRAY, "not selected"),
    ]
    for index, (gene, position, color, outcome) in enumerate(shared_rows):
        y = 3.18 + index * 0.25
        add_rank_track(8.76, y, gene, position, color)
        _add_text(slide, outcome, 11.72, y - 0.055, 0.80, 0.21, size=7.8, color=color, bold=True, align=PP_ALIGN.RIGHT, margin=0)
    shared_cutoff_x = 8.76 + 0.90 + 1.92 * 0.75
    _add_rule(slide, shared_cutoff_x, 3.10, 0.022, TEAL, 1.27)
    _add_text(slide, "selection cutoff", 11.15, 4.42, 1.20, 0.18, size=8.5, color=TEAL, italic=True, align=PP_ALIGN.RIGHT, margin=0)

    _add_text(slide, "Top-ranked genes become classifier inputs", 0.52, 4.84, 4.65, 0.28, size=16.0, color=NAVY, bold=True, margin=0)
    _add_panel(slide, 0.51, 5.18, 3.37, 1.28, fill="F4F8FA", line="D5E1E7", radius=False)
    _add_text(slide, "Rank 974 genes", 0.77, 5.45, 2.84, 0.25, size=13.0, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "real, generated, or consensus", 0.77, 5.88, 2.84, 0.22, size=10.0, color=GRAY, align=PP_ALIGN.CENTER, margin=0)
    _add_arrow(slide, 3.99, 5.68, 0.31, 0.20, MID_GRAY)

    _add_panel(slide, 4.41, 5.18, 3.37, 1.28, fill=PALE_TEAL, line="C9DFDB", radius=False)
    _add_text(slide, "Keep the top k", 4.67, 5.39, 2.84, 0.25, size=13.0, color=TEAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "k = 10, 25, 50, or 100", 4.67, 5.76, 2.84, 0.22, size=10.2, color=DARK, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "validation chooses k", 4.67, 6.06, 2.84, 0.18, size=8.8, color=GRAY, italic=True, align=PP_ALIGN.CENTER, margin=0)
    _add_arrow(slide, 7.89, 5.68, 0.31, 0.20, MID_GRAY)

    _add_panel(slide, 8.31, 5.18, 4.34, 1.28, fill="F6F7F8", line="D9DFE3", radius=False)
    _add_text(slide, "Fit FLT-GC classifier", 8.57, 5.45, 3.82, 0.25, size=13.0, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "using only the selected gene columns", 8.57, 5.88, 3.82, 0.22, size=10.0, color=DARK, align=PP_ALIGN.CENTER, margin=0)

    _add_source(slide, "Five genes are schematic examples from the 974-gene ranking. Real and generated normalized ranks are combined.")


def _slide_correlated_gene_dilution(slide):
    _add_slide_title(
        slide,
        "Interpretation",
        "Correlated genes dilute one-at-a-time importance",
        "A pathway can matter even when no single member causes a large permutation loss.",
    )

    panels = [
        (0.46, "1", "CORRELATED MODULE", BLUE, PALE_BLUE),
        (4.48, "2", "SHUFFLE ONE GENE", CORAL, PALE_CORAL),
        (8.50, "3", "CONSENSUS PANEL", TEAL, PALE_TEAL),
    ]
    for x, number, heading, color, fill in panels:
        _add_panel(slide, x, 2.03, 3.65, 3.86, fill=fill, line=color, radius=False)
        _add_text(slide, number, x + 0.20, 2.22, 0.42, 0.34, size=18, color=color, bold=True, margin=0)
        _add_text(slide, heading, x + 0.69, 2.24, 2.70, 0.29, size=13.0, color=NAVY, bold=True, margin=0)

    gene_y = [2.95, 3.65, 4.35]
    for index, y in enumerate(gene_y):
        color = [BLUE, TEAL, ORANGE][index]
        _add_panel(slide, 0.82, y, 1.05, 0.43, fill=WHITE, line=color, radius=False)
        _add_text(slide, f"Gene {chr(65 + index)}", 0.90, y + 0.09, 0.89, 0.22, size=11.2, color=color, bold=True, align=PP_ALIGN.CENTER, margin=0)
        _add_arrow(slide, 2.02, y + 0.10, 0.52, 0.23, color)
    _add_panel(slide, 2.70, 3.45, 1.00, 0.72, fill=WHITE, line=NAVY, radius=False)
    _add_text(slide, "FLT/GC\nclassifier", 2.79, 3.57, 0.82, 0.45, size=10.8, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "Similar expression lets these genes share predictive weight.", 0.78, 5.15, 3.08, 0.47, size=11.2, color=DARK, align=PP_ALIGN.CENTER, margin=0)

    for index, y in enumerate(gene_y):
        color = [BLUE, TEAL, ORANGE][index]
        _add_panel(slide, 4.84, y, 1.05, 0.43, fill=WHITE, line=color, radius=False)
        _add_text(slide, f"Gene {chr(65 + index)}", 4.92, y + 0.09, 0.89, 0.22, size=11.2, color=color, bold=True, align=PP_ALIGN.CENTER, margin=0)
        if index:
            _add_arrow(slide, 6.04, y + 0.10, 0.52, 0.23, color)
    _add_text(slide, "x", 5.20, 2.82, 0.32, 0.30, size=20, color=CORAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_panel(slide, 6.72, 3.45, 1.00, 0.72, fill=WHITE, line=NAVY, radius=False)
    _add_text(slide, "FLT/GC\nclassifier", 6.81, 3.57, 0.82, 0.45, size=10.8, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "Genes B and C retain the shared signal, so the score changes little.", 4.78, 5.15, 3.18, 0.47, size=11.2, color=DARK, align=PP_ALIGN.CENTER, margin=0)

    _add_text(slide, "Real rank", 8.86, 2.91, 0.91, 0.23, size=10.3, color=BLUE, bold=True, margin=0)
    _add_text(slide, "Generated rank", 10.30, 2.91, 1.21, 0.23, size=10.3, color=CORAL, bold=True, margin=0)
    consensus_rows = [
        ("Gene A", 0.90, 0.84),
        ("Gene B", 0.83, 0.88),
        ("Gene C", 0.78, 0.81),
        ("Gene D", 0.40, 0.32),
    ]
    for index, (gene, real_rank, generated_rank) in enumerate(consensus_rows):
        y = 3.33 + index * 0.48
        selected = index < 3
        color = TEAL if selected else MID_GRAY
        _add_text(slide, gene, 8.84, y, 0.76, 0.24, size=10.2, color=color, bold=selected, margin=0)
        _add_rule(slide, 9.66, y + 0.10, 0.90 * real_rank, BLUE, 0.075)
        _add_rule(slide, 10.96, y + 0.10, 0.90 * generated_rank, CORAL, 0.075)
        _add_text(slide, "selected" if selected else "not selected", 11.82, y, 0.62, 0.24, size=8.2, color=color, bold=selected, align=PP_ALIGN.RIGHT, margin=0)
    _add_text(slide, "Repeated support keeps the correlated genes together as one panel.", 8.82, 5.35, 3.44, 0.47, size=11.2, color=DARK, align=PP_ALIGN.CENTER, margin=0)

    _add_panel(slide, 0.47, 6.13, 12.18, 0.58, fill=NAVY, line=NAVY, radius=False)
    _add_text(
        slide,
        "Low marginal importance means a gene was replaceable in this classifier; it does not make the gene or pathway biologically irrelevant.",
        0.76,
        6.29,
        11.60,
        0.27,
        size=12.5,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    _add_source(slide, "Schematic explanation of permutation importance in a correlated expression feature space.")


def _slide_guidance_boundary(slide):
    _add_slide_title(
        slide,
        "Analysis",
        "Synthetic guidance can shift the FLT/GC boundary",
        "Opaque points train the classifier; transparent points are the same held-out real profiles in both panels.",
    )

    def add_point(x: float, y: float, color: str, *, held_out: bool):
        if not held_out:
            _add_circle(slide, x, y, 0.17, color)
            return
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x),
            Inches(y),
            Inches(0.20),
            Inches(0.20),
        )
        _set_fill(shape, color, transparency=58)
        _set_line(shape, color, 1.2)

    def add_scatter(x: float, y: float, *, guided: bool):
        plot_w = 4.75
        plot_h = 3.40
        _add_rule(slide, x, y + plot_h, plot_w, MID_GRAY, 0.022)
        _add_rule(slide, x, y, 0.022, MID_GRAY, plot_h + 0.02)
        train_flight = [(0.45, 0.52), (0.82, 2.25), (1.20, 1.35), (1.55, 2.95), (1.75, 0.70)]
        train_ground = [(2.85, 0.55), (3.15, 2.25), (3.50, 1.30), (3.90, 2.95), (4.25, 0.75)]
        held_flight = [(2.05, 1.55), (2.28, 2.55)]
        held_ground = [(2.58, 1.00), (2.80, 2.75)]
        for px, py in train_flight:
            add_point(x + px, y + py, CORAL, held_out=False)
        for px, py in train_ground:
            add_point(x + px, y + py, BLUE, held_out=False)
        for px, py in held_flight:
            add_point(x + px, y + py, CORAL, held_out=True)
        for px, py in held_ground:
            add_point(x + px, y + py, BLUE, held_out=True)
        if guided:
            _add_rule(slide, x + 2.42, y + 0.18, 0.045, TEAL, 2.98)
        else:
            _add_rule(slide, x + 1.85, y + 0.18, 0.035, GRAY, 2.98)
        _add_text(slide, "predicted FLT", x + 0.12, y + 3.52, 1.12, 0.20, size=9.2, color=CORAL, margin=0)
        _add_text(slide, "predicted GC", x + 3.48, y + 3.52, 1.12, 0.20, size=9.2, color=BLUE, align=PP_ALIGN.RIGHT, margin=0)

    _add_circle(slide, 8.79, 1.96, 0.13, CORAL)
    _add_text(slide, "FLT", 8.98, 1.92, 0.36, 0.21, size=9.2, color=GRAY, margin=0)
    _add_circle(slide, 9.50, 1.96, 0.13, BLUE)
    _add_text(slide, "GC", 9.69, 1.92, 0.34, 0.21, size=9.2, color=GRAY, margin=0)
    _add_circle(slide, 10.30, 1.96, 0.13, MID_GRAY)
    _add_text(slide, "training", 10.49, 1.92, 0.62, 0.21, size=9.2, color=GRAY, margin=0)
    held_legend = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.28), Inches(1.94), Inches(0.16), Inches(0.16))
    _set_fill(held_legend, MID_GRAY, transparency=58)
    _set_line(held_legend, MID_GRAY, 1.0)
    _add_text(slide, "held-out", 11.52, 1.92, 0.72, 0.21, size=9.2, color=GRAY, margin=0)

    _add_text(slide, "REAL-ONLY CLASSIFIER", 0.72, 2.23, 4.75, 0.27, size=13.0, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "SYNTHETIC-GUIDED CLASSIFIER", 7.85, 2.23, 4.75, 0.27, size=13.0, color=TEAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    add_scatter(0.72, 2.62, guided=False)
    add_scatter(7.85, 2.62, guided=True)
    _add_arrow(slide, 5.87, 3.88, 0.82, 0.44, MID_GRAY)
    _add_text(slide, "same held-out profiles", 5.57, 4.45, 1.42, 0.39, size=9.6, color=GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)

    _add_source(slide, "Schematic of the held-out evaluation used for each tissue-specific classifier candidate.")


def _slide_8(slide, utility_chart: Path):
    _add_slide_title(
        slide,
        "Utility",
        "Matched augmentation helped many tissues, but not all",
        "One pooled classifier declined; tissue-specific real-plus-synthetic models were compared with fixed all-gene baselines.",
    )
    _add_rule(slide, 0.48, 2.05, 3.30, GOLD, 0.035)
    _add_text(slide, "Pooled across tissues", 0.67, 2.25, 2.90, 0.35, size=17, color=NAVY, bold=True)
    labels = ["Real", "Generated", "Real + synth."]
    values = [0.754, 0.695, 0.737]
    colors = [BLUE, CORAL, TEAL]
    baseline = 6.00
    chart_height = 3.00
    chart_left = 0.96
    chart_right = 3.62
    axis_title = _add_text(
        slide,
        "Balanced accuracy",
        -0.28,
        4.35,
        1.55,
        0.28,
        size=10.5,
        color=GRAY,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    axis_title.rotation = 270
    for tick, label in [(0.0, "0"), (0.5, "0.5"), (1.0, "1.0")]:
        tick_y = baseline - tick * chart_height
        _add_rule(slide, chart_left, tick_y, chart_right - chart_left, "D8DEE2", 0.012)
        _add_text(
            slide,
            label,
            0.69,
            tick_y - 0.10,
            0.22,
            0.20,
            size=8.5,
            color=MID_GRAY,
            align=PP_ALIGN.RIGHT,
            margin=0,
        )
    for index, (label, value, color) in enumerate(zip(labels, values, colors)):
        x = 0.99 + index * 0.89
        bar_h = value * chart_height
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(baseline - bar_h), Inches(0.62), Inches(bar_h))
        _set_fill(shape, color)
        shape.line.fill.background()
        _add_text(slide, f"{value:.3f}", x - 0.08, baseline - bar_h - 0.35, 0.78, 0.28, size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, label, x - 0.13, baseline + 0.08, 0.88, 0.32, size=9.0, color=DARK, align=PP_ALIGN.CENTER, margin=0)
    _add_rule(slide, chart_left, baseline, chart_right - chart_left, "86949D", 0.02)

    _add_rule(slide, 3.98, 2.05, 0.015, "D5DDE2", 4.55)
    _add_rule(slide, 4.25, 2.05, 8.48, TEAL, 0.035)
    _add_text(slide, "Matched real + synthetic by tissue", 4.42, 2.25, 5.0, 0.34, size=17, color=NAVY, bold=True)
    _add_picture_contain(
        slide,
        utility_chart,
        4.32,
        2.65,
        8.32,
        3.70,
        alt="Balanced accuracy before and after matched real-plus-synthetic training",
    )
    _add_text(slide, "Teal passed all pooled and accession-macro metrics; coral did not.", 4.44, 6.25, 7.2, 0.27, size=11.7, color=GRAY, italic=True)
    _add_source(slide, "Matched analysis on held-out real profiles from represented accessions; all arms use the same 974 genes and ridge setting.")


def _slide_9(slide):
    _add_slide_title(
        slide,
        "Feature analysis",
        "Synthetic guidance changed ranking, not statistical evidence",
        "Blue marks real-only selection; teal marks synthetic-guided selection. Every displayed association has real OSDR support.",
    )
    _add_text(slide, "Stable selection after real-data support", 0.62, 2.02, 4.30, 0.30, size=15.5, color=NAVY, bold=True, margin=0)
    _add_circle(slide, 0.64, 2.42, 0.13, BLUE)
    _add_text(slide, "Blue: stable in real-only ranking", 0.86, 2.36, 2.42, 0.25, size=10.5, color=BLUE, bold=True, margin=0)
    _add_circle(slide, 3.43, 2.42, 0.13, TEAL)
    _add_text(slide, "Teal: stable in synthetic-guided ranking", 3.65, 2.36, 2.98, 0.25, size=10.5, color=TEAL, bold=True, margin=0)

    left = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.62), Inches(2.78), Inches(3.86), Inches(2.78))
    _set_fill(left, BLUE, transparency=82)
    _set_line(left, BLUE, 2.6)
    right = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.00), Inches(2.78), Inches(3.86), Inches(2.78))
    _set_fill(right, TEAL, transparency=80)
    _set_line(right, TEAL, 2.6)

    _add_text(slide, "34", 1.16, 3.42, 1.42, 0.48, size=29, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "real-only", 1.12, 3.93, 1.50, 0.30, size=14.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "WITHOUT GUIDANCE", 1.03, 4.31, 1.68, 0.24, size=9.2, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)

    _add_text(slide, "23", 3.16, 3.42, 1.17, 0.48, size=29, color=ORANGE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "reinforced", 3.10, 3.93, 1.29, 0.30, size=13.2, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "BOTH RANKINGS", 3.04, 4.31, 1.40, 0.24, size=9.2, color=ORANGE, bold=True, align=PP_ALIGN.CENTER, margin=0)

    _add_text(slide, "26", 4.89, 3.42, 1.50, 0.48, size=29, color=TEAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "promoted", 4.89, 3.93, 1.50, 0.30, size=14.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "GUIDANCE ONLY", 4.80, 4.31, 1.68, 0.24, size=9.2, color=TEAL, bold=True, align=PP_ALIGN.CENTER, margin=0)

    _add_rule(slide, 7.18, 2.10, 0.018, "D5DDE2", 3.98)
    _add_text(slide, "Real OSDR evidence gate", 7.53, 2.43, 4.92, 0.40, size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    evidence_rows = [
        (3.18, 0.38, "Estimate FLT vs GC inside each accession"),
        (3.82, 0.58, "Combine accession effects with a\nrandom-effects model"),
        (4.63, 0.38, "Apply BH FDR within each tissue"),
    ]
    for y, height, label in evidence_rows:
        _add_circle(slide, 7.79, y + 0.10, 0.10, BLUE)
        _add_text(
            slide,
            label,
            8.05,
            y,
            4.15,
            height,
            size=14.2,
            color=DARK,
            valign=MSO_ANCHOR.TOP,
            margin=0,
        )
    _add_panel(slide, 7.68, 5.14, 4.46, 0.78, fill=NAVY, line=NAVY, radius=False)
    _add_text(slide, "23 reinforced + 26 promoted", 7.90, 5.26, 4.02, 0.24, size=11.5, color="BFD0E1", bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "49 synthetic-informed associations", 7.88, 5.52, 4.06, 0.28, size=15.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER, margin=0)

    _add_rule(slide, 0.52, 6.26, 12.15, "D6DEE3", 0.018)
    _add_text(slide, "The 34 real-only associations remain real-data findings; synthetic guidance did not reinforce them.", 0.72, 6.39, 11.75, 0.26, size=12.6, color=GRAY, align=PP_ALIGN.CENTER)
    _add_source(slide, "Counts are tissue-gene associations after BH FDR and compatible real-effect direction; genes may recur across tissues.")


def _slide_matched_consensus_crosswalk(slide):
    _add_slide_title(
        slide,
        "Feature analysis",
        "Matched and consensus analyses answer different questions",
        "Matched importance tests synthetic contribution; consensus ranking organizes compact correlated panels.",
    )

    _add_text(slide, "BH-FDR tissue-gene associations", 0.63, 2.05, 5.75, 0.31, size=16.0, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    matched = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.72), Inches(2.62), Inches(3.78), Inches(2.70))
    _set_fill(matched, CORAL, transparency=78)
    _set_line(matched, CORAL, 2.4)
    consensus = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.93), Inches(2.62), Inches(3.78), Inches(2.70))
    _set_fill(consensus, BLUE, transparency=80)
    _set_line(consensus, BLUE, 2.4)
    _add_text(slide, "Matched all-gene", 0.82, 2.38, 2.15, 0.27, size=12.2, color=CORAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "Consensus panels", 4.42, 2.38, 2.15, 0.27, size=12.2, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "10", 1.25, 3.47, 1.20, 0.50, size=29, color=CORAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "matched only", 1.15, 4.03, 1.40, 0.28, size=11.0, color=DARK, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "11", 3.10, 3.47, 1.20, 0.50, size=29, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "both", 3.10, 4.03, 1.20, 0.28, size=11.0, color=DARK, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "38", 4.97, 3.47, 1.20, 0.50, size=29, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "consensus only", 4.81, 4.03, 1.52, 0.28, size=11.0, color=DARK, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "21 total", 1.57, 5.05, 1.55, 0.28, size=12.0, color=CORAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "49 total", 4.32, 5.05, 1.55, 0.28, size=12.0, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)

    _add_rule(slide, 7.10, 2.05, 0.015, "D5DDE2", 4.40)
    rows = [
        (2.17, "PRIMARY", "Matched all-gene", "Same genes and classifier; asks whether training source changes held-out-real utility or importance.", CORAL, PALE_CORAL),
        (3.50, "SECONDARY", "Consensus panel", "Ranks compact top-k panels; preserves groups of correlated genes that can replace one another.", BLUE, PALE_BLUE),
        (4.83, "SHARED", "Real association gate", "FLT-GC effects and BH FDR are estimated from observed OSDR profiles only.", TEAL, PALE_TEAL),
    ]
    for y, label, heading, detail, color, fill in rows:
        _add_panel(slide, 7.42, y, 5.10, 1.06, fill=fill, line=color, radius=False)
        _add_text(slide, label, 7.68, y + 0.18, 1.00, 0.23, size=9.2, color=color, bold=True, margin=0)
        _add_text(slide, heading, 8.83, y + 0.14, 3.34, 0.29, size=14.0, color=NAVY, bold=True, margin=0)
        _add_text(slide, detail, 7.68, y + 0.51, 4.49, 0.39, size=10.5, color=DARK, margin=0)

    _add_panel(slide, 0.58, 5.87, 12.02, 0.64, fill=NAVY, line=NAVY, radius=False)
    _add_text(slide, "Thymus is strongest because matched importance and the consensus cell-cycle panel converge.", 0.85, 6.05, 11.48, 0.28, size=13.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_source(slide, "Crosswalk: 21 matched associations, 49 consensus associations, and 11 supported by both analyses.")


def _slide_10(slide):
    _add_slide_title(
        slide,
        "Literature interpretation",
        "Consensus selection and literature are separate dimensions",
        "All 49 secondary consensus associations received a selection status and an independent literature interpretation.",
    )

    steps = [
        ("01", "Fix the candidate", "Gene, tissue, and FLT direction", BLUE),
        ("02", "Search", "Spaceflight and mechanism literature", TEAL),
        ("03", "Compare", "Observed result with published evidence", ORANGE),
        ("04", "Record", "Category, rationale, and source IDs", PURPLE),
    ]
    for index, (number, heading, detail, color) in enumerate(steps):
        x = 0.48 + index * 3.16
        if index:
            _add_rule(slide, x - 0.28, 2.02, 0.012, "D8DFE3", 0.95)
        _add_text(slide, number, x, 2.07, 0.40, 0.24, size=10.8, color=color, bold=True, margin=0)
        _add_text(slide, heading, x + 0.48, 2.02, 2.30, 0.30, size=14.5, color=NAVY, bold=True, margin=0)
        _add_text(slide, detail, x + 0.48, 2.40, 2.30, 0.42, size=10.8, color=GRAY, margin=0)
        _add_rule(slide, x, 2.95, 2.72, color, 0.025)

    annotations = pd.read_csv(
        PAPER_DIR / "source_data/table_s16_promoted_gene_literature_annotations.tsv",
        sep="\t",
    )
    counts = annotations["literature_classification"].value_counts().to_dict()
    selection_counts = annotations["selection_status"].value_counts().to_dict()
    direct = int(
        annotations["evidence_scope"]
        .eq("direct_same_gene_same_tissue_same_direction")
        .sum()
    )
    if len(annotations) != 49 or direct != 5:
        raise ValueError("Unexpected consensus literature inventory")

    _add_text(slide, "Selection status", 0.55, 3.35, 3.58, 0.34, size=17, color=NAVY, bold=True)
    selection_rows = [
        ("Promoted", selection_counts.get("promoted", 0), "Stable only with synthetic guidance", CORAL),
        ("Reinforced", selection_counts.get("reinforced", 0), "Stable with and without guidance", TEAL),
    ]
    for index, (label, count, detail, color) in enumerate(selection_rows):
        y = 4.02 + index * 1.05
        if index:
            _add_rule(slide, 0.55, y - 0.18, 3.58, "E0E5E8", 0.012)
        _add_text(slide, str(count), 0.55, y - 0.05, 0.68, 0.52, size=29, color=color, bold=True, margin=0)
        _add_text(slide, label, 1.38, y - 0.02, 2.65, 0.30, size=15, color=color, bold=True, margin=0)
        _add_text(slide, detail, 1.38, y + 0.31, 2.70, 0.38, size=11.5, color=DARK, margin=0)

    _add_rule(slide, 4.45, 3.33, 0.015, "D5DDE2", 2.77)
    _add_text(slide, "Literature interpretation", 4.82, 3.35, 7.52, 0.34, size=17, color=NAVY, bold=True)
    _add_text(slide, "Class", 4.84, 3.86, 1.55, 0.22, size=9.5, color=GRAY, bold=True, margin=0)
    _add_text(slide, "Count", 6.86, 3.86, 0.70, 0.22, size=9.5, color=GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "Meaning", 7.69, 3.86, 4.31, 0.22, size=9.5, color=GRAY, bold=True, margin=0)
    _add_rule(slide, 4.82, 4.13, 7.54, NAVY, 0.020)
    category_rows = [
        ("Aligning", counts.get("aligning", 0), "Direct or same-tissue process agreement", GREEN),
        ("Complementary", counts.get("complementary", 0), "Related process or mechanism", BLUE),
        ("Ambiguous", counts.get("ambiguous", 0), "Mixed or context-dependent evidence", ORANGE),
        ("Unmatched", counts.get("unmatched", 0), "No sufficiently specific match found", PURPLE),
    ]
    for index, (label, count, detail, color) in enumerate(category_rows):
        y = 4.27 + index * 0.48
        if index % 2 == 0:
            shade = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.82), Inches(y - 0.04), Inches(7.54), Inches(0.42))
            _set_fill(shade, "F5F7F8")
            shade.line.fill.background()
        _add_rule(slide, 4.86, y + 0.11, 0.12, color, 0.12)
        _add_text(slide, label, 5.10, y, 1.68, 0.34, size=12.6, color=color, bold=True, valign=MSO_ANCHOR.MIDDLE, margin=0)
        _add_text(slide, str(count), 6.87, y, 0.66, 0.34, size=14.0, color=DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
        _add_text(slide, detail, 7.69, y, 4.41, 0.34, size=11.8, color=DARK, valign=MSO_ANCHOR.MIDDLE, margin=0)

    _add_rule(slide, 0.55, 6.35, 12.02, "CED7DD", 0.018)
    _add_text(
        slide,
        "A reinforced gene can be complementary, ambiguous, or unmatched; an aligning gene can be promoted.",
        0.65,
        6.49,
        11.82,
        0.28,
        size=13.5,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_source(
        slide,
        "Table S16 records gene-level rationale, evidence scope and source IDs; Table S17 provides 33 citations with DOI/URL and data relationship.",
    )


def _inventory_gene_segments(
    rows: pd.DataFrame,
    annotation_lookup: dict[tuple[str, str, str], str],
) -> list[tuple[str, dict]]:
    class_colors = {
        "aligning": GREEN,
        "complementary": BLUE,
        "ambiguous": ORANGE,
        "unmatched": PURPLE,
    }
    ordered = rows.sort_values(["real_meta_fdr", "symbol"])
    segments: list[tuple[str, dict]] = []
    for index, row in enumerate(ordered.itertuples(index=False)):
        if index:
            segments.append((", ", {"size": 9.0, "color": MID_GRAY}))
        key = (row.analysis_scope, row.tissue, row.gene)
        literature_class = annotation_lookup.get(key)
        if literature_class not in class_colors:
            raise ValueError(f"Missing synthetic-informed annotation for {key}")
        color = class_colors[literature_class]
        segments.append(
            (
                row.symbol,
                {
                    "size": 9.0,
                    "color": color,
                    "bold": True,
                    "italic": True,
                },
            )
        )
    return segments


def _add_gene_inventory_block(
    slide,
    inventory: pd.DataFrame,
    annotation_lookup: dict[tuple[str, str, str], str],
    *,
    scope: str,
    tissue: str,
    label: str,
    x: float,
    y: float,
    w: float,
    h: float,
) -> None:
    rows = inventory[
        inventory["analysis_scope"].eq(scope) & inventory["tissue"].eq(tissue)
    ]
    if rows.empty:
        raise ValueError(f"No synthetic-informed genes for {scope}/{tissue}")

    _add_rule(slide, x, y, w, "DDE4E8", 0.018)
    count_label = "gene" if len(rows) == 1 else "genes"
    label_size = 10.8 if len(label.replace("\n", "")) > 12 else 11.8
    _add_rich_text(
        slide,
        [
            (f"{label}\n", {"size": label_size, "color": NAVY, "bold": True}),
            (f"{len(rows)} {count_label}", {"size": 9.0, "color": MID_GRAY}),
        ],
        x + 0.06,
        y + 0.07,
        1.36,
        h - 0.10,
        margin=0,
    )

    groups = []
    for direction, direction_label, direction_color in [
        ("FLT_higher", "FLT higher", CORAL),
        ("FLT_lower", "FLT lower", TEAL),
    ]:
        for selection, selection_label, selection_color in [
            ("synthetic_promoted", "Promoted", CORAL),
            ("reinforced_real_and_synthetic", "Reinforced", TEAL),
        ]:
            group_rows = rows[
                rows["flt_gc_direction"].eq(direction)
                & rows["selection_interpretation"].eq(selection)
            ]
            if group_rows.empty:
                continue
            symbol_characters = int(group_rows["symbol"].str.len().sum()) + 2 * max(0, len(group_rows) - 1)
            estimated_lines = max(1, int(np.ceil(symbol_characters / 38)))
            groups.append(
                (
                    direction_label,
                    direction_color,
                    selection_label,
                    selection_color,
                    group_rows,
                    1.0 + 0.75 * (estimated_lines - 1),
                )
            )

    available_height = h - 0.11
    weight_total = sum(group[-1] for group in groups)
    row_y = y + 0.06
    for index, (direction_label, direction_color, selection_label, selection_color, group_rows, weight) in enumerate(groups):
        row_h = available_height * weight / weight_total
        if index:
            _add_rule(slide, x + 1.44, row_y, w - 1.50, "E2E7EA", 0.010)
        _add_text(
            slide,
            direction_label,
            x + 1.47,
            row_y + 0.01,
            0.72,
            row_h - 0.02,
            size=8.0,
            color=direction_color,
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
            margin=0,
        )
        _add_text(
            slide,
            selection_label,
            x + 2.22,
            row_y + 0.01,
            0.78,
            row_h - 0.02,
            size=7.9,
            color=selection_color,
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
            margin=0,
        )
        _add_rich_text(
            slide,
            _inventory_gene_segments(group_rows, annotation_lookup),
            x + 3.04,
            row_y + 0.01,
            w - 3.12,
            row_h - 0.02,
            size=9.0,
            valign=MSO_ANCHOR.MIDDLE,
            margin=0,
        )
        row_y += row_h


def _slide_all_tissue_coverage(slide):
    _add_slide_title(
        slide,
        "Analysis coverage",
        "The screen covered all 27 completed tissue analyses",
        "The biological discussion narrows only after reporting synthetic-informed, real-only, and null outcomes.",
    )
    summary = pd.read_csv(
        PAPER_DIR / "source_data/table_s12_bh_fdr_tissue_summary.tsv",
        sep="\t",
    )
    if len(summary) != 27:
        raise ValueError("Expected 27 completed tissue analysis units")

    synthetic = summary.loc[
        summary["n_reinforced_real_and_synthetic"]
        .add(summary["n_synthetic_promoted"])
        .gt(0)
    ]
    real_only = summary.loc[
        summary["n_bh_fdr_genes"].gt(0)
        & ~summary.index.to_series().isin(synthetic.index)
    ]
    null = summary.loc[summary["n_bh_fdr_genes"].eq(0)]
    if (len(synthetic), len(real_only), len(null)) != (10, 5, 12):
        raise ValueError("Unexpected all-tissue coverage partition")

    display = {
        "skeletal_muscle": "Skeletal muscle (pooled)",
        "edl": "EDL",
    }

    def names(frame: pd.DataFrame) -> list[str]:
        return [
            display.get(value, value.replace("_", " ").title())
            for value in frame["tissue"]
        ]

    columns = [
        (
            0.43,
            "10",
            "Synthetic-informed",
            names(synthetic),
            TEAL,
            PALE_TEAL,
            "Promoted or reinforced BH-FDR gene",
        ),
        (
            4.49,
            "5",
            "Real BH-FDR only",
            names(real_only),
            BLUE,
            PALE_BLUE,
            "Real association, no synthetic-informed selection",
        ),
        (
            8.55,
            "12",
            "No BH-FDR gene",
            names(null),
            ORANGE,
            PALE_GOLD,
            "No BH-FDR gene in the 974-gene panel",
        ),
    ]
    for index, (x, count, heading, tissue_names, color, fill, detail) in enumerate(columns):
        if index:
            _add_rule(slide, x - 0.18, 2.05, 0.012, "D8DFE3", 4.58)
        _add_rule(slide, x, 2.04, 3.65, color, 0.035)
        _add_text(slide, count, x + 0.02, 2.22, 0.70, 0.55, size=31, color=color, bold=True, margin=0)
        _add_text(slide, heading, x + 0.84, 2.27, 2.72, 0.38, size=16.0, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE, margin=0)
        _add_text(slide, detail, x + 0.02, 2.92, 3.48, 0.42, size=10.5, color=GRAY, valign=MSO_ANCHOR.MIDDLE, margin=0)
        _add_rule(slide, x, 3.50, 3.58, "D7DEE2", 0.014)
        _add_text(slide, "ANALYSIS UNITS", x + 0.02, 3.66, 1.20, 0.19, size=8.8, color=GRAY, bold=True, margin=0)

        list_top = 3.96
        list_height = 2.42
        row_height = min(0.36, list_height / len(tissue_names))
        list_size = 9.8 if len(tissue_names) >= 12 else 10.4 if len(tissue_names) >= 10 else 11.2
        for index, tissue_name in enumerate(tissue_names):
            row_y = list_top + index * row_height
            _add_rule(slide, x + 0.03, row_y + (row_height - 0.045) / 2, 0.13, color, 0.045)
            _add_text(
                slide,
                tissue_name,
                x + 0.25,
                row_y,
                3.25,
                row_height,
                size=list_size,
                color=DARK,
                valign=MSO_ANCHOR.MIDDLE,
                margin=0,
            )
    _add_source(slide, "Twenty-two canonical tissues plus EDL, gastrocnemius, quadriceps, soleus and tibialis anterior; complete counts in Table S12.")


def _slide_matched_coverage(slide):
    _add_slide_title(
        slide,
        "Primary analysis coverage",
        "The matched screen covered all 27 tissue analyses",
        "Predictive utility and gene-level support are separate outcomes.",
    )
    utility = pd.read_csv(
        PAPER_DIR / "source_data/table_s18_matched_all_gene_utility.tsv",
        sep="\t",
    )
    candidates = pd.read_csv(
        PAPER_DIR / "source_data/table_s19_matched_all_gene_candidates.tsv",
        sep="\t",
    )
    utility = utility.loc[utility["arm"].eq("real_plus_generated")].copy()
    candidate_tissues = set(candidates["tissue"].astype(str))
    with_genes = utility.loc[utility["tissue"].isin(candidate_tissues)]
    utility_only = utility.loc[
        utility["joint_mean_all_metrics_nonworse"].astype(bool)
        & ~utility["tissue"].isin(candidate_tissues)
    ]
    did_not_pass = utility.loc[
        ~utility["joint_mean_all_metrics_nonworse"].astype(bool)
    ]
    if (len(with_genes), len(utility_only), len(did_not_pass)) != (4, 14, 9):
        raise ValueError("Unexpected matched all-gene coverage partition")

    display = {
        "skeletal_muscle": "Skeletal muscle (pooled)",
        "edl": "EDL",
        "brown_adipose_tissue": "Brown adipose tissue",
        "white_adipose_tissue": "White adipose tissue",
    }

    def names(frame: pd.DataFrame) -> list[str]:
        return [
            display.get(value, value.replace("_", " ").title())
            for value in frame["tissue"].astype(str)
        ]

    columns = [
        (0.43, "4", "Matched genes", names(with_genes), CORAL, PALE_CORAL, "Utility passed + BH-FDR gene importance"),
        (4.49, "14", "Utility only", names(utility_only), TEAL, PALE_TEAL, "All six metrics passed; no retained gene"),
        (8.55, "9", "Gate not passed", names(did_not_pass), BLUE, PALE_BLUE, "At least one mean metric was worse"),
    ]
    for column_index, (x, count, heading, tissue_names, color, _fill, detail) in enumerate(columns):
        if column_index:
            _add_rule(slide, x - 0.18, 2.05, 0.012, "D8DFE3", 4.58)
        _add_rule(slide, x, 2.04, 3.65, color, 0.035)
        _add_text(slide, count, x + 0.02, 2.22, 0.70, 0.55, size=31, color=color, bold=True, margin=0)
        _add_text(slide, heading, x + 0.84, 2.27, 2.72, 0.38, size=16.0, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE, margin=0)
        _add_text(slide, detail, x + 0.02, 2.92, 3.48, 0.42, size=10.5, color=GRAY, valign=MSO_ANCHOR.MIDDLE, margin=0)
        _add_rule(slide, x, 3.50, 3.58, "D7DEE2", 0.014)
        _add_text(slide, "ANALYSIS UNITS", x + 0.02, 3.66, 1.20, 0.19, size=8.8, color=GRAY, bold=True, margin=0)
        list_top = 3.96
        list_height = 2.42
        row_height = min(0.36, list_height / len(tissue_names))
        list_size = 9.6 if len(tissue_names) >= 12 else 10.3
        for row_index, tissue_name in enumerate(tissue_names):
            row_y = list_top + row_index * row_height
            _add_rule(slide, x + 0.03, row_y + (row_height - 0.045) / 2, 0.13, color, 0.045)
            _add_text(slide, tissue_name, x + 0.25, row_y, 3.25, row_height, size=list_size, color=DARK, valign=MSO_ANCHOR.MIDDLE, margin=0)
    _add_source(slide, "Real-plus-synthetic matched arm across 22 canonical tissues and five anatomical muscle groups.")


def _slide_11(slide):
    _add_slide_title(
        slide,
        "Secondary consensus inventory",
        "Ten tissue analyses contained consensus-selected genes",
        "All 49 associations passed BH FDR in real OSDR data; consensus ranking affected prioritization, not the statistical test.",
    )
    inventory = pd.read_csv(
        PAPER_DIR / "source_data/table_s10_synthetic_informed_bh_fdr_genes.tsv",
        sep="\t",
    )
    annotations = pd.read_csv(
        PAPER_DIR / "source_data/table_s16_promoted_gene_literature_annotations.tsv",
        sep="\t",
    )
    if len(inventory) != 49 or len(annotations) != 49:
        raise ValueError("Unexpected consensus gene inventory")
    if inventory[["analysis_scope", "tissue"]].drop_duplicates().shape[0] != 10:
        raise ValueError("Expected 10 tissue analyses in consensus inventory")
    annotation_lookup = {
        (row.analysis_scope, row.tissue, row.gene): row.literature_classification
        for row in annotations.itertuples(index=False)
    }

    _add_text(slide, "Rows: FLT direction + status", 0.45, 1.86, 2.42, 0.25, size=9.8, color=NAVY, bold=True)
    legend = [
        ("aligning", GREEN, 3.02, 0.84),
        ("complementary", BLUE, 4.02, 1.15),
        ("ambiguous", ORANGE, 5.37, 0.96),
        ("unmatched", PURPLE, 6.53, 0.96),
    ]
    for label, color, x, width in legend:
        _add_rule(slide, x, 1.91, 0.14, color, 0.05)
        _add_text(slide, label, x + 0.19, 1.85, width, 0.25, size=9.8, color=color, bold=True)
    _add_text(slide, "Gene color = literature | Non-empty rows shown", 8.64, 1.85, 4.20, 0.25, size=9.8, color=MID_GRAY, italic=True, align=PP_ALIGN.RIGHT)

    left_blocks = [
        ("canonical_tissue", "thymus", "Thymus", 2.24, 1.45),
        ("canonical_tissue", "spleen", "Spleen", 3.78, 0.67),
        ("canonical_tissue", "kidney", "Kidney", 4.54, 0.61),
        ("skeletal_muscle_group", "gastrocnemius", "Gastrocnemius", 5.24, 0.61),
        ("canonical_tissue", "eye", "Eye", 5.94, 0.51),
    ]
    right_blocks = [
        ("canonical_tissue", "skeletal_muscle", "Skeletal muscle\n(pooled)", 2.24, 1.25),
        ("skeletal_muscle_group", "soleus", "Soleus", 3.58, 0.68),
        ("skeletal_muscle_group", "tibialis_anterior", "Tibialis anterior", 4.35, 0.69),
        ("canonical_tissue", "adrenal_gland", "Adrenal gland", 5.13, 0.61),
        ("canonical_tissue", "skin", "Skin", 5.83, 0.55),
    ]
    for scope, tissue, label, y, h in left_blocks:
        _add_gene_inventory_block(
            slide,
            inventory,
            annotation_lookup,
            scope=scope,
            tissue=tissue,
            label=label,
            x=0.43,
            y=y,
            w=6.15,
            h=h,
        )
    for scope, tissue, label, y, h in right_blocks:
        _add_gene_inventory_block(
            slide,
            inventory,
            annotation_lookup,
            scope=scope,
            tissue=tissue,
            label=label,
            x=6.76,
            y=y,
            w=6.15,
            h=h,
        )
    _add_source(slide, "Secondary consensus set: rows report FLT direction and selection status; gene color independently reports literature interpretation.")


def _slide_12(slide):
    _add_slide_title(
        slide,
        "Thymus",
        "Thymus is strongest across both analyses",
        "Matched importance supports synthetic contribution; consensus ranking expands the correlated cell-cycle panel.",
    )
    figure = PAPER_DIR / "figures/figure_3_thymus_biology.png"
    _add_picture_contain(slide, figure, 0.35, 1.93, 8.55, 4.78, alt="Thymus gene effects and Reactome processes")
    _add_panel(slide, 9.10, 2.02, 3.82, 4.57, fill=PALE_CORAL, line="E8C9C2", radius=False)
    _add_text(slide, "MATCHED ALL-GENE", 9.43, 2.27, 2.95, 0.23, size=9.5, color=CORAL, bold=True, margin=0)
    _add_text(slide, "15", 9.43, 2.56, 1.10, 0.65, size=35, color=CORAL, bold=True)
    _add_text(slide, "BH-FDR genes with\nsynthetic-supported importance", 10.38, 2.58, 2.12, 0.62, size=12.7, color=DARK, bold=True)
    _add_text(slide, "7 promoted | 8 shared", 9.44, 3.34, 3.05, 0.30, size=14, color=TEAL, bold=True)
    _add_bullet_rows(
        slide,
        [
            "26 significant Reactome terms; mitotic cell cycle leads",
            "Nine genes overlap the 16-gene consensus panel",
            "Consensus adds correlated Cdk1, Top2a, Aurka, Ccne2, Pcna and Ccnf",
            "Lower mitotic renewal or fewer cycling thymocytes in flight",
        ],
        9.44,
        3.77,
        3.03,
        size=11.8,
        bullet_color=CORAL,
        row_h=0.61,
    )
    _add_source(slide, "Figure shows the secondary consensus panel; matched genes and Reactome crosswalk are in Tables S19-S20. Context: Gridley et al. (2013); Horie et al. (2019).")


def _slide_13(slide):
    _add_slide_title(
        slide,
        "Soleus",
        "Soleus is a coherent secondary consensus result",
        "The five-gene program is biologically interpretable, but the matched all-gene gate did not pass.",
    )
    figure = PAPER_DIR / "figures/figure_4_soleus_biology.png"
    _add_picture_contain(slide, figure, 0.34, 1.96, 9.15, 4.70, alt="Soleus gene effects and Reactome processes")
    _add_panel(slide, 9.63, 2.04, 3.25, 4.52, fill=PALE_TEAL, line="C7DDD8", radius=False)
    _add_text(slide, "5-gene consensus panel", 9.87, 2.30, 2.77, 0.62, size=17, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "reinforced in real and consensus rankings", 9.90, 2.96, 2.70, 0.36, size=10.8, color=GRAY, align=PP_ALIGN.CENTER)
    _add_panel(slide, 9.92, 3.39, 2.65, 0.74, fill=WHITE, line=CORAL, radius=False)
    _add_text(slide, "Matched gate not passed", 10.10, 3.50, 2.30, 0.25, size=12.2, color=CORAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "BA +0.013 | AUROC -0.015 | AP -0.009", 10.03, 3.79, 2.44, 0.20, size=8.5, color=DARK, align=PP_ALIGN.CENTER, margin=0)
    _add_bullet_rows(
        slide,
        [
            "Lower Bdh1, Ech1, Bnip3 and Decr1",
            "Higher Tpm1",
            "Mitochondrial turnover and fatty-acid metabolism",
            "Panel-level hypothesis, not matched synthetic contribution",
        ],
        9.94,
        4.31,
        2.54,
        size=11.7,
        bullet_color=TEAL,
        row_h=0.52,
    )
    _add_source(slide, "Literature context: Gambara et al. (2017) and Stein et al. (2002).")


def _add_additional_finding_rows(
    slide,
    rows,
    *,
    first_label: str = "Promoted",
    second_label: str = "Reinforced",
    selection_heading: str = "Selection",
):
    _add_text(slide, "Analysis unit", 0.72, 2.00, 1.70, 0.28, size=11.5, color=GRAY, bold=True)
    _add_text(slide, selection_heading, 2.48, 2.00, 1.15, 0.28, size=11.5, color=GRAY, bold=True)
    _add_text(slide, "Genes and FLT direction", 3.78, 2.00, 3.15, 0.28, size=11.5, color=GRAY, bold=True)
    _add_text(slide, "Interpretation", 7.20, 2.00, 5.10, 0.28, size=11.5, color=GRAY, bold=True)
    for index, (tissue, promoted, reinforced, interpretation, color, fill) in enumerate(rows):
        y = 2.35 + index * 1.06
        _add_panel(slide, 0.43, y, 12.44, 0.88, fill=fill, line="DDE4E8", radius=False)
        _add_rule(slide, 0.43, y, 0.08, color, 0.88)
        _add_text(slide, tissue, 0.72, y + 0.20, 1.62, 0.44, size=15.0, color=color, bold=True, valign=MSO_ANCHOR.MIDDLE)
        _add_rule(slide, 2.42, y + 0.44, 4.42, "DDE4E8", 0.012)
        _add_text(slide, first_label, 2.48, y + 0.08, 1.12, 0.25, size=10.5, color=CORAL, bold=True, valign=MSO_ANCHOR.MIDDLE)
        _add_text(slide, promoted, 3.78, y + 0.05, 3.08, 0.35, size=10.5, color=DARK, valign=MSO_ANCHOR.MIDDLE)
        _add_text(slide, second_label, 2.48, y + 0.52, 1.12, 0.25, size=10.5, color=TEAL, bold=True, valign=MSO_ANCHOR.MIDDLE)
        _add_text(slide, reinforced, 3.78, y + 0.48, 3.08, 0.35, size=10.5, color=DARK, valign=MSO_ANCHOR.MIDDLE)
        _add_text(slide, interpretation, 7.20, y + 0.10, 5.20, 0.66, size=12.2, color=DARK, valign=MSO_ANCHOR.MIDDLE)


def _slide_14(slide):
    _add_slide_title(
        slide,
        "Additional matched findings",
        "Liver, skin and spleen add narrower gene-level results",
        "Pooled muscle improves prediction, while its gene panel remains secondary consensus evidence.",
    )
    rows = [
        (
            "Liver",
            "Lower: Grb10, Ppic, H2-DMa, Gtf2a2",
            "No retained panel",
            "Four shared-importance candidates; no significant Reactome set enrichment.",
            BLUE,
            PALE_BLUE,
        ),
        (
            "Skin",
            "Higher: Plscr1",
            "Higher: Plscr1 plus broader selected pathways",
            "Matched and consensus support the same gene; broader cell-cycle and DNA-repair context is panel-level.",
            PURPLE,
            "F0ECF6",
        ),
        (
            "Spleen",
            "Higher: Loxl1",
            "Higher: Loxl1, Rai14, Ptprk, Myl9",
            "Loxl1 is the matched anchor; the four-gene adhesion and cytoskeletal panel is tentative.",
            TEAL,
            PALE_TEAL,
        ),
        (
            "Pooled muscle",
            "All six metrics passed; no gene passed the full gate",
            "12-gene interferon and sialic-acid panel",
            "Predictive utility is matched; individual gene interpretation comes from the secondary consensus analysis.",
            ORANGE,
            PALE_GOLD,
        ),
    ]
    _add_additional_finding_rows(
        slide,
        rows,
        first_label="Matched",
        second_label="Consensus",
        selection_heading="Evidence",
    )
    _add_source(slide, "Matched genes pass real-data BH FDR and synthetic-supported marginal importance; consensus rows provide secondary panel context.")


def _slide_additional_2(slide):
    _add_slide_title(
        slide,
        "Additional consensus findings",
        "Several narrower panels did not pass the matched gene gate",
        "These remain tissue-specific hypotheses from the secondary ranking analysis.",
    )
    rows = [
        (
            "Kidney",
            "Inpp4b higher",
            "Slc37a4 higher",
            "Renal phosphoinositide signaling and glucose-handling hypothesis; the pair had no shared Reactome enrichment.",
            GREEN,
            "ECF4ED",
        ),
        (
            "Adrenal gland",
            "Psmb8 lower",
            "Tspan4 lower",
            "Literature-unmatched immune, proteostasis or tissue-composition candidates; no adrenal mechanism is established.",
            MID_GRAY,
            "F2F4F5",
        ),
        (
            "Gastrocnemius",
            "Nfkbia higher; Fhl2 lower",
            "None",
            "NF-kappaB stress alignment plus an autophagy or myogenesis candidate; the two genes do not form a coherent pathway.",
            BLUE,
            PALE_BLUE,
        ),
        (
            "Tibialis anterior",
            "Cebpd higher",
            "Cdkn1a, St3gal5, Bnip3 higher",
            "Stress, cell-cycle, ganglioside-signaling and mitophagy candidates with mixed or complementary prior evidence.",
            TEAL,
            PALE_TEAL,
        ),
    ]
    _add_additional_finding_rows(slide, rows)
    _add_source(slide, "Secondary consensus results and independent literature annotations are reported in Tables S10 and S16.")


def _slide_15(slide):
    _add_slide_title(
        slide,
        "Takeaways",
        "Use matched tests for contribution; consensus for biology",
        "Matched classifiers test utility; consensus ranking explains correlated programs; biological claims still use OSDR profiles.",
    )
    columns = [
        (0.42, "1", "Generate", "Conditional DDIM produced high-fidelity profiles with near-chance real-versus-synthetic separation.", BLUE, PALE_BLUE),
        (4.44, "2", "Test", "Matched real-plus-synthetic classifiers passed all six metric checks in 18 of 27 tissue analyses and retained 21 BH-FDR associations.", TEAL, PALE_TEAL),
        (8.46, "3", "Interpret", "Consensus ranking recovered broader correlated panels. Thymus was strongest across both analyses; soleus remained secondary.", ORANGE, PALE_GOLD),
    ]
    for index, (x, number, heading, body, color, _fill) in enumerate(columns):
        if index:
            _add_rule(slide, x - 0.20, 2.10, 0.012, "D8DFE3", 3.53)
        _add_rule(slide, x, 2.10, 3.55, color, 0.035)
        _add_text(slide, f"0{number}", x + 0.02, 2.33, 0.68, 0.40, size=21, color=color, bold=True, margin=0)
        _add_text(slide, heading, x + 0.86, 2.35, 2.46, 0.36, size=19, color=NAVY, bold=True, margin=0)
        _add_text(slide, body, x + 0.03, 3.12, 3.28, 2.10, size=16, color=DARK, margin=0)


def _slide_16(slide):
    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.38), Inches(SLIDE_W), Inches(SLIDE_H - 0.38))
    _set_fill(body, NAVY)
    body.line.fill.background()
    _add_text(slide, "Thank you", 0.65, 1.22, 5.40, 0.80, size=34, color=WHITE, bold=True)
    _add_text(slide, "Questions?", 0.67, 2.10, 4.0, 0.45, size=21, color="BFD0E1")
    _add_rule(slide, 0.68, 2.84, 3.0, ORANGE, 0.05)
    _add_text(slide, "Jason Trinh", 0.68, 3.21, 3.0, 0.35, size=18, color=WHITE, bold=True)
    _add_text(slide, "jasontrinh@berkeley.edu", 0.68, 3.63, 3.5, 0.30, size=14, color="BFD0E1")
    _add_text(slide, "Mentor", 6.02, 1.48, 1.30, 0.27, size=11, color="7FA0BC", bold=True)
    _add_text(slide, "James Casaletto", 7.43, 1.45, 4.55, 0.34, size=17, color=WHITE, bold=True)
    _add_text(slide, "Program", 6.02, 2.16, 1.30, 0.27, size=11, color="7FA0BC", bold=True)
    _add_text(slide, "NASA Space Life Sciences Training Program", 7.43, 2.13, 4.82, 0.52, size=15, color=WHITE)
    _add_text(slide, "Data", 6.02, 3.04, 1.30, 0.27, size=11, color="7FA0BC", bold=True)
    _add_text(slide, "NASA OSDR | ARCHS4 | Reactome", 7.43, 3.01, 4.55, 0.34, size=15, color=WHITE)
    _add_text(slide, "Compute", 6.02, 3.73, 1.30, 0.27, size=11, color="7FA0BC", bold=True)
    _add_text(slide, "NASA Ames Research Center", 7.43, 3.70, 4.55, 0.34, size=15, color=WHITE)
    _add_text(slide, "Code and manuscript", 6.02, 4.42, 1.30, 0.40, size=11, color="7FA0BC", bold=True)
    _add_text(slide, "github.com/jasont314/nasa-mouse", 7.43, 4.40, 4.55, 0.34, size=15, color=WHITE)
    _add_text(slide, "All biological associations were tested in observed OSDR profiles.", 6.02, 5.50, 6.00, 0.40, size=13.5, color="BFD0E1", italic=True)


def _add_notes(slide, note: SlideNote) -> None:
    frame = slide.notes_slide.notes_text_frame
    frame.text = f"Target time: {note.time}\n\n{note.text}"


def _write_notes(notes: list[SlideNote]) -> None:
    lines = [
        "# SLSTP 2026 mouse spaceflight transcriptomics speaker notes",
        "",
        "Target length: 12-15 minutes. Planned speaking time: about 13 minutes.",
        "",
    ]
    for note in notes:
        lines.extend([
            f"## {note.number}. {note.title} ({note.time})",
            "",
            note.text,
            "",
        ])
    NOTES_PATH.write_text("\n".join(lines), encoding="utf-8")


def build() -> Path:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    utility_chart = _build_tissue_utility_chart()
    trajectory = _build_trajectory_crop()
    tissue_accession_pca = _build_pca_comparison_chart("tissue")
    condition_accession_pca = _build_pca_comparison_chart("condition")
    architecture_figure = ASSET_DIR / "lacan_figure1c_generator_architecture.png"
    diffusion_gif = ASSET_DIR / "anim_diffdata.gif"
    wgan_gif = ASSET_DIR / "anim_wgan.gif"
    if not architecture_figure.exists():
        raise FileNotFoundError(
            "Missing Lacan et al. architecture excerpt; see presentation/generative_slstp_2026/assets/SOURCES.md"
        )
    if not diffusion_gif.exists():
        raise FileNotFoundError(
            "Missing midpoint reverse-diffusion animation; see presentation/generative_slstp_2026/assets/SOURCES.md"
        )
    if not wgan_gif.exists():
        raise FileNotFoundError(
            "Missing midpoint WGAN animation; see presentation/generative_slstp_2026/assets/SOURCES.md"
        )
    presentation = Presentation(TEMPLATE)
    _set_title_slide(presentation.slides[0])
    _prepare_content_slide(presentation.slides[1], 2)
    builders = [
        None,
        _slide_project_scope,
        _slide_autoencoder_foundation,
        _slide_expimap_program_scores,
        _slide_expimap_tissue_results,
        _slide_why_synthetic,
        _slide_2,
        lambda slide: _slide_wgan_explainer(slide, wgan_gif),
        _slide_3,
        lambda slide: _slide_diffusion_explainer(slide, diffusion_gif),
        lambda slide: _slide_4(slide, architecture_figure),
        lambda slide: _slide_5(slide, trajectory),
        lambda slide: _slide_6(slide, tissue_accession_pca),
        lambda slide: _slide_condition_accession_pca(slide, condition_accession_pca),
        _slide_matched_classifier_design,
        _slide_guidance_mechanism,
        _slide_correlated_gene_dilution,
        lambda slide: _slide_8(slide, utility_chart),
        _slide_matched_consensus_crosswalk,
        _slide_10,
        _slide_matched_coverage,
        _slide_11,
        _slide_12,
        _slide_13,
        _slide_14,
        _slide_additional_2,
        _slide_15,
        _slide_16,
    ]
    while len(presentation.slides) < len(builders):
        number = len(presentation.slides) + 1
        slide = presentation.slides.add_slide(presentation.slide_layouts[3])
        _prepare_content_slide(slide, number)

    for index, builder in enumerate(builders):
        if builder is not None:
            builder(presentation.slides[index])

    notes = [
        SlideNote(1, "Interpretable and generative models for mouse spaceflight", "0:15", "This project uses machine learning to study mouse bulk RNA-seq from NASA spaceflight experiments. The first model asks which pathways change. The second asks whether realistic synthetic profiles can improve a tissue-specific FLT versus GC analysis."),
        SlideNote(2, "One dataset, two machine-learning questions", "0:35", "Both parts begin with the same OSDR flight and ground-control data. expiMap organizes expression into interpretable Reactome programs. The generative pipeline models the expression distribution and then tests whether generated profiles improve classification and gene ranking."),
        SlideNote(3, "Autoencoders compress expression into a latent space", "0:40", "An autoencoder takes a high-dimensional gene-expression profile, compresses it into a small set of latent variables, and reconstructs the original profile. Nearby points in latent space have similar expression. In a standard autoencoder those axes may have no clear biological meaning. expiMap constrains them with known gene programs."),
        SlideNote(4, "Program scores become pathway-level changes", "0:45", "These are observed OSDR examples from the expiMap analysis. Each sample receives a latent score for each Reactome program. The table shows project-centered mean scores for flight and ground control, followed by flight minus ground. The formal analysis estimates the change inside each project and then combines project effects. All four examples are lower in flight; a lower latent score does not by itself prove biochemical inhibition."),
        SlideNote(5, "Four tissues produced the clearest pathway patterns", "0:50", "Thymus showed lower repair, cytoskeletal, and stromal-interaction programs. Skin showed lower chromatin regulation, repair, Hedgehog, sphingolipid, and cell-junction programs. Liver showed lower MHC class II and T-cell receptor scores. Spleen combined lower T-cell receptor, neutrophil-degranulation, and C-type lectin programs."),
        SlideNote(6, "What is synthetic transcriptomics?", "0:25", "A generator learns the distribution of measured expression and samples new numeric profiles for a chosen tissue and FLT or GC context. We use those profiles to test classifiers and rank genes. They are model output, not new animals or independent biological measurements."),
        SlideNote(7, "Small studies and study effects complicate tissue comparisons", "0:30", "OSDR covers many tissues, but its 1,610 profiles are spread across 75 accessions. ARCHS4 supplies a much larger mouse reference. The model must preserve tissue and condition structure without simply learning study identity."),
        SlideNote(8, "A generator learns by competing with a critic", "0:30", "A WGAN-GP alternates between two updates. The generator turns noise and biological conditions into an expression profile. The critic compares measured and generated profiles. Their competition teaches the generator to match the observed distribution while retaining the requested tissue, flight or ground-control, accession, and material context."),
        SlideNote(8, "We built a configurable bulk RNA-seq generation pipeline", "0:45", "The pipeline can change data scope, transformation, harmonization, model, training source, and conditioning. The branch used here applies TPM, MaxAbs scaling, 974 landmarks, ARCHS4 pretraining, OSDR adaptation, and conditioning on tissue, FLT or GC, accession, and material type."),
        SlideNote(9, "Diffusion turns noise into a conditioned sample", "0:25", "This teaching animation starts from random points. During reverse diffusion, the model repeatedly predicts and removes noise while receiving the requested condition. In the RNA-seq model, the output is a vector of 974 gene values conditioned on tissue, flight or ground control, accession, and material type."),
        SlideNote(10, "DDIM matched expression and reduced separability", "0:40", "We compared WGAN-GP with the conditional diffusion model. Both matched expression well. DDIM had higher F1, adversarial accuracy close to chance, and lower distributional distance, so the remaining analysis uses DDIM."),
        SlideNote(11, "Diffusion learns tissue structure from noise", "0:25", "The same generated profiles begin as noise, develop structure by timestep 200, and reach their tissue-conditioned regions at timestep zero. All three panels share PCA axes."),
        SlideNote(12, "Tissue and study structure dominate the PCA space", "0:25", "These panels use the same locked coordinates. The left panel colors each profile by tissue, while the right colors it by OSDR accession. Circles are observed profiles and crosses are matched DDIM profiles. Both tissue and study structure are reproduced in the generated data."),
        SlideNote(13, "Flight condition is subtler than study structure", "0:25", "The left panel now colors the same profiles by flight or ground-control condition, while the right repeats the accession view. FLT and GC overlap much more than the study clusters. This is why the downstream analysis estimates FLT-GC effects within accession rather than treating the pooled separation as biology."),
        SlideNote(14, "The primary comparison changes only the training source", "0:40", "Every matched classifier uses all 974 genes, the same real-fitted scaler, the same outer split, and one regularization value selected from real training data. Real-only, synthetic-only, and real-plus-synthetic models are evaluated on the same held-out real profiles. This isolates training source within the classifier analysis."),
        SlideNote(15, "Consensus ranking is a secondary panel analysis", "0:25", "Real and generated profiles rank the same 974 genes. Combining those rankings can move a gene into or out of a compact top-k panel. This is useful for pathway interpretation, but it does not isolate training source as directly as the matched all-gene comparison."),
        SlideNote(16, "Correlated genes dilute marginal importance", "0:35", "Suppose several genes carry the same pathway signal. Ridge can divide weight among them. If I shuffle Gene A, Genes B and C remain, so held-out performance changes little and Gene A receives low permutation importance. Consensus ranking can still retain the group. Low individual importance means replaceable in this classifier, not biologically irrelevant."),
        SlideNote(17, "Matched augmentation helped many tissues, but not all", "0:35", "The pooled multi-tissue classifier declined with augmentation. In separate tissue models, real plus synthetic passed all pooled and accession-macro balanced-accuracy, AUROC, and average-precision checks in 18 of 27 analyses. The coral examples show why a balanced-accuracy gain alone is not enough when another metric declines."),
        SlideNote(18, "Matched and consensus results overlap only partly", "0:30", "The matched analysis retained 21 BH-FDR tissue-gene associations, and consensus ranking retained 49. Eleven appear in both. Matched results are primary evidence that synthetic training changes classifier behavior. Consensus-only results are secondary panel evidence."),
        SlideNote(19, "Consensus selection and literature are separate dimensions", "0:30", "For the 49 consensus associations, promoted or reinforced describes feature selection. Aligning, complementary, ambiguous, or unmatched describes the literature review. These labels answer different questions and can occur in any combination."),
        SlideNote(20, "The matched screen covered all 27 tissue analyses", "0:25", "Four tissues had both matched utility and retained BH-FDR genes. Fourteen more passed the utility gate without a retained gene. Nine failed at least one mean metric. Predictive improvement does not automatically produce a biological candidate."),
        SlideNote(21, "The secondary consensus inventory spans ten tissue analyses", "0:20", "This slide lists all 49 consensus associations. Rows separate FLT direction and selection status. Gene color gives the independent literature classification."),
        SlideNote(22, "Thymus is strongest across both analyses", "0:45", "The matched analysis retained 15 thymus genes, seven promoted after augmentation, and 26 significant Reactome terms led by mitotic cell cycle. Nine genes overlap the 16-gene consensus panel, which adds correlated cell-cycle members. Together they support lower proliferative renewal or fewer cycling thymocytes in flight."),
        SlideNote(23, "Soleus remains a secondary consensus result", "0:35", "The consensus analysis reinforces lower Bdh1, Ech1, Bnip3, and Decr1 with higher Tpm1, a coherent mitochondrial and lipid-metabolism panel. In the fixed all-gene comparison, balanced accuracy rose slightly but AUROC and average precision fell, so the matched gate did not pass."),
        SlideNote(24, "Three additional tissues have matched genes", "0:25", "Liver contributes four flight-lower shared-importance genes without pathway enrichment. Skin Plscr1 and spleen Loxl1 are supported by both analyses. Pooled muscle improves prediction in the matched classifier, but its individual gene interpretation remains consensus-level."),
        SlideNote(25, "Additional panels remain consensus-only", "0:20", "Kidney, adrenal gland, gastrocnemius, and tibialis anterior add smaller tissue-specific consensus candidates. They did not pass the matched gene gate, so they remain exploratory panel-level hypotheses."),
        SlideNote(26, "Use matched tests for contribution and consensus for programs", "0:30", "Conditional DDIM produced realistic profiles. Matched classifiers show where synthetic training changes held-out-real prediction and gene importance. Consensus ranking organizes correlated biological panels. Thymus is strongest across both; all association statistics still come from observed OSDR samples."),
        SlideNote(27, "Thank you", "0:10", "Acknowledge James Casaletto, SLSTP, NASA OSDR, ARCHS4, Reactome, and NASA Ames compute, then invite questions."),
    ]
    notes = [
        SlideNote(index, note.title, note.time, note.text)
        for index, note in enumerate(notes, start=1)
    ]
    for note, slide in zip(notes, presentation.slides):
        _add_notes(slide, note)
    _write_notes(notes)
    presentation.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    print(build())
