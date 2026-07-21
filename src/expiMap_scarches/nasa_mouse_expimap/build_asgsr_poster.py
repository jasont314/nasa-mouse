"""Build the approved-template ASGSR expiMap scientific poster."""

from __future__ import annotations

import shutil
import subprocess
import tempfile
import textwrap
from io import BytesIO
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.lines import Line2D
import numpy as np
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from .integrate_reassessed_tissues_paper import (
    MAIN_TISSUES,
    RETAINED_TERMS,
    load_retained_evidence,
    project_effects,
)


ROOT = Path(__file__).resolve().parents[3]
PAPER_DIR = ROOT / "paper/asgsr_expimap_hvg"
FIGURE_DIR = PAPER_DIR / "figures"
POSTER_DIR = PAPER_DIR / "poster"
ASSET_DIR = POSTER_DIR / "assets"
TEMPLATE_PATH = (
    ROOT
    / "assets/poster_template/00 Poster Session Student Template_Approved by Legal.pptx"
)

TITLE = (
    "Cross-mission expiMap analysis recovers established tissue responses "
    "and identifies complementary pathway shifts in mouse spaceflight "
    "transcriptomes"
)
TITLE_LINE_1 = (
    "Cross-mission expiMap analysis recovers established tissue responses"
)
TITLE_LINE_2 = (
    "and identifies complementary pathway shifts in mouse spaceflight transcriptomes"
)

SLIDE_W = 48.0
SLIDE_H = 27.0
FONT = "Arial"

WHITE = "FFFFFF"
INK = "071C3F"
BODY = "33445D"
MUTED = "697889"
HEADER = "DDECF7"
PANEL = "DCEAF4"
PANEL_LIGHT = "EEF4F8"
RULE = "AEBCC5"
NAVY = "082552"
BLUE = "236DA2"
GOLD = "9A6D1D"
GREEN = "178C6A"
ORANGE = "D97800"
THYMUS = "6A4C93"
SKIN = "C65D37"
LIVER = "13877C"
SPLEEN = "A64D67"

ALIGNED_CONTEXT_TERMS = {
    "thymus": (
        "R-MMU-69278_CELL_CYCLE_MITOTIC",
        "R-MMU-202403_TCR_SIGNALING",
    ),
    "skin": ("R-MMU-6805567_KERATINIZATION",),
    "liver": ("R-MMU-422356_REGULATION_OF_INSULIN_SECRETION",),
    "spleen": (),
}
POSTER_TERMS = {
    "thymus": (
        *ALIGNED_CONTEXT_TERMS["thymus"],
        *RETAINED_TERMS["thymus"],
    ),
    "skin": (
        *ALIGNED_CONTEXT_TERMS["skin"],
        "R-MMU-421270_CELL_CELL_JUNCTION_ORGANIZATION",
        "R-MMU-3247509_CHROMATIN_MODIFYING_ENZYMES",
        "R-MMU-5358351_SIGNALING_BY_HEDGEHOG",
        "R-MMU-428157_SPHINGOLIPID_METABOLISM",
        "R-MMU-73894_DNA_REPAIR",
    ),
    "liver": (
        *ALIGNED_CONTEXT_TERMS["liver"],
        *RETAINED_TERMS["liver"],
    ),
    "spleen": RETAINED_TERMS["spleen"],
}
CONTEXT_DISPLAY_LABELS = {
    "R-MMU-69278_CELL_CYCLE_MITOTIC": "Mitotic cell cycle",
    "R-MMU-202403_TCR_SIGNALING": "T-cell receptor signaling",
    "R-MMU-6805567_KERATINIZATION": "Keratinization",
    "R-MMU-422356_REGULATION_OF_INSULIN_SECRETION": (
        "Regulation of insulin secretion"
    ),
}


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_fill(shape, color: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)


def set_line(shape, color: str, width: float = 1.0) -> None:
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)


def keep_words_intact(paragraph) -> None:
    """Require PowerPoint to wrap Latin text only at word boundaries."""
    properties = paragraph._p.get_or_add_pPr()
    properties.set("latinLnBrk", "0")
    properties.set("eaLnBrk", "0")
    properties.set("hangingPunct", "0")


def add_text(
    slide,
    value: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float,
    color: str = BODY,
    bold: bool = False,
    italic: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.0,
    word_wrap: bool = True,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = word_wrap
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0
    keep_words_intact(paragraph)
    run = paragraph.add_run()
    run.text = value
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return shape


def add_text_lines(
    slide,
    values: tuple[str, ...],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float,
    color: str = BODY,
    bold: bool = False,
    italic: bool = False,
    align=PP_ALIGN.LEFT,
) -> list:
    """Author one complete, nonwrapping line per PowerPoint text box."""
    if not values or any("\n" in value for value in values):
        raise ValueError("Fixed text lines must be nonempty and contain no line breaks")
    line_height = h / len(values)
    if line_height < size / 72:
        raise ValueError("Fixed text lines do not have enough vertical space")
    shapes = []
    for index, value in enumerate(values):
        shape = add_text(
            slide,
            value,
            x,
            y + index * line_height,
            w,
            line_height,
            size=size,
            color=color,
            bold=bold,
            italic=italic,
            align=align,
            valign=MSO_ANCHOR.MIDDLE,
            word_wrap=False,
        )
        shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        shapes.append(shape)
    return shapes


def add_paragraphs(
    slide,
    values: list[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float,
    color: str = BODY,
    gap: float = 4.0,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    for index, value in enumerate(values):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(gap)
        paragraph.line_spacing = 1.0
        keep_words_intact(paragraph)
        run = paragraph.add_run()
        run.text = value
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = rgb(color)
    return shape


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str,
    *,
    line: str | None = None,
    line_width: float = 1.0,
    rounded: bool = False,
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        kind, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_fill(shape, fill)
    if line:
        set_line(shape, line, line_width)
    else:
        shape.line.fill.background()
    return shape


def add_panel(slide, x: float, y: float, w: float, h: float) -> None:
    add_rect(slide, x, y, w, h, PANEL)


def add_section(
    slide, label: str, x: float, y: float, w: float, *, size: float = 26
) -> None:
    add_rect(slide, x, y, w, 0.72, WHITE, line=RULE, line_width=1.1)
    add_text(
        slide,
        label,
        x + 0.15,
        y + 0.02,
        w - 0.30,
        0.66,
        size=size,
        color="111111",
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        word_wrap=False,
    )


def add_badge(
    slide,
    label: str,
    x: float,
    y: float,
    w: float,
    *,
    color: str = BLUE,
) -> None:
    shape = add_rect(
        slide,
        x,
        y,
        w,
        0.72,
        WHITE,
        line=color,
        line_width=1.2,
        rounded=True,
    )
    shape.adjustments[0] = 0.08
    add_text(
        slide,
        label,
        x + 0.08,
        y + 0.03,
        w - 0.16,
        0.63,
        size=17,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        word_wrap=False,
    )


def add_box(
    slide,
    title: str,
    subtitle: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str,
    line: str,
    title_size: float = 18,
    subtitle_size: float = 14.5,
) -> None:
    shape = add_rect(
        slide, x, y, w, h, fill, line=line, line_width=1.4, rounded=True
    )
    shape.adjustments[0] = 0.06
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.margin_left = Inches(0.10)
    frame.margin_right = Inches(0.10)
    frame.margin_top = Inches(0.05)
    frame.margin_bottom = Inches(0.04)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    first = frame.paragraphs[0]
    first.alignment = PP_ALIGN.CENTER
    first.space_after = Pt(2)
    keep_words_intact(first)
    run = first.add_run()
    run.text = title
    run.font.name = FONT
    run.font.size = Pt(title_size)
    run.font.bold = True
    run.font.color.rgb = rgb(INK)
    second = frame.add_paragraph()
    second.alignment = PP_ALIGN.CENTER
    second.space_before = Pt(0)
    second.space_after = Pt(0)
    keep_words_intact(second)
    run = second.add_run()
    run.text = subtitle
    run.font.name = FONT
    run.font.size = Pt(subtitle_size)
    run.font.color.rgb = rgb(MUTED)


def add_down_arrow(slide, center_x: float, y: float, color: str = BLUE) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        Inches(center_x - 0.18),
        Inches(y),
        Inches(0.36),
        Inches(0.34),
    )
    set_fill(shape, color)
    shape.line.fill.background()


def add_right_arrow(
    slide, x: float, y: float, w: float, h: float, color: str = BLUE
) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_fill(shape, color)
    shape.line.fill.background()


def add_connector(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: str,
    width: float = 1.1,
    dashed: bool = False,
):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    set_line(connector, color, width)
    if dashed:
        connector.line.dash_style = 2
    return connector


def set_table_cell(
    cell,
    value: str,
    *,
    size: float,
    color: str,
    bold: bool,
    center: bool,
) -> None:
    cell.text = ""
    cell.margin_left = Inches(0.20 if not center else 0.07)
    cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.03)
    cell.margin_bottom = Inches(0.03)
    cell.text_frame.word_wrap = False
    cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
    keep_words_intact(paragraph)
    run = paragraph.add_run()
    run.text = value
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_scope_table(slide, x: float, y: float, w: float, h: float) -> None:
    rows = [
        ("Thymus", "1,362", "117", "5", "387", THYMUS),
        ("Skin", "2,593", "151", "4", "319", SKIN),
        ("Liver", "5,000", "197", "9", "364", LIVER),
        ("Spleen", "6,289", "100", "5", "360", SPLEEN),
    ]
    shape = slide.shapes.add_table(5, 5, Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    widths = [3.10, 2.75, 2.35, 2.25, 2.65]
    for index, value in enumerate(widths):
        table.columns[index].width = Inches(value)
    for column, value in enumerate(("Tissue", "ARCHS4", "OSDR", "Projects", "Programs")):
        cell = table.cell(0, column)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(NAVY)
        set_table_cell(
            cell,
            value,
            size=16,
            color=WHITE,
            bold=True,
            center=column > 0,
        )
    for row_index, row in enumerate(rows, 1):
        for column, value in enumerate(row[:5]):
            cell = table.cell(row_index, column)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(WHITE if row_index % 2 else PANEL_LIGHT)
            set_table_cell(
                cell,
                value,
                size=16.5,
                color=row[5] if column == 0 else BODY,
                bold=column == 0,
                center=column > 0,
            )
    for row in table.rows:
        row.height = Inches(h / 5)


def add_nasa_logo(slide, template: Presentation) -> None:
    template_slide = template.slides[0]
    candidates = [
        shape
        for shape in template_slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        and shape.left / template.slide_width > 0.85
        and shape.top / template.slide_height < 0.25
    ]
    if not candidates:
        raise RuntimeError("NASA logo was not found in the approved template")
    logo = max(candidates, key=lambda shape: shape.width * shape.height)
    slide.shapes.add_picture(BytesIO(logo.image.blob), Inches(44.20), Inches(0.78), width=Inches(3.05))


def add_challenge_equation(slide, x: float, y: float, w: float) -> None:
    add_text(
        slide,
        "THE CENTRAL CHALLENGE",
        x,
        y,
        w,
        0.40,
        size=19,
        color=GOLD,
        bold=True,
    )
    add_text(
        slide,
        "Observed expression",
        x,
        y + 0.48,
        2.70,
        0.60,
        size=17,
        color=INK,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        word_wrap=False,
    )
    parts = [
        ("Spaceflight biology", ORANGE, 2.75),
        ("Tissue biology", NAVY, 2.20),
        ("Animal variation", MUTED, 2.35),
        ("Study / protocol", GOLD, 2.55),
    ]
    cursor = x + 2.85
    for index, (label, color, width) in enumerate(parts):
        if index:
            add_text(
                slide,
                "+",
                cursor,
                y + 0.48,
                0.38,
                0.60,
                size=22,
                color=MUTED,
                bold=True,
                align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE,
                word_wrap=False,
            )
            cursor += 0.43
        add_rect(slide, cursor, y + 0.48, width, 0.60, color, rounded=True)
        add_text(
            slide,
            label,
            cursor + 0.07,
            y + 0.50,
            width - 0.14,
            0.54,
            size=15.5,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
            word_wrap=False,
        )
        cursor += width


def add_architecture(slide, x: float, y: float, w: float, h: float) -> None:
    add_text(
        slide,
        "Pathway structure is wired into the decoder",
        x,
        y,
        w,
        0.50,
        size=22,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
        word_wrap=False,
    )
    add_text(
        slide,
        (
            "Train a tissue reference, then map flight and ground samples into "
            "the same annotated program space."
        ),
        x + 0.35,
        y + 0.52,
        w - 0.70,
        0.46,
        size=16,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        word_wrap=False,
    )

    add_text(
        slide,
        "1  REFERENCE TRAINING",
        x + 0.15,
        y + 1.03,
        4.00,
        0.35,
        size=17,
        color=BLUE,
        bold=True,
    )
    add_box(
        slide,
        "ARCHS4 reference",
        "non-spaceflight tissue counts",
        x + 0.15,
        y + 1.45,
        3.55,
        1.10,
        fill="E9F3F8",
        line=BLUE,
        title_size=18,
        subtitle_size=14,
    )
    add_right_arrow(slide, x + 3.88, y + 1.79, 0.55, 0.36, BLUE)

    model = add_rect(
        slide,
        x + 4.55,
        y + 1.28,
        10.20,
        2.65,
        WHITE,
        line=BLUE,
        line_width=1.2,
        rounded=True,
    )
    model.adjustments[0] = 0.04
    add_text(
        slide,
        "expiMap tissue model",
        x + 4.82,
        y + 1.42,
        4.40,
        0.36,
        size=17.5,
        color=INK,
        bold=True,
        word_wrap=False,
    )

    input_x = x + 4.95
    encoder_x = x + 5.48
    encoder_w = 1.48
    latent_x = x + 7.42
    decoder_x = x + 10.58
    decoder_w = 1.48
    output_x = x + 13.55
    gene_ys = [y + 2.05 + 0.42 * index for index in range(4)]
    for gene_y in gene_ys:
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(input_x),
            Inches(gene_y),
            Inches(0.24),
            Inches(0.24),
        )
        set_fill(node, "A9BBD0")
        node.line.fill.background()

    encoder = slide.shapes.add_shape(
        MSO_SHAPE.TRAPEZOID,
        Inches(encoder_x),
        Inches(y + 1.92),
        Inches(encoder_w),
        Inches(1.42),
    )
    encoder.rotation = 90
    set_fill(encoder, "DCE6F2")
    set_line(encoder, BLUE, 1.4)
    add_text(
        slide,
        "Encoder",
        encoder_x - 0.05,
        y + 2.42,
        encoder_w + 0.10,
        0.38,
        size=14,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        word_wrap=False,
    )
    for gene_y in gene_ys:
        add_connector(
            slide,
            input_x + 0.24,
            gene_y + 0.12,
            encoder_x + 0.04,
            y + 2.63,
            color="A9BBD0",
            width=0.8,
        )

    latent_rows = [
        ("DNA repair", THYMUS, y + 1.95),
        ("T-cell signaling", SPLEEN, y + 2.43),
        ("Cell junctions", SKIN, y + 2.91),
    ]
    for label, color, node_y in latent_rows:
        add_connector(
            slide,
            encoder_x + encoder_w - 0.04,
            y + 2.63,
            latent_x,
            node_y + 0.16,
            color="A9BBD0",
            width=0.9,
        )
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(latent_x),
            Inches(node_y),
            Inches(0.32),
            Inches(0.32),
        )
        set_fill(node, color)
        node.line.fill.background()
        label_box = add_rect(
            slide,
            latent_x + 0.40,
            node_y - 0.02,
            2.22,
            0.36,
            "F8FAFC",
            line=color,
            line_width=0.9,
            rounded=True,
        )
        label_box.adjustments[0] = 0.08
        add_text(
            slide,
            label,
            latent_x + 0.50,
            node_y - 0.02,
            2.02,
            0.36,
            size=15.5,
            color=color,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
            word_wrap=False,
        )

    output_ys = [y + 2.03 + 0.38 * index for index in range(4)]
    connection_map = [(0, 0), (0, 2), (1, 1), (1, 3), (2, 2), (2, 3)]
    for latent_index, output_index in connection_map:
        _, color, node_y = latent_rows[latent_index]
        add_connector(
            slide,
            latent_x + 2.62,
            node_y + 0.16,
            output_x,
            output_ys[output_index] + 0.12,
            color=color,
            width=1.2,
        )

    decoder = slide.shapes.add_shape(
        MSO_SHAPE.TRAPEZOID,
        Inches(decoder_x),
        Inches(y + 1.92),
        Inches(decoder_w),
        Inches(1.42),
    )
    decoder.rotation = 270
    decoder.fill.background()
    set_line(decoder, GREEN, 1.4)
    add_rect(
        slide,
        decoder_x + 0.27,
        y + 2.43,
        0.94,
        0.34,
        WHITE,
    )
    add_text(
        slide,
        "Decoder",
        decoder_x + 0.22,
        y + 2.42,
        1.04,
        0.38,
        size=13.5,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        word_wrap=False,
    )

    for out_y in output_ys:
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(output_x),
            Inches(out_y),
            Inches(0.24),
            Inches(0.24),
        )
        set_fill(node, "A9BBD0")
        node.line.fill.background()

    add_rect(
        slide,
        x + 10.12,
        y + 1.42,
        3.78,
        0.58,
        "E9F4EE",
        line=GREEN,
        line_width=1.1,
        rounded=True,
    )
    add_text(
        slide,
        "Reactome mask selects pathway-gene edges",
        x + 10.25,
        y + 1.48,
        3.52,
        0.44,
        size=13.2,
        color=GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        word_wrap=False,
    )
    add_down_arrow(slide, decoder_x + decoder_w / 2, y + 2.03, GREEN)

    add_text(
        slide,
        "genes",
        input_x - 0.18,
        y + 3.48,
        0.66,
        0.35,
        size=14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        word_wrap=False,
    )
    add_text(
        slide,
        "program scores",
        latent_x - 0.05,
        y + 3.48,
        2.65,
        0.35,
        size=14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        word_wrap=False,
    )
    add_text(
        slide,
        "masked decoder",
        decoder_x - 0.18,
        y + 3.48,
        decoder_w + 0.36,
        0.35,
        size=14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        word_wrap=False,
    )
    add_text(
        slide,
        "output genes",
        x + 12.85,
        y + 3.48,
        1.65,
        0.35,
        size=14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        word_wrap=False,
    )

    add_text(
        slide,
        "2  QUERY MAPPING AND CROSS-MISSION COMPARISON",
        x + 0.15,
        y + 4.08,
        8.10,
        0.36,
        size=17,
        color=ORANGE,
        bold=True,
    )
    add_box(
        slide,
        "NASA OSDR query",
        "FLT, GC, and accession",
        x + 0.15,
        y + 4.48,
        3.20,
        1.08,
        fill="FAEFE7",
        line=ORANGE,
        title_size=17.5,
        subtitle_size=14,
    )
    add_right_arrow(slide, x + 3.48, y + 4.82, 0.45, 0.36, ORANGE)
    add_box(
        slide,
        "Reference-query map",
        "scArches adaptation",
        x + 4.05,
        y + 4.48,
        3.10,
        1.08,
        fill=WHITE,
        line=ORANGE,
        title_size=17,
        subtitle_size=14,
    )
    add_right_arrow(slide, x + 7.28, y + 4.82, 0.45, 0.36, BLUE)
    add_box(
        slide,
        "Program scores",
        "posterior mean per sample",
        x + 7.85,
        y + 4.48,
        2.95,
        1.08,
        fill=WHITE,
        line=NAVY,
        title_size=17,
        subtitle_size=14,
    )
    add_right_arrow(slide, x + 10.93, y + 4.82, 0.40, 0.36, BLUE)
    add_box(
        slide,
        "Project-balanced shift",
        "equal-weight FLT - GC mean",
        x + 11.45,
        y + 4.48,
        3.30,
        1.08,
        fill=WHITE,
        line=BLUE,
        title_size=16.5,
        subtitle_size=14,
    )

    add_down_arrow(slide, x + 13.10, y + 5.60, GREEN)
    add_box(
        slide,
        "Post-score annotation",
        "literature and mission-context role",
        x + 9.10,
        y + 5.94,
        5.65,
        0.62,
        fill="EBF5EF",
        line=GREEN,
        title_size=15.5,
        subtitle_size=13.5,
    )

    add_text(
        slide,
        "~2,000 HVGs | 319-387 Reactome programs | 250-epoch query mapping",
        x + 0.20,
        y + 6.08,
        8.55,
        0.34,
        size=14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        word_wrap=False,
    )


def load_poster_evidence() -> pd.DataFrame:
    """Combine the retained core with explicitly qualified literature anchors."""
    retained = load_retained_evidence().copy()
    retained["poster_tier"] = "retained"
    retained.loc[
        retained["term"].eq("R-MMU-421270_CELL_CELL_JUNCTION_ORGANIZATION"),
        "evidence_role",
    ] = "aligned"

    context = pd.read_csv(
        PAPER_DIR / "source_data/table_s24_pathway_robustness_evidence.tsv",
        sep="\t",
    )
    context = context.loc[
        [
            row.term in ALIGNED_CONTEXT_TERMS.get(row.tissue, ())
            for row in context.itertuples(index=False)
        ]
    ].copy()
    context["display_label"] = context["term"].map(CONTEXT_DISPLAY_LABELS)
    context["evidence_role"] = "aligned"
    context["analysis_role"] = "context_only"
    context["poster_tier"] = "literature_anchor"
    context["seed_effect_median"] = context[
        ["effect_seed2020", "effect_seed2021", "effect_seed2022"]
    ].median(axis=1)
    context["seed_effect_minimum"] = context[
        ["effect_seed2020", "effect_seed2021", "effect_seed2022"]
    ].min(axis=1)
    context["seed_effect_maximum"] = context[
        ["effect_seed2020", "effect_seed2021", "effect_seed2022"]
    ].max(axis=1)

    columns = sorted(set(retained.columns).union(context.columns))
    return pd.concat(
        [retained.reindex(columns=columns), context.reindex(columns=columns)],
        ignore_index=True,
    )


def render_poster_pathway_asset(output: Path) -> None:
    """Render retained pathways and qualified aligned anchors for the poster."""
    role_colors = {
        "aligned": "#009E73",
        "complementary": "#0072B2",
        "context_sensitive": "#D55E00",
    }
    role_markers = {
        "aligned": "o",
        "complementary": "s",
        "context_sensitive": "^",
    }
    evidence = load_poster_evidence()
    style = {
        "font.family": "DejaVu Sans",
        "font.size": 14,
        "axes.labelsize": 15,
        "xtick.labelsize": 13,
        "ytick.labelsize": 15,
        "text.color": "#202629",
        "axes.labelcolor": "#202629",
        "xtick.color": "#202629",
        "ytick.color": "#202629",
    }
    with plt.rc_context(style):
        fig, axes = plt.subplots(
            2,
            2,
            figsize=(14.4, 8.0),
            layout="constrained",
        )
        for panel, ax, tissue in zip("abcd", axes.flat, MAIN_TISSUES):
            terms = POSTER_TERMS[tissue]
            subset = (
                evidence.loc[evidence["tissue"].eq(tissue)]
                .set_index("term")
                .loc[list(terms)]
                .reset_index()
            )
            points = project_effects(tissue, terms)
            positions = np.arange(len(subset))[::-1]
            for position, row in zip(positions, subset.itertuples(index=False)):
                local = points.loc[
                    points["term"].eq(row.term), "project_effect"
                ].to_numpy(dtype=float)
                ax.scatter(
                    local,
                    np.full(len(local), position),
                    s=64,
                    facecolor="white",
                    edgecolor="#707A7F",
                    linewidth=1.2,
                    zorder=2,
                )
                role = str(row.evidence_role)
                is_anchor = row.poster_tier == "literature_anchor"
                ax.hlines(
                    position,
                    float(row.seed_effect_minimum),
                    float(row.seed_effect_maximum),
                    color=role_colors[role],
                    linewidth=2.0 if is_anchor else 2.4,
                    linestyle="--" if is_anchor else "-",
                    zorder=3,
                )
                ax.scatter(
                    float(row.seed_effect_median),
                    position,
                    marker="D" if is_anchor else role_markers[role],
                    s=132,
                    facecolor="white" if is_anchor else role_colors[role],
                    edgecolor=role_colors[role] if is_anchor else "white",
                    linewidth=1.8 if is_anchor else 1.1,
                    zorder=4,
                )

            ax.axvline(0, color="#3F494E", linewidth=1.1)
            ax.set_yticks(positions)
            ax.set_yticklabels(
                [
                    textwrap.fill(
                        str(label),
                        width=29,
                        break_long_words=False,
                        break_on_hyphens=False,
                    )
                    for label in subset["display_label"]
                ]
            )
            for tick, row in zip(ax.get_yticklabels(), subset.itertuples(index=False)):
                tick.set_color(role_colors[str(row.evidence_role)])
                if row.poster_tier == "literature_anchor":
                    tick.set_fontstyle("italic")
            project_count = int(subset["expimap_n_projects"].max())
            ax.set_title(
                f"{panel}  {tissue.title()} ({project_count} projects)",
                loc="left",
                fontweight="bold",
                fontsize=18,
                pad=9,
            )
            if panel in "cd":
                ax.set_xlabel("Flight - ground pathway shift")
            ax.grid(axis="x", color="#DCE1E3", linewidth=1.0)
            ax.set_axisbelow(True)
            ax.tick_params(axis="y", length=0, pad=8)
            ax.tick_params(axis="x", width=1.0, length=4)
            for side in ("top", "right", "left"):
                ax.spines[side].set_visible(False)
            ax.spines["bottom"].set_color("#697277")
            ax.spines["bottom"].set_linewidth(1.0)

        handles = [
            Line2D(
                [0],
                [0],
                marker="o",
                linestyle="none",
                markerfacecolor="white",
                markeredgecolor="#707A7F",
                markeredgewidth=1.2,
                markersize=8,
                label="OSDR project",
            ),
            *[
                Line2D(
                    [0],
                    [0],
                    marker=role_markers[role],
                    color=role_colors[role],
                    linewidth=2.4,
                    markersize=8,
                    label=label,
                )
                for role, label in (
                    ("aligned", "Aligned and retained"),
                    ("complementary", "Complementary and retained"),
                    ("context_sensitive", "Context sensitive and retained"),
                )
            ],
            Line2D(
                [0],
                [0],
                marker="D",
                color=role_colors["aligned"],
                linestyle="--",
                markerfacecolor="white",
                markeredgecolor=role_colors["aligned"],
                markeredgewidth=1.5,
                linewidth=2.0,
                markersize=7,
                label="Literature aligned, not retained",
            ),
        ]
        fig.legend(
            handles=handles,
            loc="outside lower center",
            ncol=5,
            frameon=False,
            fontsize=12.5,
            handlelength=1.8,
            columnspacing=1.25,
        )
        fig.get_layout_engine().set(
            w_pad=0.12,
            h_pad=0.12,
            wspace=0.12,
            hspace=0.12,
        )
        output.parent.mkdir(parents=True, exist_ok=True)
        fig.savefig(
            output,
            dpi=400,
            facecolor="white",
            metadata={
                "Title": "Aligned anchors and retained expiMap pathway shifts"
            },
        )
        plt.close(fig)


def render_pdf_asset(source: Path, output: Path, dpi: int = 700) -> None:
    output.parent.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        [
            "pdftocairo",
            "-singlefile",
            "-png",
            "-r",
            str(dpi),
            str(source),
            str(output.with_suffix("")),
        ],
        check=True,
    )


def add_picture_contain(
    slide,
    path: Path,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    name: str,
    description: str,
) -> float:
    with Image.open(path) as image:
        px_w, px_h = image.size
    image_ratio = px_w / px_h
    region_ratio = w / h
    if image_ratio >= region_ratio:
        placed_w = w
        placed_h = w / image_ratio
        placed_x = x
        placed_y = y + (h - placed_h) / 2
    else:
        placed_h = h
        placed_w = h * image_ratio
        placed_x = x + (w - placed_w) / 2
        placed_y = y
    picture = slide.shapes.add_picture(
        str(path),
        Inches(placed_x),
        Inches(placed_y),
        Inches(placed_w),
        Inches(placed_h),
    )
    properties = picture._element.xpath(".//p:cNvPr")[0]
    properties.set("name", name)
    properties.set("descr", description)
    ppi = min(px_w / placed_w, px_h / placed_h)
    if ppi < 300:
        raise ValueError(f"{path.name} has only {ppi:.0f} effective ppi")
    return ppi


def validate(prs: Presentation) -> None:
    if len(prs.slides) != 1:
        raise ValueError("Poster must contain exactly one slide")
    for shape in prs.slides[0].shapes:
        if shape.left < 0 or shape.top < 0:
            raise ValueError(f"{shape.name!r} starts outside the slide")
        if shape.left + shape.width > prs.slide_width + 1:
            raise ValueError(f"{shape.name!r} exceeds the slide width")
        if shape.top + shape.height > prs.slide_height + 1:
            raise ValueError(f"{shape.name!r} exceeds the slide height")
        if shape.has_text_frame:
            normalized = " ".join(shape.text.split())
            if len(normalized) >= 60 and shape.text_frame.word_wrap:
                raise ValueError(
                    f"{shape.name!r} contains long reflowable text; use fixed lines"
                )
            for paragraph in shape.text_frame.paragraphs:
                if not paragraph.text:
                    continue
                properties = paragraph._p.get_or_add_pPr()
                if properties.get("latinLnBrk") != "0":
                    raise ValueError(
                        f"{shape.name!r} allows a Latin word to break between letters"
                    )
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        if not paragraph.text:
                            continue
                        properties = paragraph._p.get_or_add_pPr()
                        if properties.get("latinLnBrk") != "0":
                            raise ValueError(
                                "A table cell allows a Latin word to break between letters"
                            )


def render_poster(pptx_path: Path) -> tuple[Path | None, Path | None]:
    libreoffice = shutil.which("libreoffice")
    pdftocairo = shutil.which("pdftocairo")
    if not libreoffice:
        return None, None
    pdf_path = pptx_path.with_suffix(".pdf")
    pdf_path.unlink(missing_ok=True)
    with tempfile.TemporaryDirectory(prefix="asgsr_poster_lo_") as profile:
        subprocess.run(
            [
                libreoffice,
                "--headless",
                f"-env:UserInstallation={Path(profile).resolve().as_uri()}",
                "--convert-to",
                "pdf",
                "--outdir",
                str(pptx_path.parent),
                str(pptx_path),
            ],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            text=True,
        )
    if not pdf_path.exists():
        raise FileNotFoundError(f"LibreOffice did not create {pdf_path}")
    if not pdftocairo:
        return pdf_path, None
    preview = pptx_path.with_name(pptx_path.stem + "_preview.png")
    subprocess.run(
        [
            pdftocairo,
            "-singlefile",
            "-png",
            "-r",
            "100",
            str(pdf_path),
            str(preview.with_suffix("")),
        ],
        check=True,
    )
    return pdf_path, preview


def render_architecture_crop(pdf_path: Path) -> Path:
    output = ASSET_DIR / "expimap_architecture_visualization_300dpi.png"
    dpi = 300
    x, y, w, h = 16.18, 6.28, 15.05, 7.08
    subprocess.run(
        [
            "pdftocairo",
            "-singlefile",
            "-png",
            "-r",
            str(dpi),
            "-x",
            str(round(x * dpi)),
            "-y",
            str(round(y * dpi)),
            "-W",
            str(round(w * dpi)),
            "-H",
            str(round(h * dpi)),
            str(pdf_path),
            str(output.with_suffix("")),
        ],
        check=True,
    )
    return output


def build() -> tuple[Path, Path | None, Path | None, Path | None, list[float]]:
    POSTER_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    result_asset = ASSET_DIR / "figure_3_tissue_pathway_shifts_poster_400dpi.png"
    hypothesis_asset = ASSET_DIR / "figure_6_tissue_state_hypotheses_700dpi.png"
    render_poster_pathway_asset(result_asset)
    render_pdf_asset(
        FIGURE_DIR / "figure_6_tissue_state_hypotheses.pdf", hypothesis_asset
    )

    template = Presentation(TEMPLATE_PATH)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.title = TITLE
    prs.core_properties.author = "Jason Trinh"
    prs.core_properties.subject = (
        "Template-based cross-mission expiMap poster using NASA OSDR mouse RNA-seq"
    )
    prs.core_properties.keywords = (
        "spaceflight, expiMap, OSDR, ARCHS4, Reactome, transcriptomics, NASA"
    )
    prs.core_properties.comments = (
        "Layout and NASA branding reference 00 Poster Session Student "
        "Template_Approved by Legal.pptx."
    )
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(WHITE)

    # Approved-template header
    add_rect(slide, 0, 0, SLIDE_W, 5.05, HEADER)
    add_text(
        slide,
        "National Aeronautics and Space Administration",
        0.80,
        0.28,
        18.0,
        0.48,
        size=16,
        color=MUTED,
    )
    for line, line_y in ((TITLE_LINE_1, 1.00), (TITLE_LINE_2, 1.92)):
        add_text(
            slide,
            line,
            0.60,
            line_y,
            42.80,
            0.82,
            size=50,
            color="111111",
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
            word_wrap=False,
        )
    add_text(
        slide,
        (
            "Jason Trinh / University of California, Berkeley | "
            "NASA Space Life Sciences Training Program, NASA Ames Research Center"
        ),
        1.50,
        3.70,
        41.0,
        0.72,
        size=26,
        color=MUTED,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        word_wrap=False,
    )
    add_nasa_logo(slide, template)

    # Template-style panels
    left_x, left_w = 0.35, 15.15
    center_x, center_w = 15.85, 15.75
    right_x, right_w = 31.95, 15.70

    add_panel(slide, left_x, 5.35, left_w, 3.78)
    add_section(slide, "ABSTRACT", left_x + 0.28, 5.55, left_w - 0.56)
    abstract_lines = (
        "Mission differences can obscure spaceflight responses shared across studies. We mapped NASA OSDR mouse bulk RNA",
        "sequencing samples to tissue-matched ARCHS4 references. Reactome-constrained expiMap models analyzed about",
        "2,000 highly variable genes. Flight-ground shifts were balanced across projects and evaluated with gene-set",
        "enrichment, project holdouts, three complete training runs, composition proxies, and member-gene review.",
        "Literature review followed scoring. Retained programs were lower for thymic repair, skin maintenance, liver",
        "adaptive immunity, and splenic adaptive and innate immunity. Spleen showed the strongest multi-pathway result.",
    )
    add_text_lines(
        slide,
        abstract_lines,
        left_x + 0.50,
        6.48,
        left_w - 1.00,
        1.86,
        size=18.5,
        color=BODY,
    )

    add_panel(slide, left_x, 9.43, left_w, 15.25)
    add_section(slide, "INTRODUCTION", left_x + 0.28, 9.63, left_w - 0.56)
    add_text(
        slide,
        "PROJECT OBJECTIVE",
        left_x + 0.55,
        10.62,
        left_w - 1.10,
        0.42,
        size=19,
        color=GOLD,
        bold=True,
    )
    add_text(
        slide,
        "Identify gene programs that shift consistently across spaceflight missions.",
        left_x + 0.55,
        11.05,
        left_w - 1.10,
        1.10,
        size=25,
        color=INK,
        bold=True,
        word_wrap=False,
    )
    add_challenge_equation(slide, left_x + 0.55, 12.30, left_w - 1.10)
    add_text_lines(
        slide,
        (
            "Study identity and protocol can dominate an unconstrained expression representation. Tissue-matched",
            "reference mapping, accession conditioning, and equal project weighting reduce this bias without claiming",
            "to remove all mission confounding.",
        ),
        left_x + 0.55,
        13.62,
        left_w - 1.10,
        0.96,
        size=19,
        color=BODY,
    )
    add_text(
        slide,
        "FINAL ANALYSIS SCOPE",
        left_x + 0.55,
        15.32,
        left_w - 1.10,
        0.40,
        size=19,
        color=GOLD,
        bold=True,
    )
    add_scope_table(
        slide, left_x + 0.55, 15.83, left_w - 1.10, 4.05
    )
    add_text_lines(
        slide,
        ("OSDR counts are primary analysis samples. Spleen excludes one condition-strain-confounded project.",),
        left_x + 0.55,
        19.97,
        left_w - 1.10,
        0.30,
        size=14.5,
        color=MUTED,
        italic=True,
    )
    add_text(
        slide,
        "ROBUSTNESS GATE",
        left_x + 0.55,
        20.65,
        left_w - 1.10,
        0.40,
        size=19,
        color=GOLD,
        bold=True,
    )
    badge_y = 21.18
    badge_w = 4.24
    badge_gap = 0.42
    for index, label in enumerate(("ssGSEA + GSEA", "Held-out projects", "3 full trainings")):
        add_badge(
            slide,
            label,
            left_x + 0.55 + index * (badge_w + badge_gap),
            badge_y,
            badge_w,
            color=BLUE if index < 2 else THYMUS,
        )
    for index, label in enumerate(("Composition proxies", "Member genes", "Literature review")):
        add_badge(
            slide,
            label,
            left_x + 0.55 + index * (badge_w + badge_gap),
            badge_y + 0.90,
            badge_w,
            color=GREEN,
        )
    add_text_lines(
        slide,
        (
            "Complementary programs add a plausible perspective beyond the dominant phenotype in prior literature. De novo",
            "nodes were not retained in the final models.",
        ),
        left_x + 0.55,
        23.08,
        left_w - 1.10,
        0.56,
        size=16.5,
        color=BODY,
    )

    add_panel(slide, center_x, 5.35, center_w, 8.05)
    add_section(slide, "METHODS", center_x + 0.28, 5.55, center_w - 0.56)
    add_architecture(
        slide,
        center_x + 0.40,
        6.42,
        center_w - 0.80,
        6.65,
    )

    add_panel(slide, center_x, 13.70, center_w, 10.98)
    add_section(
        slide,
        "RESULTS: ALIGNED ANCHORS AND RETAINED PROGRAM SHIFTS",
        center_x + 0.28,
        13.90,
        center_w - 0.56,
        size=24,
    )
    result_ppi = add_picture_contain(
        slide,
        result_asset,
        center_x + 0.35,
        14.83,
        center_w - 0.70,
        8.82,
        name="Retained pathway shifts",
        description=(
            "Literature-aligned anchors and retained expiMap pathway shifts for "
            "thymus, skin, liver, and spleen, with project effects and "
            "three-training ranges."
        ),
    )
    add_text_lines(
        slide,
        (
            "Filled markers show 13 retained programs. Open diamonds are literature-aligned context only, not retained findings;",
            "keratinization reversed in one of three trainings. Colored ranges span three complete trainings; axes are tissue specific.",
        ),
        center_x + 0.60,
        23.62,
        center_w - 1.20,
        0.56,
        size=14.5,
        color=MUTED,
        italic=True,
        align=PP_ALIGN.CENTER,
    )

    add_panel(slide, right_x, 5.35, right_w, 13.28)
    add_section(
        slide,
        "RESULTS: BIOLOGICAL INTERPRETATION",
        right_x + 0.28,
        5.55,
        right_w - 0.56,
        size=24,
    )
    add_text(
        slide,
        "Prior evidence, observed shifts, and hypotheses to test",
        right_x + 0.50,
        6.42,
        right_w - 1.00,
        0.42,
        size=19,
        color=BLUE,
        bold=True,
        align=PP_ALIGN.CENTER,
        word_wrap=False,
    )
    hypothesis_ppi = add_picture_contain(
        slide,
        hypothesis_asset,
        right_x + 0.35,
        6.87,
        right_w - 0.70,
        11.22,
        name="Tissue-state hypotheses",
        description=(
            "Prior-literature phenotypes, observed lower expiMap program scores, "
            "and complementary hypotheses for thymus, skin, liver, and spleen."
        ),
    )

    add_panel(slide, right_x, 18.93, right_w, 3.55)
    add_section(slide, "CONCLUSIONS", right_x + 0.28, 19.13, right_w - 0.56)
    add_text_lines(
        slide,
        (
            "\u2022 Spleen was strongest: T cell receptor, neutrophil degranulation, and C\u2011type lectin",
            "  receptor programs were lower across five projects and three trainings; all had GSEA FDR <0.05.",
        ),
        right_x + 0.55,
        20.02,
        right_w - 1.10,
        0.68,
        size=17.5,
        color=BODY,
    )
    add_text_lines(
        slide,
        (
            "\u2022 Skin and thymus emphasized lower tissue maintenance programs; liver showed lower",
            "  adaptive immune signaling amid heterogeneous metabolism.",
        ),
        right_x + 0.55,
        20.80,
        right_w - 1.10,
        0.68,
        size=17.5,
        color=BODY,
    )
    add_text_lines(
        slide,
        (
            "\u2022 Bulk latent scores are testable hypotheses, not causal or cell-resolved measures",
            "  of pathway activity.",
        ),
        right_x + 0.55,
        21.50,
        right_w - 1.10,
        0.58,
        size=17.5,
        color=BODY,
    )

    add_panel(slide, right_x, 22.75, right_w, 1.93)
    add_section(
        slide,
        "REFERENCES",
        right_x + 0.28,
        22.95,
        right_w - 0.56,
        size=22,
    )
    add_text_lines(
        slide,
        (
            "1 Gebre et al., NAR 2025 (OSDR). 2 Lotfollahi et al., Nat Cell Biol 2023 (expiMap).",
            "3 Lachmann et al., Nat Commun 2018 (ARCHS4). 4 Milacic et al., NAR 2024 (Reactome).",
            "5 Horie et al., Sci Rep 2019. 6 Cope et al., Commun Med 2024.",
            "7 Beheshti et al., Sci Rep 2019. 8 Gridley et al., J Appl Physiol 2009.",
        ),
        right_x + 0.52,
        23.78,
        right_w - 1.04,
        0.78,
        size=12.5,
        color=BODY,
    )

    # Approved-template footer and acknowledgements
    add_rect(slide, center_x, 24.92, center_w + right_w + 0.05, 2.08, HEADER)
    add_text(
        slide,
        "Acknowledgements",
        center_x + 0.35,
        25.02,
        5.0,
        0.45,
        size=22,
        color=MUTED,
        bold=True,
    )
    add_text_lines(
        slide,
        (
            "Mentor: James Casaletto | NASA Space Life Sciences Training Program",
            "NASA OSDR, ARCHS4, and Reactome investigators and curation teams",
        ),
        center_x + 0.35,
        25.52,
        20.60,
        0.66,
        size=16,
        color=BODY,
    )
    add_text(
        slide,
        "jasontrinh@berkeley.edu\ngithub.com/jasont314/nasa-mouse",
        39.25,
        25.32,
        7.92,
        0.90,
        size=16,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.RIGHT,
        word_wrap=False,
    )
    add_text(
        slide,
        "www.nasa.gov",
        0.70,
        26.10,
        3.60,
        0.40,
        size=15,
        color=MUTED,
    )

    validate(prs)
    output = POSTER_DIR / "asgsr_expimap_poster.pptx"
    prs.save(output)
    pdf_path, preview_path = render_poster(output)
    architecture_path = render_architecture_crop(pdf_path) if pdf_path else None
    return (
        output,
        pdf_path,
        preview_path,
        architecture_path,
        [result_ppi, hypothesis_ppi],
    )


def run() -> None:
    output, pdf_path, preview_path, architecture_path, ppi = build()
    print(output)
    if pdf_path:
        print(pdf_path)
    if preview_path:
        print(preview_path)
    if architecture_path:
        print(architecture_path)
    print(
        "Embedded data-figure effective resolution: "
        + ", ".join(f"{value:.0f} ppi" for value in ppi)
    )


if __name__ == "__main__":
    run()
