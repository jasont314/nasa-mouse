"""Build the 2026 SLSTP presentation for the generative transcriptomics study."""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[3]
FINAL_DIR = ROOT / "presentation/final"
PACKAGE_DIR = FINAL_DIR / "source"
TEMPLATE = PACKAGE_DIR / "SLSTP_template_2026.pptx"
OUTPUT = FINAL_DIR / "SLSTP_2026_Generative_Transcriptomics.pptx"
ASSET_DIR = PACKAGE_DIR / "assets"
NOTES_PATH = FINAL_DIR / "speaker_notes.md"
PAPER_DIR = ROOT / "paper/synthetic_guided_spaceflight"
EXPIMAP_PAPER_DIR = ROOT / "paper/asgsr_expimap_hvg"

SLIDE_W = 13.333333
SLIDE_H = 7.5
FONT = "Arial"

NAVY = "082B55"
BLUE = "2D6496"
TEAL = "178681"
GREEN = "478A64"
ORANGE = "D37B00"
GOLD = "9A6B20"
CORAL = "D96552"
PURPLE = "8063A6"
DARK = "263746"
GRAY = "59697A"
MID_GRAY = "8A969E"
LIGHT = "F2F5F7"
PALE_BLUE = "E9F1F7"
PALE_TEAL = "E8F3F1"
PALE_GOLD = "F8F1E2"
PALE_CORAL = "F8ECE9"
WHITE = "FFFFFF"


@dataclass(frozen=True)
class SlideNote:
    number: int
    title: str
    time: str
    text: str


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def _diverging_color(value: float, limit: float) -> str:
    position = np.clip((float(value) / float(limit) + 1.0) / 2.0, 0.0, 1.0)
    colormap = matplotlib.colors.LinearSegmentedColormap.from_list(
        "program_score", [f"#{BLUE}", f"#{WHITE}", "#C9342F"]
    )
    return matplotlib.colors.to_hex(colormap(position)).lstrip("#").upper()


def _set_fill(shape, color: str, transparency: int | None = None) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    if transparency is not None:
        color_nodes = shape._element.spPr.xpath("./a:solidFill/a:srgbClr")
        if color_nodes:
            alpha = OxmlElement("a:alpha")
            alpha.set("val", str(int((100 - transparency) * 1000)))
            color_nodes[0].append(alpha)


def _set_line(shape, color: str, width: float = 1.0) -> None:
    shape.line.color.rgb = _rgb(color)
    shape.line.width = Pt(width)


def _add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: str = DARK,
    bold: bool = False,
    italic: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.02,
    name: str | None = None,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color)
    return shape


def _add_rich_text(
    slide,
    segments: Iterable[tuple[str, dict]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: str = DARK,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.02,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0
    for text, style in segments:
        run = paragraph.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(style.get("size", size))
        run.font.bold = style.get("bold", False)
        run.font.italic = style.get("italic", False)
        run.font.color.rgb = _rgb(style.get("color", color))
    return shape


def _add_panel(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = LIGHT,
    line: str = "DDE4E8",
    radius: bool = True,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    _set_fill(shape, fill)
    _set_line(shape, line, 0.8)
    return shape


def _add_rule(slide, x: float, y: float, w: float, color: str, height: float = 0.04):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(height)
    )
    _set_fill(shape, color)
    shape.line.fill.background()
    return shape


def _add_circle(slide, x: float, y: float, diameter: float, color: str):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x),
        Inches(y),
        Inches(diameter),
        Inches(diameter),
    )
    _set_fill(shape, color)
    shape.line.fill.background()
    return shape


def _add_data_badge(slide, label: str, x: float, y: float, color: str, diameter: float = 0.24):
    _add_circle(slide, x, y, diameter, color)
    _add_text(
        slide,
        label,
        x,
        y + 0.005,
        diameter,
        diameter - 0.01,
        size=8.2,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )


def _add_arrow(slide, x: float, y: float, w: float, h: float, color: str = MID_GRAY):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    _set_fill(shape, color)
    shape.line.fill.background()
    return shape


def _add_bullet_rows(
    slide,
    rows: list[str],
    x: float,
    y: float,
    w: float,
    *,
    size: float = 16,
    color: str = DARK,
    bullet_color: str = TEAL,
    row_h: float = 0.48,
):
    for index, row in enumerate(rows):
        top = y + index * row_h
        _add_circle(slide, x, top + 0.10, 0.10, bullet_color)
        _add_text(
            slide,
            row,
            x + 0.20,
            top,
            w - 0.20,
            row_h,
            size=size,
            color=color,
            valign=MSO_ANCHOR.TOP,
            margin=0,
        )


def _add_picture_contain(
    slide,
    path: Path,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    alt: str,
):
    with Image.open(path) as image:
        aspect = image.width / image.height
    box_aspect = w / h
    if aspect >= box_aspect:
        width = w
        height = w / aspect
        left = x
        top = y + (h - height) / 2
    else:
        height = h
        width = h * aspect
        left = x + (w - width) / 2
        top = y
    shape = slide.shapes.add_picture(
        str(path), Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.name = alt
    shape._element.nvPicPr.cNvPr.set("descr", alt)
    return shape


def _add_full_slide_image(slide, path: Path, *, alt: str):
    shape = slide.shapes.add_picture(
        str(path),
        Inches(0),
        Inches(0),
        Inches(SLIDE_W),
        Inches(SLIDE_H),
    )
    shape.name = alt
    shape._element.nvPicPr.cNvPr.set("descr", alt)
    return shape


def _add_slide_title(slide, eyebrow: str, title: str, subtitle: str | None = None):
    _add_text(
        slide,
        eyebrow.upper(),
        0.32,
        0.58,
        5.5,
        0.25,
        size=10.5,
        color=GOLD,
        bold=True,
    )
    _add_text(
        slide,
        title,
        0.32,
        0.86,
        12.5,
        0.64,
        size=27,
        color=NAVY,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    if subtitle:
        _add_text(
            slide,
            subtitle,
            0.34,
            1.48,
            12.3,
            0.34,
            size=14.5,
            color=GRAY,
            valign=MSO_ANCHOR.MIDDLE,
        )


def _prepare_content_slide(slide, number: int):
    for placeholder in slide.placeholders:
        kind = placeholder.placeholder_format.type
        if kind in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.OBJECT):
            placeholder.text = ""
        elif kind == PP_PLACEHOLDER.FOOTER:
            placeholder.text = ""
        elif kind == PP_PLACEHOLDER.SLIDE_NUMBER:
            placeholder.text = str(number)
            for paragraph in placeholder.text_frame.paragraphs:
                paragraph.space_before = Pt(0)
                paragraph.space_after = Pt(0)
                paragraph.line_spacing = 1.0
                for run in paragraph.runs:
                    run.font.name = FONT
    return slide


def _add_source(slide, text: str):
    _add_text(
        slide,
        text,
        0.36,
        7.12,
        12.0,
        0.19,
        size=7.2,
        color=MID_GRAY,
        valign=MSO_ANCHOR.MIDDLE,
    )


def _build_tissue_utility_chart() -> Path:
    data = pd.read_csv(
        PAPER_DIR / "source_data/table_s18_matched_all_gene_utility.tsv",
        sep="\t",
    )
    data = data.loc[data["arm"].eq("real_plus_generated")].copy()
    canonical_order = [
        "adrenal_gland",
        "bone",
        "bone_marrow",
        "brain",
        "brown_adipose_tissue",
        "cecum",
        "cerebellum",
        "colon",
        "eye",
        "heart",
        "hippocampus",
        "kidney",
        "liver",
        "lung",
        "mammary_gland",
        "optic_nerve",
        "retina",
        "skeletal_muscle",
        "skin",
        "spleen",
        "thymus",
        "white_adipose_tissue",
    ]
    muscle_order = [
        "edl",
        "gastrocnemius",
        "quadriceps",
        "soleus",
        "tibialis_anterior",
    ]
    order = canonical_order + muscle_order
    if len(data) != len(order) or set(data["tissue"]) != set(order):
        raise ValueError("Expected 22 canonical tissues and five muscle groups")
    data = data.set_index("tissue").loc[order].reset_index()
    display = {
        "adrenal_gland": "Adrenal gland",
        "bone_marrow": "Bone marrow",
        "brown_adipose_tissue": "Brown adipose",
        "mammary_gland": "Mammary gland",
        "optic_nerve": "Optic nerve",
        "skeletal_muscle": "Skeletal muscle, pooled",
        "white_adipose_tissue": "White adipose",
        "edl": "EDL",
        "tibialis_anterior": "Tibialis anterior",
    }
    frame = data.reset_index(drop=True)
    fig, ax = plt.subplots(figsize=(12.4, 5.55))
    fig.patch.set_alpha(0)
    ax.set_facecolor("none")
    y = np.arange(len(frame))
    real = frame["real_mean_balanced_accuracy"].to_numpy(float)
    augmented = frame["arm_mean_balanced_accuracy"].to_numpy(float)
    delta = frame["mean_delta_balanced_accuracy"].to_numpy(float)
    passed = frame["joint_mean_all_metrics_nonworse"].astype(bool).to_numpy()
    colors = np.where(passed, f"#{TEAL}", f"#{CORAL}")
    for row_y, start, end in zip(y, real, augmented):
        ax.plot(
            [start, end],
            [row_y, row_y],
            color="#AAB5BC",
            lw=1.55,
            zorder=1,
        )
    ax.scatter(
        real,
        y,
        s=27,
        color=f"#{MID_GRAY}",
        edgecolors="white",
        linewidths=0.4,
        zorder=3,
        label="Real only",
    )
    ax.scatter(
        augmented[passed],
        y[passed],
        s=39,
        color=f"#{TEAL}",
        edgecolors="white",
        linewidths=0.5,
        zorder=4,
        label="Real + synthetic, all metrics pass",
    )
    ax.scatter(
        augmented[~passed],
        y[~passed],
        s=39,
        color=f"#{CORAL}",
        edgecolors="white",
        linewidths=0.5,
        zorder=4,
        label="Real + synthetic, mixed result",
    )
    for row_y, change, color in zip(y, delta, colors):
        text = "0" if abs(change) < 0.0005 else f"{change:+.2f}"
        ax.text(
            1.055,
            row_y,
            text,
            ha="center",
            va="center",
            fontsize=7.0,
            color=color,
            weight="bold",
        )
    labels = [
        display.get(value, value.replace("_", " ").title())
        for value in frame["tissue"].astype(str)
    ]
    ax.set_yticks(y, labels)
    ax.invert_yaxis()
    ax.set_xlim(0.42, 1.10)
    ax.set_xticks(np.arange(0.5, 1.01, 0.1))
    ax.grid(axis="x", color="#DDE4E8", lw=0.75)
    ax.axvline(1.02, color="#D4DBDF", lw=0.7)
    ax.text(
        1.055,
        -1.05,
        "CHANGE",
        ha="center",
        va="center",
        fontsize=7.2,
        color=f"#{GRAY}",
        weight="bold",
    )
    muscle_positions = np.flatnonzero(frame["scope"].eq("muscle_group"))
    if len(muscle_positions):
        ax.axhline(
            float(muscle_positions[0]) - 0.5,
            color=f"#{BLUE}",
            lw=1.0,
        )
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.spines["bottom"].set_color("#AAB5BC")
    ax.tick_params(axis="y", length=0, labelsize=7.4, colors=f"#{DARK}")
    ax.tick_params(axis="x", labelsize=8.0, colors=f"#{GRAY}")
    ax.set_xlabel(
        "Balanced accuracy for FLT versus GC on held-out real samples",
        fontsize=9.3,
        color=f"#{GRAY}",
        labelpad=7,
    )
    fig.subplots_adjust(left=0.19, right=0.97, top=0.99, bottom=0.11)
    path = ASSET_DIR / "tissue_flt_gc_balanced_accuracy_all27.png"
    fig.savefig(path, dpi=240, transparent=True, bbox_inches="tight")
    plt.close(fig)
    return path


def _build_trajectory_crop() -> Path:
    source = PAPER_DIR / "figures/figure_2a_archs4_denoising_trajectory.png"
    with Image.open(source) as image:
        top = int(image.height * 0.07)
        bottom = int(image.height * 0.84)
        cropped = image.crop((0, top, image.width, bottom))
        path = ASSET_DIR / "ddim_trajectory_panels.png"
        cropped.save(path)
    return path


def _build_pca_comparison_chart(primary: str, *, include_study: bool = True) -> Path:
    if primary not in {"tissue", "condition"}:
        raise ValueError(f"Unsupported PCA comparison: {primary!r}")
    source = pd.read_csv(
        PACKAGE_DIR / "source_data/locked_pca_by_accession.tsv",
        sep="\t",
    )
    required = {"source", primary, "accession", "pc1", "pc2"}
    missing = required.difference(source.columns)
    if missing:
        raise ValueError(
            "PCA source data are missing columns: "
            + ", ".join(sorted(missing))
        )
    accessions = sorted(
        source["accession"].astype(str).unique(),
        key=lambda value: int(value.rsplit("-", 1)[-1]),
    )
    hues = (np.arange(len(accessions)) * 0.61803398875) % 1.0
    palette = {
        accession: matplotlib.colors.hsv_to_rgb((hue, 0.66, 0.82))
        for accession, hue in zip(accessions, hues)
    }
    if primary == "condition":
        primary_levels = ["flight", "ground_control"]
        primary_palette = {
            "flight": "#D96552",
            "ground_control": "#2D6496",
        }
        primary_labels = {
            "flight": "Flight",
            "ground_control": "Ground control",
        }
    else:
        primary_levels = sorted(source[primary].astype(str).unique())
        tissue_cmap = plt.get_cmap("tab20", max(len(primary_levels), 1))
        primary_palette = {
            level: tissue_cmap(index)
            for index, level in enumerate(primary_levels)
        }
        primary_labels = {
            level: level.replace("_", " ").title()
            for level in primary_levels
        }
    x_pad = 0.04 * (source["pc1"].max() - source["pc1"].min())
    y_pad = 0.06 * (source["pc2"].max() - source["pc2"].min())
    x_limits = (source["pc1"].min() - x_pad, source["pc1"].max() + x_pad)
    y_limits = (source["pc2"].min() - y_pad, source["pc2"].max() + y_pad)

    if include_study:
        figure, axes = plt.subplots(1, 2, figsize=(12.4, 4.45), sharex=True, sharey=True)
    else:
        figure, axis = plt.subplots(1, 1, figsize=(6.0, 4.45))
        axes = [axis]
    figure.patch.set_alpha(0)

    def add_grouped_points(
        axis,
        *,
        column: str,
        levels: list[str],
        colors: dict[str, object],
        labels: dict[str, str],
        legend: bool,
    ) -> None:
        for level in levels:
            group = source.loc[source[column].astype(str).eq(level)]
            real = group.loc[group["source"].eq("real")]
            synthetic = group.loc[group["source"].eq("synthetic")]
            axis.scatter(
                real["pc1"],
                real["pc2"],
                s=24,
                marker="o",
                color=colors[level],
                alpha=0.42,
                linewidths=0,
                edgecolors="none",
            )
            axis.scatter(
                synthetic["pc1"],
                synthetic["pc2"],
                s=28,
                marker="x",
                color=colors[level],
                alpha=0.82,
                linewidths=1.0,
                label=labels[level],
            )
        if legend:
            axis.legend(
                frameon=False,
                fontsize=6.0 if column == "tissue" else 8.5,
                ncol=2 if column == "tissue" else 1,
                loc="upper right",
                borderaxespad=0.25,
                handletextpad=0.3,
                columnspacing=0.7,
            )

    add_grouped_points(
        axes[0],
        column=primary,
        levels=primary_levels,
        colors=primary_palette,
        labels=primary_labels,
        legend=True,
    )
    if include_study:
        add_grouped_points(
            axes[1],
            column="accession",
            levels=accessions,
            colors=palette,
            labels={accession: accession for accession in accessions},
            legend=False,
        )
    titles = {
        "tissue": "Colored by tissue",
        "condition": "Colored by condition",
    }
    panel_titles = [titles[primary], "Colored by study"] if include_study else [titles[primary]]
    for axis, title in zip(axes, panel_titles):
        axis.set_facecolor("none")
        axis.set_title(title, fontsize=14, color="#082B55", weight="bold", pad=9)
        axis.set_xlabel("PC1", fontsize=11, color="#263746")
        axis.set_xlim(x_limits)
        axis.set_ylim(y_limits)
        axis.grid(color="#DDE4E8", linewidth=0.75, alpha=0.78)
        axis.tick_params(labelsize=8.5, colors="#59697A")
        axis.spines[["top", "right"]].set_visible(False)
        axis.spines[["left", "bottom"]].set_color("#AEB9BF")
    axes[0].set_ylabel("PC2", fontsize=11, color="#263746")
    if include_study:
        figure.subplots_adjust(left=0.06, right=0.99, top=0.90, bottom=0.13, wspace=0.16)
        suffix = "vs_accession"
    else:
        figure.subplots_adjust(left=0.07, right=0.985, top=0.89, bottom=0.14)
        suffix = "only"
    path = ASSET_DIR / f"locked_pca_{primary}_{suffix}.png"
    figure.savefig(path, dpi=240, transparent=True)
    plt.close(figure)
    return path


def _set_title_slide(slide):
    main = next(
        placeholder
        for placeholder in slide.placeholders
        if placeholder.placeholder_format.type == PP_PLACEHOLDER.BODY
    )
    main.text_frame.clear()
    main.text_frame.word_wrap = True
    main.text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
    paragraph = main.text_frame.paragraphs[0]
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0
    run = paragraph.add_run()
    run.text = "Interpretable and generative models\nfor mouse spaceflight"
    run.font.name = FONT
    run.font.size = Pt(29)
    run.font.bold = True
    run.font.color.rgb = _rgb(WHITE)

    subtitle = next(
        placeholder
        for placeholder in slide.placeholders
        if placeholder.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE
    )
    subtitle.text_frame.clear()
    subtitle.text_frame.word_wrap = True
    entries = [
        ("Jason Trinh", 21, True, WHITE),
        ("EECS | UC Berkeley", 14, False, WHITE),
        ("", 6, False, WHITE),
        ("Mentor: James Casaletto | August 2026", 13, False, WHITE),
    ]
    for index, (text, size, bold, color) in enumerate(entries):
        paragraph = subtitle.text_frame.paragraphs[0] if index == 0 else subtitle.text_frame.add_paragraph()
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(2)
        paragraph.line_spacing = 1.0
        run = paragraph.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color)


def _slide_project_scope(slide):
    _add_slide_title(
        slide,
        "Project",
        "One dataset, two machine-learning questions",
        "We used mouse bulk RNA-seq to study how spaceflight changes genes, pathways, and tissue state.",
    )

    _add_panel(slide, 0.55, 2.05, 12.24, 0.76, fill=NAVY, line=NAVY, radius=False)
    _add_text(slide, "NASA OSDR API", 0.86, 2.20, 1.78, 0.31, size=16, color=WHITE, bold=True, margin=0)
    _add_text(
        slide,
        "Mus musculus bulk RNA-seq | spaceflight (FLT) and ground control (GC) | multiple missions and tissues",
        2.72,
        2.18,
        9.58,
        0.34,
        size=13.2,
        color="DCE7F2",
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )

    panels = [
        (
            0.55,
            "01",
            "Which biological programs change?",
            "expiMap",
            "Train a pathway-constrained reference on ARCHS4, then map OSDR samples into mouse Reactome programs.",
            "Output: tissue-specific pathway shifts",
            BLUE,
            PALE_BLUE,
        ),
        (
            6.83,
            "02",
            "Can synthetic profiles sharpen the analysis?",
            "Conditional generation",
            "Generate condition-specific expression profiles by tissue, study, and FLT or GC condition.",
            "Output: validated profiles and gene ranking",
            TEAL,
            PALE_TEAL,
        ),
    ]
    for x, number, question, method, body, output, color, fill in panels:
        _add_panel(slide, x, 3.22, 5.96, 3.20, fill=fill, line=fill, radius=False)
        _add_text(slide, number, x + 0.24, 3.48, 0.45, 0.25, size=11, color=color, bold=True, margin=0)
        _add_text(slide, question, x + 0.82, 3.40, 4.76, 0.50, size=17.2, color=NAVY, bold=True, margin=0)
        _add_rule(slide, x + 0.25, 4.12, 5.43, color, 0.025)
        _add_text(slide, method, x + 0.25, 4.37, 5.30, 0.34, size=15.4, color=color, bold=True, margin=0)
        _add_text(slide, body, x + 0.25, 4.82, 5.30, 0.66, size=13.0, color=DARK, margin=0)
        _add_text(slide, output, x + 0.25, 5.78, 5.30, 0.31, size=12.3, color=NAVY, bold=True, margin=0)
    _add_source(slide, "Data source: NASA Open Science Data Repository Biological Data API.")


def _slide_autoencoder_foundation(slide):
    _add_slide_title(
        slide,
        "Neural networks",
        "Autoencoders compress thousands of genes into a few features",
        "The encoder summarizes a gene profile; the decoder learns enough structure to reconstruct it.",
    )

    _add_text(slide, "AUTOENCODER", 0.58, 2.05, 1.44, 0.24, size=10.2, color=BLUE, bold=True, margin=0)
    _add_panel(slide, 0.50, 2.36, 7.38, 3.88, fill="F7F9FA", line="D9E1E5", radius=False)

    def add_trapezoid(points, x, y, fill, line):
        builder = slide.shapes.build_freeform(
            points[0][0],
            points[0][1],
            scale=Inches(0.01),
        )
        builder.add_line_segments(points[1:], close=True)
        shape = builder.convert_to_shape(Inches(x), Inches(y))
        _set_fill(shape, fill)
        _set_line(shape, line, 1.0)

    diagram_dx = -0.45
    add_trapezoid(
        [(0, 0), (120, 25), (120, 255), (0, 280)],
        3.15 + diagram_dx,
        2.66,
        "EAF1F5",
        "B8CCD8",
    )
    add_trapezoid(
        [(0, 25), (120, 0), (120, 280), (0, 255)],
        5.54 + diagram_dx,
        2.66,
        "E8F3F1",
        "B6D3CD",
    )

    gene_values = [0.58, 1.07, 0.73, 1.32, 0.90, 0.46]
    for index, height in enumerate(gene_values):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.64 + diagram_dx + index * 0.19),
            Inches(4.72 - height),
            Inches(0.11),
            Inches(height),
        )
        _set_fill(bar, BLUE if index % 2 == 0 else "8FB2CC")
        bar.line.fill.background()
    _add_rule(slide, 1.56 + diagram_dx, 4.72, 1.20, "AAB8C1", 0.018)
    _add_text(slide, "gene profile", 1.49 + diagram_dx, 4.96, 1.34, 0.28, size=10.8, color=GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_arrow(slide, 2.86 + diagram_dx, 3.86, 0.34, 0.25, MID_GRAY)

    encoder_positions = [
        [(3.43 + diagram_dx, 3.08 + row * (2.02 / 4)) for row in range(5)],
        [(3.93 + diagram_dx, 3.08 + row * (2.02 / 2)) for row in range(3)],
    ]
    latent_positions = [(4.98 + diagram_dx, y + 0.10) for y in [3.44, 3.88, 4.32]]
    decoder_positions = [
        [(5.96 + diagram_dx, 3.08 + row * (2.02 / 2)) for row in range(3)],
        [(6.46 + diagram_dx, 3.08 + row * (2.02 / 4)) for row in range(5)],
    ]

    def connect_layers(source, target):
        for x1, y1 in source:
            for x2, y2 in target:
                line = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    Inches(x1),
                    Inches(y1),
                    Inches(x2),
                    Inches(y2),
                )
                _set_line(line, "C5D1D8", 0.65)

    connect_layers(encoder_positions[0], encoder_positions[1])
    connect_layers(encoder_positions[1], latent_positions)
    connect_layers(latent_positions, decoder_positions[0])
    connect_layers(decoder_positions[0], decoder_positions[1])

    for col_index, positions in enumerate(encoder_positions):
        for center_x, center_y in positions:
            _add_circle(slide, center_x - 0.08, center_y - 0.08, 0.16, "6D97B5" if col_index == 0 else TEAL)
    _add_text(slide, "encoder", 3.18 + diagram_dx, 5.35, 1.05, 0.24, size=10.2, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)

    for index, (center_x, center_y) in enumerate(latent_positions):
        _add_circle(slide, center_x - 0.10, center_y - 0.10, 0.20, ORANGE if index == 1 else GOLD)
    _add_text(slide, "compressed\nfeatures", 4.51 + diagram_dx, 4.88, 0.94, 0.48, size=9.8, color=ORANGE, bold=True, align=PP_ALIGN.CENTER, margin=0)

    for col_index, positions in enumerate(decoder_positions):
        for center_x, center_y in positions:
            _add_circle(slide, center_x - 0.08, center_y - 0.08, 0.16, TEAL if col_index == 0 else "6D97B5")
    _add_text(slide, "decoder", 5.70 + diagram_dx, 5.35, 1.05, 0.24, size=10.2, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_arrow(slide, 6.82 + diagram_dx, 3.86, 0.30, 0.25, MID_GRAY)
    reconstructed_values = [0.56, 1.03, 0.76, 1.28, 0.86, 0.49]
    for index, height in enumerate(reconstructed_values):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(7.12 + diagram_dx + index * 0.11),
            Inches(4.72 - height),
            Inches(0.065),
            Inches(height),
        )
        _set_fill(bar, BLUE if index % 2 == 0 else "8FB2CC")
        bar.line.fill.background()
    _add_rule(slide, 7.06 + diagram_dx, 4.72, 0.72, "AAB8C1", 0.018)
    _add_text(slide, "output profile\n(reconstruction)", 6.88 + diagram_dx, 4.92, 1.08, 0.48, size=9.5, color=GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)

    _add_text(slide, "COMPRESSED SPACE", 8.30, 2.05, 1.78, 0.24, size=10.2, color=TEAL, bold=True, margin=0)
    _add_panel(slide, 8.22, 2.36, 4.58, 3.88, fill=PALE_TEAL, line="C9DFDB", radius=False)
    _add_rule(slide, 8.72, 5.55, 3.43, MID_GRAY, 0.018)
    _add_rule(slide, 8.72, 2.86, 0.018, MID_GRAY, 2.71)
    clusters = [
        (9.33, 4.63, BLUE),
        (10.47, 3.66, ORANGE),
        (11.48, 4.78, TEAL),
    ]
    offsets = [(-0.19, -0.07), (0.08, -0.16), (0.20, 0.05), (-0.06, 0.18), (0.02, 0.02)]
    for cx, cy, color in clusters:
        for dx, dy in offsets:
            _add_circle(slide, cx + dx, cy + dy, 0.13, color)
    _add_text(slide, "profiles with similar expression occupy nearby regions", 8.55, 5.76, 3.90, 0.24, size=10.5, color=GRAY, align=PP_ALIGN.CENTER, margin=0)

    _add_source(slide, "Concept adapted from the midpoint presentation. expiMap: Lotfollahi et al., Nature Cell Biology (2023).")


def _slide_expimap_program_scores(slide, output_heatmap: Path):
    _add_slide_title(
        slide,
        "expiMap scores",
        "Program scores summarize pathway changes within one tissue",
    )
    examples = [
        ("Extracellular matrix", 0.47, 0.08),
        ("Cell-cycle control", 0.42, 0.10),
        ("Immune signaling", 0.15, 0.13),
        ("Oxidative metabolism", -0.06, -0.03),
        ("DNA repair", -0.10, 0.34),
    ]

    _add_panel(slide, 0.47, 2.04, 8.05, 4.55, fill="F7F9FA", line="DDE4E8", radius=False)
    headers = [
        (3.35, "FLT score", DARK),
        (5.15, "GC score", DARK),
        (6.95, "FLT - GC", DARK),
    ]
    for x, label, color in headers:
        _add_text(slide, label, x, 2.47, 1.10, 0.26, size=9.8, color=color, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "-", 4.66, 2.47, 0.22, 0.26, size=13, color=MID_GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "=", 6.46, 2.47, 0.22, 0.26, size=12, color=MID_GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)

    for index, (program, flight, ground_control) in enumerate(examples):
        y = 2.86 + index * 0.63
        if index:
            _add_rule(slide, 0.77, y - 0.12, 7.45, "E1E6E9", 0.012)
        _add_rule(slide, 0.77, y + 0.05, 0.06, TEAL, 0.48)
        _add_text(slide, program, 1.02, y - 0.01, 2.10, 0.48, size=10.4, color=DARK, bold=True, valign=MSO_ANCHOR.MIDDLE, margin=0)
        values = [flight, ground_control, flight - ground_control]
        for x, value in zip((3.35, 5.15, 6.95), values):
            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x),
                Inches(y - 0.03),
                Inches(1.10),
                Inches(0.48),
            )
            _set_fill(cell, _diverging_color(value, 0.55))
            _set_line(cell, "D5DDE2", 0.6)
            _add_text(
                slide,
                f"{value:+.2f}",
                x,
                y + 0.06,
                1.10,
                0.24,
                size=10.2,
                color=WHITE if abs(value) > 0.32 else DARK,
                bold=True,
                align=PP_ALIGN.CENTER,
                margin=0,
            )
        _add_text(slide, "-", 4.66, y + 0.04, 0.22, 0.26, size=13, color=MID_GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)
        _add_text(slide, "=", 6.46, y + 0.04, 0.22, 0.26, size=12, color=MID_GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)

    legend_values = [(-1.0, "lower"), (0.0, "little change"), (1.0, "higher")]
    for index, (value, label) in enumerate(legend_values):
        x = 1.42 + index * 2.16
        square = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(6.18), Inches(0.25), Inches(0.18)
        )
        _set_fill(square, _diverging_color(value, 1.0))
        _set_line(square, "D5DDE2", 0.6)
        _add_text(slide, label, x + 0.31, 6.15, 1.05, 0.23, size=8.8, color=GRAY, margin=0)

    _add_panel(slide, 8.82, 2.04, 4.03, 4.55, fill="F7F9FA", line="DDE4E8", radius=False)
    _add_text(
        slide,
        "OUTPUT · PROGRAMS × SAMPLES",
        9.12,
        2.23,
        3.42,
        0.25,
        size=10.2,
        color=BLUE,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    _add_picture_contain(
        slide,
        output_heatmap,
        9.02,
        2.50,
        3.62,
        3.94,
        alt="Thymus expiMap output heatmap with Reactome programs as rows and samples as columns",
    )
    _add_source(slide, "Example values, ordered from higher to lower FLT - GC. Literature review appears next.")


def _slide_study_effect_umap(slide, umap_path: Path):
    _add_slide_title(
        slide,
        "The complication",
        "Study identity dominates the expression structure",
        "In EDL muscle, the two studies separate while FLT and GC overlap within each study.",
    )
    _add_picture_contain(
        slide,
        umap_path,
        1.15,
        1.82,
        11.05,
        5.15,
        alt="UMAP of EDL profiles colored by OSDR study and shaped by FLT or GC condition",
    )
    _add_source(slide, "Midpoint slide 7; EDL profiles from OSD-665 and OSD-99.")


def _slide_expimap_literature_review(
    slide,
    raw_heatmap: Path,
    reviewed_heatmap: Path,
):
    _add_slide_title(
        slide,
        "expiMap interpretation",
        "Each program, reviewed against prior literature",
        "The pathway scores stay fixed; the pathway-name colors summarize the literature comparison.",
    )

    _add_text(slide, "1  PATHWAY SCORES", 0.48, 1.98, 2.25, 0.24, size=10.3, color=BLUE, bold=True, margin=0)
    _add_picture_contain(
        slide,
        raw_heatmap,
        0.40,
        2.23,
        12.45,
        1.35,
        alt="Thymus expiMap pathway score heatmap before literature annotation",
    )

    _add_text(slide, "2  LITERATURE LABELS", 0.48, 3.77, 2.45, 0.24, size=10.3, color=TEAL, bold=True, margin=0)
    _add_picture_contain(
        slide,
        reviewed_heatmap,
        0.40,
        4.02,
        12.45,
        1.35,
        alt="The same thymus pathway score heatmap with literature-colored pathway names",
    )

    legend = [
        ("Aligned", "34A36F"),
        ("Complementary", "327AB5"),
        ("Uncertain", "E68D20"),
        ("Conflicting", "C63A2B"),
        ("Low effect", "96A1B3"),
    ]
    start_x = 1.03
    item_w = 2.28
    for index, (label, color) in enumerate(legend):
        x = start_x + index * item_w
        _add_rule(slide, x, 5.73, 0.25, color, 0.20)
        _add_text(slide, label, x + 0.35, 5.69, 1.77, 0.28, size=10.3, color=DARK, bold=True, margin=0)

    _add_text(
        slide,
        "Name color records the literature assessment; heatmap color records the program score.",
        1.02,
        6.20,
        11.30,
        0.30,
        size=11.5,
        color=GRAY,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    _add_source(slide, "Heatmaps reproduced from midpoint slides 11 and 12.")


def _slide_expimap_tissue_results(slide):
    _add_slide_title(
        slide,
        "expiMap results",
        "Five tissues showed recurring pathway patterns",
        "FLT and GC were compared within each study before the study-level changes were combined.",
    )
    rows = [
        (
            "Thymus",
            "117 samples | 5 projects",
            "LOWER IN FLIGHT",
            "DNA repair; RHOA cytoskeletal cycle; lymphoid-stromal interactions",
            "Known involution may also involve lower repair, motility, and niche coordination.",
            PURPLE,
            "F1EDF7",
        ),
        (
            "Skin",
            "151 samples | 4 projects",
            "LOWER IN FLIGHT",
            "Chromatin regulation; DNA repair; Hedgehog; sphingolipids; cell junctions",
            "Barrier injury may include a broader loss of tissue maintenance and coordination.",
            ORANGE,
            PALE_GOLD,
        ),
        (
            "Liver",
            "197 samples | 9 projects",
            "LOWER IN FLIGHT",
            "MHC class II antigen presentation; T-cell receptor signaling",
            "Metabolic heterogeneity coexisted with lower adaptive immune communication.",
            TEAL,
            PALE_TEAL,
        ),
        (
            "Spleen",
            "100 samples | 5 projects",
            "LOWER IN FLIGHT",
            "T-cell receptor signaling; neutrophil degranulation; C-type lectin signaling",
            "The most consistent multi-pathway result joined adaptive and innate immune changes.",
            CORAL,
            PALE_CORAL,
        ),
        (
            "Kidney",
            "135 samples | 6 projects",
            "HIGHER IN FLIGHT",
            "ECM proteoglycans; WNT signaling; IGF transport and uptake",
            "Matrix and growth-factor programs rose together, suggesting a renal remodeling response.",
            BLUE,
            PALE_BLUE,
        ),
    ]
    for index, (tissue, scope, direction, pathways, interpretation, color, fill) in enumerate(rows):
        y = 2.00 + index * 0.90
        _add_panel(slide, 0.48, y, 12.38, 0.80, fill=fill, line=fill, radius=False)
        _add_rule(slide, 0.48, y, 0.075, color, 0.80)
        _add_text(slide, tissue, 0.76, y + 0.11, 1.42, 0.28, size=16.0, color=color, bold=True, margin=0)
        _add_text(slide, scope, 0.76, y + 0.47, 1.58, 0.18, size=9.0, color=GRAY, margin=0)
        _add_text(slide, direction, 2.54, y + 0.10, 1.34, 0.20, size=8.6, color=CORAL if direction.startswith("HIGHER") else BLUE, bold=True, margin=0)
        _add_text(slide, pathways, 2.54, y + 0.33, 4.31, 0.34, size=11.2, color=DARK, bold=True, valign=MSO_ANCHOR.MIDDLE, margin=0)
        _add_rule(slide, 7.02, y + 0.10, 0.012, "D3DBDF", 0.58)
        _add_text(slide, interpretation, 7.30, y + 0.10, 5.18, 0.58, size=11.7, color=NAVY, valign=MSO_ANCHOR.MIDDLE, margin=0)
    _add_source(slide, "Source: expiMap manuscript Figures 3, 4, and 6. Program-score direction does not by itself prove biochemical activation or inhibition.")


def _slide_expimap_evidence(slide):
    _add_slide_title(
        slide,
        "expiMap evidence",
        "Cross-checks narrowed the biological interpretation",
        "Blue cells mark directional support. Conventional GSEA FDR is reported separately.",
    )
    retained = pd.read_csv(
        EXPIMAP_PAPER_DIR / "source_data/table_2_retained_pathway_evidence.tsv",
        sep="\t",
    )
    evidence = pd.read_csv(
        EXPIMAP_PAPER_DIR / "source_data/table_s24_pathway_robustness_evidence.tsv",
        sep="\t",
    )
    kidney_spleen_evidence = pd.read_csv(
        EXPIMAP_PAPER_DIR / "source_data/table_s27_kidney_spleen_pathway_evidence.tsv",
        sep="\t",
    )
    retained = retained.loc[retained["analysis_role"].eq("main")]
    columns = [
        "tissue",
        "term",
        "ssgsea_direction_support",
        "preranked_gsea_direction_support",
        "heldout_direction_support",
        "seed_direction_support",
        "composition_proxy_support",
    ]
    evidence = pd.concat(
        [evidence[columns], kidney_spleen_evidence[columns]],
        ignore_index=True,
    ).drop_duplicates(["tissue", "term"], keep="last")
    data = retained.merge(evidence, on=["tissue", "term"], how="left", validate="one_to_one")
    order = {"thymus": 0, "skin": 1, "liver": 2, "spleen": 3}
    data = data.assign(_order=data["tissue"].map(order)).sort_values(["_order", "display_label"])
    if len(data) != 13 or data[columns[2:]].isna().any().any():
        raise ValueError("Unexpected retained expiMap evidence table")

    label_x = 0.48
    label_w = 3.10
    check_x = 3.78
    cell_w = 0.60
    fdr_x = 6.93
    headers = ["ssGSEA", "GSEA", "Held-out", "3 seeds", "Composition"]
    _add_text(slide, "Tissue and pathway", label_x, 2.00, label_w, 0.24, size=9.5, color=GRAY, bold=True, margin=0)
    for index, header in enumerate(headers):
        _add_text(slide, header, check_x + index * cell_w - 0.05, 1.98, cell_w + 0.10, 0.32, size=8.2, color=GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "GSEA FDR", fdr_x, 2.00, 0.80, 0.24, size=8.8, color=GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_rule(slide, label_x, 2.34, 7.28, NAVY, 0.022)

    tissue_colors = {"thymus": PURPLE, "skin": ORANGE, "liver": TEAL, "spleen": CORAL}
    row_h = 0.31
    start_y = 2.46
    last_tissue = None
    for index, row in enumerate(data.itertuples(index=False)):
        y = start_y + index * row_h
        if index % 2 == 0:
            shade = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(label_x), Inches(y - 0.015), Inches(7.28), Inches(row_h - 0.01))
            _set_fill(shade, "F6F8F9")
            shade.line.fill.background()
        if last_tissue is not None and row.tissue != last_tissue:
            _add_rule(slide, label_x, y - 0.045, 7.28, "BBC7CE", 0.018)
        color = tissue_colors[row.tissue]
        label = f"{row.tissue.title()}: {row.display_label}"
        _add_rule(slide, label_x + 0.03, y + 0.105, 0.12, color, 0.045)
        _add_text(slide, label, label_x + 0.23, y + 0.035, label_w - 0.20, 0.23, size=8.8, color=DARK, margin=0)
        values = [
            row.ssgsea_direction_support,
            row.preranked_gsea_direction_support,
            row.heldout_direction_support,
            row.seed_direction_support,
            row.composition_proxy_support,
        ]
        for check_index, supported in enumerate(values):
            x = check_x + check_index * cell_w + 0.13
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + 0.025), Inches(0.34), Inches(0.24))
            _set_fill(box, BLUE if bool(supported) else "E7ECEF")
            box.line.fill.background()
            _add_text(slide, "+" if bool(supported) else "-", x, y + 0.025, 0.34, 0.24, size=9.0, color=WHITE if bool(supported) else MID_GRAY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
        fdr = float(row.gsea_fdr)
        fdr_label = "<0.001" if fdr < 0.001 else f"{fdr:.3f}"
        _add_text(slide, fdr_label, fdr_x, y + 0.035, 0.80, 0.22, size=8.8, color=NAVY if fdr < 0.05 else GRAY, bold=fdr < 0.05, align=PP_ALIGN.CENTER, margin=0)
        last_tissue = row.tissue

    _add_rule(slide, 8.02, 2.08, 0.015, "D5DDE2", 4.38)
    _add_text(slide, "WHAT REMAINED", 8.38, 2.10, 1.62, 0.24, size=10.2, color=TEAL, bold=True, margin=0)
    summaries = [
        ("Spleen", "All three programs passed GSEA FDR < 0.05 and all five directional checks.", CORAL),
        ("Skin", "Chromatin, DNA repair, and cell-junction programs passed GSEA FDR < 0.05.", ORANGE),
        ("Thymus", "DNA repair passed GSEA FDR < 0.001; the cytoskeletal and niche programs had directional support.", PURPLE),
        ("Liver", "MHC II and T-cell receptor scores were lower, but conventional FDR was 0.121 and 0.051.", TEAL),
    ]
    for index, (tissue, text, color) in enumerate(summaries):
        y = 2.59 + index * 0.90
        if index:
            _add_rule(slide, 8.38, y - 0.16, 4.02, "E0E5E8", 0.012)
        _add_text(slide, tissue, 8.38, y, 1.12, 0.26, size=13.5, color=color, bold=True, margin=0)
        _add_text(slide, text, 9.56, y - 0.02, 2.84, 0.63, size=10.8, color=DARK, margin=0)
    _add_panel(slide, 8.37, 6.19, 4.05, 0.46, fill=PALE_BLUE, line="C9DCE9", radius=False)
    _add_text(slide, "Spleen was the strongest multi-pathway result.", 8.58, 6.29, 3.62, 0.24, size=11.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_source(slide, "Source: expiMap manuscript Table 2 and Table S24. Held-out checks are internal project-wise validation.")


def _slide_why_synthetic(slide):
    _add_slide_title(
        slide,
        "Background",
        "What is synthetic transcriptomics?",
        "A generator learns patterns in measured RNA-seq data, then samples new expression vectors under chosen conditions.",
    )

    columns = [
        (
            0.50,
            "01",
            "Observed profiles",
            "Measured gene expression from real mouse tissue",
            BLUE,
            PALE_BLUE,
        ),
        (
            4.62,
            "02",
            "Learn expression patterns",
            "Capture gene relationships and variation by tissue and FLT/GC condition",
            ORANGE,
            PALE_GOLD,
        ),
        (
            8.74,
            "03",
            "Synthetic profiles",
            "Sample new numeric expression vectors for a selected context",
            TEAL,
            PALE_TEAL,
        ),
    ]
    for index, (x, number, heading, body, color, fill) in enumerate(columns):
        if index:
            _add_arrow(slide, x - 0.47, 4.00, 0.34, 0.28, MID_GRAY)
        _add_text(slide, number, x, 2.10, 0.42, 0.25, size=11, color=color, bold=True, margin=0)
        _add_text(slide, heading, x + 0.50, 2.04, 3.08, 0.38, size=17.2, color=NAVY, bold=True, margin=0)
        _add_panel(slide, x, 2.58, 3.58, 3.72, fill=fill, line=fill, radius=False)

        if index in (0, 2):
            cell_colors = [color, "AFC3D4", color, "D8E2E8", color]
            for row in range(5):
                for col in range(6):
                    cell_color = cell_colors[(row * 2 + col + index) % len(cell_colors)]
                    square = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(x + 0.36 + col * 0.40),
                        Inches(3.03 + row * 0.31),
                        Inches(0.29),
                        Inches(0.20),
                    )
                    _set_fill(square, cell_color)
                    square.line.fill.background()
            _add_text(slide, "genes", x + 0.36, 4.72, 2.29, 0.20, size=9.5, color=GRAY, align=PP_ALIGN.CENTER, margin=0)
            _add_text(slide, "profiles", x + 2.74, 3.62, 0.54, 0.22, size=9.5, color=GRAY, align=PP_ALIGN.CENTER, margin=0)
        else:
            node_positions = [
                (x + 0.36, 3.72),
                (x + 1.03, 3.20),
                (x + 1.03, 4.28),
                (x + 1.76, 3.74),
                (x + 2.48, 3.20),
                (x + 2.48, 4.28),
                (x + 3.08, 3.74),
            ]
            connections = [(0, 1), (0, 2), (1, 3), (2, 3), (3, 4), (3, 5), (4, 6), (5, 6)]
            for start, end in connections:
                x1, y1 = node_positions[start]
                x2, y2 = node_positions[end]
                line = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    Inches(x1 + 0.08),
                    Inches(y1 + 0.08),
                    Inches(x2 + 0.08),
                    Inches(y2 + 0.08),
                )
                _set_line(line, "C6B07A", 1.5)
            for node_index, (node_x, node_y) in enumerate(node_positions):
                _add_circle(slide, node_x, node_y, 0.20, ORANGE if node_index in (0, 3, 6) else GOLD)

        _add_rule(slide, x + 0.28, 5.08, 3.02, "D3DCE1", 0.018)
        _add_text(slide, body, x + 0.30, 5.30, 2.98, 0.72, size=14.2, color=DARK, valign=MSO_ANCHOR.MIDDLE, margin=0)


def _slide_scientific_objective(slide):
    _add_slide_title(
        slide,
        "Study goals",
        "Match tissue distributions, then test FLT versus GC biology",
        "The first goal validates the generator. The second asks whether synthetic data improves a real-data analysis.",
    )

    _add_text(slide, "01", 0.58, 2.08, 0.42, 0.26, size=11, color=TEAL, bold=True, margin=0)
    _add_text(slide, "Reproduce tissue structure", 1.08, 2.02, 4.50, 0.38, size=18, color=NAVY, bold=True, margin=0)
    _add_text(
        slide,
        "Synthetic bulk RNA-seq should occupy the same tissue-defined expression space as real profiles.",
        0.58,
        2.49,
        5.58,
        0.56,
        size=13.4,
        color=DARK,
        margin=0,
    )

    _add_rule(slide, 0.83, 4.58, 4.92, MID_GRAY, 0.018)
    _add_rule(slide, 0.83, 3.14, 0.018, MID_GRAY, 1.46)

    def add_cross(x, y, color):
        first = _add_rule(slide, x, y + 0.06, 0.18, color, 0.025)
        first.rotation = 45
        second = _add_rule(slide, x, y + 0.06, 0.18, color, 0.025)
        second.rotation = -45

    clusters = [
        (1.52, 3.98, TEAL, "Liver"),
        (3.12, 3.43, ORANGE, "Muscle"),
        (4.77, 4.03, PURPLE, "Thymus"),
    ]
    offsets = [(-0.23, -0.10), (0.05, -0.19), (0.22, 0.02), (-0.08, 0.18)]
    for center_x, center_y, color, label in clusters:
        for point_index, (dx, dy) in enumerate(offsets):
            _add_circle(slide, center_x + dx, center_y + dy, 0.13, color)
            add_cross(center_x + dx + 0.10, center_y + dy + (0.08 if point_index % 2 else -0.03), color)
        _add_text(slide, label, center_x - 0.37, center_y + 0.40, 0.98, 0.22, size=9.2, color=GRAY, align=PP_ALIGN.CENTER, margin=0)
    _add_circle(slide, 0.93, 3.05, 0.12, NAVY)
    _add_text(slide, "Real", 1.10, 3.02, 0.55, 0.20, size=9.2, color=GRAY, margin=0)
    add_cross(1.70, 3.03, NAVY)
    _add_text(slide, "Generated", 1.95, 3.02, 0.90, 0.20, size=9.2, color=GRAY, margin=0)

    _add_rule(slide, 6.50, 2.05, 0.015, "D6DEE3", 2.72)

    _add_text(slide, "02", 6.86, 2.08, 0.42, 0.26, size=11, color=CORAL, bold=True, margin=0)
    _add_text(slide, "Preserve the FLT/GC difference", 7.36, 2.02, 4.96, 0.38, size=18, color=NAVY, bold=True, margin=0)
    _add_text(
        slide,
        "Within each tissue, generated profiles should retain the smaller condition signal.",
        6.86,
        2.49,
        5.60,
        0.56,
        size=13.4,
        color=DARK,
        margin=0,
    )
    _add_text(slide, "FLT", 7.18, 3.24, 0.72, 0.28, size=14, color=CORAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "GC", 9.08, 3.24, 0.72, 0.28, size=14, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    flt_points = [(7.18, 3.73), (7.52, 3.58), (7.79, 3.81), (7.37, 4.04)]
    gc_points = [(9.09, 3.73), (9.43, 3.58), (9.70, 3.81), (9.28, 4.04)]
    for x, y in flt_points:
        _add_circle(slide, x, y, 0.17, CORAL)
    for x, y in gc_points:
        _add_circle(slide, x, y, 0.17, BLUE)
    _add_text(slide, "vs", 8.38, 3.75, 0.42, 0.24, size=11.5, color=MID_GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_arrow(slide, 9.98, 3.67, 0.31, 0.28, MID_GRAY)
    _add_panel(slide, 10.49, 3.30, 2.00, 1.14, fill=PALE_BLUE, line="D4E2EC", radius=False)
    _add_text(slide, "OUTPUT", 10.68, 3.47, 1.63, 0.18, size=9.2, color=BLUE, bold=True, margin=0)
    _add_text(
        slide,
        "Tissue-specific\ngenes and pathways",
        10.68,
        3.75,
        1.63,
        0.54,
        size=12.2,
        color=NAVY,
        bold=True,
        margin=0,
    )

    _add_rule(slide, 0.60, 5.12, 12.08, TEAL, 0.035)
    _add_text(slide, "Model choice", 0.62, 5.43, 1.55, 0.30, size=16.5, color=TEAL, bold=True, margin=0)
    _add_text(
        slide,
        "Use the pipeline that matches tissue structure and performs best on held-out real FLT/GC samples.",
        2.18,
        5.34,
        7.40,
        0.58,
        size=15,
        color=NAVY,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )
    _add_text(slide, "Gene effects and significance use observed OSDR data.", 9.88, 5.42, 2.57, 0.46, size=11.5, color=GRAY, margin=0)


def _slide_2(slide):
    _add_slide_title(
        slide,
        "Question",
        "Small studies and study effects complicate tissue comparisons",
        "Can a generator help rank spaceflight signal without pretending it creates new animals?",
    )
    _add_text(slide, "NASA OSDR", 0.62, 2.04, 2.20, 0.30, size=17, color=BLUE, bold=True)
    _add_text(slide, "Observed spaceflight cohort", 0.62, 2.37, 2.80, 0.24, size=11.5, color=GRAY)
    _add_text(slide, "1,610", 0.58, 2.64, 2.25, 0.69, size=42, color=NAVY, bold=True)
    _add_text(slide, "profiles", 0.64, 3.29, 1.15, 0.25, size=12.5, color=GRAY)
    _add_text(slide, "75", 3.01, 2.72, 1.20, 0.54, size=31, color=NAVY, bold=True)
    _add_text(slide, "OSDR studies", 3.04, 3.26, 1.40, 0.25, size=12.5, color=GRAY)

    bar_x = 0.64
    bar_y = 4.03
    bar_w = 5.24
    flt_w = bar_w * 835 / 1610
    _add_text(slide, "835 FLT", bar_x, 3.65, 1.10, 0.26, size=13.5, color=CORAL, bold=True)
    _add_text(slide, "775 GC", bar_x + bar_w - 1.10, 3.65, 1.10, 0.26, size=13.5, color=BLUE, bold=True, align=PP_ALIGN.RIGHT)
    _add_rule(slide, bar_x, bar_y, flt_w, CORAL, 0.12)
    _add_rule(slide, bar_x + flt_w, bar_y, bar_w - flt_w, BLUE, 0.12)

    _add_rule(slide, 6.49, 2.02, 0.018, "CBD4DA", 2.21)
    _add_text(slide, "ARCHS4 mouse", 6.88, 2.04, 2.40, 0.30, size=17, color=TEAL, bold=True)
    _add_text(slide, "Reference pretraining cohort", 6.88, 2.37, 2.80, 0.24, size=11.5, color=GRAY)
    _add_text(slide, "997,515", 6.84, 2.64, 2.55, 0.69, size=39, color=NAVY, bold=True)
    _add_text(slide, "profiles screened", 6.90, 3.29, 1.75, 0.25, size=12.5, color=GRAY)
    _add_arrow(slide, 9.26, 2.88, 0.42, 0.30, MID_GRAY)
    _add_text(slide, "17,244", 9.86, 2.66, 1.70, 0.54, size=29, color=NAVY, bold=True)
    _add_text(slide, "selected", 9.89, 3.24, 1.10, 0.25, size=12.5, color=GRAY)
    _add_text(slide, "20 tissues", 11.42, 2.73, 1.25, 0.36, size=17, color=TEAL, bold=True)

    _add_rule(slide, 0.62, 4.57, 12.08, "D6DEE3", 0.018)
    _add_text(slide, "Study effects can look like spaceflight biology.", 0.62, 4.82, 11.80, 0.42, size=22, color=NAVY, bold=True)
    challenge_rows = [
        ("01", "Preserve tissue and FLT/GC structure."),
        ("02", "Avoid memorizing individual profiles."),
        ("03", "Test biological effects in observed OSDR samples."),
    ]
    for index, (number, text) in enumerate(challenge_rows):
        y = 5.42 + index * 0.42
        _add_text(slide, number, 0.64, y, 0.42, 0.26, size=10.5, color=ORANGE, bold=True, margin=0)
        _add_text(slide, text, 1.17, y - 0.03, 10.95, 0.31, size=14.5, color=DARK, margin=0)
    _add_source(slide, "Sources: NASA OSDR Biological Data API; ARCHS4 mouse v2.5.")


def _add_midpoint_generator_slide(
    slide,
    tile_prefix: str,
    overlays: list[tuple[Path, tuple[int, int, int, int], str]],
    slide_number: int,
) -> None:
    source_width = 18_288_000
    source_height = 10_287_000
    tile_boxes = [
        (0, 0, 9_144_000, 5_143_500),
        (9_124_950, 0, 9_163_050, 5_143_500),
        (0, 5_124_450, 9_144_000, 5_162_550),
        (9_124_950, 5_124_450, 9_163_050, 5_162_550),
    ]

    def scaled_box(box: tuple[int, int, int, int]) -> tuple[float, float, float, float]:
        left, top, width, height = box
        return (
            left / source_width * SLIDE_W,
            top / source_height * SLIDE_H,
            width / source_width * SLIDE_W,
            height / source_height * SLIDE_H,
        )

    for index, box in enumerate(tile_boxes, start=1):
        path = ASSET_DIR / f"image-{tile_prefix}-{index}.jpeg"
        left, top, width, height = scaled_box(box)
        picture = slide.shapes.add_picture(
            str(path), Inches(left), Inches(top), Inches(width), Inches(height)
        )
        picture.name = f"Midpoint slide {tile_prefix} tile {index}"
        picture._element.nvPicPr.cNvPr.set(
            "descr", f"Original midpoint slide {tile_prefix}, quadrant {index}"
        )

    for path, box, alt in overlays:
        left, top, width, height = scaled_box(box)
        picture = slide.shapes.add_picture(
            str(path), Inches(left), Inches(top), Inches(width), Inches(height)
        )
        picture.name = path.name
        picture._element.nvPicPr.cNvPr.set("descr", alt)

    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(12.54),
        Inches(0.10),
        Inches(0.52),
        Inches(0.30),
    )
    _set_fill(badge, "082768")
    badge.line.fill.background()
    _add_text(
        slide,
        str(slide_number),
        12.54,
        0.135,
        0.52,
        0.21,
        size=11.5,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )


def _slide_expimap_method(slide) -> None:
    _add_midpoint_generator_slide(slide, "10", [], 7)

    cover = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.28),
        Inches(2.14),
        Inches(8.05),
        Inches(5.36),
    )
    _set_fill(cover, WHITE)
    cover.line.fill.background()

    _add_panel(slide, 6.10, 2.66, 6.28, 3.46, fill="F2F5F7", line="E4E9EC", radius=True)
    _add_text(slide, "WHY IT'S INTERPRETABLE", 6.47, 2.93, 4.55, 0.29, size=15.0, color=NAVY, bold=True, margin=0)
    rows = [
        "Each latent node is one Reactome program (~1,140).",
        "A soft mask lets a stray gene help only when the data support it.",
        "Compare FLT and GC one program at a time.",
    ]
    _add_bullet_rows(
        slide,
        rows,
        6.49,
        3.43,
        5.42,
        size=15.0,
        color=DARK,
        bullet_color=BLUE,
        row_h=0.76,
    )


def _slide_wgan_explainer(slide, wgan_gif: Path, slide_number: int):
    _add_midpoint_generator_slide(
        slide,
        "21",
        [
            (
                wgan_gif,
                (10_434_638, 3_146_108, 5_895_023, 6_550_343),
                "Animated WGAN comparison between real OSDR and generated expression profiles",
            )
        ],
        slide_number,
    )


def _slide_3(slide):
    _add_slide_title(
        slide,
        "Model development",
        "Building the RNA-seq generator",
        "We compared how to process the data, handle study differences, train the model, and tell it what to generate.",
    )

    def option_tile(label, x, y, w, color, selected=False, *, size=9.1):
        if selected:
            tile = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x),
                Inches(y),
                Inches(w),
                Inches(0.36),
            )
            _set_fill(tile, WHITE)
            _set_line(tile, color, 1.8)
            _add_rule(slide, x, y, 0.035, color, 0.36)
        _add_text(
            slide,
            label,
            x + 0.05,
            y + 0.055,
            w - 0.10,
            0.23,
            size=size,
            color=NAVY if selected else MID_GRAY,
            bold=selected,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
            margin=0,
        )

    def segmented_row(x, y, label, options, color):
        _add_text(slide, label, x + 0.16, y, 2.02, 0.18, size=8.4, color=GRAY, bold=True, margin=0)
        gap = 0.05
        available = 2.03
        tile_width = (available - gap * (len(options) - 1)) / len(options)
        for index, (option, selected) in enumerate(options):
            option_tile(
                option,
                x + 0.16 + index * (tile_width + gap),
                y + 0.22,
                tile_width,
                color,
                selected,
                size=8.2 if len(options) == 3 else 8.8,
            )

    def stage_panel(x, number, title, color, fill):
        _add_text(slide, number, x + 0.14, 2.15, 0.34, 0.24, size=11.5, color=MID_GRAY, bold=True, margin=0)
        _add_text(slide, title, x + 0.56, 2.11, 1.64, 0.32, size=14.2, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)
        _add_rule(slide, x + 0.14, 2.55, 2.05, color, 0.022)

    xs = [0.35, 2.90, 5.45, 8.00, 10.55]
    for separator_x in [2.79, 5.34, 7.89, 10.44]:
        _add_rule(slide, separator_x, 2.03, 0.012, "D9E0E4", 3.24)

    stage_panel(xs[0], "01", "Data scope", BLUE, WHITE)
    segmented_row(xs[0], 2.72, "Sources", [("OSDR API", True), ("ARCHS4", True)], BLUE)
    segmented_row(xs[0], 3.48, "Studies", [("Single", False), ("Multiple", True)], BLUE)
    segmented_row(xs[0], 4.24, "Tissues", [("Per tissue", False), ("All tissues", True)], BLUE)

    stage_panel(xs[1], "02", "Processing", TEAL, WHITE)
    segmented_row(xs[1], 2.72, "Expression", [("Raw", False), ("CPM", False), ("TPM", True)], TEAL)
    segmented_row(xs[1], 3.48, "Scaling", [("None", False), ("Z-score", False), ("MaxAbs", True)], TEAL)
    _add_text(slide, "Features", xs[1] + 0.16, 4.24, 2.02, 0.18, size=8.4, color=GRAY, bold=True, margin=0)
    feature_options = [
        ("All genes", False),
        ("HVGs", False),
        ("Reactome", False),
        ("L1000 map", True),
    ]
    for index, (label, selected) in enumerate(feature_options):
        option_tile(
            label,
            xs[1] + 0.16 + (index % 2) * 1.04,
            4.46 + (index // 2) * 0.40,
            0.99,
            TEAL,
            selected,
            size=7.7,
        )

    stage_panel(xs[2], "03", "Study effects", TEAL, WHITE)
    _add_text(slide, "Alternative methods", xs[2] + 0.16, 2.72, 2.02, 0.18, size=8.4, color=GRAY, bold=True, margin=0)
    for index, (label, selected) in enumerate([
        ("None", True),
        ("Within-study z-score", False),
        ("ComBat / MBatch", False),
        ("MOBER", False),
    ]):
        option_tile(label, xs[2] + 0.16, 2.94 + index * 0.50, 2.03, TEAL, selected, size=9.2)

    stage_panel(xs[3], "04", "Model training", ORANGE, WHITE)
    segmented_row(xs[3], 2.72, "Generator", [("WGAN-GP", False), ("DDIM", True)], ORANGE)
    _add_text(slide, "Training source", xs[3] + 0.16, 3.48, 2.02, 0.18, size=8.4, color=GRAY, bold=True, margin=0)
    for index, (label, selected) in enumerate([
        ("OSDR only", False),
        ("ARCHS4 only", False),
        ("ARCHS4 then OSDR", True),
    ]):
        option_tile(label, xs[3] + 0.16, 3.70 + index * 0.50, 2.03, ORANGE, selected, size=9.2)

    stage_panel(xs[4], "05", "Conditioning", TEAL, WHITE)
    _add_text(slide, "Model inputs", xs[4] + 0.16, 2.72, 2.02, 0.18, size=8.4, color=GRAY, bold=True, margin=0)
    for index, (label, selected) in enumerate([
        ("Tissue", True),
        ("FLT / GC", True),
        ("Study", True),
        ("Material type", True),
        ("Sex / age", False),
    ]):
        option_tile(label, xs[4] + 0.16, 2.92 + index * 0.42, 2.03, TEAL, selected, size=8.9)

    _add_panel(slide, 0.35, 5.57, 12.55, 1.23, fill="F7F9FA", line="F7F9FA", radius=False)
    _add_rule(slide, 0.35, 5.57, 12.55, "C9D3D9", 0.018)
    _add_text(slide, "Downstream", 0.64, 5.80, 1.42, 0.25, size=10.8, color=NAVY, bold=True)
    _add_text(slide, "configuration", 0.64, 6.17, 1.35, 0.24, size=9.6, color=GRAY)

    def selected_step(x, width, heading, value, color):
        _add_panel(slide, x, 5.72, width, 0.49, fill=WHITE, line=color, radius=False)
        _add_text(slide, heading, x + 0.10, 5.78, width - 0.20, 0.12, size=7.1, color=color, bold=True, margin=0)
        _add_text(slide, value, x + 0.10, 5.94, width - 0.20, 0.19, size=9.1, color=DARK, bold=True, align=PP_ALIGN.CENTER, margin=0)

    selected_path = [
        (2.10, 2.15, "PROCESS", "TPM / MaxAbs / 974 landmarks", TEAL),
        (4.65, 1.85, "PRETRAIN", "ARCHS4", TEAL),
        (6.90, 1.85, "ADAPT", "OSDR", BLUE),
        (9.15, 2.45, "GENERATE", "Conditional DDIM", ORANGE),
    ]
    for index, (x, width, heading, value, color) in enumerate(selected_path):
        selected_step(x, width, heading, value, color)
        if index < len(selected_path) - 1:
            _add_arrow(slide, x + width + 0.07, 5.86, 0.23, 0.20, MID_GRAY)

    _add_text(slide, "CONDITION", 2.10, 6.37, 0.72, 0.16, size=7.3, color=TEAL, bold=True, margin=0)
    condition_tiles = [
        (2.88, 0.90, "Tissue"),
        (3.85, 1.00, "FLT / GC"),
        (4.92, 1.14, "Study"),
        (6.13, 1.38, "Material type"),
    ]
    for x, width, label in condition_tiles:
        option_tile(label, x, 6.27, width, TEAL, True, size=7.8)

    _add_text(slide, "HARMONIZE", 7.76, 6.37, 0.82, 0.16, size=7.3, color=TEAL, bold=True, margin=0)
    option_tile("None", 8.65, 6.27, 0.78, TEAL, True, size=7.8)
    _add_text(slide, "SCOPE", 9.70, 6.37, 0.46, 0.16, size=7.3, color=BLUE, bold=True, margin=0)
    option_tile("All tissues", 10.24, 6.27, 1.20, BLUE, True, size=7.8)
    _add_source(slide, "Pipeline evaluated with OSDR and ARCHS4. Generator designs follow Vinas et al. (2022) and Lacan et al. (2026).")


def _slide_5(slide, trajectory: Path):
    _add_slide_title(
        slide,
        "Diffusion",
        "Diffusion learns tissue structure from noise",
        "The same PCA axes follow 1,024 generated profiles through reverse diffusion.",
    )
    _add_picture_contain(
        slide,
        trajectory,
        0.34,
        1.78,
        12.64,
        4.72,
        alt="DDIM reverse trajectory at timesteps 1000, 200 and 0",
    )
    _add_text(slide, "t = 1000: noise", 0.78, 6.50, 2.5, 0.28, size=13, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "t = 200: partial structure", 5.23, 6.50, 2.9, 0.28, size=13, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "t = 0: tissue-conditioned profiles", 9.15, 6.50, 3.4, 0.28, size=13, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    _add_source(slide, "Gray points are real ARCHS4 profiles; colors identify generated tissue conditions. PC1 and PC2 limits are shared across panels.")


def _slide_diffusion_explainer(
    slide,
    diffusion_image_gif: Path,
    diffusion_data_gif: Path,
    slide_number: int,
):
    _add_midpoint_generator_slide(
        slide,
        "20",
        [
            (
                diffusion_image_gif,
                (3_076_575, 5_393_055, 3_657_600, 2_057_400),
                "Animated image denoising example",
            ),
            (
                diffusion_data_gif,
                (10_487_978, 3_146_108, 5_788_343, 6_550_343),
                "Animated FLT- and GC-conditioned expression diffusion example",
            ),
        ],
        slide_number,
    )


def _slide_6(slide, tissue_accession_pca: Path):
    _add_slide_title(
        slide,
        "Diffusion output",
        "Tissue and study structure dominate the PCA space",
        "We therefore compare FLT and GC within tissues and account for study.",
    )
    _add_picture_contain(
        slide,
        tissue_accession_pca,
        0.38,
        1.93,
        12.52,
        4.38,
        alt="PCA of real and generated profiles colored by tissue and study",
    )
    _add_text(
        slide,
        "Left: tissue. Right: study. Circles are real samples; crosses are generated samples.",
        0.76,
        6.45,
        11.80,
        0.26,
        size=11.5,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )


def _slide_condition_pca(slide, condition_pca: Path):
    _add_slide_title(
        slide,
        "Diffusion output",
        "FLT and GC overlap in the global PCA view",
        "The condition signal is not clear when all tissues and studies are combined.",
    )
    _add_picture_contain(
        slide,
        condition_pca,
        0.62,
        1.93,
        12.08,
        4.38,
        alt="PCA of real and generated profiles colored by FLT and ground-control condition",
    )
    _add_text(
        slide,
        "Color: FLT or GC. Circles are real samples; crosses are generated samples.",
        0.76,
        6.45,
        11.80,
        0.26,
        size=11.5,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )


def _slide_4(slide, architecture_figure: Path):
    _add_slide_title(
        slide,
        "Validation",
        "Diffusion best reproduced the measured expression distribution",
        "Both models reproduced gene expression. Diffusion samples were harder to distinguish from measured samples.",
    )
    _add_text(slide, "Architecture from Lacan et al.", 0.52, 2.08, 7.12, 0.32, size=16.5, color=NAVY, bold=True)
    _add_picture_contain(
        slide,
        architecture_figure,
        0.20,
        2.39,
        7.48,
        2.73,
        alt="Lacan et al. Figure 1C residual diffusion generator architecture",
    )
    _add_text(slide, "The network repeatedly predicts and removes noise using tissue and time information.", 0.56, 5.17, 6.90, 0.28, size=11.3, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "For OSDR, the model also receives FLT/GC, study, and sample-material information.", 0.56, 5.51, 6.90, 0.34, size=10.6, color=GRAY, align=PP_ALIGN.CENTER)

    _add_rule(slide, 7.83, 2.05, 0.015, "D5DDE2", 4.32)
    _add_text(slide, "Metric", 8.04, 2.17, 1.42, 0.28, size=10.0, color=GRAY, bold=True)
    _add_text(slide, "WGAN-GP", 9.57, 2.17, 0.85, 0.28, size=10.0, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "DDIM", 10.52, 2.17, 0.76, 0.28, size=10.0, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "Read", 11.43, 2.17, 1.10, 0.28, size=10.0, color=GRAY, bold=True)
    _add_rule(slide, 8.04, 2.56, 4.58, NAVY, 0.022)
    metrics = [
        ("Expression corr.", "0.976", "0.974", "higher"),
        ("Coverage (F1)", "0.985", "0.997", "higher"),
        ("Real-vs-synth. acc.", "0.636", "0.475", "0.5 is ideal"),
        ("Distribution distance", "0.144", "0.074", "lower"),
    ]
    for index, (label, wgan, ddim, reading) in enumerate(metrics):
        y = 2.80 + index * 0.59
        if index % 2 == 0:
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.02), Inches(y - 0.07), Inches(4.61), Inches(0.50))
            _set_fill(shape, "F5F7F8")
            shape.line.fill.background()
        _add_text(slide, label, 8.10, y, 1.39, 0.32, size=10.3, color=DARK, bold=index >= 2)
        _add_text(slide, wgan, 9.57, y, 0.85, 0.32, size=12.6, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, ddim, 10.52, y, 0.76, 0.32, size=12.6, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, reading, 11.43, y, 1.10, 0.32, size=10.2, color=GRAY)

    _add_rule(slide, 0.55, 5.91, 7.03, "D5DDE2", 0.015)
    _add_text(slide, "Tissue identity retained", 0.55, 6.08, 1.90, 0.24, size=10.4, color=BLUE, bold=True, margin=0)
    _add_text(slide, "Real 0.781  |  Synthetic 0.781", 4.30, 6.06, 3.12, 0.27, size=10.4, color=DARK, align=PP_ALIGN.RIGHT, margin=0)
    _add_text(slide, "Balanced accuracy is preserved in generated tissue labels.", 0.55, 6.39, 6.87, 0.25, size=9.7, color=GRAY, align=PP_ALIGN.CENTER, margin=0)
    _add_rule(slide, 8.04, 5.36, 0.055, TEAL, 0.78)
    _add_text(slide, "Use diffusion", 8.29, 5.40, 1.55, 0.31, size=16, color=NAVY, bold=True, margin=0)
    _add_text(slide, "Better coverage and closer to the real distribution.", 9.94, 5.38, 2.57, 0.47, size=10.8, color=DARK, valign=MSO_ANCHOR.MIDDLE, margin=0)
    _add_text(slide, "A score near 0.5 means a prediction model cannot reliably tell measured and generated samples apart.", 8.29, 5.97, 4.22, 0.40, size=9.1, color=GRAY, margin=0)
    _add_source(slide, "Architecture: Lacan et al. (2026), Fig. 1C, doi:10.1186/s12859-026-06470-8. Each model was evaluated on its own fixed test split.")


def _slide_7(slide):
    _add_slide_title(
        slide,
        "Analysis",
        "Five arms separate gene ranking from classifier fitting",
        "Guided arms use synthetic profiles to help choose genes; they do not treat generated profiles as new animals.",
    )
    _add_text(slide, "INPUTS", 0.47, 2.00, 0.66, 0.22, size=8.8, color=GRAY, bold=True, margin=0)
    _add_data_badge(slide, "R", 1.20, 1.98, BLUE)
    _add_text(slide, "real OSDR", 1.52, 1.98, 1.18, 0.24, size=10.8, color=BLUE, bold=True, margin=0)
    _add_data_badge(slide, "S", 2.91, 1.98, CORAL)
    _add_text(slide, "matched DDIM", 3.23, 1.98, 1.42, 0.24, size=10.8, color=CORAL, bold=True, margin=0)
    _add_text(slide, "GUIDED", 6.10, 1.99, 0.70, 0.22, size=8.8, color=TEAL, bold=True, margin=0)
    _add_text(slide, "R + S help rank genes", 6.85, 1.97, 2.05, 0.25, size=10.5, color=DARK, margin=0)
    _add_text(slide, "5%", 9.41, 1.99, 0.38, 0.22, size=8.8, color=ORANGE, bold=True, margin=0)
    _add_text(slide, "S carries 5% of total fit weight", 9.84, 1.97, 2.78, 0.25, size=10.5, color=DARK, margin=0)

    headers = [
        ("ANALYSIS ARM", 0.66, 2.18),
        ("GENE RANKING", 3.07, 2.90),
        ("CLASSIFIER FIT", 6.43, 2.95),
        ("WHAT IT TESTS", 9.82, 2.72),
    ]
    for heading, x, width in headers:
        _add_text(slide, heading, x, 2.38, width, 0.22, size=9.0, color=GRAY, bold=True, margin=0)
    _add_rule(slide, 0.47, 2.64, 12.18, NAVY, 0.022)

    arms = [
        ("Real only", BLUE, ("R",), "Rank from real", ("R",), "Fit real", "Baseline"),
        ("Generated only", CORAL, ("S",), "Rank from generated", ("S",), "Fit generated", "Synthetic-only stress test"),
        ("Real + generated", TEAL, ("R", "S"), "Consensus rank", ("R", "S"), "Equal total weight", "Direct augmentation"),
        ("Guided: real fit", TEAL, ("R", "S"), "Consensus rank", ("R",), "Fit real only", "Feature guidance only"),
        ("Guided: 5% synthetic", ORANGE, ("R", "S"), "Consensus rank", ("R", "S"), "S = 5% weight", "Guidance + light regularization"),
    ]

    def add_sources(tokens: tuple[str, ...], x: float, y: float) -> float:
        if len(tokens) == 1:
            color = BLUE if tokens[0] == "R" else CORAL
            _add_data_badge(slide, tokens[0], x, y, color)
            return x + 0.34
        _add_data_badge(slide, "R", x, y, BLUE)
        _add_text(slide, "+", x + 0.27, y + 0.01, 0.18, 0.20, size=9.5, color=GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)
        _add_data_badge(slide, "S", x + 0.47, y, CORAL)
        return x + 0.82

    for index, (arm, color, rank_tokens, rank_label, fit_tokens, fit_label, purpose) in enumerate(arms):
        top = 2.73 + index * 0.57
        if index % 2 == 0:
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.47), Inches(top), Inches(12.18), Inches(0.50))
            _set_fill(shape, "F4F6F7")
            shape.line.fill.background()
        _add_rule(slide, 0.47, top, 0.055, color, 0.50)
        _add_text(slide, arm, 0.67, top + 0.08, 2.17, 0.33, size=11.2, color=color, bold=True, valign=MSO_ANCHOR.MIDDLE, margin=0)

        rank_end = add_sources(rank_tokens, 3.08, top + 0.13)
        _add_arrow(slide, rank_end, top + 0.17, 0.26, 0.14, MID_GRAY)
        _add_text(slide, rank_label, rank_end + 0.37, top + 0.08, 1.78, 0.33, size=10.4, color=DARK, bold=True, valign=MSO_ANCHOR.MIDDLE, margin=0)

        fit_end = add_sources(fit_tokens, 6.45, top + 0.13)
        _add_arrow(slide, fit_end, top + 0.17, 0.26, 0.14, MID_GRAY)
        _add_text(slide, fit_label, fit_end + 0.37, top + 0.08, 1.71, 0.33, size=10.4, color=DARK, bold=True, valign=MSO_ANCHOR.MIDDLE, margin=0)
        _add_text(slide, purpose, 9.83, top + 0.08, 2.70, 0.33, size=10.5, color=DARK, valign=MSO_ANCHOR.MIDDLE, margin=0)

    _add_rule(slide, 0.47, 5.64, 12.18, "CDD6DC", 0.018)
    _add_text(slide, "All five arms", 0.62, 5.83, 1.36, 0.27, size=12.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_arrow(slide, 2.13, 5.87, 0.34, 0.19, MID_GRAY)
    _add_text(slide, "Held-out real profiles", 2.66, 5.82, 2.00, 0.30, size=12.5, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_arrow(slide, 4.83, 5.87, 0.34, 0.19, MID_GRAY)
    _add_text(slide, "BA  |  AUROC  |  AP", 5.37, 5.83, 1.88, 0.27, size=12.2, color=DARK, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_arrow(slide, 7.41, 5.87, 0.34, 0.19, MID_GRAY)
    _add_text(slide, "Choose an eligible arm within each tissue", 7.92, 5.80, 4.40, 0.34, size=12.8, color=TEAL, bold=True, align=PP_ALIGN.CENTER, margin=0)

    _add_panel(slide, 0.46, 6.25, 12.22, 0.58, fill=NAVY, line=NAVY, radius=False)
    _add_text(slide, "Association test", 0.74, 6.40, 1.56, 0.25, size=13.5, color=WHITE, bold=True, margin=0)
    _add_text(slide, "FLT vs GC effects and BH FDR use real OSDR profiles only.", 2.47, 6.39, 6.40, 0.27, size=12.9, color="DCE7F2", margin=0)
    _add_text(slide, "Animal n stays unchanged.", 9.19, 6.39, 3.13, 0.27, size=11.5, color="FFD69A", bold=True, align=PP_ALIGN.RIGHT, margin=0)


def _slide_feature_importance_venn(slide):
    _add_slide_title(
        slide,
        "Synthetic analysis",
        "Compare what each classifier finds important",
    )

    _add_text(slide, "Real-only classifier", 3.07, 1.74, 3.30, 0.30, size=13.5, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "Real + synthetic classifier", 6.76, 1.74, 3.72, 0.30, size=13.5, color=TEAL, bold=True, align=PP_ALIGN.CENTER, margin=0)

    real_set = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(1.72),
        Inches(2.02),
        Inches(6.00),
        Inches(4.35),
    )
    _set_fill(real_set, BLUE, transparency=82)
    _set_line(real_set, BLUE, 2.4)
    synthetic_set = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(5.62),
        Inches(2.02),
        Inches(6.00),
        Inches(4.35),
    )
    _set_fill(synthetic_set, TEAL, transparency=82)
    _set_line(synthetic_set, TEAL, 2.4)

    _add_text(slide, "REAL ONLY", 2.25, 3.35, 3.08, 0.35, size=17.0, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "Important without", 2.25, 3.96, 3.08, 0.30, size=17.0, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "synthetic samples", 2.25, 4.30, 3.08, 0.30, size=17.0, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "PROMOTED", 8.12, 3.35, 2.85, 0.35, size=17.0, color=TEAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "Important after adding", 8.00, 3.96, 3.10, 0.30, size=15.6, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "synthetic samples", 8.00, 4.30, 3.10, 0.30, size=15.6, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)

    _add_text(slide, "REINFORCED", 5.83, 3.47, 1.67, 0.35, size=14.0, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "Important", 5.88, 4.04, 1.57, 0.27, size=13.2, color=DARK, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "in both", 5.88, 4.34, 1.57, 0.27, size=13.2, color=DARK, bold=True, align=PP_ALIGN.CENTER, margin=0)


def _slide_synthetic_interpretation_map(slide):
    _add_slide_title(
        slide,
        "Synthetic analysis",
        "We compared feature importance at three levels",
        "Each analysis compares the real-only and real + synthetic classifiers.",
    )

    panels = [
        (0.48, BLUE, PALE_BLUE, "1", "INDIVIDUAL GENES", "Permutation + SHAP"),
        (4.47, TEAL, PALE_TEAL, "2", "REACTOME GROUPS", "Grouped permutation + grouped SHAP"),
        (8.46, ORANGE, PALE_GOLD, "3", "CONSENSUS RANKING", "Compact gene panels"),
    ]
    for x, color, fill, number, heading, method in panels:
        _add_panel(slide, x, 2.14, 3.72, 4.47, fill=fill, line=color, radius=False)
        _add_rule(slide, x, 2.14, 3.72, color, 0.055)
        _add_circle(slide, x + 0.25, 2.40, 0.34, color)
        _add_text(slide, number, x + 0.25, 2.465, 0.34, 0.18, size=10.0, color=WHITE, bold=True, align=PP_ALIGN.CENTER, margin=0)
        _add_text(slide, heading, x + 0.72, 2.37, 2.64, 0.28, size=13.2, color=color, bold=True, margin=0)
        _add_text(slide, method, x + 0.28, 2.82, 3.16, 0.28, size=11.3, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)

    # Gene-level view: perturb one feature and attribute its prediction contribution.
    gene_x = 0.92
    gene_y = 3.35
    gene_heights = [0.26, 0.52, 0.38, 0.72, 0.44]
    for index, height in enumerate(gene_heights):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(gene_x + index * 0.28),
            Inches(gene_y + 0.75 - height),
            Inches(0.15),
            Inches(height),
        )
        _set_fill(bar, CORAL if index == 3 else BLUE)
        bar.line.fill.background()
    _add_text(slide, "shuffle", 2.34, 3.42, 0.58, 0.22, size=9.4, color=CORAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_arrow(slide, 2.30, 3.64, 0.56, 0.18, CORAL)
    _add_text(slide, "Permutation: does prediction weaken?", 0.78, 4.46, 3.10, 0.27, size=11.0, color=DARK, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "SHAP: does the gene push toward FLT or GC?", 0.75, 4.88, 3.18, 0.44, size=10.3, color=GRAY, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "OUTPUT · RANKED GENES", 0.87, 6.18, 2.95, 0.22, size=9.4, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)

    # Group-level view: related Reactome genes move together.
    group_colors = [BLUE, TEAL, ORANGE, BLUE, TEAL, ORANGE]
    group_positions = [(4.95, 3.49), (5.37, 3.32), (5.79, 3.53), (5.16, 3.88), (5.59, 3.81), (6.01, 3.91)]
    for (x, y), color in zip(group_positions, group_colors):
        _add_circle(slide, x, y, 0.20, color)
    _add_arrow(slide, 6.35, 3.65, 0.42, 0.22, MID_GRAY)
    _add_panel(slide, 6.92, 3.35, 0.88, 0.72, fill=WHITE, line=TEAL, radius=False)
    _add_text(slide, "one\npathway", 7.00, 3.48, 0.72, 0.42, size=9.7, color=TEAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "Move related genes together", 4.79, 4.46, 3.06, 0.27, size=11.0, color=DARK, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "Test correlated genes as one group", 4.74, 4.88, 3.16, 0.44, size=10.3, color=GRAY, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "OUTPUT · RANKED PATHWAYS", 4.82, 6.18, 3.00, 0.22, size=9.4, color=TEAL, bold=True, align=PP_ALIGN.CENTER, margin=0)

    # Consensus view: combine measured and generated rankings into compact panels.
    rank_rows = [("Gene 1", 0.86, 0.74), ("Gene 2", 0.68, 0.71), ("Gene 3", 0.47, 0.82)]
    for index, (label, real_rank, synthetic_rank) in enumerate(rank_rows):
        y = 3.32 + index * 0.40
        _add_text(slide, label, 8.78, y, 0.65, 0.22, size=9.2, color=DARK, bold=True, margin=0)
        _add_rule(slide, 9.51, y + 0.09, real_rank, BLUE, 0.07)
        _add_rule(slide, 10.52, y + 0.09, synthetic_rank, CORAL, 0.07)
    _add_text(slide, "real", 9.50, 3.09, 0.88, 0.19, size=8.8, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "generated", 10.45, 3.09, 0.95, 0.19, size=8.8, color=CORAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_arrow(slide, 11.42, 3.65, 0.24, 0.22, MID_GRAY)
    _add_text(slide, "top-k", 11.70, 3.60, 0.38, 0.24, size=9.4, color=ORANGE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "Combine measured and generated ranks", 8.72, 4.46, 3.42, 0.27, size=11.0, color=DARK, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "Test small panels of the highest-ranked genes", 8.70, 4.88, 3.48, 0.44, size=10.3, color=GRAY, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "OUTPUT · COMPACT GENE PANELS", 8.75, 6.18, 3.38, 0.22, size=9.4, color=ORANGE, bold=True, align=PP_ALIGN.CENTER, margin=0)

def _slide_matched_classifier_design(slide):
    _add_slide_title(
        slide,
        "Primary analysis",
        "Does adding synthetic data improve FLT vs GC prediction?",
        "Train the same tissue-specific classifier two ways, then compare both versions on real samples.",
    )

    rows = [
        (
            2.24,
            "01",
            "TRAIN ON REAL DATA",
            "Real profiles",
            ("R",),
            BLUE,
            PALE_BLUE,
        ),
        (
            4.15,
            "02",
            "TRAIN ON REAL + SYNTHETIC DATA",
            "Real + synthetic profiles",
            ("R", "S"),
            TEAL,
            PALE_TEAL,
        ),
    ]
    for y, number, heading, data_label, tokens, color, fill in rows:
        _add_panel(slide, 0.63, y, 7.55, 1.44, fill=fill, line=fill, radius=False)
        _add_rule(slide, 0.63, y, 0.075, color, 1.44)
        _add_text(slide, number, 0.96, y + 0.25, 0.42, 0.24, size=10.2, color=color, bold=True, margin=0)
        _add_text(slide, heading, 1.47, y + 0.22, 5.80, 0.30, size=14.5, color=NAVY, bold=True, margin=0)
        if len(tokens) == 1:
            _add_data_badge(slide, "R", 1.48, y + 0.78, BLUE, diameter=0.34)
            label_x = 1.98
        else:
            _add_data_badge(slide, "R", 1.48, y + 0.78, BLUE, diameter=0.34)
            _add_text(slide, "+", 1.87, y + 0.82, 0.18, 0.20, size=12.0, color=GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)
            _add_data_badge(slide, "S", 2.10, y + 0.78, CORAL, diameter=0.34)
            label_x = 2.60
        _add_text(slide, data_label, label_x, y + 0.81, 2.16, 0.24, size=11.5, color=DARK, margin=0)
        _add_arrow(slide, 4.62, y + 0.76, 0.56, 0.28, MID_GRAY)
        _add_panel(slide, 5.42, y + 0.51, 2.28, 0.78, fill=WHITE, line=color, radius=False)
        _add_text(slide, "FLT / GC classifier", 5.61, y + 0.75, 1.90, 0.28, size=13.0, color=color, bold=True, align=PP_ALIGN.CENTER, margin=0)
        _add_arrow(slide, 7.82, y + 0.76, 0.58, 0.28, MID_GRAY)

    _add_panel(slide, 8.60, 2.24, 4.06, 3.35, fill=WHITE, line=ORANGE, radius=False)
    _add_rule(slide, 8.60, 2.24, 4.06, ORANGE, 0.06)
    _add_text(slide, "TEST BOTH", 8.96, 2.66, 3.34, 0.28, size=14.0, color=ORANGE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "Same real samples", 8.96, 3.18, 3.34, 0.34, size=18.0, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "FLT", 9.42, 3.92, 1.00, 0.30, size=15.0, color=CORAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "GC", 10.84, 3.92, 1.00, 0.30, size=15.0, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_circle(slide, 9.78, 4.36, 0.22, CORAL)
    _add_circle(slide, 10.14, 4.23, 0.22, CORAL)
    _add_circle(slide, 10.22, 4.61, 0.22, CORAL)
    _add_circle(slide, 11.20, 4.31, 0.22, BLUE)
    _add_circle(slide, 11.56, 4.48, 0.22, BLUE)
    _add_circle(slide, 11.18, 4.67, 0.22, BLUE)
    _add_text(slide, "Compare prediction accuracy", 8.96, 5.10, 3.34, 0.28, size=12.3, color=DARK, bold=True, align=PP_ALIGN.CENTER, margin=0)

    _add_source(slide, "Both classifiers are evaluated on real OSDR profiles.")


def _slide_all_gene_candidate_filter(slide):
    _add_slide_title(
        slide,
        "Analysis 1 | Individual genes",
        "How 974 scored genes became 21 candidates",
        "Permutation and SHAP score every gene. The final list keeps genes supported by prediction and measured OSDR data.",
    )

    _add_panel(slide, 0.48, 1.96, 12.17, 0.58, fill="F2F5F7", line="D3DDE2", radius=False)
    _add_text(
        slide,
        "START: ALL 974 GENES REMAIN IN EACH CLASSIFIER",
        0.75,
        2.13,
        11.63,
        0.24,
        size=12.4,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
    )

    gates = [
        (
            0.48,
            "1",
            "TISSUE MODEL",
            "Synthetic training holds or improves prediction",
            "Checked on unseen measured samples across repeated splits",
            BLUE,
            PALE_BLUE,
        ),
        (
            3.55,
            "2",
            "GENE IMPORTANCE",
            "Shuffle one gene in unseen real samples",
            "Keep genes whose removal repeatedly weakens prediction",
            CORAL,
            PALE_CORAL,
        ),
        (
            6.62,
            "3",
            "DIRECTION",
            "Coefficient and SHAP point the same way",
            "Direction agrees with the measured FLT-GC change",
            ORANGE,
            "FFF3E2",
        ),
        (
            9.69,
            "4",
            "MEASURED ASSOCIATION",
            "The gene changes between FLT and GC in real OSDR samples",
            "Measured-data evidence is required; generated samples do not count",
            TEAL,
            PALE_TEAL,
        ),
    ]
    for x, number, heading, main, detail, color, fill in gates:
        _add_panel(slide, x, 2.79, 2.66, 2.30, fill=fill, line=color, radius=False)
        _add_circle(slide, x + 0.19, 3.02, 0.34, color)
        _add_text(
            slide,
            number,
            x + 0.19,
            3.075,
            0.34,
            0.20,
            size=10.2,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            margin=0,
        )
        _add_text(slide, heading, x + 0.65, 3.03, 1.76, 0.25, size=10.6, color=color, bold=True, margin=0)
        _add_rule(slide, x + 0.20, 3.43, 2.26, color, 0.022)
        _add_text(
            slide,
            main,
            x + 0.22,
            3.63,
            2.22,
            0.55,
            size=12.0,
            color=NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
            margin=0,
        )
        _add_text(
            slide,
            detail,
            x + 0.22,
            4.32,
            2.22,
            0.55,
            size=9.4,
            color=DARK,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
            margin=0,
        )
        _add_rule(slide, x + 1.32, 5.09, 0.018, color, 0.34)

    _add_rule(slide, 1.80, 5.42, 9.21, "A8B7C0", 0.018)
    down_arrow = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        Inches(6.48),
        Inches(5.40),
        Inches(0.38),
        Inches(0.46),
    )
    _set_fill(down_arrow, TEAL)
    down_arrow.line.fill.background()

    _add_panel(slide, 3.29, 5.91, 6.75, 0.66, fill=NAVY, line=NAVY, radius=False)
    _add_text(slide, "ALL FOUR CHECKS", 3.57, 6.12, 2.08, 0.24, size=11.3, color="BFD0E1", bold=True, margin=0)
    _add_text(slide, "21 tissue-gene associations", 5.80, 6.08, 3.87, 0.30, size=16.2, color=WHITE, bold=True, align=PP_ALIGN.CENTER, margin=0)

    _add_text(
        slide,
        "These four checks are specific to individual genes; pathway and panel analyses use their own evidence screens.",
        0.63,
        6.73,
        12.03,
        0.28,
        size=10.8,
        color=GRAY,
        italic=True,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    _add_source(slide, "Exact repeat, effect-direction, and measured-data thresholds are provided in the speaker notes.")


def _slide_guidance_mechanism(slide):
    _add_slide_title(
        slide,
        "Analysis 3 | Compact panels",
        "Combine real and synthetic rankings to choose a small gene panel",
        "Genes supported by both rankings rise; only the top genes enter the FLT-GC model.",
    )

    def add_rank_track(x: float, y: float, label: str, position: float, color: str):
        track_x = x + 0.90
        track_w = 1.92
        _add_text(slide, label, x, y - 0.055, 0.78, 0.23, size=10.4, color=DARK, bold=True, margin=0)
        _add_rule(slide, track_x, y + 0.035, track_w, "C9D2D8", 0.025)
        _add_circle(slide, track_x + track_w * position - 0.075, y - 0.025, 0.15, color)

    panel_y = 2.05
    panel_h = 2.60
    _add_panel(slide, 0.48, panel_y, 3.61, panel_h, fill=PALE_BLUE, line="C9DCE9", radius=False)
    _add_text(slide, "1  REAL-DATA RANKING", 0.72, 2.27, 2.35, 0.25, size=13.1, color=BLUE, bold=True, margin=0)
    _add_text(slide, "Rank all 974 genes using real OSDR", 0.72, 2.61, 3.02, 0.23, size=10.4, color=DARK, margin=0)
    _add_text(slide, "lower-ranked", 1.59, 2.96, 0.82, 0.18, size=8.5, color=GRAY, margin=0)
    _add_text(slide, "top-ranked", 3.02, 2.96, 0.72, 0.18, size=8.5, color=GRAY, align=PP_ALIGN.RIGHT, margin=0)
    example_genes = ["Gene 1", "Gene 2", "Gene 3", "Gene 4", "Gene 5"]
    real_positions = [0.90, 0.82, 0.72, 0.56, 0.35]
    for index, (gene, position) in enumerate(zip(example_genes, real_positions)):
        add_rank_track(0.72, 3.18 + index * 0.25, gene, position, BLUE)

    _add_arrow(slide, 4.16, 2.93, 0.27, 0.20, MID_GRAY)
    _add_panel(slide, 4.50, panel_y, 3.61, panel_h, fill=PALE_CORAL, line="E6CEC8", radius=False)
    _add_text(slide, "2  SYNTHETIC-DATA RANKING", 4.74, 2.27, 3.02, 0.25, size=12.3, color=CORAL, bold=True, margin=0)
    _add_text(slide, "Rank the same 974 genes using generated profiles", 4.74, 2.61, 3.07, 0.23, size=10.4, color=DARK, margin=0)
    _add_text(slide, "lower-ranked", 5.61, 2.96, 0.82, 0.18, size=8.5, color=GRAY, margin=0)
    _add_text(slide, "top-ranked", 7.04, 2.96, 0.72, 0.18, size=8.5, color=GRAY, align=PP_ALIGN.RIGHT, margin=0)
    generated_positions = [0.88, 0.45, 0.86, 0.68, 0.32]
    for index, (gene, position) in enumerate(zip(example_genes, generated_positions)):
        add_rank_track(4.74, 3.18 + index * 0.25, gene, position, CORAL)
    _add_arrow(slide, 8.18, 2.93, 0.27, 0.20, MID_GRAY)
    _add_panel(slide, 8.52, panel_y, 4.13, panel_h, fill=PALE_TEAL, line="C9DFDB", radius=False)
    _add_text(slide, "3  COMBINED RANKING", 8.76, 2.27, 2.55, 0.25, size=13.4, color=TEAL, bold=True, margin=0)
    _add_text(slide, "Agreement moves genes upward", 8.76, 2.61, 2.82, 0.23, size=10.4, color=DARK, margin=0)
    _add_text(slide, "lower-ranked", 9.66, 2.96, 0.82, 0.18, size=8.5, color=GRAY, margin=0)
    _add_text(slide, "top-ranked", 11.61, 2.96, 0.72, 0.18, size=8.5, color=GRAY, align=PP_ALIGN.RIGHT, margin=0)
    shared_rows = [
        ("Gene 1", 0.89, TEAL, "selected"),
        ("Gene 2", 0.61, MID_GRAY, "not selected"),
        ("Gene 3", 0.79, TEAL, "selected"),
        ("Gene 4", 0.62, MID_GRAY, "not selected"),
        ("Gene 5", 0.33, MID_GRAY, "not selected"),
    ]
    for index, (gene, position, color, outcome) in enumerate(shared_rows):
        y = 3.18 + index * 0.25
        add_rank_track(8.76, y, gene, position, color)
        _add_text(slide, outcome, 11.72, y - 0.055, 0.80, 0.21, size=7.8, color=color, bold=True, align=PP_ALIGN.RIGHT, margin=0)
    shared_cutoff_x = 8.76 + 0.90 + 1.92 * 0.75
    _add_rule(slide, shared_cutoff_x, 3.10, 0.022, TEAL, 1.27)
    _add_text(slide, "selection cutoff", 11.15, 4.42, 1.20, 0.18, size=8.5, color=TEAL, italic=True, align=PP_ALIGN.RIGHT, margin=0)

    _add_text(slide, "Top-ranked genes enter a repeated panel test", 0.52, 4.84, 5.20, 0.28, size=16.0, color=NAVY, bold=True, margin=0)
    lower_steps = [
        (0.51, 2.62, "Rank 974 genes", "real, synthetic, or combined", BLUE, "F4F8FA", "D5E1E7"),
        (3.57, 2.62, "Keep top-ranked genes", "test several compact panel sizes", TEAL, PALE_TEAL, "C9DFDB"),
        (6.63, 2.62, "Repeat panel tests", "use the same real train/test design", NAVY, "F6F7F8", "D9DFE3"),
        (9.69, 2.96, "Retain stable genes", "recur across splits, keep direction,\nand pass measured-data support", ORANGE, "FFF3E2", "E8CFAC"),
    ]
    for index, (x, width, heading, detail, color, fill, line) in enumerate(lower_steps):
        _add_panel(slide, x, 5.18, width, 1.28, fill=fill, line=line, radius=False)
        _add_text(slide, heading, x + 0.15, 5.40, width - 0.30, 0.27, size=12.0, color=color, bold=True, align=PP_ALIGN.CENTER, margin=0)
        _add_text(slide, detail, x + 0.15, 5.80, width - 0.30, 0.45, size=9.2, color=DARK, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
        if index < len(lower_steps) - 1:
            _add_arrow(slide, x + width + 0.09, 5.69, 0.25, 0.18, MID_GRAY)

    _add_source(slide, "Five genes illustrate ranking. Exact panel sizes and stability thresholds are provided in the speaker notes.")


def _slide_grouped_pathway_analysis(slide):
    _add_slide_title(
        slide,
        "Analysis 2 | Reactome groups",
        "Move a pathway together to reveal shared predictive signal",
        "Grouped permutation asks whether the pathway matters; grouped SHAP shows whether it points toward FLT or GC.",
    )

    panels = [
        (0.48, 3.61, "A", "GROUP RELATED GENES", BLUE, PALE_BLUE),
        (4.50, 3.61, "B", "GROUPED PERMUTATION", CORAL, PALE_CORAL),
        (8.52, 4.13, "C", "GROUPED SHAP", TEAL, PALE_TEAL),
    ]
    for x, width, label, heading, color, fill in panels:
        _add_panel(slide, x, 2.03, width, 2.74, fill=fill, line=color, radius=False)
        _add_text(slide, label, x + 0.21, 2.22, 0.38, 0.30, size=16.0, color=color, bold=True, margin=0)
        _add_text(slide, heading, x + 0.65, 2.23, width - 0.86, 0.28, size=12.6, color=NAVY, bold=True, margin=0)

    # A: Reactome supplies predefined groups so correlated genes move together.
    gene_colors = [BLUE, TEAL, ORANGE]
    for index, color in enumerate(gene_colors):
        y = 2.84 + index * 0.47
        _add_panel(slide, 0.84, y, 0.90, 0.34, fill=WHITE, line=color, radius=False)
        _add_text(slide, f"Gene {chr(65 + index)}", 0.91, y + 0.07, 0.76, 0.19, size=9.4, color=color, bold=True, align=PP_ALIGN.CENTER, margin=0)
        _add_arrow(slide, 1.90, y + 0.075, 0.37, 0.18, color)
    _add_panel(slide, 2.42, 3.17, 1.28, 0.72, fill=WHITE, line=BLUE, radius=False)
    _add_text(slide, "one Reactome\npathway", 2.53, 3.31, 1.06, 0.40, size=10.2, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "Related genes can share signal, so test them as one unit.", 0.82, 4.26, 2.95, 0.34, size=10.3, color=DARK, align=PP_ALIGN.CENTER, margin=0)

    # B: use one shared row permutation for every member of the pathway.
    _add_panel(slide, 4.84, 3.05, 1.20, 0.72, fill=WHITE, line=CORAL, radius=False)
    _add_text(slide, "shuffle the\nwhole pathway", 4.95, 3.18, 0.98, 0.42, size=9.8, color=CORAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_arrow(slide, 6.16, 3.29, 0.34, 0.20, MID_GRAY)
    _add_panel(slide, 6.63, 3.05, 1.03, 0.72, fill=WHITE, line=NAVY, radius=False)
    _add_text(slide, "FLT/GC\nmodel", 6.72, 3.18, 0.85, 0.42, size=10.0, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "Does prediction weaken?", 4.86, 3.96, 2.76, 0.25, size=11.0, color=CORAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "A repeatable drop means the pathway carries predictive information.", 4.82, 4.28, 2.98, 0.34, size=10.0, color=DARK, align=PP_ALIGN.CENTER, margin=0)

    # C: sum member-gene SHAP values and read the pathway-level direction.
    _add_panel(slide, 8.86, 3.05, 1.05, 0.72, fill=WHITE, line=TEAL, radius=False)
    _add_text(slide, "pathway\ngenes", 8.96, 3.18, 0.85, 0.42, size=10.0, color=TEAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_arrow(slide, 10.02, 3.29, 0.28, 0.20, MID_GRAY)
    _add_panel(slide, 10.40, 3.05, 1.03, 0.72, fill=WHITE, line=TEAL, radius=False)
    _add_text(slide, "sum member\nSHAP values", 10.49, 3.18, 0.85, 0.42, size=9.4, color=TEAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_arrow(slide, 11.54, 3.29, 0.25, 0.20, MID_GRAY)
    _add_text(slide, "FLT or GC", 11.87, 3.26, 0.54, 0.24, size=9.4, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "Which condition does the pathway support?", 8.87, 3.96, 3.37, 0.25, size=11.0, color=TEAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "Its direction should agree with the measured pathway change.", 8.84, 4.28, 3.48, 0.34, size=10.0, color=DARK, align=PP_ALIGN.CENTER, margin=0)

    _add_text(slide, "KEEP A PATHWAY WHEN", 0.60, 4.98, 2.12, 0.24, size=11.3, color=NAVY, bold=True, margin=0)
    _add_panel(slide, 0.48, 5.24, 12.17, 0.96, fill="F2F5F7", line="D3DDE2", radius=False)
    gates = [
        (0.71, "1", "Prediction holds or improves"),
        (3.66, "2", "Pathway permutation\nrepeatedly weakens prediction"),
        (6.65, "3", "SHAP direction matches\nthe measured effect"),
        (9.66, "4", "Pathway change appears\nin measured OSDR"),
    ]
    for index, (x, number, label) in enumerate(gates):
        if index:
            _add_rule(slide, x - 0.26, 5.43, 0.012, "C7D1D7", 0.57)
        _add_circle(slide, x, 5.50, 0.32, TEAL)
        _add_text(slide, number, x, 5.55, 0.32, 0.18, size=9.2, color=WHITE, bold=True, align=PP_ALIGN.CENTER, margin=0)
        _add_text(slide, label, x + 0.43, 5.43, 2.08, 0.46, size=10.5, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE, margin=0)

    _add_panel(slide, 3.50, 6.37, 6.33, 0.45, fill=NAVY, line=NAVY, radius=False)
    _add_text(slide, "ALL FOUR CHECKS  |  RETAINED PATHWAY", 3.76, 6.49, 5.81, 0.22, size=11.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_source(slide, "Exact repeat and measured-data thresholds are provided in the speaker notes.")


def _slide_guidance_boundary(slide):
    _add_slide_title(
        slide,
        "Analysis",
        "Synthetic guidance can shift the FLT/GC boundary",
        "Opaque points train the classifier; transparent points are the same held-out real profiles in both panels.",
    )

    def add_point(x: float, y: float, color: str, *, held_out: bool):
        if not held_out:
            _add_circle(slide, x, y, 0.17, color)
            return
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x),
            Inches(y),
            Inches(0.20),
            Inches(0.20),
        )
        _set_fill(shape, color, transparency=58)
        _set_line(shape, color, 1.2)

    def add_scatter(x: float, y: float, *, guided: bool):
        plot_w = 4.75
        plot_h = 3.40
        _add_rule(slide, x, y + plot_h, plot_w, MID_GRAY, 0.022)
        _add_rule(slide, x, y, 0.022, MID_GRAY, plot_h + 0.02)
        train_flight = [(0.45, 0.52), (0.82, 2.25), (1.20, 1.35), (1.55, 2.95), (1.75, 0.70)]
        train_ground = [(2.85, 0.55), (3.15, 2.25), (3.50, 1.30), (3.90, 2.95), (4.25, 0.75)]
        held_flight = [(2.05, 1.55), (2.28, 2.55)]
        held_ground = [(2.58, 1.00), (2.80, 2.75)]
        for px, py in train_flight:
            add_point(x + px, y + py, CORAL, held_out=False)
        for px, py in train_ground:
            add_point(x + px, y + py, BLUE, held_out=False)
        for px, py in held_flight:
            add_point(x + px, y + py, CORAL, held_out=True)
        for px, py in held_ground:
            add_point(x + px, y + py, BLUE, held_out=True)
        if guided:
            _add_rule(slide, x + 2.42, y + 0.18, 0.045, TEAL, 2.98)
        else:
            _add_rule(slide, x + 1.85, y + 0.18, 0.035, GRAY, 2.98)
        _add_text(slide, "predicted FLT", x + 0.12, y + 3.52, 1.12, 0.20, size=9.2, color=CORAL, margin=0)
        _add_text(slide, "predicted GC", x + 3.48, y + 3.52, 1.12, 0.20, size=9.2, color=BLUE, align=PP_ALIGN.RIGHT, margin=0)

    _add_circle(slide, 8.79, 1.96, 0.13, CORAL)
    _add_text(slide, "FLT", 8.98, 1.92, 0.36, 0.21, size=9.2, color=GRAY, margin=0)
    _add_circle(slide, 9.50, 1.96, 0.13, BLUE)
    _add_text(slide, "GC", 9.69, 1.92, 0.34, 0.21, size=9.2, color=GRAY, margin=0)
    _add_circle(slide, 10.30, 1.96, 0.13, MID_GRAY)
    _add_text(slide, "training", 10.49, 1.92, 0.62, 0.21, size=9.2, color=GRAY, margin=0)
    held_legend = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.28), Inches(1.94), Inches(0.16), Inches(0.16))
    _set_fill(held_legend, MID_GRAY, transparency=58)
    _set_line(held_legend, MID_GRAY, 1.0)
    _add_text(slide, "held-out", 11.52, 1.92, 0.72, 0.21, size=9.2, color=GRAY, margin=0)

    _add_text(slide, "REAL-ONLY CLASSIFIER", 0.72, 2.23, 4.75, 0.27, size=13.0, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "SYNTHETIC-GUIDED CLASSIFIER", 7.85, 2.23, 4.75, 0.27, size=13.0, color=TEAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    add_scatter(0.72, 2.62, guided=False)
    add_scatter(7.85, 2.62, guided=True)
    _add_arrow(slide, 5.87, 3.88, 0.82, 0.44, MID_GRAY)
    _add_text(slide, "same held-out profiles", 5.57, 4.45, 1.42, 0.39, size=9.6, color=GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)

    _add_source(slide, "Schematic of the held-out evaluation used for each tissue-specific classifier candidate.")


def _slide_8(slide, utility_chart: Path):
    _add_slide_title(
        slide,
        "Prediction",
        "Real + synthetic vs real-only balanced accuracy across 27 tissues",
    )
    _add_picture_contain(
        slide,
        utility_chart,
        0.38,
        1.52,
        12.55,
        5.12,
        alt="Changes in FLT versus GC classifier performance for all 27 tissue analyses",
    )
    legend_items = [
        (3.15, MID_GRAY, "Real only", 0.78),
        (5.02, TEAL, "All six metrics passed", 1.55),
        (7.95, CORAL, "At least one declined", 1.62),
    ]
    for x, color, label, width in legend_items:
        _add_circle(slide, x, 6.77, 0.11, color)
        _add_text(
            slide,
            label,
            x + 0.17,
            6.735,
            width,
            0.20,
            size=8.9,
            color=color,
            bold=True,
            margin=0,
        )
    _add_source(
        slide,
        "Outcome: FLT versus ground control within each tissue. All versions used the same 974 genes and prediction settings.",
    )


def _slide_9(slide):
    _add_slide_title(
        slide,
        "Feature analysis",
        "Synthetic guidance changed ranking, not statistical evidence",
        "Blue marks real-only selection; teal marks synthetic-guided selection. Every displayed association has real OSDR support.",
    )
    _add_text(slide, "Stable selection after real-data support", 0.62, 2.02, 4.30, 0.30, size=15.5, color=NAVY, bold=True, margin=0)
    _add_circle(slide, 0.64, 2.42, 0.13, BLUE)
    _add_text(slide, "Blue: stable in real-only ranking", 0.86, 2.36, 2.42, 0.25, size=10.5, color=BLUE, bold=True, margin=0)
    _add_circle(slide, 3.43, 2.42, 0.13, TEAL)
    _add_text(slide, "Teal: stable in synthetic-guided ranking", 3.65, 2.36, 2.98, 0.25, size=10.5, color=TEAL, bold=True, margin=0)

    left = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.62), Inches(2.78), Inches(3.86), Inches(2.78))
    _set_fill(left, BLUE, transparency=82)
    _set_line(left, BLUE, 2.6)
    right = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.00), Inches(2.78), Inches(3.86), Inches(2.78))
    _set_fill(right, TEAL, transparency=80)
    _set_line(right, TEAL, 2.6)

    _add_text(slide, "34", 1.16, 3.42, 1.42, 0.48, size=29, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "real-only", 1.12, 3.93, 1.50, 0.30, size=14.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "WITHOUT GUIDANCE", 1.03, 4.31, 1.68, 0.24, size=9.2, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)

    _add_text(slide, "23", 3.16, 3.42, 1.17, 0.48, size=29, color=ORANGE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "reinforced", 3.10, 3.93, 1.29, 0.30, size=13.2, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "BOTH RANKINGS", 3.04, 4.31, 1.40, 0.24, size=9.2, color=ORANGE, bold=True, align=PP_ALIGN.CENTER, margin=0)

    _add_text(slide, "26", 4.89, 3.42, 1.50, 0.48, size=29, color=TEAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "promoted", 4.89, 3.93, 1.50, 0.30, size=14.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "GUIDANCE ONLY", 4.80, 4.31, 1.68, 0.24, size=9.2, color=TEAL, bold=True, align=PP_ALIGN.CENTER, margin=0)

    _add_rule(slide, 7.18, 2.10, 0.018, "D5DDE2", 3.98)
    _add_text(slide, "Real OSDR evidence gate", 7.53, 2.43, 4.92, 0.40, size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    evidence_rows = [
        (3.18, 0.38, "Estimate FLT vs GC inside each accession"),
        (3.82, 0.58, "Combine accession effects with a\nrandom-effects model"),
        (4.63, 0.38, "Apply BH FDR within each tissue"),
    ]
    for y, height, label in evidence_rows:
        _add_circle(slide, 7.79, y + 0.10, 0.10, BLUE)
        _add_text(
            slide,
            label,
            8.05,
            y,
            4.15,
            height,
            size=14.2,
            color=DARK,
            valign=MSO_ANCHOR.TOP,
            margin=0,
        )
    _add_panel(slide, 7.68, 5.14, 4.46, 0.78, fill=NAVY, line=NAVY, radius=False)
    _add_text(slide, "23 reinforced + 26 promoted", 7.90, 5.26, 4.02, 0.24, size=11.5, color="BFD0E1", bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "49 synthetic-informed associations", 7.88, 5.52, 4.06, 0.28, size=15.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER, margin=0)

    _add_rule(slide, 0.52, 6.26, 12.15, "D6DEE3", 0.018)
    _add_text(slide, "The 34 real-only associations remain real-data findings; synthetic guidance did not reinforce them.", 0.72, 6.39, 11.75, 0.26, size=12.6, color=GRAY, align=PP_ALIGN.CENTER)
    _add_source(slide, "Counts are tissue-gene associations after BH FDR and compatible real-effect direction; genes may recur across tissues.")


def _slide_matched_consensus_crosswalk(slide):
    _add_slide_title(
        slide,
        "Feature analysis",
        "All-gene models and compact panels serve different purposes",
        "Both start from gene changes supported in measured OSDR samples.",
    )

    _add_text(slide, "Genes supported in measured OSDR data", 0.63, 2.05, 5.75, 0.31, size=14.0, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    matched = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.72), Inches(2.62), Inches(3.78), Inches(2.70))
    _set_fill(matched, CORAL, transparency=78)
    _set_line(matched, CORAL, 2.4)
    consensus = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.93), Inches(2.62), Inches(3.78), Inches(2.70))
    _set_fill(consensus, BLUE, transparency=80)
    _set_line(consensus, BLUE, 2.4)
    _add_text(slide, "All-gene model", 0.72, 2.38, 2.45, 0.27, size=11.7, color=CORAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "Combined gene panels", 4.32, 2.38, 2.35, 0.27, size=11.7, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "10", 1.25, 3.47, 1.20, 0.50, size=29, color=CORAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "all-gene only", 1.15, 4.03, 1.40, 0.28, size=11.0, color=DARK, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "11", 3.10, 3.47, 1.20, 0.50, size=29, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "both", 3.10, 4.03, 1.20, 0.28, size=11.0, color=DARK, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "38", 4.97, 3.47, 1.20, 0.50, size=29, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "panel only", 4.81, 4.03, 1.52, 0.28, size=11.0, color=DARK, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "21 total", 1.57, 5.05, 1.55, 0.28, size=12.0, color=CORAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "49 total", 4.32, 5.05, 1.55, 0.28, size=12.0, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)

    _add_rule(slide, 7.10, 2.05, 0.015, "D5DDE2", 4.40)
    rows = [
        (2.17, "PRIMARY", "All-gene model", "Only the training data change. Prediction and gene importance are tested on unseen real samples.", CORAL, PALE_CORAL),
        (3.50, "SECONDARY", "Combined gene panel", "Combines real and synthetic rankings to retain small groups of correlated genes.", BLUE, PALE_BLUE),
        (4.83, "SHARED", "Measured-data support", "FLT-GC differences are checked in measured OSDR samples, not generated profiles.", TEAL, PALE_TEAL),
    ]
    for y, label, heading, detail, color, fill in rows:
        _add_panel(slide, 7.42, y, 5.10, 1.06, fill=fill, line=color, radius=False)
        _add_text(slide, label, 7.68, y + 0.18, 1.00, 0.23, size=9.2, color=color, bold=True, margin=0)
        _add_text(slide, heading, 8.83, y + 0.14, 3.34, 0.29, size=14.0, color=NAVY, bold=True, margin=0)
        _add_text(slide, detail, 7.68, y + 0.51, 4.49, 0.39, size=10.5, color=DARK, margin=0)

    _add_panel(slide, 0.58, 5.87, 12.02, 0.64, fill=NAVY, line=NAVY, radius=False)
    _add_text(slide, "Thymus is strongest because both gene-level methods identify the same cell-cycle signal.", 0.85, 6.05, 11.48, 0.28, size=13.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_source(slide, "Overlap: 21 all-gene associations, 49 combined-panel associations, and 11 found by both gene-level methods.")


def _load_retained_candidate_inventory() -> tuple[pd.DataFrame, pd.DataFrame, pd.DataFrame]:
    source_dir = PAPER_DIR / "source_data"

    matched_raw = pd.read_csv(
        source_dir / "table_s19_matched_all_gene_candidates.tsv",
        sep="\t",
    )
    matched = (
        matched_raw.groupby(["analysis_scope", "tissue", "gene"], as_index=False)
        .agg(
            symbol=("symbol", "first"),
            promoted=(
                "pattern",
                lambda values: values.astype(str)
                .str.contains("synthetic_promoted_real_transfer")
                .any(),
            ),
        )
    )
    matched["selection_status"] = np.where(
        matched["promoted"], "promoted", "reinforced"
    )

    grouped = pd.read_csv(
        source_dir / "table_s23_grouped_pathway_literature_annotations.tsv",
        sep="\t",
    ).drop_duplicates(["scope", "tissue", "term"])
    grouped["selection_status"] = np.where(
        grouped["group_importance_patterns"]
        .astype(str)
        .str.contains("synthetic_promoted_group"),
        "promoted",
        "reinforced",
    )

    compact = pd.read_csv(
        source_dir / "table_s10_synthetic_informed_bh_fdr_genes.tsv",
        sep="\t",
    ).copy()
    compact["selection_status"] = compact["selection_interpretation"].map(
        {
            "synthetic_promoted": "promoted",
            "reinforced_real_and_synthetic": "reinforced",
        }
    )

    expected = {
        "matched": (len(matched), int(matched["promoted"].sum())),
        "grouped": (
            len(grouped),
            int(grouped["selection_status"].eq("promoted").sum()),
        ),
        "compact": (
            len(compact),
            int(compact["selection_status"].eq("promoted").sum()),
        ),
    }
    if expected != {
        "matched": (21, 7),
        "grouped": (10, 4),
        "compact": (49, 26),
    }:
        raise ValueError(f"Unexpected retained-candidate inventory: {expected}")
    if compact["selection_status"].isna().any():
        raise ValueError("Unknown compact-panel selection status")
    return matched, grouped, compact


def _slide_retained_outputs(slide):
    _add_slide_title(
        slide,
        "Results",
        "The three analyses retained different types of results",
        "Every retained candidate was also supported by the real OSDR data.",
    )
    matched, grouped, compact = _load_retained_candidate_inventory()

    methods = [
        {
            "x": 0.48,
            "color": CORAL,
            "heading": "Individual genes",
            "method": "Permutation + SHAP",
            "total": len(matched),
            "promoted": int(matched["selection_status"].eq("promoted").sum()),
            "reinforced": int(matched["selection_status"].eq("reinforced").sum()),
            "unit": "associations in 4 tissues",
            "rows": [("Thymus", "15"), ("Liver", "4"), ("Skin", "1"), ("Spleen", "1")],
        },
        {
            "x": 4.49,
            "color": GOLD,
            "heading": "Reactome groups",
            "method": "Grouped permutation + SHAP",
            "total": len(grouped),
            "promoted": int(grouped["selection_status"].eq("promoted").sum()),
            "reinforced": int(grouped["selection_status"].eq("reinforced").sum()),
            "unit": "pathways in 3 tissues",
            "rows": [("Thymus", "7"), ("Skin", "2"), ("Spleen", "1")],
        },
        {
            "x": 8.50,
            "color": BLUE,
            "heading": "Consensus ranking",
            "method": "Compact gene panels",
            "total": len(compact),
            "promoted": int(compact["selection_status"].eq("promoted").sum()),
            "reinforced": int(compact["selection_status"].eq("reinforced").sum()),
            "unit": "associations in 10 analyses",
            "rows": [("Thymus", "16"), ("Pooled muscle", "12"), ("Soleus", "5"), ("Seven others", "16")],
        },
    ]

    for index, method in enumerate(methods):
        x = method["x"]
        if index:
            _add_rule(slide, x - 0.22, 2.02, 0.012, "D9E0E4", 4.63)
        _add_rule(slide, x, 2.02, 3.56, method["color"], 0.045)
        _add_text(slide, method["heading"], x, 2.20, 3.56, 0.34, size=18.5, color=NAVY, bold=True, margin=0)
        _add_text(slide, method["method"], x, 2.58, 3.56, 0.24, size=10.5, color=GRAY, margin=0)
        _add_text(slide, str(method["total"]), x, 2.96, 1.08, 0.62, size=35, color=method["color"], bold=True, margin=0)
        _add_text(slide, method["unit"], x + 1.04, 3.06, 2.45, 0.36, size=12.5, color=DARK, bold=True, valign=MSO_ANCHOR.MIDDLE, margin=0)

        _add_rule(slide, x, 3.76, 0.13, CORAL, 0.13)
        _add_text(slide, f"{method['promoted']} promoted", x + 0.22, 3.69, 1.43, 0.27, size=11.2, color=CORAL, bold=True, margin=0)
        _add_rule(slide, x + 1.80, 3.76, 0.13, TEAL, 0.13)
        _add_text(slide, f"{method['reinforced']} reinforced", x + 2.02, 3.69, 1.54, 0.27, size=11.2, color=TEAL, bold=True, margin=0)

        _add_text(slide, "RETAINED BY TISSUE", x, 4.28, 2.10, 0.20, size=8.8, color=GRAY, bold=True, margin=0)
        _add_rule(slide, x, 4.55, 3.56, "D8E0E4", 0.014)
        for row_index, (label, count) in enumerate(method["rows"]):
            y = 4.70 + row_index * 0.43
            _add_text(slide, label, x + 0.02, y, 2.70, 0.29, size=11.8, color=DARK, margin=0)
            _add_text(slide, count, x + 2.84, y, 0.62, 0.29, size=12.2, color=method["color"], bold=True, align=PP_ALIGN.RIGHT, margin=0)
            if row_index < len(method["rows"]) - 1:
                _add_rule(slide, x + 0.02, y + 0.34, 3.43, "E4E9EC", 0.010)

    _add_source(slide, "Promoted: retained only after adding synthetic information. Reinforced: retained in both analyses.")


def _count_candidates_by_tissue(frame: pd.DataFrame) -> dict[str, tuple[int, int, int]]:
    counts = {}
    for tissue, rows in frame.groupby("tissue"):
        promoted = int(rows["selection_status"].eq("promoted").sum())
        reinforced = int(rows["selection_status"].eq("reinforced").sum())
        counts[str(tissue)] = (len(rows), promoted, reinforced)
    return counts


def _slide_tissue_method_convergence(slide):
    _add_slide_title(
        slide,
        "Results",
        "All three analyses found candidates in thymus, skin and spleen",
    )
    matched, grouped, compact = _load_retained_candidate_inventory()
    method_counts = [
        _count_candidates_by_tissue(matched),
        _count_candidates_by_tissue(grouped),
        _count_candidates_by_tissue(compact),
    ]
    rows = [
        ("thymus", "Thymus", True),
        ("skin", "Skin", True),
        ("spleen", "Spleen", True),
        ("liver", "Liver", False),
        ("skeletal_muscle", "Pooled muscle", False),
        ("soleus", "Soleus", False),
        ("kidney", "Kidney", False),
        ("tibialis_anterior", "Tibialis anterior", False),
        ("gastrocnemius", "Gastrocnemius", False),
        ("adrenal_gland", "Adrenal gland", False),
        ("eye", "Eye", False),
    ]
    method_headers = [
        ("Individual genes", "Permutation + SHAP", CORAL, PALE_CORAL),
        ("Reactome groups", "Grouped permutation + SHAP", GOLD, PALE_GOLD),
        ("Consensus ranking", "Compact gene panels", BLUE, PALE_BLUE),
    ]
    tissue_x = 0.55
    column_x = [3.06, 6.29, 9.52]
    column_w = 2.96

    _add_text(slide, "TISSUE", tissue_x, 2.01, 2.18, 0.21, size=9.0, color=GRAY, bold=True, margin=0)
    for x, (heading, method, color, _fill) in zip(column_x, method_headers):
        _add_rule(slide, x, 1.98, column_w, color, 0.035)
        _add_text(slide, heading, x, 2.09, column_w, 0.25, size=12.0, color=NAVY, bold=True, align=PP_ALIGN.CENTER, margin=0)
        _add_text(slide, method, x, 2.35, column_w, 0.20, size=8.7, color=GRAY, align=PP_ALIGN.CENTER, margin=0)

    row_top = 2.66
    row_h = 0.36
    for row_index, (tissue, label, convergent) in enumerate(rows):
        y = row_top + row_index * row_h
        if convergent:
            background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.47), Inches(y), Inches(12.05), Inches(row_h - 0.02))
            _set_fill(background, "EDF5F3")
            background.line.fill.background()
            _add_rule(slide, 0.47, y, 0.055, TEAL, row_h - 0.02)
        else:
            background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.47), Inches(y), Inches(12.05), Inches(row_h - 0.02))
            _set_fill(background, "F7F9FA" if row_index % 2 else WHITE)
            background.line.fill.background()
        if row_index == 3:
            _add_rule(slide, 0.47, y - 0.045, 12.05, "C9D4DA", 0.018)
        _add_text(slide, label, tissue_x, y + 0.03, 2.20, row_h - 0.07, size=10.5, color=NAVY if convergent else DARK, bold=convergent, valign=MSO_ANCHOR.MIDDLE, margin=0)

        for method_index, (x, (_, _, color, fill)) in enumerate(zip(column_x, method_headers)):
            values = method_counts[method_index].get(tissue)
            if values is None:
                _add_text(slide, "-", x, y + 0.02, column_w, row_h - 0.06, size=11.0, color="B3BDC3", align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
                continue
            total, promoted, reinforced = values
            _add_panel(slide, x, y + 0.035, column_w, row_h - 0.09, fill=fill, line=fill, radius=False)
            _add_text(
                slide,
                f"{total} | {promoted} promoted, {reinforced} reinforced",
                x + 0.04,
                y + 0.04,
                column_w - 0.08,
                row_h - 0.10,
                size=9.3,
                color=color,
                bold=True,
                align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE,
                margin=0,
            )

    _add_text(slide, "FOUND BY ALL THREE", 0.60, 6.77, 2.18, 0.20, size=8.5, color=TEAL, bold=True, margin=0)
    _add_text(slide, "Results below the divider came from one analysis.", 3.05, 6.73, 9.38, 0.26, size=10.8, color=GRAY, align=PP_ALIGN.RIGHT, margin=0)
    _add_source(slide, "Counts refer to tissue-feature associations. Genes and pathways are different evidence units.")


def _add_annotation_bar(slide, label: str, counts: list[int], total: int, y: float):
    colors = [GREEN, BLUE, ORANGE, PURPLE]
    bar_x = 3.25
    bar_w = 8.10
    _add_text(slide, label, 0.60, y + 0.04, 2.12, 0.33, size=13.5, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE, margin=0)
    _add_text(slide, f"n = {total}", 2.42, y + 0.05, 0.67, 0.31, size=10.0, color=GRAY, align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE, margin=0)
    cursor = bar_x
    for count, color in zip(counts, colors):
        if count == 0:
            continue
        width = bar_w * count / total
        segment = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cursor), Inches(y), Inches(width), Inches(0.43))
        _set_fill(segment, color)
        segment.line.fill.background()
        if width >= 0.30:
            _add_text(slide, str(count), cursor, y + 0.02, width, 0.36, size=10.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
        cursor += width


def _slide_literature_annotation_summary(slide):
    _add_slide_title(
        slide,
        "Literature review",
        "Feature selection and literature review answer separate questions",
        "Selection compares classifiers. Literature review compares each candidate with prior evidence.",
    )

    process_rows = [
        (2.05, "MODEL COMPARISON", "Real-only importance", "Real + synthetic importance", "Promoted or reinforced", BLUE),
        (2.68, "EVIDENCE REVIEW", "Candidate + tissue + direction", "Search prior studies", "Evidence category", GOLD),
    ]
    for y, label, first, second, third, color in process_rows:
        _add_text(slide, label, 0.55, y + 0.12, 1.54, 0.20, size=8.7, color=color, bold=True, margin=0)
        positions = [(2.18, first, 2.72), (5.51, second, 2.55), (8.67, third, 3.24)]
        for index, (x, text, width) in enumerate(positions):
            _add_panel(slide, x, y, width, 0.48, fill=WHITE, line=color, radius=False)
            _add_text(slide, text, x + 0.08, y + 0.07, width - 0.16, 0.32, size=10.6, color=DARK, bold=index == 2, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
            if index < len(positions) - 1:
                next_x = positions[index + 1][0]
                _add_arrow(slide, x + width + 0.13, y + 0.15, next_x - (x + width) - 0.26, 0.18, color)

    _add_text(slide, "Literature classifications among retained candidates", 0.58, 3.43, 6.90, 0.31, size=16.5, color=NAVY, bold=True, margin=0)
    _add_annotation_bar(slide, "Individual genes", [9, 9, 1, 2], 21, 3.91)
    _add_annotation_bar(slide, "Reactome groups", [6, 2, 2, 0], 10, 4.60)
    _add_annotation_bar(slide, "Consensus ranking", [22, 19, 4, 4], 49, 5.29)

    legend = [
        ("Aligning", GREEN, 2.18),
        ("Complementary", BLUE, 4.13),
        ("Ambiguous", ORANGE, 6.65),
        ("Unmatched", PURPLE, 8.68),
    ]
    for label, color, x in legend:
        _add_rule(slide, x, 6.13, 0.16, color, 0.12)
        _add_text(slide, label, x + 0.24, 6.06, 1.70, 0.27, size=10.8, color=color, bold=True, margin=0)
    _add_source(slide, "Annotations used the fixed tissue, FLT-GC direction and candidate identity; detailed rationale and sources are in Tables S16, S22 and S23.")


def _slide_thymus_integrated_result(slide):
    _add_slide_title(
        slide,
        "Thymus",
        "Thymus: all three analyses point to lower cell-cycle activity in flight",
        "Genes, Reactome pathways and the compact panel all show the same direction.",
    )
    figure = PAPER_DIR / "figures/figure_3_thymus_biology.png"
    _add_picture_contain(slide, figure, 0.34, 1.98, 8.35, 4.74, alt="Thymus FLT-GC gene effects and shared cell-cycle processes")
    _add_panel(slide, 0.34, 1.98, 8.35, 0.56, fill=WHITE, line=WHITE, radius=False)
    _add_text(slide, "Real OSDR effects and shared Reactome processes", 0.57, 2.18, 7.89, 0.23, size=11.2, color=DARK, bold=True, margin=0)

    _add_rule(slide, 8.92, 2.02, 0.015, "D2DCE1", 4.66)
    method_rows = [
        (2.07, "INDIVIDUAL GENES", "15", "7 promoted | 8 reinforced", "Most retained genes are FLT lower.", CORAL),
        (3.12, "REACTOME GROUPS", "7", "2 promoted | 5 reinforced", "Six pathways describe mitosis, DNA replication or APC/C control.", GOLD),
        (4.35, "CONSENSUS RANKING", "16", "13 promoted | 3 reinforced", "Thirteen genes are FLT lower, including CDK1, TOP2A and AURKA.", BLUE),
    ]
    for y, label, count, status, detail, color in method_rows:
        _add_text(slide, label, 9.19, y, 2.92, 0.20, size=8.8, color=color, bold=True, margin=0)
        _add_text(slide, count, 9.17, y + 0.27, 0.70, 0.47, size=28, color=color, bold=True, margin=0)
        _add_text(slide, status, 9.91, y + 0.29, 2.58, 0.25, size=10.8, color=DARK, bold=True, margin=0)
        _add_text(slide, detail, 9.91, y + 0.56, 2.58, 0.42, size=10.3, color=GRAY, margin=0)
        if y < 4.0:
            _add_rule(slide, 9.18, y + 0.95, 3.30, "E0E6E9", 0.012)

    _add_panel(slide, 9.15, 5.56, 3.38, 1.03, fill=NAVY, line=NAVY, radius=False)
    _add_text(slide, "BIOLOGICAL READING", 9.40, 5.72, 2.88, 0.20, size=8.4, color="AFC6DC", bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "Lower thymic proliferative activity", 9.40, 5.98, 2.88, 0.40, size=13.0, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
    _add_source(slide, "Individual genes: Table S19. Reactome groups: Table S23. Consensus ranking: Table S10.")


def _add_tissue_method_row(slide, x: float, y: float, w: float, label: str, text: str, color: str):
    _add_text(slide, label, x, y, 1.42, 0.22, size=8.7, color=color, bold=True, margin=0)
    _add_text(slide, text, x + 1.49, y - 0.02, w - 1.49, 0.44, size=11.3, color=DARK, valign=MSO_ANCHOR.TOP, margin=0)
    _add_rule(slide, x, y + 0.49, w, "E0E6E9", 0.011)


def _slide_skin_spleen_integrated_result(slide):
    _add_slide_title(
        slide,
        "Biological interpretation",
        "Skin and spleen show different flight-associated patterns",
        "Skin: regulated cell death  |  Spleen: structural remodeling",
    )
    _add_rule(slide, 6.66, 2.02, 0.015, "D3DCE1", 4.60)

    _add_rule(slide, 0.50, 2.02, 5.83, PURPLE, 0.045)
    _add_text(slide, "SKIN", 0.50, 2.19, 1.04, 0.22, size=9.2, color=PURPLE, bold=True, margin=0)
    _add_text(slide, "Interferon-linked regulated cell death", 0.50, 2.50, 5.83, 0.42, size=18.0, color=NAVY, bold=True, margin=0)
    _add_tissue_method_row(slide, 0.53, 3.10, 5.72, "GENE EVIDENCE", "PLSCR1 was higher in flight and retained by two analyses", CORAL)
    _add_tissue_method_row(slide, 0.53, 3.73, 5.72, "PATHWAY EVIDENCE", "Two necroptosis pathways were higher in flight", GOLD)
    _add_tissue_method_row(slide, 0.53, 4.36, 5.72, "CONTEXT", "PLSCR1 is interferon inducible", BLUE)
    _add_panel(slide, 0.53, 5.22, 5.72, 1.03, fill="F2EEF7", line="D9CEE8", radius=False)
    _add_text(slide, "COMPLEMENTARY HYPOTHESIS", 0.79, 5.42, 5.20, 0.20, size=8.8, color=PURPLE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "Direct spaceflight evidence for skin necroptosis is limited", 0.79, 5.67, 5.20, 0.34, size=12.2, color=DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)

    _add_rule(slide, 6.99, 2.02, 5.84, TEAL, 0.045)
    _add_text(slide, "SPLEEN", 6.99, 2.19, 1.30, 0.22, size=9.2, color=TEAL, bold=True, margin=0)
    _add_text(slide, "Matrix, adhesion and cytoskeletal remodeling", 6.99, 2.50, 5.83, 0.42, size=18.0, color=NAVY, bold=True, margin=0)
    _add_tissue_method_row(slide, 7.02, 3.10, 5.72, "GENE EVIDENCE", "LOXL1 was retained by two analyses; RAI14, PTPRK and MYL9 were also flight higher", CORAL)
    _add_tissue_method_row(slide, 7.02, 3.73, 5.72, "PATHWAY EVIDENCE", "No shared structural Reactome pathway", GOLD)
    _add_tissue_method_row(slide, 7.02, 4.36, 5.72, "CONTEXT", "The genes span matrix, adhesion and cytoskeletal functions", BLUE)
    _add_panel(slide, 7.02, 5.22, 5.72, 1.03, fill=PALE_BLUE, line="CADAE7", radius=False)
    _add_text(slide, "COMPLEMENTARY HYPOTHESIS", 7.28, 5.42, 5.20, 0.20, size=8.8, color=BLUE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "Matrix, adhesion and cytoskeletal remodeling", 7.28, 5.67, 5.20, 0.34, size=12.2, color=DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
    _add_source(slide, "All displayed gene results were higher in flight. Direct spaceflight evidence for both complementary hypotheses was limited.")


def _slide_method_specific_hypotheses(slide):
    _add_slide_title(
        slide,
        "Follow-up findings",
        "Other tissue results came from one analysis",
        "The other two analyses did not retain matching genes or pathways.",
    )
    columns = [
        {
            "x": 0.48,
            "color": CORAL,
            "tissue": "Liver",
            "method": "INDIVIDUAL GENES",
            "count": "4 FLT-lower genes",
            "genes": "GRB10, PPIC, H2-DMA, GTF2A2",
            "reading": "The genes helped FLT-GC prediction but did not form a shared Reactome pathway.",
        },
        {
            "x": 4.49,
            "color": GOLD,
            "tissue": "Soleus",
            "method": "CONSENSUS RANKING",
            "count": "5 reinforced genes",
            "genes": "Lower: BDH1, ECH1, BNIP3, DECR1\nHigher: TPM1",
            "reading": "The panel points to mitochondrial turnover and fatty-acid metabolism; prediction benefit was mixed.",
        },
        {
            "x": 8.50,
            "color": BLUE,
            "tissue": "Pooled muscle",
            "method": "CONSENSUS RANKING",
            "count": "12 genes: 4 promoted, 8 reinforced",
            "genes": "Interferon, stress-response and sialic-acid-related candidates",
            "reading": "Prediction improved after adding synthetic samples, but the individual-gene analysis did not retain a matching candidate.",
        },
    ]
    for index, column in enumerate(columns):
        x = column["x"]
        if index:
            _add_rule(slide, x - 0.22, 2.03, 0.012, "D9E0E4", 4.55)
        _add_rule(slide, x, 2.03, 3.56, column["color"], 0.045)
        _add_text(slide, column["tissue"], x, 2.23, 3.56, 0.38, size=21, color=NAVY, bold=True, margin=0)
        _add_text(slide, column["method"], x, 2.72, 3.56, 0.22, size=8.8, color=column["color"], bold=True, margin=0)
        count_size = 12.7 if len(column["count"]) > 24 else 15.2
        _add_text(slide, column["count"], x, 3.13, 3.56, 0.36, size=count_size, color=column["color"], bold=True, margin=0)
        _add_rule(slide, x, 3.62, 3.56, "DDE4E8", 0.014)
        _add_text(slide, "RETAINED FEATURES", x, 3.86, 2.10, 0.20, size=8.5, color=GRAY, bold=True, margin=0)
        _add_text(slide, column["genes"], x, 4.16, 3.50, 0.74, size=12.4, color=DARK, bold=True, margin=0)
        _add_text(slide, "INTERPRETATION", x, 5.15, 1.55, 0.20, size=8.5, color=GRAY, bold=True, margin=0)
        _add_text(slide, column["reading"], x, 5.45, 3.50, 0.91, size=11.8, color=DARK, margin=0)
    _add_source(slide, "Liver: Table S19. Soleus and pooled muscle: Table S10.")


def _slide_10(slide):
    _add_slide_title(
        slide,
        "Literature interpretation",
        "How a gene was selected and what the literature says are separate",
        "We labeled all 49 combined-panel associations by selection pattern and by prior biological evidence.",
    )

    steps = [
        ("01", "Fix the candidate", "Gene, tissue, and FLT direction", BLUE),
        ("02", "Search", "Spaceflight and mechanism literature", TEAL),
        ("03", "Compare", "Observed result with published evidence", ORANGE),
        ("04", "Record", "Category, rationale, and source IDs", PURPLE),
    ]
    for index, (number, heading, detail, color) in enumerate(steps):
        x = 0.48 + index * 3.16
        if index:
            _add_rule(slide, x - 0.28, 2.02, 0.012, "D8DFE3", 0.95)
        _add_text(slide, number, x, 2.07, 0.40, 0.24, size=10.8, color=color, bold=True, margin=0)
        _add_text(slide, heading, x + 0.48, 2.02, 2.30, 0.30, size=14.5, color=NAVY, bold=True, margin=0)
        _add_text(slide, detail, x + 0.48, 2.40, 2.30, 0.42, size=10.8, color=GRAY, margin=0)
        _add_rule(slide, x, 2.95, 2.72, color, 0.025)

    annotations = pd.read_csv(
        PAPER_DIR / "source_data/table_s16_promoted_gene_literature_annotations.tsv",
        sep="\t",
    )
    counts = annotations["literature_classification"].value_counts().to_dict()
    selection_counts = annotations["selection_status"].value_counts().to_dict()
    direct = int(
        annotations["evidence_scope"]
        .eq("direct_same_gene_same_tissue_same_direction")
        .sum()
    )
    if len(annotations) != 49 or direct != 5:
        raise ValueError("Unexpected consensus literature inventory")

    _add_text(slide, "Selection status", 0.55, 3.35, 3.58, 0.34, size=17, color=NAVY, bold=True)
    selection_rows = [
        ("Promoted", selection_counts.get("promoted", 0), "Selected only after adding the synthetic ranking", CORAL),
        ("Reinforced", selection_counts.get("reinforced", 0), "Selected by both real-only and combined rankings", TEAL),
    ]
    for index, (label, count, detail, color) in enumerate(selection_rows):
        y = 4.02 + index * 1.05
        if index:
            _add_rule(slide, 0.55, y - 0.18, 3.58, "E0E5E8", 0.012)
        _add_text(slide, str(count), 0.55, y - 0.05, 0.68, 0.52, size=29, color=color, bold=True, margin=0)
        _add_text(slide, label, 1.38, y - 0.02, 2.65, 0.30, size=15, color=color, bold=True, margin=0)
        _add_text(slide, detail, 1.38, y + 0.31, 2.70, 0.38, size=11.5, color=DARK, margin=0)

    _add_rule(slide, 4.45, 3.33, 0.015, "D5DDE2", 2.77)
    _add_text(slide, "Literature interpretation", 4.82, 3.35, 7.52, 0.34, size=17, color=NAVY, bold=True)
    _add_text(slide, "Class", 4.84, 3.86, 1.55, 0.22, size=9.5, color=GRAY, bold=True, margin=0)
    _add_text(slide, "Count", 6.86, 3.86, 0.70, 0.22, size=9.5, color=GRAY, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "Meaning", 7.69, 3.86, 4.31, 0.22, size=9.5, color=GRAY, bold=True, margin=0)
    _add_rule(slide, 4.82, 4.13, 7.54, NAVY, 0.020)
    category_rows = [
        ("Aligning", counts.get("aligning", 0), "Direct or same-tissue process agreement", GREEN),
        ("Complementary", counts.get("complementary", 0), "Related process or mechanism", BLUE),
        ("Ambiguous", counts.get("ambiguous", 0), "Mixed or context-dependent evidence", ORANGE),
        ("Unmatched", counts.get("unmatched", 0), "No sufficiently specific match found", PURPLE),
    ]
    for index, (label, count, detail, color) in enumerate(category_rows):
        y = 4.27 + index * 0.48
        if index % 2 == 0:
            shade = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.82), Inches(y - 0.04), Inches(7.54), Inches(0.42))
            _set_fill(shade, "F5F7F8")
            shade.line.fill.background()
        _add_rule(slide, 4.86, y + 0.11, 0.12, color, 0.12)
        _add_text(slide, label, 5.10, y, 1.68, 0.34, size=12.6, color=color, bold=True, valign=MSO_ANCHOR.MIDDLE, margin=0)
        _add_text(slide, str(count), 6.87, y, 0.66, 0.34, size=14.0, color=DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
        _add_text(slide, detail, 7.69, y, 4.41, 0.34, size=11.8, color=DARK, valign=MSO_ANCHOR.MIDDLE, margin=0)

    _add_rule(slide, 0.55, 6.35, 12.02, "CED7DD", 0.018)
    _add_text(
        slide,
        "Selection describes our analysis. The literature label describes previous biological evidence.",
        0.65,
        6.49,
        11.82,
        0.28,
        size=13.5,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_source(
        slide,
        "Detailed gene-level rationale and literature sources are provided in Supplementary Tables S16-S17.",
    )


def _inventory_gene_segments(
    rows: pd.DataFrame,
    annotation_lookup: dict[tuple[str, str, str], str],
) -> list[tuple[str, dict]]:
    class_colors = {
        "aligning": GREEN,
        "complementary": BLUE,
        "ambiguous": ORANGE,
        "unmatched": PURPLE,
    }
    ordered = rows.sort_values(["real_meta_fdr", "symbol"])
    segments: list[tuple[str, dict]] = []
    for index, row in enumerate(ordered.itertuples(index=False)):
        if index:
            segments.append((", ", {"size": 9.0, "color": MID_GRAY}))
        key = (row.analysis_scope, row.tissue, row.gene)
        literature_class = annotation_lookup.get(key)
        if literature_class not in class_colors:
            raise ValueError(f"Missing synthetic-informed annotation for {key}")
        color = class_colors[literature_class]
        segments.append(
            (
                row.symbol,
                {
                    "size": 9.0,
                    "color": color,
                    "bold": True,
                    "italic": True,
                },
            )
        )
    return segments


def _add_gene_inventory_block(
    slide,
    inventory: pd.DataFrame,
    annotation_lookup: dict[tuple[str, str, str], str],
    *,
    scope: str,
    tissue: str,
    label: str,
    x: float,
    y: float,
    w: float,
    h: float,
) -> None:
    rows = inventory[
        inventory["analysis_scope"].eq(scope) & inventory["tissue"].eq(tissue)
    ]
    if rows.empty:
        raise ValueError(f"No synthetic-informed genes for {scope}/{tissue}")

    _add_rule(slide, x, y, w, "DDE4E8", 0.018)
    count_label = "gene" if len(rows) == 1 else "genes"
    label_size = 10.8 if len(label.replace("\n", "")) > 12 else 11.8
    _add_rich_text(
        slide,
        [
            (f"{label}\n", {"size": label_size, "color": NAVY, "bold": True}),
            (f"{len(rows)} {count_label}", {"size": 9.0, "color": MID_GRAY}),
        ],
        x + 0.06,
        y + 0.07,
        1.36,
        h - 0.10,
        margin=0,
    )

    groups = []
    for direction, direction_label, direction_color in [
        ("FLT_higher", "FLT higher", CORAL),
        ("FLT_lower", "FLT lower", TEAL),
    ]:
        for selection, selection_label, selection_color in [
            ("synthetic_promoted", "Promoted", CORAL),
            ("reinforced_real_and_synthetic", "Reinforced", TEAL),
        ]:
            group_rows = rows[
                rows["flt_gc_direction"].eq(direction)
                & rows["selection_interpretation"].eq(selection)
            ]
            if group_rows.empty:
                continue
            symbol_characters = int(group_rows["symbol"].str.len().sum()) + 2 * max(0, len(group_rows) - 1)
            estimated_lines = max(1, int(np.ceil(symbol_characters / 38)))
            groups.append(
                (
                    direction_label,
                    direction_color,
                    selection_label,
                    selection_color,
                    group_rows,
                    1.0 + 0.75 * (estimated_lines - 1),
                )
            )

    available_height = h - 0.11
    weight_total = sum(group[-1] for group in groups)
    row_y = y + 0.06
    for index, (direction_label, direction_color, selection_label, selection_color, group_rows, weight) in enumerate(groups):
        row_h = available_height * weight / weight_total
        if index:
            _add_rule(slide, x + 1.44, row_y, w - 1.50, "E2E7EA", 0.010)
        _add_text(
            slide,
            direction_label,
            x + 1.47,
            row_y + 0.01,
            0.72,
            row_h - 0.02,
            size=8.0,
            color=direction_color,
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
            margin=0,
        )
        _add_text(
            slide,
            selection_label,
            x + 2.22,
            row_y + 0.01,
            0.78,
            row_h - 0.02,
            size=7.9,
            color=selection_color,
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
            margin=0,
        )
        _add_rich_text(
            slide,
            _inventory_gene_segments(group_rows, annotation_lookup),
            x + 3.04,
            row_y + 0.01,
            w - 3.12,
            row_h - 0.02,
            size=9.0,
            valign=MSO_ANCHOR.MIDDLE,
            margin=0,
        )
        row_y += row_h


def _slide_all_tissue_coverage(slide):
    _add_slide_title(
        slide,
        "Analysis coverage",
        "The screen covered all 27 completed tissue analyses",
        "The biological discussion narrows only after reporting synthetic-informed, real-only, and null outcomes.",
    )
    summary = pd.read_csv(
        PAPER_DIR / "source_data/table_s12_bh_fdr_tissue_summary.tsv",
        sep="\t",
    )
    if len(summary) != 27:
        raise ValueError("Expected 27 completed tissue analysis units")

    synthetic = summary.loc[
        summary["n_reinforced_real_and_synthetic"]
        .add(summary["n_synthetic_promoted"])
        .gt(0)
    ]
    real_only = summary.loc[
        summary["n_bh_fdr_genes"].gt(0)
        & ~summary.index.to_series().isin(synthetic.index)
    ]
    null = summary.loc[summary["n_bh_fdr_genes"].eq(0)]
    if (len(synthetic), len(real_only), len(null)) != (10, 5, 12):
        raise ValueError("Unexpected all-tissue coverage partition")

    display = {
        "skeletal_muscle": "Skeletal muscle (pooled)",
        "edl": "EDL",
    }

    def names(frame: pd.DataFrame) -> list[str]:
        return [
            display.get(value, value.replace("_", " ").title())
            for value in frame["tissue"]
        ]

    columns = [
        (
            0.43,
            "10",
            "Synthetic-informed",
            names(synthetic),
            TEAL,
            PALE_TEAL,
            "Promoted or reinforced BH-FDR gene",
        ),
        (
            4.49,
            "5",
            "Real BH-FDR only",
            names(real_only),
            BLUE,
            PALE_BLUE,
            "Real association, no synthetic-informed selection",
        ),
        (
            8.55,
            "12",
            "No BH-FDR gene",
            names(null),
            ORANGE,
            PALE_GOLD,
            "No BH-FDR gene in the 974-gene panel",
        ),
    ]
    for index, (x, count, heading, tissue_names, color, fill, detail) in enumerate(columns):
        if index:
            _add_rule(slide, x - 0.18, 2.05, 0.012, "D8DFE3", 4.58)
        _add_rule(slide, x, 2.04, 3.65, color, 0.035)
        _add_text(slide, count, x + 0.02, 2.22, 0.70, 0.55, size=31, color=color, bold=True, margin=0)
        _add_text(slide, heading, x + 0.84, 2.27, 2.72, 0.38, size=16.0, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE, margin=0)
        _add_text(slide, detail, x + 0.02, 2.92, 3.48, 0.42, size=10.5, color=GRAY, valign=MSO_ANCHOR.MIDDLE, margin=0)
        _add_rule(slide, x, 3.50, 3.58, "D7DEE2", 0.014)
        _add_text(slide, "TISSUES / MUSCLE GROUPS", x + 0.02, 3.66, 1.80, 0.19, size=8.3, color=GRAY, bold=True, margin=0)

        list_top = 3.96
        list_height = 2.42
        row_height = min(0.36, list_height / len(tissue_names))
        list_size = 9.8 if len(tissue_names) >= 12 else 10.4 if len(tissue_names) >= 10 else 11.2
        for index, tissue_name in enumerate(tissue_names):
            row_y = list_top + index * row_height
            _add_rule(slide, x + 0.03, row_y + (row_height - 0.045) / 2, 0.13, color, 0.045)
            _add_text(
                slide,
                tissue_name,
                x + 0.25,
                row_y,
                3.25,
                row_height,
                size=list_size,
                color=DARK,
                valign=MSO_ANCHOR.MIDDLE,
                margin=0,
            )
    _add_source(slide, "Twenty-two canonical tissues plus EDL, gastrocnemius, quadriceps, soleus and tibialis anterior; complete counts in Table S12.")


def _slide_matched_coverage(slide):
    _add_slide_title(
        slide,
        "Tissue coverage",
        "We compared the same training strategies in all 27 tissue analyses",
        "Better FLT-GC prediction and interpretable genes did not always occur together.",
    )
    utility = pd.read_csv(
        PAPER_DIR / "source_data/table_s18_matched_all_gene_utility.tsv",
        sep="\t",
    )
    candidates = pd.read_csv(
        PAPER_DIR / "source_data/table_s19_matched_all_gene_candidates.tsv",
        sep="\t",
    )
    utility = utility.loc[utility["arm"].eq("real_plus_generated")].copy()
    candidate_tissues = set(candidates["tissue"].astype(str))
    with_genes = utility.loc[utility["tissue"].isin(candidate_tissues)]
    utility_only = utility.loc[
        utility["joint_mean_all_metrics_nonworse"].astype(bool)
        & ~utility["tissue"].isin(candidate_tissues)
    ]
    did_not_pass = utility.loc[
        ~utility["joint_mean_all_metrics_nonworse"].astype(bool)
    ]
    if (len(with_genes), len(utility_only), len(did_not_pass)) != (4, 14, 9):
        raise ValueError("Unexpected matched all-gene coverage partition")

    display = {
        "skeletal_muscle": "Skeletal muscle (pooled)",
        "edl": "EDL",
        "brown_adipose_tissue": "Brown adipose tissue",
        "white_adipose_tissue": "White adipose tissue",
    }

    def names(frame: pd.DataFrame) -> list[str]:
        return [
            display.get(value, value.replace("_", " ").title())
            for value in frame["tissue"].astype(str)
        ]

    columns = [
        (0.43, "4", "Prediction + genes", names(with_genes), CORAL, PALE_CORAL, "Prediction held or improved; at least one measured-data gene helped the model"),
        (4.49, "14", "Prediction only", names(utility_only), TEAL, PALE_TEAL, "Prediction held or improved; no individual gene passed the interpretation screen"),
        (8.55, "9", "Mixed result", names(did_not_pass), BLUE, PALE_BLUE, "At least one prediction measure worsened"),
    ]
    for column_index, (x, count, heading, tissue_names, color, _fill, detail) in enumerate(columns):
        if column_index:
            _add_rule(slide, x - 0.18, 2.05, 0.012, "D8DFE3", 4.58)
        _add_rule(slide, x, 2.04, 3.65, color, 0.035)
        _add_text(slide, count, x + 0.02, 2.22, 0.70, 0.55, size=31, color=color, bold=True, margin=0)
        _add_text(slide, heading, x + 0.84, 2.27, 2.72, 0.38, size=16.0, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE, margin=0)
        _add_text(slide, detail, x + 0.02, 2.92, 3.48, 0.42, size=10.5, color=GRAY, valign=MSO_ANCHOR.MIDDLE, margin=0)
        _add_rule(slide, x, 3.50, 3.58, "D7DEE2", 0.014)
        _add_text(slide, "TISSUES / MUSCLE GROUPS", x + 0.02, 3.66, 1.80, 0.19, size=8.3, color=GRAY, bold=True, margin=0)
        list_top = 3.96
        list_height = 2.42
        row_height = min(0.36, list_height / len(tissue_names))
        list_size = 9.6 if len(tissue_names) >= 12 else 10.3
        for row_index, tissue_name in enumerate(tissue_names):
            row_y = list_top + row_index * row_height
            _add_rule(slide, x + 0.03, row_y + (row_height - 0.045) / 2, 0.13, color, 0.045)
            _add_text(slide, tissue_name, x + 0.25, row_y, 3.25, row_height, size=list_size, color=DARK, valign=MSO_ANCHOR.MIDDLE, margin=0)
    _add_source(slide, "Real-plus-synthetic comparison across 22 tissue categories and five anatomical muscle groups.")


def _slide_11(slide):
    _add_slide_title(
        slide,
        "Combined-ranking gene inventory",
        "Ten tissue analyses had genes prioritized by the combined ranking",
        "All 49 associations were supported in measured OSDR data. Synthetic data changed their priority, not the biological comparison.",
    )
    inventory = pd.read_csv(
        PAPER_DIR / "source_data/table_s10_synthetic_informed_bh_fdr_genes.tsv",
        sep="\t",
    )
    annotations = pd.read_csv(
        PAPER_DIR / "source_data/table_s16_promoted_gene_literature_annotations.tsv",
        sep="\t",
    )
    if len(inventory) != 49 or len(annotations) != 49:
        raise ValueError("Unexpected consensus gene inventory")
    if inventory[["analysis_scope", "tissue"]].drop_duplicates().shape[0] != 10:
        raise ValueError("Expected 10 tissue analyses in consensus inventory")
    annotation_lookup = {
        (row.analysis_scope, row.tissue, row.gene): row.literature_classification
        for row in annotations.itertuples(index=False)
    }

    _add_text(slide, "Rows: FLT direction + selection status", 0.45, 1.86, 2.66, 0.25, size=9.5, color=NAVY, bold=True)
    legend = [
        ("aligning", GREEN, 3.02, 0.84),
        ("complementary", BLUE, 4.02, 1.15),
        ("ambiguous", ORANGE, 5.37, 0.96),
        ("unmatched", PURPLE, 6.53, 0.96),
    ]
    for label, color, x, width in legend:
        _add_rule(slide, x, 1.91, 0.14, color, 0.05)
        _add_text(slide, label, x + 0.19, 1.85, width, 0.25, size=9.8, color=color, bold=True)
    _add_text(slide, "Gene color = literature | Non-empty rows shown", 8.64, 1.85, 4.20, 0.25, size=9.8, color=MID_GRAY, italic=True, align=PP_ALIGN.RIGHT)

    left_blocks = [
        ("canonical_tissue", "thymus", "Thymus", 2.24, 1.45),
        ("canonical_tissue", "spleen", "Spleen", 3.78, 0.67),
        ("canonical_tissue", "kidney", "Kidney", 4.54, 0.61),
        ("skeletal_muscle_group", "gastrocnemius", "Gastrocnemius", 5.24, 0.61),
        ("canonical_tissue", "eye", "Eye", 5.94, 0.51),
    ]
    right_blocks = [
        ("canonical_tissue", "skeletal_muscle", "Skeletal muscle\n(pooled)", 2.24, 1.25),
        ("skeletal_muscle_group", "soleus", "Soleus", 3.58, 0.68),
        ("skeletal_muscle_group", "tibialis_anterior", "Tibialis anterior", 4.35, 0.69),
        ("canonical_tissue", "adrenal_gland", "Adrenal gland", 5.13, 0.61),
        ("canonical_tissue", "skin", "Skin", 5.83, 0.55),
    ]
    for scope, tissue, label, y, h in left_blocks:
        _add_gene_inventory_block(
            slide,
            inventory,
            annotation_lookup,
            scope=scope,
            tissue=tissue,
            label=label,
            x=0.43,
            y=y,
            w=6.15,
            h=h,
        )
    for scope, tissue, label, y, h in right_blocks:
        _add_gene_inventory_block(
            slide,
            inventory,
            annotation_lookup,
            scope=scope,
            tissue=tissue,
            label=label,
            x=6.76,
            y=y,
            w=6.15,
            h=h,
        )
    _add_source(slide, "Rows show FLT direction and selection status. Gene color separately shows the literature interpretation.")


def _slide_12(slide):
    _add_slide_title(
        slide,
        "Thymus",
        "Thymus is strongest across all three analyses",
        "Genes, Reactome groups, and the compact panel converge on lower cell-cycle activity in flight.",
    )
    figure = PAPER_DIR / "figures/figure_3_thymus_biology.png"
    _add_picture_contain(slide, figure, 0.35, 1.93, 8.55, 4.78, alt="Thymus gene effects and Reactome processes")
    _add_panel(slide, 0.35, 1.93, 8.55, 0.48, fill=WHITE, line=WHITE, radius=False)
    _add_text(
        slide,
        "Thymus genes point to lower cell-cycle activity in flight",
        0.58,
        2.08,
        8.02,
        0.24,
        size=11.5,
        color=DARK,
        bold=True,
        margin=0,
    )
    _add_panel(slide, 9.10, 2.02, 3.82, 4.57, fill=PALE_CORAL, line="E8C9C2", radius=False)
    _add_text(slide, "ALL-GENE MODEL", 9.43, 2.27, 2.95, 0.23, size=9.2, color=CORAL, bold=True, margin=0)
    _add_text(slide, "15", 9.43, 2.56, 1.10, 0.65, size=35, color=CORAL, bold=True)
    _add_text(slide, "genes that also helped\nFLT-GC prediction", 10.38, 2.58, 2.12, 0.62, size=12.7, color=DARK, bold=True)
    _add_text(slide, "7 synthetic-added | 8 shared", 9.44, 3.34, 3.05, 0.30, size=13.2, color=TEAL, bold=True)
    _add_bullet_rows(
        slide,
        [
            "A broad Reactome pattern led by mitotic cell cycle",
            "Nine genes also appear in the 16-gene combined panel",
            "The combined ranking adds Cdk1, Top2a, Aurka, Ccne2, Pcna and Ccnf",
            "Lower mitotic renewal or fewer cycling thymocytes in flight",
        ],
        9.44,
        3.77,
        3.03,
        size=11.8,
        bullet_color=CORAL,
        row_h=0.61,
    )
    _add_source(slide, "Figure shows the combined panel; controlled-test genes and Reactome results are in Tables S19-S20. Context: Gridley et al. (2013); Horie et al. (2019).")


def _slide_13(slide):
    _add_slide_title(
        slide,
        "Soleus",
        "Soleus has a coherent secondary gene panel",
        "The five genes form a biological story, but synthetic profiles did not improve prediction consistently.",
    )
    figure = PAPER_DIR / "figures/figure_4_soleus_biology.png"
    _add_picture_contain(slide, figure, 0.34, 1.96, 9.15, 4.70, alt="Soleus gene effects and Reactome processes")
    _add_panel(slide, 9.63, 2.04, 3.25, 4.52, fill=PALE_TEAL, line="C7DDD8", radius=False)
    _add_text(slide, "5-gene combined panel", 9.87, 2.30, 2.77, 0.62, size=17, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "selected by both real and combined rankings", 9.90, 2.96, 2.70, 0.36, size=10.4, color=GRAY, align=PP_ALIGN.CENTER)
    _add_panel(slide, 9.92, 3.39, 2.65, 0.74, fill=WHITE, line=CORAL, radius=False)
    _add_text(slide, "Synthetic benefit varied", 10.03, 3.53, 2.44, 0.20, size=10.8, color=CORAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
    _add_text(slide, "Treat as a follow-up hypothesis", 10.03, 3.80, 2.44, 0.20, size=9.0, color=DARK, align=PP_ALIGN.CENTER, margin=0)
    _add_bullet_rows(
        slide,
        [
            "Lower Bdh1, Ech1, Bnip3 and Decr1",
            "Higher Tpm1",
            "Mitochondrial turnover and fatty-acid metabolism",
            "Panel-level hypothesis; no confirmed synthetic benefit",
        ],
        9.94,
        4.31,
        2.54,
        size=11.7,
        bullet_color=TEAL,
        row_h=0.52,
    )
    _add_source(slide, "Literature context: Gambara et al. (2017) and Stein et al. (2002).")


def _add_additional_finding_rows(
    slide,
    rows,
    *,
    first_label: str = "Promoted",
    second_label: str = "Reinforced",
    selection_heading: str = "Selection",
):
    _add_text(slide, "Analysis unit", 0.72, 2.00, 1.70, 0.28, size=11.5, color=GRAY, bold=True)
    _add_text(slide, selection_heading, 2.48, 2.00, 1.35, 0.28, size=11.2, color=GRAY, bold=True)
    _add_text(slide, "Genes and FLT direction", 3.93, 2.00, 3.00, 0.28, size=11.5, color=GRAY, bold=True)
    _add_text(slide, "Interpretation", 7.20, 2.00, 5.10, 0.28, size=11.5, color=GRAY, bold=True)
    for index, (tissue, promoted, reinforced, interpretation, color, fill) in enumerate(rows):
        y = 2.35 + index * 1.06
        _add_panel(slide, 0.43, y, 12.44, 0.88, fill=fill, line="DDE4E8", radius=False)
        _add_rule(slide, 0.43, y, 0.08, color, 0.88)
        _add_text(slide, tissue, 0.72, y + 0.20, 1.62, 0.44, size=15.0, color=color, bold=True, valign=MSO_ANCHOR.MIDDLE)
        _add_rule(slide, 2.42, y + 0.44, 4.42, "DDE4E8", 0.012)
        _add_text(slide, first_label, 2.48, y + 0.08, 1.33, 0.25, size=9.8, color=CORAL, bold=True, valign=MSO_ANCHOR.MIDDLE)
        _add_text(slide, promoted, 3.93, y + 0.05, 2.93, 0.35, size=10.3, color=DARK, valign=MSO_ANCHOR.MIDDLE)
        _add_text(slide, second_label, 2.48, y + 0.52, 1.33, 0.25, size=9.8, color=TEAL, bold=True, valign=MSO_ANCHOR.MIDDLE)
        _add_text(slide, reinforced, 3.93, y + 0.48, 2.93, 0.35, size=10.3, color=DARK, valign=MSO_ANCHOR.MIDDLE)
        _add_text(slide, interpretation, 7.20, y + 0.10, 5.20, 0.66, size=12.2, color=DARK, valign=MSO_ANCHOR.MIDDLE)


def _slide_14(slide):
    _add_slide_title(
        slide,
        "Additional tissue findings",
        "Liver, skin and spleen add narrower gene-level results",
        "Pooled muscle also benefited from synthetic training, while its gene panel remains a follow-up result.",
    )
    rows = [
        (
            "Liver",
            "Lower: Grb10, Ppic, H2-DMa, Gtf2a2",
            "No retained panel",
            "Four genes helped the model, but they did not form a clear shared Reactome pathway.",
            BLUE,
            PALE_BLUE,
        ),
        (
            "Skin",
            "Higher: Plscr1",
            "Higher: Plscr1 plus broader selected pathways",
            "Both gene-level analyses support Plscr1; broader cell-cycle and DNA-repair context comes from the gene panel.",
            PURPLE,
            "F0ECF6",
        ),
        (
            "Spleen",
            "Higher: Loxl1",
            "Higher: Loxl1, Rai14, Ptprk, Myl9",
            "Loxl1 appears in the controlled test; the four-gene adhesion and cytoskeletal panel is tentative.",
            TEAL,
            PALE_TEAL,
        ),
        (
            "Pooled muscle",
            "Prediction held or improved; no single gene passed the interpretation screen",
            "12-gene interferon and sialic-acid panel",
            "Adding synthetic samples improves prediction; individual genes come from the secondary ranking analysis.",
            ORANGE,
            PALE_GOLD,
        ),
    ]
    _add_additional_finding_rows(
        slide,
        rows,
        first_label="All-gene test",
        second_label="Combined ranking",
        selection_heading="Evidence source",
    )
    _add_source(slide, "All-gene candidates are supported in measured OSDR data and help prediction. Combined-ranking rows add panel context.")


def _slide_additional_2(slide):
    _add_slide_title(
        slide,
        "Additional combined-ranking findings",
        "Several smaller panels did not show a clear synthetic benefit",
        "These remain tissue-specific hypotheses from the secondary ranking analysis.",
    )
    rows = [
        (
            "Kidney",
            "Inpp4b higher",
            "Slc37a4 higher",
            "Renal phosphoinositide signaling and glucose-handling hypothesis; the pair had no shared Reactome enrichment.",
            GREEN,
            "ECF4ED",
        ),
        (
            "Adrenal gland",
            "Psmb8 lower",
            "Tspan4 lower",
            "No close prior match was found for these immune, proteostasis, or tissue-composition candidates.",
            MID_GRAY,
            "F2F4F5",
        ),
        (
            "Gastrocnemius",
            "Nfkbia higher; Fhl2 lower",
            "None",
            "NF-kappaB stress alignment plus an autophagy or myogenesis candidate; the two genes do not form a coherent pathway.",
            BLUE,
            PALE_BLUE,
        ),
        (
            "Tibialis anterior",
            "Cebpd higher",
            "Cdkn1a, St3gal5, Bnip3 higher",
            "Stress, cell-cycle, ganglioside-signaling and mitophagy candidates with mixed or complementary prior evidence.",
            TEAL,
            PALE_TEAL,
        ),
    ]
    _add_additional_finding_rows(slide, rows)
    _add_source(slide, "Secondary combined-ranking results and independent literature annotations are reported in Tables S10 and S16.")


def _slide_15(slide):
    _add_slide_title(
        slide,
        "Takeaways",
        "What the models can and cannot tell us",
    )
    columns = [
        (0.42, "1", "Computational evidence", "expiMap and the generative models summarize patterns in the measured RNA-seq data. They do not reveal mechanism directly.", BLUE, PALE_BLUE),
        (4.44, "2", "Candidates to follow up", "Agreement among genes, pathways, models, and prior studies narrows the list of tissues and targets.", TEAL, PALE_TEAL),
        (8.46, "3", "Experimental testing", "Wet-lab experiments must confirm the biological changes and distinguish regulation from shifts in cell composition.", ORANGE, PALE_GOLD),
    ]
    for index, (x, number, heading, body, color, _fill) in enumerate(columns):
        if index:
            _add_rule(slide, x - 0.20, 2.10, 0.012, "D8DFE3", 3.53)
        _add_rule(slide, x, 2.10, 3.55, color, 0.035)
        _add_text(slide, f"0{number}", x + 0.02, 2.33, 0.68, 0.40, size=21, color=color, bold=True, margin=0)
        _add_text(slide, heading, x + 0.86, 2.35, 2.46, 0.36, size=19, color=NAVY, bold=True, margin=0)
        _add_text(slide, body, x + 0.03, 3.12, 3.28, 2.10, size=16, color=DARK, margin=0)


def _slide_16(slide):
    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.38), Inches(SLIDE_W), Inches(SLIDE_H - 0.38))
    _set_fill(body, NAVY)
    body.line.fill.background()
    _add_text(slide, "Thank you", 0.65, 1.22, 5.40, 0.80, size=34, color=WHITE, bold=True)
    _add_text(slide, "Questions?", 0.67, 2.10, 4.0, 0.45, size=21, color="BFD0E1")
    _add_rule(slide, 0.68, 2.84, 3.0, ORANGE, 0.05)
    _add_text(slide, "Jason Trinh", 0.68, 3.21, 3.0, 0.35, size=18, color=WHITE, bold=True)
    _add_text(slide, "jasontrinh@berkeley.edu", 0.68, 3.63, 3.5, 0.30, size=14, color="BFD0E1")
    _add_text(slide, "Mentor", 6.02, 1.48, 1.30, 0.27, size=11, color="7FA0BC", bold=True)
    _add_text(slide, "James Casaletto", 7.43, 1.45, 4.55, 0.34, size=17, color=WHITE, bold=True)
    _add_text(slide, "Program", 6.02, 2.16, 1.30, 0.27, size=11, color="7FA0BC", bold=True)
    _add_text(slide, "NASA Space Life Sciences Training Program", 7.43, 2.13, 4.82, 0.52, size=15, color=WHITE)
    _add_text(slide, "Data", 6.02, 3.04, 1.30, 0.27, size=11, color="7FA0BC", bold=True)
    _add_text(slide, "NASA OSDR | ARCHS4 | Reactome", 7.43, 3.01, 4.55, 0.34, size=15, color=WHITE)
    _add_text(slide, "AI tools", 6.02, 3.73, 1.30, 0.27, size=11, color="7FA0BC", bold=True)
    _add_text(slide, "ChatGPT | Claude", 7.43, 3.70, 4.55, 0.34, size=15, color=WHITE)
    _add_text(slide, "Code and manuscript", 6.02, 4.42, 1.30, 0.40, size=11, color="7FA0BC", bold=True)
    _add_text(slide, "github.com/jasont314/nasa-mouse", 7.43, 4.40, 4.55, 0.34, size=15, color=WHITE)
    _add_text(slide, "All biological associations were tested in observed OSDR profiles.", 6.02, 5.50, 6.00, 0.40, size=13.5, color="BFD0E1", italic=True)


def _add_notes(slide, note: SlideNote) -> None:
    frame = slide.notes_slide.notes_text_frame
    frame.text = f"Target time: {note.time}\n\n{note.text}"


def _write_notes(notes: list[SlideNote]) -> None:
    total_seconds = sum(
        int(minutes) * 60 + int(seconds)
        for note in notes
        for minutes, seconds in [note.time.split(":", maxsplit=1)]
    )
    lines = [
        "# SLSTP 2026 mouse spaceflight transcriptomics speaker notes",
        "",
        (
            "Target length: 12-15 minutes. Suggested pacing totals "
            f"{total_seconds // 60}:{total_seconds % 60:02d}."
        ),
        "",
    ]
    for note in notes:
        lines.extend([
            f"## {note.number}. {note.title} ({note.time})",
            "",
            note.text,
            "",
        ])
    NOTES_PATH.write_text("\n".join(lines), encoding="utf-8")


def build() -> Path:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    utility_chart = _build_tissue_utility_chart()
    trajectory = _build_trajectory_crop()
    tissue_accession_pca = _build_pca_comparison_chart("tissue")
    condition_pca = _build_pca_comparison_chart("condition", include_study=False)
    architecture_figure = ASSET_DIR / "lacan_figure1c_generator_architecture.png"
    diffusion_gif = ASSET_DIR / "anim_diffdata.gif"
    diffusion_image_gif = ASSET_DIR / "anim_diffimg.gif"
    mober_gif = ASSET_DIR / "anim_mober.gif"
    wgan_gif = ASSET_DIR / "anim_wgan.gif"
    midpoint_slide2 = ASSET_DIR / "midpoint_slide2.png"
    midpoint_slide7_umap = ASSET_DIR / "midpoint_slide7_umap.png"
    midpoint_slide10_output_heatmap = ASSET_DIR / "midpoint_slide10_output_heatmap.png"
    midpoint_slide11_heatmap = ASSET_DIR / "midpoint_slide11_heatmap.png"
    midpoint_slide12_heatmap = ASSET_DIR / "midpoint_slide12_heatmap.png"
    midpoint_tiles = [
        ASSET_DIR / f"image-{slide_number}-{tile}.jpeg"
        for slide_number in (8, 9, 10, 20, 21)
        for tile in range(1, 5)
    ]
    required_assets = [
        architecture_figure,
        diffusion_gif,
        diffusion_image_gif,
        mober_gif,
        wgan_gif,
        midpoint_slide2,
        midpoint_slide7_umap,
        midpoint_slide10_output_heatmap,
        midpoint_slide11_heatmap,
        midpoint_slide12_heatmap,
        *midpoint_tiles,
    ]
    missing_assets = [path for path in required_assets if not path.exists()]
    if missing_assets:
        raise FileNotFoundError(
            "Missing presentation assets: "
            + ", ".join(path.name for path in missing_assets)
            + ". See presentation/final/source/assets/SOURCES.md"
        )
    presentation = Presentation(TEMPLATE)
    _set_title_slide(presentation.slides[0])
    _prepare_content_slide(presentation.slides[1], 2)
    builders = [
        None,
        lambda slide: _add_full_slide_image(
            slide,
            midpoint_slide2,
            alt="Original midpoint presentation slide 2: project objective and two project goals",
        ),
        _slide_autoencoder_foundation,
        lambda slide: _slide_study_effect_umap(slide, midpoint_slide7_umap),
        lambda slide: _add_midpoint_generator_slide(
            slide,
            "8",
            [
                (
                    mober_gif,
                    (9_637_395, 3_690_938, 7_490_460, 6_119_813),
                    "Animated MOBER example aligning three study distributions",
                )
            ],
            5,
        ),
        lambda slide: _add_midpoint_generator_slide(slide, "9", [], 6),
        _slide_expimap_method,
        lambda slide: _slide_expimap_program_scores(
            slide,
            midpoint_slide10_output_heatmap,
        ),
        lambda slide: _slide_expimap_literature_review(
            slide,
            midpoint_slide11_heatmap,
            midpoint_slide12_heatmap,
        ),
        _slide_expimap_tissue_results,
        _slide_why_synthetic,
        lambda slide: _slide_diffusion_explainer(
            slide, diffusion_image_gif, diffusion_gif, 12
        ),
        lambda slide: _slide_wgan_explainer(slide, wgan_gif, 13),
        _slide_3,
        lambda slide: _slide_4(slide, architecture_figure),
        lambda slide: _slide_5(slide, trajectory),
        lambda slide: _slide_condition_pca(slide, condition_pca),
        lambda slide: _slide_6(slide, tissue_accession_pca),
        _slide_matched_classifier_design,
        lambda slide: _slide_8(slide, utility_chart),
        _slide_feature_importance_venn,
        _slide_synthetic_interpretation_map,
        _slide_retained_outputs,
        _slide_tissue_method_convergence,
        _slide_literature_annotation_summary,
        _slide_thymus_integrated_result,
        _slide_skin_spleen_integrated_result,
        _slide_15,
        _slide_16,
    ]
    while len(presentation.slides) < len(builders):
        number = len(presentation.slides) + 1
        slide = presentation.slides.add_slide(presentation.slide_layouts[3])
        _prepare_content_slide(slide, number)

    for index, builder in enumerate(builders):
        if builder is not None:
            builder(presentation.slides[index])

    notes = [
        SlideNote(1, "Interpretable and generative models for mouse spaceflight", "0:15", "This project uses machine learning to study mouse bulk RNA-seq from NASA spaceflight experiments. One model asks which pathways change. The other asks whether realistic synthetic samples can improve a tissue-specific comparison of flight and ground control."),
        SlideNote(2, "Learn how spaceflight changes living systems", "0:35", "The project has two connected goals. First, identify which genes, pathways, and biological systems differ between flight and ground-control mice. Second, learn meaningful expression patterns and generate realistic profiles that can support that comparison."),
        SlideNote(3, "Autoencoders compress thousands of genes into a few features", "0:40", "An autoencoder compresses thousands of gene measurements into a smaller set of features and then reconstructs the original profile. The output bars are close to, but not exactly the same as, the input bars. Samples with similar compressed profiles sit near one another. MOBER uses this compressed space to reduce study identity, while expiMap connects its features to known Reactome gene programs."),
        SlideNote(4, "Study identity dominates the expression structure", "0:25", "This is EDL skeletal muscle from two OSDR studies. Color marks the study and shape marks flight or ground control. The two colors form separate UMAP clusters, while triangles and circles overlap within each cluster. Study identity is the strongest visible structure."),
        SlideNote(4, "MOBER tries to remove study identity", "0:35", "Study differences can be larger than the biological effect. MOBER combines an autoencoder with a source discriminator that tries to identify the originating study. The encoder makes study labels harder to predict while still reconstructing the expression profiles. The animation starts with separated studies and moves them into a shared distribution."),
        SlideNote(5, "MOBER reduces study separation, but FLT and GC still overlap", "0:30", "This UMAP compares the same two muscle studies before and after MOBER. Before correction, study identity dominates the layout. After MOBER, samples from the studies mix more closely, but flight and ground control still overlap. Removing visible study structure does not automatically reveal a strong spaceflight axis."),
        SlideNote(6, "expiMap assigns each latent feature to a known pathway", "0:35", "GLARE and MOBER learn their latent structure from expression. expiMap instead connects each latent node to a known mouse Reactome pathway through a masked decoder. That lets us compare flight and ground control one named biological program at a time. The next slide shows how those program scores are summarized."),
        SlideNote(4, "Program scores summarize pathway changes within one tissue", "0:45", "The table is an illustrative example for one tissue, not a project result. Each sample receives one score per Reactome program. We compare the average flight and ground-control scores and subtract GC from FLT. The rows are ordered from higher red shifts to lower blue shifts. The heatmap at right shows the complete program-by-sample output that these comparisons summarize."),
        SlideNote(5, "Each program was compared with prior literature", "0:30", "The top heatmap shows the pathway scores alone. The lower heatmap shows the same scores after each pathway name was reviewed against spaceflight literature. Green agrees with prior work, blue adds a related or complementary interpretation, orange is uncertain, red conflicts, and gray has little effect. The score values do not change; the colors describe the evidence attached to each pathway."),
        SlideNote(5, "Five tissues showed recurring pathway patterns", "0:55", "Thymus showed lower repair, cytoskeletal, and stromal-interaction programs. Skin showed lower chromatin regulation, repair, Hedgehog, sphingolipid, and cell-junction programs. Liver showed lower MHC class II and T-cell receptor scores. Spleen combined lower T-cell receptor, neutrophil-degranulation, and C-type lectin programs. Kidney showed higher ECM proteoglycan, WNT, and IGF-transport programs, suggesting structural and growth-factor remodeling."),
        SlideNote(6, "What is synthetic transcriptomics?", "0:25", "A generator learns patterns from measured RNA-seq and creates new numeric expression profiles for a chosen tissue and FLT or GC condition. These profiles are useful model outputs, but they are not new animals or independent biological measurements."),
        SlideNote(8, "Diffusion: denoise from noise", "0:30", "The left animation shows diffusion for images: start with random noise and remove it step by step until an image appears. The right side applies the same idea to gene expression. The phrase data manifold on the slide means the region occupied by real expression profiles."),
        SlideNote(9, "Conditional WGAN-GP: generator versus critic", "0:30", "A GAN trains two networks against each other. The generator makes expression profiles. The critic, which is similar to a discriminator, scores how much each profile resembles measured data. Tissue, FLT or GC, study, and sample material tell the generator what type of profile to make."),
        SlideNote(10, "Building the RNA-seq generator", "0:45", "We compared different ways to prepare the data, handle study differences, train the generator, and specify what it should make. Feature-space options included all shared genes, highly variable genes, Reactome genes, and a mouse mapping of the human L1000 landmark panel. The selected diffusion model uses the 974-gene L1000 mapping with TPM and ARCHS4-training MaxAbs scaling. It first learns broad mouse tissue patterns from ARCHS4, then adapts to OSDR."),
        SlideNote(11, "Diffusion best reproduced the measured expression distribution", "0:45", "We compared WGAN-GP with diffusion. Correlation measures gene-expression agreement. Coverage F1 asks whether generated samples cover the same regions as measured samples. Real-versus-synthetic accuracy asks whether a prediction model can tell the two apart, where 0.5 is ideal. Distribution distance measures overall separation, so lower is better. Diffusion had better coverage, near-chance discrimination, and the smaller distance, so we used it for the biological analysis."),
        SlideNote(12, "Diffusion learns tissue structure from noise", "0:25", "The same 1,024 generated samples begin as noise, develop tissue structure by step 200, and reach their final expression profiles at step zero. PCA compresses the 974 gene values to two axes so the movement can be plotted."),
        SlideNote(13, "FLT and GC overlap in the global PCA view", "0:25", "We first ask whether flight and ground-control profiles separate when all tissues and studies are viewed together. They overlap across the first two principal components, so there is no clear global condition axis. This motivates a closer look at the tissue and study structure."),
        SlideNote(14, "Tissue and study structure dominate the PCA space", "0:30", "Recoloring the same profiles by tissue and study reveals much stronger structure. The left panel shows tissue and the right shows study. Circles are real samples and crosses are generated samples. We therefore compare flight and ground control within each tissue and account for study in the downstream analysis."),
        SlideNote(15, "Does adding synthetic data improve FLT vs GC prediction?", "0:30", "We train the same tissue-specific classifier twice: once with real profiles and once with real plus synthetic profiles. Both versions are tested on the same real samples. This directly asks whether adding synthetic profiles improves flight-versus-ground-control prediction."),
        SlideNote(19, "Real + synthetic vs real-only balanced accuracy across 27 tissues", "0:45", "Each line compares real-only training with real-plus-synthetic training for flight-versus-ground-control classification in one tissue analysis. All scores come from held-out real OSDR samples. Balanced accuracy averages flight sensitivity and ground-control specificity. Teal points identify the 18 analyses where balanced accuracy, AUROC, and average precision all held or improved both overall and after giving each study equal weight. Coral marks the nine mixed results where at least one measure declined."),
        SlideNote(20, "Compare what each classifier finds important", "0:30", "Both classifiers see the same 974 genes and are evaluated on the same real samples. Features important in both classifiers are reinforced. Features that become important only after adding synthetic samples are promoted."),
        SlideNote(15, "We compared feature importance at three levels", "0:35", "Each analysis asks which features support flight-versus-ground-control prediction. Individual-gene permutation and SHAP score one gene at a time. Grouped permutation and grouped SHAP score a Reactome pathway together. Consensus ranking compares real and generated gene rankings and tests compact panels."),
        SlideNote(21, "The three analyses retained different types of results", "0:35", "The individual-gene analysis retained 21 associations. The grouped analysis retained 10 Reactome pathways, and consensus ranking retained 49 gene associations in compact panels. Promoted means the feature became important after synthetic samples were added. Reinforced means it was important with and without synthetic samples."),
        SlideNote(22, "All three analyses found candidates in thymus, skin and spleen", "0:40", "Thymus, skin, and spleen have retained candidates in every column. This agreement is why the biological interpretation focuses on these three tissues. The full analysis retains the single-method candidates for follow-up."),
        SlideNote(23, "Feature selection and literature review answer separate questions", "0:35", "The model comparison labels a candidate as promoted or reinforced. The literature review then fixes the candidate, tissue, and flight direction and asks whether prior work agrees, supports a related mechanism, is ambiguous, or offers no close match."),
        SlideNote(24, "Thymus: all three analyses point to lower cell-cycle activity in flight", "0:45", "Thymus has the clearest agreement. Fifteen individual genes, seven Reactome groups, and a 16-gene compact panel point toward lower mitotic activity in flight. This could reflect reduced proliferative renewal or a smaller fraction of cycling thymocytes. The expression data cannot distinguish those explanations."),
        SlideNote(25, "Skin and spleen show different flight-associated patterns", "0:45", "In skin, flight-higher PLSCR1 and two flight-higher necroptosis pathways suggest an interferon-linked regulated cell-death response. This is a complementary hypothesis because direct spaceflight evidence for skin necroptosis is limited. In spleen, flight-higher LOXL1, RAI14, PTPRK, and MYL9 suggest extracellular-matrix, adhesion, and cytoskeletal remodeling. This is also complementary because prior work supports the component mechanisms but not the same spleen-flight gene directions. The four genes did not form a significant shared Reactome pathway."),
        SlideNote(28, "What the models can and cannot tell us", "0:30", "The models summarize patterns in the RNA-seq data and narrow the list of tissues, pathways, and genes to examine. They do not establish mechanism. Wet-lab experiments must confirm the biological changes and distinguish gene regulation from shifts in cell composition."),
        SlideNote(29, "Thank you", "0:10", "Acknowledge James Casaletto, SLSTP, NASA OSDR, ARCHS4, Reactome, ChatGPT, and Claude, then invite questions."),
    ]
    final_times = (
        "0:10", "0:30", "0:35", "0:20", "0:30", "0:25", "0:30",
        "0:35", "0:25", "0:45", "0:20", "0:25", "0:25", "0:35",
        "0:35", "0:20", "0:20", "0:25", "0:25", "0:35", "0:25",
        "0:30", "0:30", "0:35", "0:30", "0:40", "0:40", "0:25",
        "0:10",
    )
    if len(notes) != len(final_times):
        raise ValueError("Presentation note count does not match final pacing plan")
    notes = [
        SlideNote(index, note.title, final_times[index - 1], note.text)
        for index, note in enumerate(notes, start=1)
    ]
    for note, slide in zip(notes, presentation.slides):
        _add_notes(slide, note)
    _write_notes(notes)
    presentation.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    print(build())
